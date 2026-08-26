import io
import uuid

from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from config.database import Base
from exports.executive_group_summary_pptx_builder import build_executive_group_summary_pptx
from models.models import User
from models.report_schemas import ReportConfiguration
from models.schemas import EvaluationData, PerformanceRecord
from services.report_service import ReportService


def _payload():
    driver = {
        "key": "response_time",
        "label": "Response Time",
        "actual": 62,
        "target": 30,
        "unit": "minutes",
        "direction": "lower_better",
        "achievement_pct": 48.4,
        "baseline_actual": 70,
        "mom": 4.0,
        "weighted_impact": 25.8,
        "status": "Requires action",
        "affected_scope": ["Offshore EGY", "Marketing"],
        "affected_regions": ["Offshore EGY"],
        "affected_teams": ["Marketing"],
    }
    group = {
        "name": "Marketing",
        "scope_type": "Team",
        "scope_label": "Offshore EGY",
        "score": 74.5,
        "baseline": 72.0,
        "movement": 2.5,
        "target": 100.0,
        "gap_to_target": 25.5,
        "status": "Watch",
        "record_count": 2,
        "below_threshold": 1,
        "kpis": [driver],
        "action_status": "Owner / due unresolved",
    }
    return {
        "period_label": "July 2026",
        "scope_label": "All authorized regions and teams",
        "overall": {
            "score": 82.1,
            "baseline": 80.0,
            "movement": 2.1,
            "target": 100.0,
            "status": "Watch",
        },
        "regions": [{**group, "name": "Offshore EGY", "scope_type": "Region"}],
        "teams": [group],
        "kpis": [driver],
        "drivers": [driver],
        "trend": [
            {"label": "June 2026", "score": 80.0, "target": 100.0, "record_count": 2},
            {"label": "July 2026", "score": 82.1, "target": 100.0, "record_count": 2},
        ],
        "history_count": 2,
        "best_period": {"label": "July 2026", "score": 82.1},
        "worst_period": {"label": "June 2026", "score": 80.0},
        "net_movement": 2.1,
        "actions": [{
            "is_proposed": True,
            "action_type": "Proposed management action",
            "action_text": "Validate response time workflow.",
            "team": "Marketing",
            "region": "EGY",
            "scope": "Offshore EGY, Marketing",
            "linked_kpi_key": "response_time",
            "owner": "Owner needed",
            "due_date": "Due date needed",
            "status": "Proposed",
            "priority": "High",
            "success_metric": "Move Response Time toward 30 min",
            "employee_name": "Authorized group scope",
        }],
        "raw_actions": [],
        "evidence": [],
        "evidence_rows": [[
            "Offshore EGY, Marketing",
            "Marketing",
            "Response Time",
            "62 min",
            "30 min",
            "48.4%",
            "25.8%",
            "KPI Evidence — Investigation Required",
        ]],
    }


def _text(presentation: Presentation) -> str:
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def test_executive_builder_uses_approved_story_order_and_group_safe_content():
    presentation = Presentation(io.BytesIO(build_executive_group_summary_pptx("July 2026", _payload())))

    assert len(presentation.slides) == 10
    text = _text(presentation)
    for title in (
        "Executive Group Performance Overview",
        "Region / Team Priority Map",
        "Performance by Region / Team",
        "Performance Trend & KPI Health",
        "Driver → Region / Team",
        "KPI Evidence & Root-Cause State",
        "Executive Commitment Tracker",
        "CEO Management Summary / Next Review",
        "Authorized Group Appendix",
    ):
        assert title in text
    assert "Response Time" in text
    assert "Lower is better" in text
    assert "Employee M1" not in text
    assert "Employee M2" not in text


def test_executive_builder_marks_missing_comparison_as_new():
    payload = _payload()
    payload["overall"]["movement"] = None
    payload["regions"][0]["movement"] = None
    payload["teams"][0]["movement"] = None
    payload["drivers"][0]["movement"] = None
    payload["kpis"][0]["movement"] = None

    text = _text(Presentation(io.BytesIO(build_executive_group_summary_pptx("July 2026", payload))))

    assert "MoM: NEW" in text


class _RecordService:
    def __init__(self, records):
        self.records = records

    def list_records(self, **filters):
        return [
            record
            for record in self.records
            if all(value is None or str(getattr(record, key, "")) == str(value) for key, value in filters.items())
        ]


def _record(employee_id, team, region, month, score, response_time):
    return PerformanceRecord(
        id=f"{employee_id}_{month}",
        employee_id=employee_id,
        employee_name=f"Employee {employee_id}",
        team=team,
        position="Analyst",
        region=region,
        performance_level="Employee",
        year=2026,
        month=month,
        status="Below" if score < 90 else "Meets",
        evaluation=EvaluationData(score=score, grade="B"),
        kpi_values=[{
            "kpi_key": "response_time",
            "label": "Response Time",
            "unit": "minutes",
            "direction": "lower_better",
            "actual_value": response_time,
            "target_value": 30,
            "achievement_ratio": 30 / response_time,
            "weight_applied": 1.0,
        }],
    )


def test_executive_report_is_live_comparable_direction_aware_and_integrated():
    engine = create_engine(
        "sqlite:///:memory:",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    Base.metadata.create_all(engine)
    session = sessionmaker(bind=engine)()
    user = User(
        id=uuid.uuid4(),
        username="executive-test",
        email="executive-test@example.com",
        password_hash="not-used",
        role="Admin",
    )
    session.add(user)
    session.commit()
    records = [
        _record("M1", "Marketing", "EGY", "June", 70, 75),
        _record("M2", "Marketing", "EGY", "June", 80, 65),
        _record("M1", "Marketing", "EGY", "July", 76, 60),
        _record("M2", "Marketing", "EGY", "July", 83, 64),
        _record("O1", "Operations", "UAE", "June", 95, 20),
        _record("O1", "Operations", "UAE", "July", 93, 25),
    ]
    service = ReportService(session, _RecordService(records))
    service.actions.list_active = lambda: []
    scope = {
        "user": user,
        "user_id": str(user.id),
        "role": "Admin",
        "accessible_teams": [],
        "accessible_team_levels": [],
        "is_general_manager": True,
        "legacy_unscoped": False,
    }
    configuration = ReportConfiguration(
        report_type="executive_group_summary",
        report_name="Executive Group Summary",
        start_month="July",
        start_year=2026,
        comparison_month="June",
        comparison_year=2026,
        included_sections=["summary", "details"],
        output_format="pptx",
    )

    payload = service._executive_group_presentation_data(configuration, scope)

    assert payload["overall"]["score"] == 84.0
    assert payload["overall"]["baseline"] == 81.66666666666667
    assert [(row["name"], row["score"]) for row in payload["regions"]] == [
        ("Offshore EGY", 79.5),
        ("UAE", 93.0),
    ]
    driver = payload["drivers"][0]
    assert driver["label"] == "Response Time"
    assert driver["direction"] == "lower_better"
    assert driver["actual"] > driver["target"]
    assert driver["weighted_impact"] > 0
    assert payload["actions"]
    assert payload["actions"][0]["owner"] == "Owner needed"
    assert payload["actions"][0]["due_date"] == "Due date needed"

    report = service.generate(configuration, scope)
    presentation = Presentation(io.BytesIO(bytes(report.file_data)))
    text = _text(presentation)
    assert report.file_name == "Executive_Group_Summary.pptx"
    assert report.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert "July 2026" in text
    assert "June 2026" in text
    assert "Response Time" in text
    assert "Lower is better" in text
    assert "Employee M1" not in text
    assert "Employee M2" not in text
