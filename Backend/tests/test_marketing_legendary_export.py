import io
import re

import pytest
from pptx import Presentation

from exports.marketing_legendary_pptx_builder import build_marketing_legendary_pptx
from services.insights_report_service import build_insights_snapshot


def _record(employee_id, month, score, kpis, *, name=None, position="Media Buyer"):
    return {
        "employee_id": employee_id,
        "employee_name": name or f"Employee {employee_id}",
        "team": "Marketing",
        "position": position,
        "region": "EGY",
        "year": 2026,
        "month": month,
        "score": score,
        "kpis": kpis,
    }


def _kpi(key, label, actual, target, weight, *, direction="higher_better", unit="%"):
    return {
        "kpi_key": key,
        "label": label,
        "actual_value": actual,
        "target_value": target,
        "weight_applied": weight,
        "direction": direction,
        "unit": unit,
    }


def _fixture_data():
    june_nadia = _record(
        "E1",
        "June",
        55,
        [
            _kpi("response_time", "Response Time", 90, 30, 0.40, direction="lower_better", unit="minutes"),
            _kpi("qualified_leads", "Qualified Leads", 60, 100, 0.40),
        ],
        name="Nadia Hassan",
    )
    june_omar = _record(
        "E2",
        "June",
        90,
        [
            _kpi("response_time", "Response Time", 40, 30, 0.40, direction="lower_better", unit="minutes"),
            _kpi("qualified_leads", "Qualified Leads", 110, 100, 0.40),
        ],
        name="Omar Ali",
        position="Marketing Manager",
    )
    july_nadia = _record(
        "E1",
        "July",
        60,
        [
            _kpi("response_time", "Response Time", 75, 30, 0.40, direction="lower_better", unit="minutes"),
            _kpi("qualified_leads", "Qualified Leads", 70, 100, 0.40),
            _kpi("fresh_campaigns", "Fresh Campaigns", 2, 5, 0.20, unit="count"),
        ],
        name="Nadia Hassan",
    )
    july_omar = _record(
        "E2",
        "July",
        94,
        [
            _kpi("response_time", "Response Time", 25, 30, 0.40, direction="lower_better", unit="minutes"),
            _kpi("qualified_leads", "Qualified Leads", 120, 100, 0.40),
            _kpi("fresh_campaigns", "Fresh Campaigns", 5, 5, 0.20, unit="count"),
        ],
        name="Omar Ali",
        position="Marketing Manager",
    )
    return {
        "scope_label": "Marketing",
        "filters": {"team": "Marketing"},
        "records": [july_nadia, july_omar],
        "selected_records": [july_nadia, july_omar],
        "history": [
            {"key": "2026-06", "label": "June 2026", "records": [june_nadia, june_omar]},
            {"key": "2026-07", "label": "July 2026", "records": [july_nadia, july_omar]},
        ],
        "actions": [
            {
                "action_type": "Coaching",
                "action_text": "Review response workflow",
                "linked_kpi_key": "response_time",
                "employee_id": "E1",
                "owner": None,
                "due_date": None,
                "status": "Open",
                "success_metric": "Response time at target",
                "root_cause_note": "Queue handoff needs review",
                "evidence_reference": "Ticket sample",
            }
        ],
    }


def _presentation(deck):
    return Presentation(io.BytesIO(deck))


def _deck_text(deck):
    presentation = _presentation(deck)
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def test_legendary_marketing_export_uses_dynamic_story_order_and_values():
    data = _fixture_data()
    deck = build_marketing_legendary_pptx("July 2026", data)
    presentation = _presentation(deck)
    text = _deck_text(deck)

    assert len(presentation.slides) >= 14
    expected_order = [
        "Marketing Team",
        "Monthly Performance Overview",
        "Performance by Role",
        "Performance Trend (2 periods)",
        "KPI Health Overview",
        "Performance Driver Impact",
        "Driver → Role → Employee",
        "Employee Performance Priority",
        "Employee Detail: Nadia Hassan",
        "Root Cause / Evidence",
        "Corrective Action Tracker",
        "Management Summary / Next Review",
        "Appendix: Full Employee Scorecard",
        "Appendix: Full KPI Reference",
    ]
    slide_text = [
        " | ".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in presentation.slides
    ]
    cursor = 0
    for expected in expected_order:
        index = next(index for index in range(cursor, len(slide_text)) if expected in slide_text[index])
        cursor = index + 1

    role_slide_text = [shape.text for shape in presentation.slides[2].shapes if hasattr(shape, "text")]
    assert any("Media Buyer" in value and "Nadia Hassan" in value for value in role_slide_text)
    assert any("Marketing Manager" in value and "Omar Ali" in value for value in role_slide_text)

    assert "July 2026" in text
    assert "Nadia Hassan" in text
    assert "Omar Ali" in text
    assert "Bahy Hamed Amer" not in text
    assert "Ziad sayed hassan" not in text
    assert "Response Time" in text
    assert "75 min" in text
    assert "30 min" in text
    assert "Lower is better" in text
    assert "NEW" in text
    assert "Owner needed" in text
    assert "Due date needed" in text
    assert "Confirmed Root Cause" in text
    assert "Investigation Required" in text
    assert not re.search(r"\d+(?:\.\d+)?\s*pp\b", text.casefold())
    assert not re.search(r"\d+(?:\.\d+)?\s*pts?\b", text.casefold())


def test_legendary_snapshot_preserves_direction_baseline_mom_and_loss_order():
    data = _fixture_data()
    from exports.marketing_legendary_pptx_builder import _prepare_snapshot

    payload = _prepare_snapshot("July 2026", data)
    assert payload["latest_period_label"] == "July 2026"
    assert payload["baseline_period_label"] == "June 2026"
    assert payload["baseline_score"] == 72.5
    assert payload["movement"] == 4.5

    response = next(row for row in payload["kpis"] if row["label"] == "Response Time")
    fresh = next(row for row in payload["kpis"] if row["label"] == "Fresh Campaigns")
    assert response["direction"] == "lower_better"
    assert response["achievement_pct"] == pytest.approx(60.0)
    assert response["baseline_actual"] == 65.0
    assert response["mom"] is not None
    assert fresh["mom"] is None
    assert payload["kpis"][0]["weighted_impact"] >= payload["kpis"][1]["weighted_impact"]


def test_legendary_export_is_not_july_specific_when_period_and_names_change():
    data = _fixture_data()
    for record in data["records"]:
        record["month"] = "September"
        record["employee_name"] = f"September {record['employee_name']}"
    for item in data["history"]:
        item["label"] = item["label"].replace("July", "September")
        for record in item["records"]:
            if record["month"] == "July":
                record["month"] = "September"
    deck = build_marketing_legendary_pptx("September 2026", data)
    text = _deck_text(deck)

    assert "September 2026" in text
    assert "September Nadia Hassan" in text
    assert "Marketing_PMS_Jul_2026_Legendary_Report" not in text
