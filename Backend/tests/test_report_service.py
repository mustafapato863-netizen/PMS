import io
import uuid
import zipfile
from types import SimpleNamespace

import pytest
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from config.database import Base
from models.models import GeneratedReport, SavedReportTemplate, User
from models.report_schemas import ReportConfiguration
from models.schemas import EvaluationData, PerformanceRecord
from services.insights_report_service import build_insights_snapshot
from services.report_service import ReportAccessError, ReportNotFoundError, ReportService, ReportValidationError


class StubRecordService:
    def __init__(self, records):
        self.records = records

    def list_records(self, **filters):
        return [
            record
            for record in self.records
            if all(value is None or str(getattr(record, key)) == str(value) for key, value in filters.items())
        ]


class RecordingRecordService:
    def __init__(self, records):
        self.records = records
        self.calls = []

    def list_records(self, **filters):
        self.calls.append(filters)
        return self.records


def _record(employee_id: str, team: str, level: str, month: str = "June") -> PerformanceRecord:
    return PerformanceRecord(
        id=f"{employee_id}_2026_{month}",
        employee_id=employee_id,
        employee_name=f"Employee {employee_id}",
        team=team,
        position="Analyst",
        region="EGY",
        performance_level=level,
        year=2026,
        month=month,
        status="Meets",
        evaluation=EvaluationData(score=91.5, grade="B"),
        kpi_values=[
            {
                "kpi_key": "quality",
                "label": "Quality",
                "actual_value": 95,
                "target_value": 90,
                "achievement_ratio": 1.055,
                "contribution": 0.25,
            }
        ],
    )


@pytest.fixture()
def db():
    engine = create_engine(
        "sqlite:///:memory:",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    Base.metadata.create_all(
        engine,
        tables=[User.__table__, GeneratedReport.__table__, SavedReportTemplate.__table__],
    )
    session = sessionmaker(bind=engine)()
    user = User(
        id=uuid.uuid4(),
        username="report-admin",
        email="reports@example.com",
        password_hash="not-used",
        role="Admin",
    )
    session.add(user)
    session.commit()
    yield session, user
    session.close()


def _admin_scope(user: User) -> dict:
    return {
        "user": user,
        "user_id": str(user.id),
        "role": "Admin",
        "employee_id": user.employee_id,
        "accessible_teams": [],
        "accessible_team_levels": [],
        "is_general_manager": True,
        "legacy_unscoped": False,
    }


def _configuration(**overrides) -> ReportConfiguration:
    values = {
        "report_type": "team",
        "report_name": "June Team Report",
        "start_month": "June",
        "start_year": 2026,
        "team": "Marketing",
        "included_sections": ["summary", "kpi_breakdown", "details"],
    }
    values.update(overrides)
    return ReportConfiguration(**values)


def test_options_are_restricted_to_manager_team_and_level(db):
    session, _user = db
    service = ReportService(
        session,
        StubRecordService([
            _record("EMP1", "Marketing", "Employee"),
            _record("EMP2", "Marketing", "Managerial"),
            _record("EMP3", "Sales", "Employee"),
        ]),
    )
    scope = {
        "role": "Manager",
        "accessible_teams": ["Marketing"],
        "accessible_team_levels": [("Marketing", "Managerial")],
        "is_general_manager": False,
        "legacy_unscoped": False,
    }

    options = service.options(scope)

    assert options["teams"] == ["Marketing"]
    assert options["performance_levels"] == ["Managerial"]
    assert [employee["id"] for employee in options["employees"]] == ["EMP2"]
    assert options["can_export"] is True


def test_admin_options_expose_all_report_levels_even_when_snapshot_has_employee_rows_only(db):
    session, user = db
    service = ReportService(
        session,
        StubRecordService([_record("EMP1", "Marketing", "Employee")]),
    )


def _management_analysis_record(employee_id: str = "MGR1", month: str = "June") -> dict:
    return {
        "id": f"{employee_id}_2026_{month}",
        "employee_id": employee_id,
        "employee_name": "Manager One",
        "team": "Marketing",
        "position": "Marketing Manager",
        "region": "EGY",
        "performance_level": "Managerial",
        "year": 2026,
        "month": month,
        "status": "Below",
        "evaluation": {"score": 86.0, "grade": "B"},
        "kpi_values": [{
            "kpi_key": "manager_quality",
            "label": "Manager Quality",
            "unit": "%",
            "direction": "higher_better",
            "actual_value": 86,
            "target_value": 100,
            "achievement_ratio": 0.86,
            "weight_applied": 1.0,
            "contribution": 0.86,
        }],
    }

    options = service.options(_admin_scope(user))

    assert options["performance_levels"] == ["Employee", "Managerial", "Corporate"]


def test_all_level_configuration_aggregates_the_full_team_scope(db):
    session, user = db
    service = ReportService(
        session,
        StubRecordService([
            _record("EMP1", "Marketing", "Employee"),
            _record("EMP2", "Marketing", "Employee"),
        ]),
    )
    service.record_service.records[1].evaluation.score = 80.0
    service.record_service.records[0].evaluation.score = 90.0

    configuration = _configuration(report_type="team_marketing", performance_level="All")

    assert configuration.performance_level is None
    data = service._performance_data(configuration, _admin_scope(user))

    assert data.summary["employee_count"] == 2
    assert data.summary["average_score"] == 85.0


def test_all_level_report_merges_management_bsc_rows_but_explicit_level_stays_scoped(db, monkeypatch):
    session, user = db
    manager_row = _management_analysis_record()
    monkeypatch.setattr(
        "services.report_service.ManagementBSCService.list_analysis_records",
        lambda _service: [manager_row],
    )
    service = ReportService(
        session,
        StubRecordService([_record("EMP1", "Marketing", "Employee")]),
    )
    scope = _admin_scope(user)

    all_records = service._performance_records(
        _configuration(report_type="team_marketing", performance_level="All"),
        scope,
    )
    assert {record.performance_level for record in all_records} == {"Employee", "Managerial"}
    assert {record.employee_id for record in all_records} == {"EMP1", "MGR1"}

    employee_records = service._performance_records(
        _configuration(report_type="team_marketing", performance_level="Employee"),
        scope,
    )
    assert [record.employee_id for record in employee_records] == ["EMP1"]

    managerial_records = service._performance_records(
        _configuration(report_type="team_marketing", performance_level="Managerial"),
        scope,
    )
    assert [record.employee_id for record in managerial_records] == ["MGR1"]


def test_preview_uses_strict_period_without_falling_back(db):
    session, user = db
    service = ReportService(session, StubRecordService([_record("EMP1", "Marketing", "Employee", "May")]))

    with pytest.raises(ReportNotFoundError):
        service.preview(_configuration(), _admin_scope(user))


def test_generate_persists_real_pptx_and_selected_sections(db):
    session, user = db
    service = ReportService(session, StubRecordService([_record("EMP1", "Marketing", "Employee")]))

    report = service.generate(_configuration(), _admin_scope(user))

    assert report.record_count == 1
    assert report.output_format == "pptx"
    assert report.file_name == "June_Team_Report.pptx"
    assert session.query(GeneratedReport).count() == 1
    with zipfile.ZipFile(io.BytesIO(report.file_data)) as archive:
        assert "ppt/presentation.xml" in archive.namelist()
        assert "ppt/slides/slide1.xml" in archive.namelist()
        assert archive.read("ppt/slides/slide1.xml").startswith(b"<?xml")


def test_marketing_pptx_uses_selected_period_records_and_actions(db):
    session, user = db
    marketing_record = _record("EMP1", "Marketing", "Employee", "July")
    marketing_record.year = 2026
    marketing_record.position = "Media Buyer"
    marketing_record.evaluation.score = 61.0
    marketing_record.kpi_values = [{
        "kpi_key": "mb_leads",
        "label": "Leads",
        "unit": "count",
        "direction": "higher_better",
        "actual_value": 35,
        "target_value": 100,
        "achievement_ratio": 0.35,
        "weight_applied": 0.30,
    }]
    service = ReportService(session, StubRecordService([marketing_record]))
    service.actions.list_active = lambda: []

    report = service.generate(
        _configuration(
            report_type="team_marketing",
            report_name="July Marketing Summary",
            start_month="July",
            start_year=2026,
            team="Marketing",
            performance_level="Employee",
            included_sections=["summary", "details"],
        ),
        _admin_scope(user),
    )

    assert report.output_format == "pptx"
    assert report.file_name == "July_Marketing_Summary.pptx"
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(report.file_data))
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "July 2026" in slide_text
    assert "Leads" in slide_text
    assert "61.0%" in slide_text
    assert "June 2026" not in slide_text


def test_marketing_export_forwards_scope_and_comparison_filters(db):
    session, user = db
    current = _record("EMP1", "Marketing", "Employee", "June")
    previous = _record("EMP1", "Marketing", "Employee", "April")
    record_service = RecordingRecordService([current, previous])
    service = ReportService(session, record_service)
    service.actions.list_active = lambda: []

    configuration = _configuration(
        report_type="team_marketing",
        report_name="June Marketing Filter Check",
        start_month="June",
        start_year=2026,
        comparison_month="April",
        comparison_year=2026,
        region="EGY",
        team="Marketing",
        position="Analyst",
        performance_level="Employee",
        employee_id="EMP1",
        grade="B",
        status="Meets",
        kpi="quality",
        included_sections=["summary", "details"],
    )

    payload = service._marketing_presentation_data(configuration, _admin_scope(user))

    assert record_service.calls
    assert record_service.calls[0] == {
        "team": "Marketing",
        "employee_id": "EMP1",
        "grade": "B",
        "status": "Meets",
        "performance_level": "Employee",
        "position": "Analyst",
        "region": "EGY",
        "kpi": "quality",
    }
    assert payload["filters"]["comparison_period"] == "2026-04"
    assert payload["scope_label"] == (
        "Region EGY | Team Marketing | Level Employee | Position Analyst | "
        "Employee EMP1 | Grade B | Status Meets | KPI quality"
    )


def test_marketing_export_scopes_actions_to_the_same_filters_as_records(db):
    session, user = db
    good = _record("EMP1", "Marketing", "Employee", "June")
    good.position = "Analyst"

    wrong_grade = _record("EMP2", "Marketing", "Employee", "June")
    wrong_grade.position = "Analyst"
    wrong_grade.evaluation.grade = "A"

    wrong_status = _record("EMP3", "Marketing", "Employee", "June")
    wrong_status.position = "Analyst"
    wrong_status.status = "Below"

    team = SimpleNamespace(display_name="Marketing", name="Marketing", region="EGY")

    def action(employee_id: str, *, linked_kpi: str = "quality", region: str = "EGY"):
        employee = SimpleNamespace(
            employee_id=employee_id,
            name=f"Employee {employee_id}",
            region=region,
            position_name="Analyst",
            performance_level="Employee",
        )
        return SimpleNamespace(
            year=2026,
            month="June",
            team=team,
            employee=employee,
            action_type="Coaching",
            action_text="Review the KPI evidence.",
            root_cause_note="",
            status="Open",
            priority="High",
            linked_kpi_key=linked_kpi,
            owner=None,
            due_date=None,
            plan=None,
            evidence_reference="",
        )

    service = ReportService(
        session,
        RecordingRecordService([good, wrong_grade, wrong_status]),
    )
    service.actions.list_active = lambda: [
        action("EMP1"),
        action("EMP2"),
        action("EMP3"),
        action("EMP1", linked_kpi="other"),
        action("EMP1", region="UAE"),
    ]

    configuration = _configuration(
        report_type="team_marketing",
        report_name="June Marketing Action Filter Check",
        start_month="June",
        start_year=2026,
        region="Egypt",
        team="Marketing",
        position="Analyst",
        performance_level="Employee",
        employee_id=None,
        grade="B",
        status="Meets",
        kpi="quality",
    )

    payload = service._marketing_presentation_data(configuration, _admin_scope(user))

    assert [row["employee_id"] for row in payload["actions"]] == ["EMP1"]


def test_report_configuration_validates_comparison_period_pair():
    with pytest.raises(ValueError, match="Comparison month and comparison year"):
        _configuration(comparison_month="April")


def test_marketing_snapshot_uses_the_explicit_comparison_period():
    def row(month: str, score: float) -> dict:
        return {
            "employee_id": "EMP1",
            "employee_name": "Employee EMP1",
            "team": "Marketing",
            "year": 2026,
            "month": month,
            "score": score,
            "kpis": [],
        }

    snapshot = build_insights_snapshot({
        "period_label": "June 2026",
        "filters": {"comparison_period": "2026-04"},
        "selected_records": [row("June", 80)],
        "history": [
            {"key": "2026-04", "label": "April 2026", "records": [row("April", 50)]},
            {"key": "2026-05", "label": "May 2026", "records": [row("May", 70)]},
            {"key": "2026-06", "label": "June 2026", "records": [row("June", 80)]},
        ],
    })

    assert snapshot["comparison_period_label"] == "April 2026"
    assert snapshot["movement"] == pytest.approx(30)


def test_insights_pptx_uses_filtered_period_history_and_current_scope(db):
    session, user = db
    current = _record("EMP1", "Marketing", "Employee", "June")
    previous = _record("EMP1", "Marketing", "Employee", "May")
    current.evaluation.score = 72.0
    previous.evaluation.score = 64.0
    current.kpi_values = [{
        "kpi_key": "quality",
        "label": "Quality",
        "unit": "%",
        "direction": "higher_better",
        "actual_value": 72,
        "target_value": 100,
        "achievement_ratio": 0.72,
        "weight_applied": 0.50,
    }]
    previous.kpi_values = [{
        "kpi_key": "quality",
        "label": "Quality",
        "unit": "%",
        "direction": "higher_better",
        "actual_value": 64,
        "target_value": 100,
        "achievement_ratio": 0.64,
        "weight_applied": 0.50,
    }]
    service = ReportService(session, StubRecordService([previous, current]))
    service.actions.list_active = lambda: []

    report = service.generate(
        _configuration(
            report_type="insights",
            report_name="June Insights",
            start_month="June",
            start_year=2026,
            team="Marketing",
            performance_level="Employee",
            included_sections=["summary", "team_breakdown", "kpi_breakdown", "details"],
            output_format="pptx",
        ),
        _admin_scope(user),
    )

    assert report.output_format == "pptx"
    assert report.file_name == "June_Insights.pptx"
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(report.file_data))
    assert len(presentation.slides) == 8
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "EXECUTIVE DECISION SNAPSHOT" in slide_text
    assert "June 2026" in slide_text
    assert "Quality" in slide_text
    assert "May 2026" in slide_text
    assert "PEOPLE FOCUS" in slide_text
    assert "KPI signal only" in slide_text
    assert "Proposed" in slide_text
    assert "points" not in slide_text.casefold()
    assert "pts" not in slide_text.casefold()


def test_marketing_pptx_uses_configured_volume_rollup_and_average_score(db):
    session, user = db
    first = _record("EMP1", "Marketing", "Employee", "June")
    second = _record("EMP2", "Marketing", "Employee", "June")
    first.position = second.position = "Graphic Designer"
    first.evaluation.score = 60.0
    second.evaluation.score = 90.0
    for record, actual in ((first, 150), (second, 100)):
        record.kpi_values = [{
            "kpi_key": "gd_on_schedule",
            "label": "Projects delivered on schedule",
            "unit": "count",
            "direction": "higher_better",
            "actual_value": actual,
            "target_value": 100,
            "achievement_ratio": actual / 100,
            "weight_applied": 0.40,
        }]
    service = ReportService(session, StubRecordService([first, second]))
    service.actions.list_active = lambda: []

    report = service.generate(
        _configuration(
            report_type="team_marketing",
            report_name="June Marketing Rollup",
            start_month="June",
            start_year=2026,
            team="Marketing",
            performance_level="Employee",
            included_sections=["summary", "details"],
        ),
        _admin_scope(user),
    )

    from pptx import Presentation

    presentation = Presentation(io.BytesIO(report.file_data))
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Monthly Performance Overview" in slide_text
    assert "CURRENT MONTHLY PERFORMANCE" in slide_text
    assert "75.0%" in slide_text
    assert "250" in slide_text
    assert "125" not in slide_text


def test_marketing_pptx_balances_performance_by_employee_across_periods(db):
    first_may = _record("EMP1", "Marketing", "Employee", "May")
    first_june = _record("EMP1", "Marketing", "Employee", "June")
    second_june = _record("EMP2", "Marketing", "Employee", "June")
    for record in (first_may, first_june, second_june):
        record.position = "Graphic Designer"
        record.kpi_values = [{
            "kpi_key": "gd_on_schedule",
            "label": "Projects delivered on schedule",
            "unit": "count",
            "direction": "higher_better",
            "actual_value": 100,
            "target_value": 100,
            "achievement_ratio": 1,
            "weight_applied": 0.40,
        }]
    first_may.evaluation.score = 60.0
    first_june.evaluation.score = 80.0
    second_june.evaluation.score = 90.0

    service = ReportService(
        db[0],
        StubRecordService([first_may, first_june, second_june]),
    )
    service.actions.list_active = lambda: []

    report = service.generate(
        _configuration(
            report_type="team_marketing",
            report_name="May To June Marketing Rollup",
            start_month="May",
            start_year=2026,
            end_month="June",
            end_year=2026,
            team="Marketing",
            performance_level="Employee",
            included_sections=["summary", "details"],
        ),
        _admin_scope(db[1]),
    )

    from pptx import Presentation

    presentation = Presentation(io.BytesIO(report.file_data))
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    # Multi-month Marketing exports average the canonical period-level scores
    # (May 60% and June 85%), rather than pooling employees across periods.
    assert "Performance Trend (2 periods)" in slide_text
    assert "72.5%" in slide_text
    assert "May 2026" in slide_text
    assert "June 2026" in slide_text


def test_generate_pdf_export_uses_pdf_content_type(db):
    session, user = db
    service = ReportService(
        session,
        StubRecordService([_record("EMP1", "Marketing", "Employee")]),
    )

    report = service.generate(_configuration(output_format="pdf"), _admin_scope(user))

    assert report.output_format == "pdf"
    assert report.file_name == "June_Team_Report.pdf"
    assert report.content_type == "application/pdf"
    assert report.file_data.startswith(b"%PDF-1.4")


def test_generate_excel_export_is_a_real_workbook_and_idempotency_lookup_is_stable(db):
    session, user = db
    service = ReportService(
        session,
        StubRecordService([_record("EMP1", "Marketing", "Employee")]),
    )

    report = service.generate(
        _configuration(output_format="excel"),
        _admin_scope(user),
        idempotency_key="user-1:reports-center-excel",
    )

    assert report.output_format == "excel"
    assert report.file_name == "June_Team_Report.xlsx"
    assert report.content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    with zipfile.ZipFile(io.BytesIO(report.file_data)) as archive:
        assert "[Content_Types].xml" in archive.namelist()
        assert "xl/workbook.xml" in archive.namelist()
    found = service.find_generated_by_idempotency(_admin_scope(user), "user-1:reports-center-excel")
    assert found is not None
    assert found.id == report.id


def test_generate_rejects_team_outside_scope(db):
    session, user = db
    service = ReportService(session, StubRecordService([_record("EMP1", "Marketing", "Employee")]))
    scope = _admin_scope(user) | {
        "role": "Manager",
        "is_general_manager": False,
        "accessible_teams": ["Sales"],
    }

    with pytest.raises(ReportAccessError):
        service.generate(_configuration(), scope)


def test_generate_rolls_back_when_persistence_fails(db, monkeypatch):
    session, user = db
    service = ReportService(session, StubRecordService([_record("EMP1", "Marketing", "Employee")]))

    def fail_after_add(report):
        session.add(report)
        raise RuntimeError("storage unavailable")

    monkeypatch.setattr(service.reports, "add_generated", fail_after_add)

    with pytest.raises(RuntimeError, match="storage unavailable"):
        service.generate(_configuration(), _admin_scope(user))

    assert session.query(GeneratedReport).count() == 0


def test_delete_generated_removes_the_authorized_report(db):
    session, user = db
    service = ReportService(session, StubRecordService([_record("EMP1", "Marketing", "Employee")]))
    report = service.generate(_configuration(), _admin_scope(user))

    deleted = service.delete_generated(str(report.id), _admin_scope(user))

    assert deleted == {"id": str(report.id), "name": "June Team Report"}
    assert session.query(GeneratedReport).count() == 0


def test_delete_generated_many_removes_selected_reports_atomically(db):
    session, user = db
    service = ReportService(session, StubRecordService([_record("EMP1", "Marketing", "Employee")]))
    first = service.generate(_configuration(report_name="June Team Report"), _admin_scope(user))
    second = service.generate(_configuration(report_name="June Marketing Report"), _admin_scope(user))

    deleted = service.delete_generated_many([str(first.id), str(second.id), str(first.id)], _admin_scope(user))

    assert deleted == [
        {"id": str(first.id), "name": "June Team Report"},
        {"id": str(second.id), "name": "June Marketing Report"},
    ]
    assert session.query(GeneratedReport).count() == 0


def test_delete_generated_many_validates_all_reports_before_deleting(db):
    session, user = db
    service = ReportService(session, StubRecordService([_record("EMP1", "Marketing", "Employee")]))
    report = service.generate(_configuration(), _admin_scope(user))

    with pytest.raises(ReportNotFoundError, match="not found"):
        service.delete_generated_many([str(report.id), str(uuid.uuid4())], _admin_scope(user))

    assert session.query(GeneratedReport).count() == 1


def test_delete_generated_rejects_an_unknown_report(db):
    session, user = db
    service = ReportService(session, StubRecordService([]))

    with pytest.raises(ReportNotFoundError, match="not found"):
        service.delete_generated(str(uuid.uuid4()), _admin_scope(user))


def test_saved_template_contains_configuration_only_and_rejects_duplicates(db):
    session, user = db
    service = ReportService(session, StubRecordService([_record("EMP1", "Marketing", "Employee")]))
    configuration = _configuration()

    saved = service.save_template("Monthly Marketing", configuration, _admin_scope(user))

    assert saved.configuration["team"] == "Marketing"
    assert "file_data" not in saved.configuration
    with pytest.raises(ReportValidationError, match="already exists"):
        service.save_template("Monthly Marketing", configuration, _admin_scope(user))
