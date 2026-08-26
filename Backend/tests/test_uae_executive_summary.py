import io
import uuid

from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from config.database import Base
from exports.uae_executive_summary_pptx_builder import build_uae_executive_summary_pptx
from models.models import GeneratedReport, User
from models.report_schemas import ReportConfiguration
from services.report_service import ReportService
from services.uae_executive_summary import build_uae_executive_summary_contract


def _record(employee_id, team, month, score, *, actual=None, direction="higher_better", kpi="Quality", name=None):
    actual = score if actual is None else actual
    return {
        "employee_id": employee_id,
        "employee_name": name or f"Employee {employee_id}",
        "team": team,
        "position": "Analyst",
        "region": "UAE",
        "performance_level": "Employee",
        "year": 2026,
        "month": month,
        "status": "Below" if score < 70 else "Meets",
        "evaluation": {"score": score, "grade": "B"},
        "kpi_values": [{
            "kpi_key": kpi.casefold().replace(" ", "_"),
            "label": kpi,
            "unit": "minutes" if direction == "lower_better" else "%",
            "direction": direction,
            "actual_value": actual,
            "target_value": 30 if direction == "lower_better" else 90,
            "weight_applied": 50,
        }],
    }


def test_contract_uses_one_population_and_ranks_all_departments():
    records = [
        _record("E1", "CSR", "June", 60),
        _record("E2", "PMO", "June", 95),
        _record("E3", "CRM", "June", 82),
        _record("E1", "CSR", "May", 70),
        _record("E2", "PMO", "May", 90),
        _record("E3", "CRM", "May", 82),
    ]
    payload = build_uae_executive_summary_contract(
        records,
        current_period=(2026, 6),
        known_departments=["CSR", "PMO", "CRM", "Marketing"],
        aggregate_only=True,
    )

    assert payload["active_departments"] == payload["cover"]["active_departments"] == 3
    assert [row["department"] for row in payload["departments"]] == ["CSR", "CRM", "PMO"]
    assert payload["no_current_departments"] == ["Marketing"]
    assert payload["employees_requiring_action"] == 1
    assert payload["departments"][0]["priority_components"]
    assert payload["departments"][0]["priority_score"] > payload["departments"][-1]["priority_score"]


def test_contract_direction_aware_loss_and_new_baseline():
    records = [
        _record("E1", "RCM", "June", 80, actual=60, direction="lower_better", kpi="Response Time"),
        _record("E1", "RCM", "May", 80, actual=60, direction="lower_better", kpi="Response Time"),
        _record("E2", "RCM", "June", 80, actual=80, direction="higher_better", kpi="Quality"),
    ]
    payload = build_uae_executive_summary_contract(records, current_period=(2026, 6), aggregate_only=True)
    response = next(row for row in payload["kpis"] if row["label"] == "Response Time")
    assert response["direction"] == "lower_better"
    assert response["status"] == "On Track"
    assert response["weighted_loss"] == 0
    assert payload["mom"] == 0  # baseline uses only the like-for-like E1 population
    assert all("Employee E1" not in str(row) for row in payload["employee_priority_rows"])

    new_payload = build_uae_executive_summary_contract(
        [_record("E1", "RCM", "June", 80)],
        current_period=(2026, 6),
        aggregate_only=True,
    )
    assert new_payload["mom"] is None


def test_contract_action_ownership_metrics_and_redaction():
    payload = build_uae_executive_summary_contract(
        [_record("E1", "CSR", "June", 60, name="Should Not Leak")],
        current_period=(2026, 6),
        actions=[{
            "employee_id": "E1",
            "team": "CSR",
            "year": 2026,
            "month": "June",
            "linked_kpi_key": "quality",
            "action_text": "Validate quality workflow",
            "status": "Open",
            "priority": "High",
            "owner": None,
            "due_date": None,
            "success_metric": None,
        }],
        aggregate_only=True,
    )
    assert payload["open_action_count"] == 1
    assert payload["missing_owner_count"] == 1
    assert payload["missing_due_date_count"] == 1
    assert payload["high_priority_unresolved"] == 1
    assert "Should Not Leak" not in str(payload)


def test_builder_has_concise_story_and_no_employee_names_for_aggregate_scope():
    payload = build_uae_executive_summary_contract(
        [_record("E1", "CSR", "June", 60, name="Should Not Leak")],
        current_period=(2026, 6),
        known_departments=["CSR"],
        aggregate_only=True,
    )
    presentation = Presentation(io.BytesIO(build_uae_executive_summary_pptx("June 2026", payload)))
    text = "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert 10 <= len(presentation.slides) <= 12
    for title in (
        "UAE Executive Overview",
        "UAE Department Performance Ranking",
        "UAE Risk and Opportunity Map",
        "UAE Performance Trend",
        "Largest KPI Performance Drivers",
        "Driver → Department → Function → Impact",
        "Action Ownership and Execution Risk",
        "Management Decisions Required",
        "Next Review Commitments",
    ):
        assert title in text
    assert "Should Not Leak" not in text


class _RecordService:
    def __init__(self, records):
        self.records = records

    def list_records(self, **filters):
        return self.records


def test_catalogue_route_persistence_and_filename_for_new_type():
    engine = create_engine("sqlite:///:memory:", connect_args={"check_same_thread": False}, poolclass=StaticPool)
    Base.metadata.create_all(engine, tables=[User.__table__, GeneratedReport.__table__])
    session = sessionmaker(bind=engine)()
    user = User(
        id=uuid.uuid4(),
        username="uae-executive-test",
        email="uae-executive-test@example.com",
        password_hash="not-used",
        role="Admin",
    )
    session.add(user)
    session.commit()
    service = ReportService(session, _RecordService([_record("E1", "CSR", "June", 60)]))
    service.actions.list_active = lambda: []
    scope = {
        "user": user,
        "user_id": str(user.id),
        "role": "Admin",
        "accessible_teams": [],
        "accessible_team_levels": [],
        "is_general_manager": True,
        "legacy_unscoped": False,
    }
    template = next(item for item in service.templates() if item["type"] == "uae_executive_summary")
    assert template["name"] == "UAE Executive Performance Summary"
    assert template["formats"] == ["pptx"]
    configuration = ReportConfiguration(
        report_type="uae_executive_summary",
        report_name="June UAE Executive Performance Summary",
        start_month="June",
        start_year=2026,
        region="UAE",
        included_sections=["summary", "details"],
        output_format="pptx",
    )
    report = service.generate(configuration, scope)
    assert report.file_name == "June_UAE_Executive_Performance_Summary.pptx"
    assert report.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert session.query(GeneratedReport).count() == 1
    assert service.get_download(str(report.id), scope).file_data == report.file_data
