import io

from pptx import Presentation
from pptx.dml.color import RGBColor

from exports.marketing_pptx_builder import TEMPLATE_PATH, _fmt_native, _kpi_label, _primary_kpi, build_marketing_pptx


def _record(employee_id, month, score, kpis, *, name=None, position="Media Buyer"):
    return {
        "employee_id": employee_id,
        "employee_name": name or f"Employee {employee_id}",
        "team": "Marketing",
        "position": position,
        "region": "EGY",
        "year": 2026,
        "month": month,
        "score": score,
        "kpis": kpis,
    }


def _kpi(key, label, actual, target, weight, *, direction="higher_better", unit="%"):
    return {
        "kpi_key": key,
        "label": label,
        "actual_value": actual,
        "target_value": target,
        "weight_applied": weight,
        "direction": direction,
        "unit": unit,
    }


def _deck_text(deck):
    presentation = Presentation(io.BytesIO(deck))
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def test_marketing_native_percentage_formats_fraction_and_normalized_values():
    assert _fmt_native(0.195, "%") == "19.5%"
    assert _fmt_native(19.5, "%") == "19.5%"


def test_marketing_action_and_kpi_labels_use_configured_display_names():
    payload = {"kpis": [{"key": "sms_response_time", "label": "Response Time"}]}
    assert _kpi_label(payload, "sms_response_time") == "Response Time"


def test_marketing_export_uses_supplied_template_and_replaces_sample_rows():
    assert TEMPLATE_PATH.exists()
    kpi = _kpi("response_time", "Response Time", 75, 30, 0.5, direction="lower_better", unit="minutes")
    july = _record("E1", "July", 60, [kpi], name="Nadia Hassan")
    june = _record("E1", "June", 50, [dict(kpi, actual_value=90)], name="Nadia Hassan")

    deck = build_marketing_pptx(
        "July 2026",
        {
            "scope_label": "Marketing | EGY",
            "selected_records": [july],
            "records": [july],
            "history": [
                {"key": "2026-06", "label": "June 2026", "records": [june]},
                {"key": "2026-07", "label": "July 2026", "records": [july]},
            ],
            "actions": [],
        },
    )

    presentation = Presentation(io.BytesIO(deck))
    text = _deck_text(deck)

    assert len(presentation.slides) == 12
    assert "Marketing Team Performance Overview" in text
    assert "KPI Loss Breakdown" in text
    assert "Response Time" in text
    assert "75 min" in text
    assert "30 min" in text
    assert "40.0%" in text
    assert "Achievement gap 60.0%" in text
    assert "Weighted loss" in text
    assert "Primary KPI | weighted loss" in text
    assert "Top KPI | weighted loss" in text
    assert "Next Review Commitments" in text
    assert "Nadia Hassan" in text
    assert "Bahy Hamed Amer" not in text
    assert "Ziad sayed hassan" not in text
    assert "Priority" in text
    assert "A%endix" not in text
    assert "pts" not in text.casefold()
    assert "pp" not in {word.casefold() for word in text.split()}


def test_marketing_primary_kpi_is_deterministic_when_all_kpis_are_on_track():
    person = {
        "leading_kpi": {"label": "Unstable source order", "weighted_impact": 0},
        "kpis": [
            {"label": "Second configured KPI", "weighted_impact": 0, "definition": {"display_order": 2}},
            {"label": "First configured KPI", "weighted_impact": 0, "definition": {"display_order": 1}},
        ],
    }

    assert _primary_kpi(person)["label"] == "First configured KPI"


def test_marketing_export_keeps_unavailable_mom_neutral():
    record = _record(
        "E1",
        "July",
        72,
        [_kpi("quality", "Quality", 72, 100, 0.5)],
    )
    deck = build_marketing_pptx("July 2026", {"records": [record], "selected_records": [record]})
    slide = Presentation(io.BytesIO(deck)).slides[0]

    movement_card = next(shape for shape in slide.shapes if shape.name == "Shape 11")
    assert movement_card.fill.fore_color.rgb == RGBColor(245, 246, 248)


def test_marketing_weighted_loss_bars_are_proportional():
    record = _record(
        "E1",
        "July",
        60,
        [
            _kpi("high", "High weight", 50, 100, 0.8),
            _kpi("low", "Low weight", 50, 100, 0.2),
        ],
    )
    deck = build_marketing_pptx("July 2026", {"records": [record], "selected_records": [record]})
    slide = Presentation(io.BytesIO(deck)).slides[1]
    widths = {
        shape.name: shape.width
        for shape in slide.shapes
        if shape.name in {"Shape 35", "Shape 40"}
    }

    assert widths["Shape 35"] > widths["Shape 40"]
