from __future__ import annotations

import io

from pptx import Presentation

from exports.insights_pptx_builder import build_insights_pptx
from services.insights_report_service import build_insights_snapshot, kpi_achievement


def _record(
    employee_id: str,
    month: str,
    score: float,
    *,
    team: str = "Marketing",
    kpis: list[dict] | None = None,
    year: int = 2026,
) -> dict:
    return {
        "employee_id": employee_id,
        "employee_name": f"Employee {employee_id}",
        "team": team,
        "position": "Analyst",
        "year": year,
        "month": month,
        "score": score,
        "kpis": kpis or [],
    }


def _kpi(key: str, label: str, actual: float, target: float, weight: float, *, direction: str = "higher_better", unit: str = "%") -> dict:
    return {
        "kpi_key": key,
        "label": label,
        "actual_value": actual,
        "target_value": target,
        "weight_applied": weight,
        "direction": direction,
        "unit": unit,
    }


def _report(records: list[dict], actions: list[dict] | None = None, *, scope: str = "Marketing") -> dict:
    by_period: dict[tuple[int, int], list[dict]] = {}
    month_numbers = {
        "January": 1,
        "February": 2,
        "March": 3,
        "April": 4,
        "May": 5,
        "June": 6,
        "July": 7,
    }
    for record in records:
        period = (record["year"], month_numbers[record["month"]])
        by_period.setdefault(period, []).append(record)
    history = [
        {
            "key": f"{year}-{month:02d}",
            "label": f"{next(name for name, value in month_numbers.items() if value == month)} {year}",
            "records": rows,
        }
        for (year, month), rows in sorted(by_period.items())
    ]
    return {
        "period_label": "Selected periods",
        "scope_label": scope,
        "filters": {},
        "selected_records": records,
        "history": history,
        "actions": actions or [],
    }


def _deck_text(deck: bytes, slide_number: int | None = None) -> str:
    presentation = Presentation(io.BytesIO(deck))
    slides = presentation.slides if slide_number is None else [presentation.slides[slide_number - 1]]
    return "\n".join(
        shape.text
        for slide in slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def test_achievement_respects_direction_and_caps_overachievement():
    assert kpi_achievement(80, 100, "higher_better") == 0.8
    assert kpi_achievement(75, 30, "lower_better") == 0.4
    assert kpi_achievement(30, 30, "lower_better") == 1.0
    assert kpi_achievement(150, 100, "higher_better") == 1.0


def test_lower_better_driver_keeps_native_units_and_direction_visible_in_deck():
    records = [_record(
        "E1",
        "June",
        60,
        kpis=[_kpi("response_time", "Response Time", 75, 30, 0.5, direction="lower_better", unit="minutes")],
    )]
    report = _report(records)
    snapshot = build_insights_snapshot(report)
    assert snapshot["driver"]["achievement_pct"] == 40.0
    assert snapshot["driver"]["gap"] == -60.0
    assert snapshot["driver"]["actual"] == 75.0
    assert snapshot["driver"]["target"] == 30.0

    deck = build_insights_pptx("June 2026", report)
    presentation = Presentation(io.BytesIO(deck))
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Lower is better" in slide_text
    assert "75 minutes" in slide_text
    assert "30 minutes" in slide_text
    assert "40.0%" in slide_text


def test_missing_target_is_data_quality_and_never_becomes_a_leading_driver():
    records = [_record("E1", "June", 80, kpis=[_kpi("quality", "Quality", 80, 0, 0.5)])]
    snapshot = build_insights_snapshot(_report(records))

    assert snapshot["kpis"][0]["status"] == "Data quality"
    assert snapshot["kpis"][0]["achievement"] is None
    assert snapshot["driver"] is None


def test_multi_month_performance_averages_period_level_kpi_values():
    records = [
        _record("E1", "May", 50, kpis=[_kpi("quality", "Quality", 50, 100, 0.5)]),
        _record("E1", "June", 75, kpis=[_kpi("quality", "Quality", 90, 120, 0.5)]),
    ]
    snapshot = build_insights_snapshot(_report(records))

    assert snapshot["overall_score"] == 62.5
    assert snapshot["trend"][0]["label"] == "May 2026"
    assert snapshot["trend"][1]["label"] == "June 2026"
    # (50 / 100 + 90 / 120) / 2, not 140 / 220.
    assert round(snapshot["kpis"][0]["achievement_pct"], 2) == 62.5
    assert snapshot["movement"] == 25


def test_one_measured_period_has_no_trend_and_uses_current_period_snapshot():
    records = [_record("E1", "June", 72, kpis=[_kpi("quality", "Quality", 72, 100, 0.5)])]
    snapshot = build_insights_snapshot(_report(records))

    assert snapshot["history_count"] == 1
    assert snapshot["trend_status"] == "Trend unavailable"
    assert snapshot["movement_status"] == "single_period"
    assert "only one measured period" in snapshot["trend_headline"]
    assert snapshot["movement"] is None

    deck = build_insights_pptx("June 2026", _report(records))
    page_two = _deck_text(deck, 2)
    assert "CURRENT-PERIOD SNAPSHOT" in page_two
    assert "OVERALL PERFORMANCE TREND" not in page_two
    assert "Trend unavailable" in page_two


def test_two_periods_show_movement_without_claiming_a_sustained_trend():
    records = [
        _record("E1", "May", 60, kpis=[_kpi("quality", "Quality", 60, 100, 0.5)]),
        _record("E1", "June", 72, kpis=[_kpi("quality", "Quality", 72, 100, 0.5)]),
    ]
    snapshot = build_insights_snapshot(_report(records))

    assert snapshot["history_count"] == 2
    assert snapshot["trend_status"] == "Movement only"
    assert snapshot["trend_sustained"] is False
    assert snapshot["movement"] == 12
    assert "sustained trend not established" in snapshot["trend_headline"]
    assert "improving" not in _deck_text(build_insights_pptx("June 2026", _report(records)), 2).casefold()


def test_three_periods_can_classify_a_sustained_improving_trend():
    records = [
        _record("E1", "April", 60, kpis=[_kpi("quality", "Quality", 60, 100, 0.5)]),
        _record("E1", "May", 72, kpis=[_kpi("quality", "Quality", 72, 100, 0.5)]),
        _record("E1", "June", 84, kpis=[_kpi("quality", "Quality", 84, 100, 0.5)]),
    ]
    snapshot = build_insights_snapshot(_report(records))

    assert snapshot["history_count"] == 3
    assert snapshot["trend_status"] == "Improving"
    assert snapshot["trend_sustained"] is True
    assert snapshot["trend"][0]["label"] == "April 2026"
    assert snapshot["trend"][-1]["label"] == "June 2026"


def test_single_team_uses_health_snapshot_and_multi_team_adds_comparison_page():
    one_team = [_record("E1", "June", 72, kpis=[_kpi("quality", "Quality", 72, 100, 0.5)])]
    single_snapshot = build_insights_snapshot(_report(one_team))
    assert single_snapshot["team_health"]["employee_count"] == 1
    assert single_snapshot["team_health"]["employees_below_target"] == 1
    single_deck = build_insights_pptx("June 2026", _report(one_team))
    assert len(Presentation(io.BytesIO(single_deck)).slides) == 8
    assert "TEAM HEALTH SNAPSHOT" in _deck_text(single_deck, 2)
    assert "TEAM COMPARISON" not in _deck_text(single_deck)

    multi_team = [
        _record("E1", "June", 72, team="Marketing", kpis=[_kpi("quality", "Quality", 72, 100, 0.5)]),
        _record("E2", "June", 88, team="Sales", kpis=[_kpi("quality", "Quality", 88, 100, 0.5)]),
    ]
    multi_deck = build_insights_pptx("June 2026", _report(multi_team, scope="Marketing | Sales"))
    assert len(Presentation(io.BytesIO(multi_deck)).slides) == 9
    assert "TEAM COMPARISON" in _deck_text(multi_deck)


def test_weighted_loss_bars_keep_exact_proportional_widths():
    records = [_record(
        "E1",
        "June",
        60,
        kpis=[
            _kpi("driver_a", "Driver A", 40, 100, 0.6),
            _kpi("driver_b", "Driver B", 60, 100, 0.4),
        ],
    )]
    deck = build_insights_pptx("June 2026", _report(records))
    slide = Presentation(io.BytesIO(deck)).slides[2]
    widths = {
        shape.name: shape.width
        for shape in slide.shapes
        if shape.name.startswith("weighted-loss-bar-")
    }
    assert widths["weighted-loss-bar-36.0000"] > widths["weighted-loss-bar-16.0000"]


def test_weighted_driver_ranking_beats_raw_achievement_ranking():
    records = [
        _record(
            "E1",
            "June",
            60,
            kpis=[
                _kpi("high_weight", "High weight", 80, 100, 0.8),
                _kpi("low_weight", "Low weight", 50, 100, 0.2),
            ],
        )
    ]
    snapshot = build_insights_snapshot(_report(records))

    assert snapshot["driver"]["key"] == "high_weight"
    assert round(snapshot["driver"]["weighted_impact"], 2) == 16.0
    assert snapshot["kpis"][1]["key"] == "low_weight"


def test_same_display_kpi_keys_merge_weighted_loss_and_match_alias_actions():
    records = [
        _record(
            "E1",
            "June",
            72,
            kpis=[_kpi("sms_response_time", "Response Time", 75, 30, 0.4, direction="lower_better", unit="min")],
        ),
        _record(
            "M1",
            "June",
            88,
            kpis=[_kpi("response_time", "Response Time", 75, 30, 0.15, direction="lower_better", unit="min")],
        ),
    ]
    actions = [{
        "team": "Marketing",
        "action_text": "Review response queue",
        "status": "Open",
        "linked_kpi_key": "sms_response_time",
    }]

    snapshot = build_insights_snapshot(_report(records, actions))

    response_rows = [row for row in snapshot["kpis"] if row["label"] == "Response Time"]
    assert len(response_rows) == 1
    assert set(response_rows[0]["source_keys"]) == {"response_time", "sms_response_time"}
    assert round(response_rows[0]["weighted_impact"], 1) == 33.0
    assert snapshot["driver"]["label"] == "Response Time"
    assert snapshot["proposed_action_count"] == 0
    assert "30.0 min" in snapshot["actions"][0]["success_metric_display"]


def test_uncovered_high_impact_kpi_becomes_explicit_proposal():
    records = [_record(
        "E1",
        "June",
        60,
        kpis=[
            _kpi("quality", "Quality", 60, 100, 0.6),
            _kpi("volume", "Volume", 50, 100, 0.4),
        ],
    )]
    actions = [{
        "employee_id": "E1",
        "team": "Marketing",
        "action_text": "Review quality",
        "status": "Open",
        "linked_kpi_key": "quality",
    }]

    snapshot = build_insights_snapshot(_report(records, actions))

    assert snapshot["proposed_action_count"] == 1
    proposal = next(row for row in snapshot["actions"] if row["source_display"] == "Proposed")
    assert proposal["linked_kpi_key"] == "volume"
    assert proposal["owner_display"] == "Owner needed"
    assert proposal["due_date_display"] == "Due date needed"
    assert "configured KPI target" in proposal["success_metric_display"]


def test_employee_impact_is_ranked_by_weighted_loss():
    records = [
        _record("E1", "June", 90, kpis=[_kpi("quality", "Quality", 50, 100, 0.8)]),
        _record("E2", "June", 60, kpis=[_kpi("quality", "Quality", 80, 100, 0.8)]),
    ]
    snapshot = build_insights_snapshot(_report(records))

    assert [row["employee_id"] for row in snapshot["people"]] == ["E1", "E2"]
    assert snapshot["people"][0]["weighted_loss_pct"] > snapshot["people"][1]["weighted_loss_pct"]


def test_root_cause_and_action_missing_fields_are_explicit():
    records = [_record("E1", "June", 60, kpis=[_kpi("quality", "Quality", 60, 100, 0.5)])]
    actions = [{
        "employee_id": "E1",
        "employee_name": "Employee E1",
        "team": "Marketing",
        "action_text": "Review quality workflow",
        "status": "Open",
        "linked_kpi_key": "quality",
        "root_cause_note": "Handoff delay is being investigated",
    }]
    snapshot = build_insights_snapshot(_report(records, actions))

    assert snapshot["root_cause_rows"][0]["evidence_status"] == "Evidence recorded — cause pending confirmation"
    assert snapshot["root_cause_rows"][0]["recorded_root_cause"] == "Handoff delay is being investigated"
    assert snapshot["actions"][0]["owner_display"] == "Owner needed"
    assert snapshot["actions"][0]["due_date_display"] == "Due date needed"
    assert snapshot["actions"][0]["success_metric_display"] == "Reach configured KPI target: 100.0%"


def test_confirmed_root_cause_requires_recorded_evidence_reference():
    records = [_record("E1", "June", 60, kpis=[_kpi("quality", "Quality", 60, 100, 0.5)])]
    actions = [{
        "employee_id": "E1",
        "team": "Marketing",
        "action_text": "Review quality workflow",
        "status": "Open",
        "linked_kpi_key": "quality",
        "root_cause_note": "Handoff delay confirmed by workflow review",
        "evidence_reference": "case-123",
    }]
    snapshot = build_insights_snapshot(_report(records, actions))
    assert snapshot["root_cause_rows"][0]["evidence_status"] == "Confirmed cause"


def test_missing_root_cause_requires_validation_and_direction_aware_escalation():
    records = [_record(
        "E1",
        "June",
        60,
        kpis=[_kpi("response_time", "Response Time", 75, 30, 0.5, direction="lower_better", unit="minutes")],
    )]
    snapshot = build_insights_snapshot(_report(records))

    assert snapshot["root_cause_rows"][0]["evidence_status"] == "KPI signal only"
    assert snapshot["next_review"]["root_cause_requirement"] == "Validate the operational cause before assigning corrective action."
    assert "remains above target" in snapshot["next_review"]["escalation_rule"]


def test_missing_recorded_action_creates_transparent_proposal_using_configured_target():
    records = [_record("E1", "June", 60, kpis=[_kpi("quality", "Quality", 60, 100, 0.5)])]
    snapshot = build_insights_snapshot(_report(records))

    assert snapshot["recorded_action_count"] == 0
    assert snapshot["proposed_action_count"] == 1
    assert snapshot["actions"][0]["source_display"] == "Proposed"
    assert snapshot["actions"][0]["status_display"] == "Proposed"
    assert "100.0%" in snapshot["actions"][0]["success_metric_display"]


def test_one_team_and_multi_team_scope_are_reported_without_losing_rows():
    records = [
        _record("E1", "June", 80, team="Marketing", kpis=[_kpi("quality", "Quality", 80, 100, 0.5)]),
        _record("E2", "June", 90, team="Sales", kpis=[_kpi("quality", "Quality", 90, 100, 0.5)]),
    ]
    snapshot = build_insights_snapshot(_report(records, scope="Marketing | Sales"))

    assert snapshot["team_count"] == 2
    assert {row["name"] for row in snapshot["teams"]} == {"Marketing", "Sales"}
    assert snapshot["population_size"] == 2


def test_employee_detail_appendix_is_added_only_when_primary_table_overflows():
    records = [
        _record(f"E{index}", "June", 60 + index, kpis=[_kpi("quality", "Quality", 40 + index, 100, 0.5)])
        for index in range(1, 6)
    ]
    snapshot = build_insights_snapshot(_report(records))
    deck = build_insights_pptx("June 2026", _report(records))

    assert len(snapshot["people"]) == 5
    presentation = Presentation(io.BytesIO(deck))
    assert len(presentation.slides) == 9
    assert "EMPLOYEE DETAIL APPENDIX" in _deck_text(deck, 7)
    assert "Additional weighted-loss contributors require follow-through" in _deck_text(deck, 7)


def test_empty_result_is_safe_and_has_no_fabricated_driver_or_action():
    snapshot = build_insights_snapshot({"scope_label": "Marketing", "actions": []})

    assert snapshot["kpis"] == []
    assert snapshot["driver"] is None
    assert snapshot["actions"] == []
    assert snapshot["warnings"]


def test_generated_deck_has_eight_pages_and_consistent_vocabulary():
    records = [_record("E1", "June", 72, kpis=[_kpi("quality", "Quality", 72, 100, 0.5)])]
    deck = build_insights_pptx("June 2026", _report(records))
    presentation = Presentation(io.BytesIO(deck))
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )

    assert len(presentation.slides) == 8
    assert "EXECUTIVE DECISION SNAPSHOT" in slide_text
    assert "PERFORMANCE LOSS BREAKDOWN" in slide_text
    assert "ROOT CAUSE EVIDENCE" in slide_text
    assert "ACTION PLAN" in slide_text
    assert "NEXT REVIEW QUESTION" in slide_text
    assert "points" not in slide_text.casefold()
    assert "pts" not in slide_text.casefold()
    assert "pp" not in {word.casefold() for word in slide_text.split()}
