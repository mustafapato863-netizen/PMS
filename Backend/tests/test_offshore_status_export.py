import io
import uuid

from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from config.database import Base
from models.models import User
from models.report_schemas import ReportConfiguration
from services.offshore_status_report import (
    COMPLETED,
    NOT_MEASURED,
    PENDING,
    build_offshore_status_payload,
)
from services.report_service import ReportService
from exports.offshore_status_pptx_builder import TEMPLATE_PATH, build_offshore_status_pptx


def _row(team, position, employee_name, employee_id, kpi_status, position_status="Current"):
    return {
        "team": team,
        "position": position,
        "employee_name": employee_name,
        "employee_id": employee_id,
        "raw_data": {
            "KPI Status": kpi_status,
            "Position Status": position_status,
        },
    }


def _records():
    return [
        _row("Marketing", "Campaign Specialist", "Bob Example", "M1", PENDING),
        _row("Marketing", "Campaign Specialist", "Carla Example", "M2", PENDING),
        _row("Marketing", "Content Executive", "Dina Example", "M3", NOT_MEASURED),
        _row("Operations", "Operations Manager", "Omar Example", "O1", COMPLETED),
        _row("HR", "HR Officer", "Vacant", "", NOT_MEASURED, "Vacant"),
    ]


def _presentation_payload():
    def kpi(key, label, actual, target, weight, direction="higher_better", unit="%"):
        return {
            "kpi_key": key,
            "label": label,
            "actual_value": actual,
            "target_value": target,
            "weight_applied": weight,
            "direction": direction,
            "unit": unit,
        }

    current = [
        {
            "employee_id": "M1",
            "employee_name": "Bob Example",
            "team": "Marketing",
            "position": "Campaign Specialist",
            "year": 2026,
            "month": "July",
            "score": 65,
            "grade": "C",
            "kpis": [
                kpi("response_time", "Response Time", 75, 30, 0.5, "lower_better", "minutes"),
                kpi("lead_quality", "Lead Quality", 68, 90, 0.5),
            ],
        },
        {
            "employee_id": "M2",
            "employee_name": "Carla Example",
            "team": "Marketing",
            "position": "Content Executive",
            "year": 2026,
            "month": "July",
            "score": 91,
            "grade": "A",
            "kpis": [
                kpi("response_time", "Response Time", 22, 30, 0.5, "lower_better", "minutes"),
                kpi("lead_quality", "Lead Quality", 92, 90, 0.5),
            ],
        },
    ]
    previous = [{**row, "month": "June", "score": row["score"] - 3} for row in current]
    return {
        "period_label": "July 2026",
        "scope_label": "Offshore EGY",
        "filters": {"region": "Egypt"},
        "records": current,
        "selected_records": current,
        "history": [
            {"key": "2026-06", "label": "June 2026", "records": previous},
            {"key": "2026-07", "label": "July 2026", "records": current},
        ],
        "actions": [],
    }


def test_payload_uses_explicit_statuses_and_keeps_marketing_names():
    payload = build_offshore_status_payload(
        _records(),
        period_label="July 2026",
        region_label="Offshore EGY",
    )

    assert payload["status_counts"] == {
        COMPLETED: 1,
        PENDING: 1,
        NOT_MEASURED: 2,
    }
    assert payload["position_status_counts"]["Vacant"] == 1

    campaign = next(row for row in payload["detail_rows"] if row["position"] == "Campaign Specialist")
    assert campaign["employee_names"] == ["Bob Example", "Carla Example"]
    assert campaign["count"] == 2


def test_builder_is_source_deck_shaped_and_contains_dynamic_marketing_names():
    payload = _presentation_payload()
    presentation = Presentation(io.BytesIO(build_offshore_status_pptx("July 2026", payload)))

    assert TEMPLATE_PATH.name == "Offshore_Departments_Performance_Review.pptx"
    assert len(presentation.slides) >= 20
    assert round(presentation.slide_width / 914400, 3) == 13.333
    assert round(presentation.slide_height / 914400, 3) == 7.5
    text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Offshore Departments" in text
    assert "Combined Monthly Performance Overview" in text
    assert "Marketing Performance by Role" in text
    assert "Marketing Root Cause Evidence" in text
    assert "Combined Management Summary / Next Review" in text
    assert "Offshore PMS KPIs Dashboard" not in text
    assert "Bob Example" in text
    assert "Carla Example" in text
    assert "July 2026" in text
    assert "Lower is better" in text


class _RecordService:
    def __init__(self, records):
        self.records = records

    def list_records(self, **filters):
        return [
            record
            for record in self.records
            if all(
                value is None or str(getattr(record, key, "")) == str(value)
                for key, value in filters.items()
            )
        ]


def _scope(user):
    return {
        "user": user,
        "user_id": str(user.id),
        "role": "Admin",
        "accessible_teams": [],
        "accessible_team_levels": [],
        "is_general_manager": True,
        "legacy_unscoped": False,
    }


def test_monthly_egypt_pptx_uses_offshore_source_route():
    engine = create_engine(
        "sqlite:///:memory:",
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
    )
    # The corrected route reuses the shared Marketing payload builder, which
    # legitimately reads actions and management/team scope tables even when
    # the fixture has no rows in them.
    Base.metadata.create_all(engine)
    session = sessionmaker(bind=engine)()
    user = User(
        id=uuid.uuid4(),
        username="offshore-test",
        email="offshore-test@example.com",
        password_hash="not-used",
        role="Admin",
    )
    session.add(user)
    session.commit()

    try:
        from models.schemas import EvaluationData, PerformanceRecord

        records = [
            PerformanceRecord(
                id="M1_2026_July",
                employee_id="M1",
                employee_name="Bob Example",
                team="Marketing",
                position="Campaign Specialist",
                region="EGY",
                performance_level="Employee",
                year=2026,
                month="July",
                status="Meets",
                evaluation=EvaluationData(score=91, grade="A"),
                raw_data={"KPI Status": PENDING, "Position Status": "Current"},
            )
        ]
        service = ReportService(session, _RecordService(records))
        configuration = ReportConfiguration(
            report_type="monthly_egypt",
            report_name="July Offshore Dashboard",
            start_month="July",
            start_year=2026,
            region="Egypt",
            output_format="pptx",
            included_sections=["summary", "details"],
        )

        report = service.generate(configuration, _scope(user))
        presentation = Presentation(io.BytesIO(bytes(report.file_data)))
        text = "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )

        assert report.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        assert report.file_name.endswith(".pptx")
        assert "Offshore Departments" in text
        assert "Combined Monthly Performance Overview" in text
        assert "Offshore PMS KPIs Dashboard" not in text
        assert "Bob Example" in text

        uae_record = records[0].model_copy(
            update={"id": "M1_2026_July_UAE", "region": "UAE"}
        )
        uae_service = ReportService(session, _RecordService([uae_record]))
        uae_configuration = configuration.model_copy(
            update={
                "report_type": "monthly_uae",
                "report_name": "July UAE Dashboard",
                "region": "United Arab Emirates",
            }
        )
        uae_report = uae_service.generate(uae_configuration, _scope(user))
        uae_presentation = Presentation(io.BytesIO(bytes(uae_report.file_data)))
        uae_text = "\n".join(
            shape.text
            for slide in uae_presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        assert "Offshore Departments" in uae_text
        assert "Combined Monthly Performance Overview" in uae_text
        assert "UAE PMS KPIs Dashboard" not in uae_text
    finally:
        session.close()
