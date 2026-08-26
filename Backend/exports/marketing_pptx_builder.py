"""Build the Marketing Main Score + MoM PowerPoint from a filtered snapshot.

The Reports page uses this builder for the ``team_marketing`` template.  The
approved Marketing deck is kept as an editable source template and this module
only replaces its inherited text, bars, badges, and table rows.  All report
math is delegated to :func:`services.insights_report_service.build_insights_snapshot`,
which in turn delegates KPI aggregation and direction-aware scoring to the
canonical KPI utilities.
"""

from __future__ import annotations

import io
import re
from collections import defaultdict
from pathlib import Path
from statistics import mean
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from services.insights_report_service import build_insights_snapshot


SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
TEMPLATE_PATH = Path(__file__).with_name("templates") / "Marketing_PMS_July_2026_Final_PerfectDesign.pptx"

# These values intentionally mirror the supplied Marketing source deck.
NAVY = RGBColor(27, 42, 74)
MUTED = RGBColor(88, 105, 132)
FAINT = RGBColor(103, 117, 142)
BLUE = RGBColor(53, 106, 230)
GREEN = RGBColor(34, 160, 90)
RED = RGBColor(229, 72, 77)
AMBER = RGBColor(232, 145, 0)
PURPLE = RGBColor(105, 84, 190)
NEUTRAL = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(228, 231, 236)
PALE_BLUE = RGBColor(238, 244, 255)
PALE_GREEN = RGBColor(232, 248, 239)
PALE_RED = RGBColor(253, 237, 237)
PALE_AMBER = RGBColor(254, 243, 224)
PALE_NEUTRAL = RGBColor(245, 246, 248)
TRACK = RGBColor(228, 231, 236)

MONTHS = {
    name: index
    for index, name in enumerate(
        (
            "January",
            "February",
            "March",
            "April",
            "May",
            "June",
            "July",
            "August",
            "September",
            "October",
            "November",
            "December",
        ),
        1,
    )
}


def _number(value: Any) -> float | None:
    if value is None or value == "":
        return None
    try:
        number = float(value)
    except (TypeError, ValueError):
        return None
    return number if number == number else None


def _clean(value: Any, fallback: str = "") -> str:
    if value is None:
        return fallback
    text = " ".join(str(value).split())
    text = (
        text.replace("â€”", "-")
        .replace("â€“", "-")
        .replace("â€¢", "|")
        .replace("â€™", "'")
        .replace("A%endix", "Appendix")
    )
    return text or fallback


def _report_team_label(payload: dict[str, Any]) -> str:
    """Resolve the displayed team from the same filtered snapshot as the deck."""
    filters = payload.get("filters") or {}
    configured_team = _clean(filters.get("team"))
    if configured_team:
        return configured_team
    records = payload.get("records") or payload.get("selected_records") or []
    teams = sorted({_clean(record.get("team")) for record in records if _clean(record.get("team"))})
    return teams[0] if len(teams) == 1 else "Selected"


def _fit(value: Any, limit: int = 58) -> str:
    """Keep labels readable without leaving unexplained ellipses in a deck."""

    text = _clean(value)
    if len(text) <= limit:
        return text
    words = text.split()
    result = ""
    for word in words:
        candidate = f"{result} {word}".strip()
        if len(candidate) > limit:
            break
        result = candidate
    return result or text[:limit]


def _fmt_percent(value: Any, *, signed: bool = False, fallback: str = "N/A") -> str:
    number = _number(value)
    if number is None:
        return fallback
    return f"{number:+.1f}%" if signed else f"{number:.1f}%"


def _fmt_native(value: Any, unit: Any = "") -> str:
    number = _number(value)
    if number is None:
        return "N/A"
    normalized = _clean(unit).casefold()
    if normalized in {"%", "percent", "percentage"}:
        # The canonical KPI payload stores percentage actuals/targets as
        # fractions (for example 0.195 = 19.5%).  Preserve already-normalized
        # percentage values from integrations that provide 19.5 directly.
        if abs(number) <= 1:
            number *= 100
        return _fmt_percent(number)
    if normalized in {"minutes", "minute", "min"}:
        return f"{number:,.0f} min" if number.is_integer() else f"{number:,.1f} min"
    if normalized in {"seconds", "second", "sec"}:
        return f"{number:,.0f} sec" if number.is_integer() else f"{number:,.1f} sec"
    if normalized in {"count", "number", "visits"}:
        return f"{number:,.0f}" if number.is_integer() else f"{number:,.1f}"
    if number.is_integer():
        formatted = f"{number:,.0f}"
    else:
        formatted = f"{number:,.2f}".rstrip("0").rstrip(".")
    return f"{formatted} {_clean(unit)}".strip()


def _kpi_gap_text(kpi: dict[str, Any]) -> str:
    label = _clean(kpi.get("label"), "No KPI gap")
    achievement = _fmt_percent(kpi.get("achievement_pct"))
    loss = _fmt_percent(kpi.get("weighted_impact"))
    return f"{label} — {achievement} | loss {loss}"


def _period_key(record: dict[str, Any]) -> tuple[int, int] | None:
    try:
        year = int(record.get("year"))
    except (TypeError, ValueError):
        return None
    raw_month = record.get("month")
    try:
        month = int(raw_month)
    except (TypeError, ValueError):
        month = MONTHS.get(_clean(raw_month).title())
    return (year, month) if month else None


def _history_from_records(records: list[dict[str, Any]], period_label: str) -> list[dict[str, Any]]:
    grouped: dict[tuple[int, int], list[dict[str, Any]]] = defaultdict(list)
    for record in records:
        period = _period_key(record)
        if period:
            grouped[period].append(record)
    if not grouped:
        return [{"key": period_label, "label": period_label, "records": records}] if records else []
    labels = {number: name for name, number in MONTHS.items()}
    return [
        {
            "key": f"{year}-{month:02d}",
            "label": f"{labels.get(month, month)} {year}",
            "records": rows,
        }
        for (year, month), rows in sorted(grouped.items())
    ]


def _prepare_payload(period_label: str, report_data: dict[str, Any] | None) -> dict[str, Any]:
    """Create the one snapshot used by every slide in the template."""

    data = dict(report_data or {})
    records = [dict(record) for record in (data.get("records") or [])]
    selected_records = [dict(record) for record in (data.get("selected_records") or records)]
    history = list(data.get("history") or [])
    if not history:
        history = _history_from_records(selected_records, period_label)
    snapshot_input = {
        "period_label": data.get("period_label") or period_label,
        "scope_label": data.get("scope_label") or "Marketing",
        "filters": dict(data.get("filters") or {}),
        "aggregate_only": bool(data.get("aggregate_only")),
        "kpi_definitions": list(data.get("kpi_definitions") or []),
        "records": records,
        "selected_records": selected_records,
        "history": history,
        "actions": list(data.get("actions") or []),
    }
    return build_insights_snapshot(snapshot_input)


def _shape_text(shape: Any) -> str:
    return _clean(getattr(shape, "text", ""))


def _has_text(shape: Any) -> bool:
    return bool(getattr(shape, "has_text_frame", False))


def _find_shape(slide: Any, name: str) -> Any | None:
    return next((shape for shape in slide.shapes if shape.name == name), None)


def _set_shape_text(shape: Any | None, value: Any, *, color: RGBColor | None = None, bold: bool | None = None, size: float | None = None) -> None:
    if shape is None:
        return
    frame = getattr(shape, "text_frame", None)
    if frame is None:
        return
    old_runs = [run for paragraph in frame.paragraphs for run in paragraph.runs]
    old_font = old_runs[0].font if old_runs else None
    old_name = old_font.name if old_font else None
    old_size = old_font.size if old_font else None
    old_bold = old_font.bold if old_font else None
    shape.text = "" if value is None else str(value)
    frame.word_wrap = True
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            if old_name:
                run.font.name = old_name
            if size is not None:
                from pptx.util import Pt

                run.font.size = Pt(size)
            elif old_size:
                run.font.size = old_size
            if bold is not None:
                run.font.bold = bold
            elif old_bold is not None:
                run.font.bold = old_bold
            if color is not None:
                run.font.color.rgb = color


def _set_named_text(slide: Any, name: str, value: Any, **style: Any) -> None:
    _set_shape_text(_find_shape(slide, name), value, **style)


def _set_fill(shape: Any | None, fill: RGBColor, line: RGBColor | None = None) -> None:
    if shape is None:
        return
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    except (AttributeError, TypeError):
        return
    if line is not None:
        try:
            shape.line.color.rgb = line
        except (AttributeError, TypeError):
            pass


def _set_bar(shape: Any | None, width: float, color: RGBColor) -> None:
    if shape is None:
        return
    shape.width = max(1, int(width))
    _set_fill(shape, color, color)


def _inches(value: Any) -> float:
    return float(value) / 914400


def _text_groups(slide: Any, *, top_min: float, top_max: float, threshold: float = 0.12) -> list[list[Any]]:
    candidates = [
        shape
        for shape in slide.shapes
        if _has_text(shape)
        and _shape_text(shape)
        and top_min <= _inches(shape.top) <= top_max
    ]
    candidates.sort(key=lambda shape: (_inches(shape.top), _inches(shape.left)))
    groups: list[list[Any]] = []
    for shape in candidates:
        top = _inches(shape.top)
        if not groups or top - groups[-1][0] > threshold:
            groups.append([top, shape])
        else:
            groups[-1].append(shape)
    return [sorted(group[1:], key=lambda shape: _inches(shape.left)) for group in groups]


def _status_color(status: Any) -> RGBColor:
    value = _clean(status).casefold()
    if value in {"critical", "below target", "below_target"}:
        return RED
    if value in {"at risk", "watch", "pending", "proposed"}:
        return AMBER
    if value in {"on track", "improving", "recorded"}:
        return GREEN
    return BLUE


def _status_fill(status: Any) -> RGBColor:
    value = _clean(status).casefold()
    if value in {"critical", "below target", "below_target"}:
        return PALE_RED
    if value in {"at risk", "watch", "pending", "proposed"}:
        return PALE_AMBER
    if value in {"on track", "improving", "recorded"}:
        return PALE_GREEN
    return PALE_NEUTRAL


def _movement_text(value: Any) -> str:
    return _fmt_percent(value, signed=True)


def _movement_color(value: Any) -> RGBColor:
    movement = _number(value)
    if movement is None or movement == 0:
        return NEUTRAL
    return GREEN if movement > 0 else RED


def _movement_fill(value: Any) -> RGBColor:
    movement = _number(value)
    if movement is None or movement == 0:
        return PALE_NEUTRAL
    return PALE_GREEN if movement > 0 else PALE_RED


def _person_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    if not payload.get("people_visible"):
        return []
    return [dict(row) for row in (payload.get("all_people") or [])]


def _primary_kpi(person: dict[str, Any]) -> dict[str, Any]:
    """Return a deterministic primary KPI for tie cases.

    The canonical snapshot already identifies the leading loss KPI.  When a
    person is fully on track, several KPIs can have zero loss and the source
    order is not a meaningful ranking.  Use configured display order (then
    label) for those watch/appendix fields so identical filtered snapshots do
    not produce different text across exports.
    """

    kpis = [dict(kpi) for kpi in (person.get("kpis") or []) if isinstance(kpi, dict)]
    positive = [kpi for kpi in kpis if (_number(kpi.get("weighted_impact")) or 0) > 0]
    candidates = positive or kpis
    if not candidates:
        return {}

    def sort_key(kpi: dict[str, Any]) -> tuple[Any, ...]:
        definition = kpi.get("definition") if isinstance(kpi.get("definition"), dict) else {}
        display_order = _number(kpi.get("display_order"))
        if display_order is None:
            display_order = _number(definition.get("display_order"))
        return (
            -(_number(kpi.get("weighted_impact")) or 0) if positive else 0,
            display_order if display_order is not None else 9999,
            _clean(kpi.get("label")).casefold(),
        )

    return min(candidates, key=sort_key)


def _role_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    people = _person_rows(payload)
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for person in people:
        team = _clean(person.get("team"))
        position = _clean(person.get("position"), "Unassigned")
        key = f"{team} / {position}" if team and len({p.get("team") for p in people}) > 1 else position
        grouped[key].append(person)
    rows = []
    for label, values in grouped.items():
        scores = [_number(row.get("score")) for row in values]
        movements = [_number(row.get("movement")) for row in values]
        movements = [value for value in movements if value is not None]
        rows.append({
            "label": label,
            "score": mean([value for value in scores if value is not None]) if any(value is not None for value in scores) else None,
            "movement": mean(movements) if movements else None,
            "count": len(values),
        })
    return sorted(rows, key=lambda row: (_number(row.get("score")) is None, -(_number(row.get("score")) or 0)))


def _representative_for_kpi(payload: dict[str, Any], kpi_key: str) -> dict[str, Any] | None:
    for person in _person_rows(payload):
        for kpi in person.get("kpis") or []:
            if _clean(kpi.get("key")) == _clean(kpi_key):
                return person
    return None


def _action_focus(person: dict[str, Any], kpi: dict[str, Any] | None) -> str:
    label = _clean((kpi or {}).get("label"), "the KPI")
    actions = person.get("actions") or []
    if actions:
        action = actions[0]
        source = _clean(action.get("source_display"), "Recorded")
        return f"{source}: {label}"
    return f"Action needed: {label}"


def _short_person_name(value: Any) -> str:
    """Return an intentional compact name for small insight lists."""

    words = _clean(value).split()
    return words[0] if words else "Unknown"


def _kpi_label(payload: dict[str, Any], key: Any) -> str:
    """Resolve a KPI key to the configured display label."""

    normalized = _clean(key).casefold()
    for row in payload.get("kpis") or []:
        source_keys = {_clean(row.get("key")).casefold()}
        source_keys.update(_clean(source_key).casefold() for source_key in row.get("source_keys") or [])
        if normalized in source_keys:
            return _clean(row.get("label"), _clean(key, "Scope KPI"))
    return _clean(key, "Scope KPI")


def _momentum(value: Any) -> str:
    movement = _number(value)
    if movement is None:
        return "History unavailable"
    if movement > 1:
        return "Improving"
    if movement < -1:
        return "Declining"
    return "Stable"


def _latest_period(payload: dict[str, Any]) -> str:
    return _clean(payload.get("latest_period_label") or payload.get("period_label"), "Selected period")


def _footer(slide: Any, period: str, page: int) -> None:
    footer = next((shape for shape in slide.shapes if _shape_text(shape).startswith("Marketing PMS")), None)
    title = _shape_text(_find_shape(slide, "Text 1")) or f"Page {page}"
    footer_text = f"Marketing PMS | {period} | {title}"
    _set_shape_text(footer, footer_text, color=FAINT, size=5.5 if len(footer_text) > 58 else None)


def _set_header(slide: Any, title: str, subtitle: str, period: str) -> None:
    _set_named_text(slide, "Text 1", title, color=WHITE, bold=True)
    _set_named_text(slide, "Text 2", subtitle, color=RGBColor(220, 228, 242))
    _set_named_text(slide, "Text 3", period, color=WHITE, bold=True)


def _clear_named(slide: Any, names: list[str]) -> None:
    for name in names:
        _set_named_text(slide, name, "")


def _movement_display(value: Any) -> str:
    return "NEW" if _number(value) is None else _movement_text(value)


def _movement_display_color(value: Any) -> RGBColor:
    return NEUTRAL if _number(value) is None else _movement_color(value)


def _status_label(person: dict[str, Any] | None) -> str:
    score = _number((person or {}).get("score"))
    if score is None:
        return "DATA QUALITY"
    if score < 70:
        return "BELOW"
    if score < 85:
        return "WATCH"
    return "ON TRACK"


def _status_key(person: dict[str, Any] | None) -> str:
    label = _status_label(person)
    if label == "BELOW":
        return "Critical"
    if label == "WATCH":
        return "At risk"
    if label == "ON TRACK":
        return "On track"
    return "Data quality"


def _row_shapes(slide: Any, top: float, *, tolerance: float = 0.12) -> list[Any]:
    return [
        shape
        for shape in slide.shapes
        if not _shape_text(shape)
        and abs(_inches(shape.top) - top) <= tolerance
        and _inches(shape.height) > 0.03
    ]


def _set_table_row_background(slide: Any, top: float, fill: RGBColor, *, status: str | None = None) -> None:
    for shape in _row_shapes(slide, top):
        left = _inches(shape.left)
        width = _inches(shape.width)
        height = _inches(shape.height)
        if status and left >= 9.9 and width <= 1.2 and height <= 0.25:
            _set_fill(shape, _status_fill(status), _status_fill(status))
        else:
            _set_fill(shape, fill, fill)


def _hide_table_row(slide: Any, top: float) -> None:
    """Remove an unused inherited row without deleting the source frame."""

    for shape in slide.shapes:
        if _shape_text(shape) or abs(_inches(shape.top) - top) > 0.12:
            continue
        if _inches(shape.height) > 0.03:
            shape.height = 1
        elif _inches(shape.width) > 0.2:
            shape.width = 1
        try:
            shape.line.color.rgb = WHITE
        except (AttributeError, TypeError):
            pass
        try:
            shape.line.fill.background()
        except (AttributeError, TypeError):
            pass
        _set_fill(shape, WHITE, WHITE)


def _person_kpi_count(payload: dict[str, Any], key: str) -> int:
    return sum(
        1
        for person in _person_rows(payload)
        if any(_clean(kpi.get("key")).casefold() == _clean(key).casefold() for kpi in person.get("kpis") or [])
    )


def _evidence_detail_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    root_by_key: dict[str, dict[str, Any]] = {}
    for row in payload.get("root_cause_rows") or []:
        for key in [row.get("key"), *(row.get("source_keys") or [])]:
            if _clean(key):
                root_by_key[_clean(key).casefold()] = row
    rows: list[dict[str, Any]] = []
    for person in _person_rows(payload):
        for kpi in person.get("kpis") or []:
            if (_number(kpi.get("weighted_impact")) or 0) <= 0:
                continue
            root = root_by_key.get(_clean(kpi.get("key")).casefold(), {})
            rows.append(
                {
                    **dict(kpi),
                    "employee_name": person.get("name"),
                    "employee_id": person.get("employee_id"),
                    "position": person.get("position"),
                    "evidence_status": root.get("evidence_status", "KPI signal only"),
                    "recorded_root_cause": root.get("recorded_root_cause", "No operational cause recorded"),
                    "required_validation": root.get("required_validation", ""),
                    "owner_display": root.get("owner_display", "Owner needed"),
                }
            )
    rows.sort(key=lambda row: (-(_number(row.get("weighted_impact")) or 0), _clean(row.get("employee_name"))))
    if rows:
        return rows
    return [dict(row) for row in (payload.get("root_cause_rows") or [])]


def _evidence_display(row: dict[str, Any]) -> str:
    state = _clean(row.get("evidence_status"), "KPI signal only")
    cause = _clean(row.get("recorded_root_cause"))
    if state == "Confirmed cause":
        display = f"Confirmed cause: {cause or 'Recorded cause available'}"
    elif state.startswith("Evidence recorded"):
        display = f"Evidence recorded — cause pending confirmation: {cause or 'validation required'}"
    elif state == "KPI signal only":
        display = "KPI signal — cause not confirmed; validate workflow, staffing, handoff, system, or adherence."
    else:
        display = "No evidence recorded; cause validation required."
    owner = _clean(row.get("owner_display"))
    if owner and owner != "Owner needed":
        display += f" Owner: {owner}."
    return display


def _person_action_focus(person: dict[str, Any]) -> str:
    actions = person.get("actions") or []
    if actions:
        action = actions[0]
        return f"Action focus: {_clean(action.get('action_display'), 'Recorded action text needed')}"
    leading = _primary_kpi(person)
    if leading:
        return f"Action focus: KPI signal only — validate {_clean(leading.get('label'), 'the leading KPI')} before assigning action."
    return "Action focus: no measurable KPI loss is available."


def _action_priority(action: dict[str, Any]) -> str:
    return _clean(action.get("priority"), "Priority needed")


def _action_row_values(payload: dict[str, Any], action: dict[str, Any]) -> list[str]:
    source = _clean(action.get("source_display"), "Recorded")
    kpi_label = _kpi_label(payload, action.get("linked_kpi_key"))
    scope = _clean(action.get("employee_name") or action.get("team"), "Selected scope")
    execution = (
        f"Status: {_clean(action.get('status_display'), 'Status not recorded')}. "
        f"{_clean(action.get('action_display'), 'Action text needed')} Scope: {scope}."
    )
    success = (
        f"{_clean(action.get('success_metric_display'), 'Success metric needed')}. "
        f"Evidence: {_clean(action.get('evidence_display'), 'No evidence reference')}."
    )
    return [
        _action_priority(action),
        f"{source}: {kpi_label}",
        _clean(action.get("owner_display"), "Owner needed"),
        execution,
        success,
        _clean(action.get("due_date_display"), "Due date needed"),
    ]


def _slide_one(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = prs.slides[0]
    period = _latest_period(payload)
    team_label = _report_team_label(payload)
    people = _person_rows(payload)
    overall = payload.get("overall_score")
    movement = payload.get("movement")
    driver = payload.get("driver") or {}
    decline = min(
        (person for person in people if _number(person.get("movement")) is not None),
        key=lambda person: _number(person.get("movement")) or 0,
        default=None,
    )
    top = max(
        (person for person in people if _number(person.get("score")) is not None),
        key=lambda person: _number(person.get("score")) or 0,
        default=None,
    )
    below_70 = sum(1 for person in people if (_number(person.get("score")) or 0) < 70)
    _set_header(slide, f"{team_label} Team Performance Overview", "Executive Team Summary — Main Score + Month-over-Month Change", period)

    cards = [
        ("Text 8", "Text 9", "Text 10", "Shape 6", f"{team_label.upper()} MAIN SCORE", _fmt_percent(overall), f"Target 100.0% | achievement gap {_fmt_percent(payload.get('gap_to_target'))}", BLUE, PALE_BLUE, RGBColor(37, 99, 235)),
        ("Text 13", "Text 14", "Text 15", "Shape 11", "COMPARABLE MOM CHANGE", _movement_display(movement), f"{_clean(payload.get('comparison_period_label'), 'History unavailable')} → {period}" if movement is not None else "Movement unavailable", _movement_display_color(movement), _movement_fill(movement), _movement_display_color(movement)),
        ("Text 18", "Text 19", "Text 20", "Shape 16", "EMPLOYEES BELOW 70%", str(below_70), f"{below_70} employee(s) require attention", RED if below_70 else GREEN, PALE_RED if below_70 else PALE_GREEN, RED if below_70 else GREEN),
        ("Text 23", "Text 24", "Text 25", "Shape 21", "TOP CURRENT SCORE", _fmt_percent(top.get("score") if top else None), _fit(top.get("name") if top else "No employee-level data", 25), GREEN if top else NEUTRAL, PALE_GREEN if top else PALE_NEUTRAL, GREEN if top else NEUTRAL),
        ("Text 28", "Text 29", "Text 30", "Shape 26", "LARGEST DECLINE", _movement_display(decline.get("movement") if decline else None), _fit(decline.get("name") if decline else "Movement unavailable", 25), _movement_display_color(decline.get("movement") if decline else None), _movement_fill(decline.get("movement") if decline else None), _movement_display_color(decline.get("movement") if decline else None)),
    ]
    for label_name, value_name, detail_name, card_name, label, value, detail, color, fill, rail in cards:
        _set_named_text(slide, label_name, label)
        _set_named_text(slide, value_name, value, color=color, bold=True, size=15 if len(value) > 13 else None)
        _set_named_text(slide, detail_name, detail, color=color, size=6.8 if len(detail) > 25 else None)
        _set_fill(_find_shape(slide, card_name), fill, fill)
        rail_shape = _find_shape(slide, f"Shape {int(card_name.split()[1]) + 1}")
        _set_fill(rail_shape, rail, rail)

    overall_number = _number(overall)
    if overall_number is None:
        overall_context = "Overall performance is unavailable for the selected scope."
    elif overall_number < 100:
        overall_context = f"{_fmt_percent(overall_number)} overall performance remains {100 - overall_number:.1f}% below target."
    elif overall_number == 100:
        overall_context = "100.0% overall performance is at target."
    else:
        overall_context = f"{_fmt_percent(overall_number)} overall performance is above target."
    if driver and overall_number is not None:
        alert = (
            f"Priority alert: {_clean(driver.get('label'), 'The leading KPI')} creates the largest weighted performance loss at {_fmt_percent(driver.get('weighted_impact'))}. "
            f"{overall_context} Ownership is required before the next review."
        )
    elif overall_number is not None:
        alert = f"Priority alert: no measurable KPI loss is available for the selected {team_label} scope. {overall_context} Confirm KPI coverage before assigning corrective action."
    else:
        alert = f"Priority alert: no measurable KPI loss or overall score is available for the selected {team_label} scope."
    _set_named_text(slide, "Text 32", alert, color=NAVY, bold=True, size=7.6 if len(alert) > 150 else None)
    _set_fill(_find_shape(slide, "Shape 31"), PALE_RED if driver else PALE_BLUE, RED if driver else BLUE)

    _set_named_text(slide, "Text 34", "Role Performance — Main Score & MoM")
    role_specs = [
        ("Text 35", "Shape 36", "Shape 37", "Text 38", "Text 39"),
        ("Text 40", "Shape 41", "Shape 42", "Text 43", "Text 44"),
        ("Text 45", "Shape 46", "Shape 47", "Text 48", "Text 49"),
        ("Text 50", "Shape 51", "Shape 52", "Text 53", "Text 54"),
        ("Text 55", "Shape 56", "Shape 57", "Text 58", "Text 59"),
        ("Text 60", "Shape 61", "Shape 62", "Text 63", "Text 64"),
        ("Text 65", "Shape 66", "Shape 67", "Text 68", "Text 69"),
        ("Text 70", "Shape 71", "Shape 72", "Text 73", "Text 74"),
    ]
    roles = _role_rows(payload)[:8]
    if not roles and not payload.get("people_visible"):
        roles = [{"label": "Aggregate-only view", "score": overall, "movement": movement}]
    for index, (name_name, track_name, bar_name, score_name, move_name) in enumerate(role_specs):
        row = roles[index] if index < len(roles) else None
        if not row:
            _clear_named(slide, [name_name, score_name, move_name])
            for shape_name in (track_name, bar_name):
                shape = _find_shape(slide, shape_name)
                if shape:
                    shape.width = 1
                    _set_fill(shape, WHITE, WHITE)
            continue
        score = _number(row.get("score"))
        status = "On track" if score is not None and score >= 85 else "At risk"
        role_label = _fit(row.get("label"), 28)
        _set_named_text(slide, name_name, role_label, color=MUTED, size=6.8 if len(role_label) > 23 else None)
        _set_named_text(slide, score_name, _fmt_percent(score), color=_status_color(status), bold=True)
        _set_named_text(slide, move_name, _movement_display(row.get("movement")), color=_movement_display_color(row.get("movement")), bold=True, size=6.6)
        track = _find_shape(slide, track_name)
        if track:
            track.width = int(3.4 * 914400)
            _set_fill(track, TRACK, TRACK)
        _set_bar(_find_shape(slide, bar_name), 3.4 * 914400 * max(0.0, min(1.0, (score or 0) / 100)), _status_color(status))

    if roles and len(roles) < len(role_specs):
        coverage_index = len(roles)
        name_name, track_name, bar_name, score_name, move_name = role_specs[coverage_index]
        record_count = len(people) if people else int(payload.get("selected_record_count") or 0)
        _set_named_text(slide, name_name, "Coverage in scope", color=MUTED, size=6.7)
        _set_named_text(slide, score_name, f"{len(roles)} roles", color=BLUE, bold=True, size=5.8)
        _set_named_text(slide, move_name, f"{record_count} records", color=BLUE, bold=True, size=5.6)
        for shape_name in (track_name, bar_name):
            shape = _find_shape(slide, shape_name)
            if shape:
                shape.width = 1
                _set_fill(shape, WHITE, WHITE)

    _set_named_text(slide, "Text 76", "Key Insights & Next Steps")
    insight_rows = [
        ("Text 78", f"{_clean(decline.get('name'), 'No measured decline') if decline else 'Trend unavailable'} is the largest urgency signal", f"{_fmt_percent(decline.get('score'))} score, {_movement_display(decline.get('movement'))} movement." if decline else "Use the next measured period to establish direction."),
        ("Text 80", f"{_clean(driver.get('label'), 'No leading KPI')} is the leading loss driver" if driver else "No measured KPI loss is available", f"Weighted performance loss {_fmt_percent(driver.get('weighted_impact'))}." if driver else "No action order can be derived from the selected data."),
        ("Text 82", "Decision before next review", "Assign an owner, review cadence, and success measure." if any(_clean(action.get("owner_display")) == "Owner needed" or _clean(action.get("due_date_display")) == "Due date needed" for action in payload.get("actions") or []) else "Confirm the owner, cadence, and evidence that will prove improvement."),
    ]
    for text_name, title, detail in insight_rows:
        text_shape = _find_shape(slide, text_name)
        if text_shape is None:
            continue
        # The source boxes each contain one text frame. Keep the title/detail
        # hierarchy in the inherited frame by using two short lines.
        _set_shape_text(text_shape, f"{_fit(title, 48)}\n{_fit(detail, 72)}", color=NAVY if text_name != "Text 82" else AMBER, bold=False, size=7.1 if len(detail) > 58 else 7.6)
    _set_fill(_find_shape(slide, "Shape 81"), PALE_AMBER, PALE_AMBER)
    _footer(slide, period, 1)


def _slide_two(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = prs.slides[1]
    period = _latest_period(payload)
    team_label = _report_team_label(payload)
    kpis = list(payload.get("kpis") or [])[:9]
    driver = payload.get("driver") or {}
    paid_media = next((row for row in kpis if any(token in _clean(row.get("label")).casefold() for token in ("cr", "conversion", "cpv", "cpl", "leads", "revenue")) and row is not driver), None)
    if paid_media is None and len(kpis) > 1:
        paid_media = kpis[1]
    gap_count = sum(int(round(_number(row.get("affected_count")) or 0)) for row in kpis)
    _set_header(slide, "2. KPI Loss Breakdown", "Which KPI gaps are costing the most score points?", period)
    cards = [
        ("Text 8", "Text 9", "Text 10", "Shape 6", "MAIN SCORE", _fmt_percent(payload.get("overall_score")), f"{team_label} overall", BLUE, PALE_BLUE, "Shape 7"),
        ("Text 13", "Text 14", "Text 15", "Shape 11", "MOM CHANGE", _movement_display(payload.get("movement")), "Comparable employee base" if payload.get("movement") is not None else "History unavailable", _movement_display_color(payload.get("movement")), _movement_fill(payload.get("movement")), "Shape 12"),
        ("Text 18", "Text 19", "Text 20", "Shape 16", "TOP LOSS KPI", _clean(driver.get("label"), "No driver"), _fmt_percent(driver.get("weighted_impact")) + " weighted loss" if driver else "No measured loss", RED if driver else NEUTRAL, PALE_RED if driver else PALE_NEUTRAL, "Shape 17"),
        ("Text 23", "Text 24", "Text 25", "Shape 21", "TOP PAID MEDIA LOSS", _clean(paid_media.get("label"), "No paid-media KPI") if paid_media else "No paid-media KPI", _fmt_percent(paid_media.get("weighted_impact")) + " weighted loss" if paid_media else "No measured loss", AMBER if paid_media else NEUTRAL, PALE_AMBER if paid_media else PALE_NEUTRAL, "Shape 22"),
        ("Text 28", "Text 29", "Text 30", "Shape 26", "LARGEST GAP COUNT", str(gap_count), "Below-target KPI gaps", AMBER if gap_count else GREEN, PALE_AMBER if gap_count else PALE_GREEN, "Shape 27"),
    ]
    for label_name, value_name, detail_name, card_name, label, value, detail, color, fill, rail_name in cards:
        _set_named_text(slide, label_name, label)
        _set_named_text(slide, value_name, value, color=color, bold=True, size=13 if len(value) > 14 else None)
        _set_named_text(slide, detail_name, detail, color=color, size=6.8 if len(detail) > 25 else None)
        _set_fill(_find_shape(slide, card_name), fill, fill)
        _set_fill(_find_shape(slide, rail_name), color, color)

    _set_named_text(slide, "Text 32", "Weighted score loss by KPI")
    row_specs = [
        ("Text 33", "Shape 34", "Shape 35", "Text 36", "Text 37"),
        ("Text 38", "Shape 39", "Shape 40", "Text 41", "Text 42"),
        ("Text 43", "Shape 44", "Shape 45", "Text 46", "Text 47"),
        ("Text 48", "Shape 49", "Shape 50", "Text 51", "Text 52"),
        ("Text 53", "Shape 54", "Shape 55", "Text 56", "Text 57"),
        ("Text 58", "Shape 59", "Shape 60", "Text 61", "Text 62"),
        ("Text 63", "Shape 64", "Shape 65", "Text 66", "Text 67"),
        ("Text 68", "Shape 69", "Shape 70", "Text 71", "Text 72"),
        ("Text 73", "Shape 74", "Shape 75", "Text 76", "Text 77"),
    ]
    max_loss = max([_number(row.get("weighted_impact")) or 0 for row in kpis] or [0])
    for index, (label_name, track_name, bar_name, value_name, affected_name) in enumerate(row_specs):
        kpi = kpis[index] if index < len(kpis) else None
        if not kpi:
            _clear_named(slide, [label_name, value_name, affected_name])
            for shape_name in (track_name, bar_name):
                shape = _find_shape(slide, shape_name)
                if shape:
                    shape.width = 1
                    _set_fill(shape, WHITE, WHITE)
            continue
        loss = _number(kpi.get("weighted_impact")) or 0
        color = RED if index == 0 and loss > 0 else AMBER if loss > 0 else BLUE
        kpi_label = _fit(kpi.get("label"), 25)
        _set_named_text(slide, label_name, kpi_label, color=NAVY, bold=index == 0, size=6.8 if len(kpi_label) > 21 else None)
        _set_named_text(slide, value_name, _fmt_percent(loss), color=color, bold=True)
        affected = int(round(_number(kpi.get("affected_count")) or 0))
        denominator = _person_kpi_count(payload, _clean(kpi.get("key"))) or len(_person_rows(payload))
        affected_text = f"{affected}/{denominator} <70%" if denominator else "No employee detail"
        _set_named_text(slide, affected_name, affected_text, color=FAINT, size=6.4)
        track = _find_shape(slide, track_name)
        if track:
            track.width = int(5.55 * 914400)
            _set_fill(track, TRACK, TRACK)
        _set_bar(_find_shape(slide, bar_name), 5.55 * 914400 * (loss / max_loss) if max_loss > 0 else 1, color)
    _set_named_text(slide, "Text 79", "Reading rule: use weighted loss to choose the action order. Do not treat every KPI gap equally; solve the highest score-impact drivers first.", color=NAVY, bold=True, size=8.0)
    _footer(slide, period, 2)


def _slide_three(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = prs.slides[2]
    period = _latest_period(payload)
    people = _person_rows(payload)[:13]
    _set_header(slide, "3. Employee Performance Priority Ranking", "Ranked by weighted performance loss; score shows severity and MoM shows urgency.", period)
    _set_named_text(slide, "Text 18", "Primary KPI | weighted loss", color=FAINT, bold=True, size=5.8)
    groups = _text_groups(slide, top_min=1.45, top_max=6.2, threshold=0.14)[:13]
    for index in range(13):
        group = groups[index] if index < len(groups) else []
        fields = sorted(group, key=lambda shape: _inches(shape.left))
        row = people[index] if index < len(people) else None
        if not row or len(fields) < 7:
            for field in fields:
                _set_shape_text(field, "")
            if group:
                _set_table_row_background(slide, _inches(group[0].top), WHITE)
            continue
        status = _status_key(row)
        status_label = _status_label(row)
        score = _number(row.get("score"))
        kpi = _primary_kpi(row)
        values = [
            str(index + 1),
            _clean(row.get("name"), "Unknown employee"),
            _fit(row.get("position"), 27),
            _fmt_percent(score),
            _movement_display(row.get("movement")),
            _kpi_gap_text(kpi) if kpi else "No KPI gap",
            status_label,
        ]
        for field, value in zip(fields[:7], values):
            color = _status_color(status) if field in fields[3:5] or field is fields[6] else RED if field is fields[5] and kpi and (_number(kpi.get("weighted_impact")) or 0) > 0 else None
            size = 5.0 if field is fields[5] and len(value) > 36 else 5.4 if field is fields[5] and len(value) > 28 else 6.1 if len(value) > 28 else None
            _set_shape_text(field, value, color=color, bold=field is fields[1] or field is fields[6], size=size)
        _set_table_row_background(slide, _inches(fields[0].top), PALE_RED if status_label == "BELOW" else PALE_NEUTRAL if index % 2 else WHITE, status=status)
    message = f"Immediate focus: rows are ranked by weighted loss contribution. The bottom {sum(1 for row in people if _status_label(row) == 'BELOW')} require the fastest review."
    _set_named_text(slide, "Text 230", message, color=NAVY, bold=True, size=8.0)
    _footer(slide, period, 3)


def _slide_four(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = prs.slides[3]
    period = _latest_period(payload)
    _set_header(slide, "4. Root Cause Evidence", "Evidence first: actual result, target, achievement, score loss, and observed operating issue.", period)
    headers = [("Text 8", "Employee"), ("Text 10", "Position"), ("Text 12", "KPI"), ("Text 14", "Actual"), ("Text 16", "Target"), ("Text 18", "Ach."), ("Text 20", "Weighted loss"), ("Text 22", "Evidence / operating issue")]
    weighted_header = _find_shape(slide, "Text 20")
    evidence_header = _find_shape(slide, "Text 22")
    if weighted_header and evidence_header:
        weighted_header.width = int(0.82 * 914400)
        evidence_header.left = int(8.1 * 914400)
        evidence_header.width = int(3.92 * 914400)
    for name, value in headers:
        _set_named_text(slide, name, value, color=FAINT, bold=True, size=5.8 if name == "Text 20" else None)
    evidence_rows = _evidence_detail_rows(payload)[:8]
    groups = _text_groups(slide, top_min=1.4, top_max=5.8, threshold=0.14)[:8]
    for index in range(8):
        group = groups[index] if index < len(groups) else []
        fields = sorted(group, key=lambda shape: _inches(shape.left))
        row = evidence_rows[index] if index < len(evidence_rows) else None
        if not row or len(fields) < 8:
            for field in fields:
                _set_shape_text(field, "")
            if group:
                _set_table_row_background(slide, _inches(group[0].top), WHITE)
            continue
        values = [
            _clean(row.get("employee_name") or "Selected scope", "Selected scope"),
            _clean(row.get("position") or "Aggregate scope", "Aggregate scope"),
            _fit(row.get("label"), 20),
            _fmt_native(row.get("actual"), row.get("unit")),
            _fmt_native(row.get("target"), row.get("unit")),
            _fmt_percent(row.get("achievement_pct")),
            _fmt_percent(row.get("weighted_impact")),
            _evidence_display(row),
        ]
        for field, value in zip(fields[:8], values):
            color = RED if field in fields[5:7] and (_number(row.get("weighted_impact")) or 0) > 0 else None
            size = 5.1 if field is fields[1] and len(value) > 18 else 5.55 if field is fields[7] else 6.0 if len(value) > 20 else None
            _set_shape_text(field, value, color=color, bold=field is fields[2], size=size)
        _set_table_row_background(slide, _inches(fields[0].top), PALE_NEUTRAL if index % 2 else WHITE)
    note = "Important: this slide separates evidence from assumption. The KPI identifies where to investigate, but the operational cause has not been confirmed unless recorded in the system."
    _set_named_text(slide, "Text 160", note, color=NAVY, bold=True, size=7.8)
    _footer(slide, period, 4)


def _slide_five(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = prs.slides[4]
    period = _latest_period(payload)
    actions = list(payload.get("actions") or [])[:4]
    _set_header(slide, "5. Action Plan", "Prioritized by score impact. Owners, execution, success measure, and timing are explicit.", period)
    table_frame = _find_shape(slide, "Shape 6")
    if table_frame:
        table_frame.height = int((0.45 + max(1, len(actions)) * 0.92) * 914400)
    _set_named_text(slide, "Text 8", "Priority", color=FAINT, bold=True, size=5.6)
    groups = _text_groups(slide, top_min=1.45, top_max=5.8, threshold=0.18)[:4]
    for index in range(4):
        group = groups[index] if index < len(groups) else []
        fields = sorted(group, key=lambda shape: _inches(shape.left))
        action = actions[index] if index < len(actions) else None
        if not action or len(fields) < 6:
            for field in fields:
                _set_shape_text(field, "")
            if group:
                _hide_table_row(slide, _inches(group[0].top))
            continue
        values = _action_row_values(payload, action)
        source = _clean(action.get("source_display"), "Recorded")
        for field, value in zip(fields[:6], values):
            color = AMBER if source == "Proposed" or value.endswith("needed") else None
            _set_shape_text(field, value, color=color, bold=field in fields[:2], size=5.15 if field is fields[0] else 5.7 if len(value) > 68 else 6.0 if len(value) > 45 else None)
        row_fill = PALE_AMBER if source == "Proposed" else PALE_NEUTRAL if index % 2 else WHITE
        _set_table_row_background(slide, _inches(fields[0].top), row_fill)
    # The source frame includes zero-height divider line shapes that have no
    # text group of their own.  Hide them explicitly when the filtered scope
    # contains fewer actions so PowerPoint cannot render ghost row dividers.
    if len(actions) < 4:
        for shape in slide.shapes:
            if _shape_text(shape):
                continue
            if 2.0 <= _inches(shape.top) <= 5.4 and _inches(shape.height) <= 0.03:
                _hide_table_row(slide, _inches(shape.top))
    missing = [
        _clean(action.get(field))
        for action in actions
        for field in ("owner_display", "due_date_display", "success_metric_display")
        if _clean(action.get(field)).endswith("needed")
    ]
    if not actions:
        governance = "Decision needed: no recorded or proposed action is available. Assign an owner, due date or cadence, and success metric."
        governance_fill = PALE_AMBER
    elif missing:
        governance = "Decision needed before the next review: " + ", ".join(dict.fromkeys(missing)) + ". Proposed actions are not recorded actions."
        governance_fill = PALE_AMBER
    else:
        governance = "Governance: daily priority check, twice-weekly owner review, and end-of-period validation on both Main Score and MoM Change."
        governance_fill = PALE_GREEN
    _set_fill(_find_shape(slide, "Shape 71"), governance_fill, governance_fill)
    _set_named_text(slide, "Text 72", governance, color=NAVY, bold=True, size=7.7)
    _footer(slide, period, 5)


_DETAIL_CARD_SPECS = [
    {"card": "Shape 18", "rail": "Shape 19", "kpi": "Text 20", "achievement": "Text 21", "loss_bg": "Shape 22", "loss": "Text 23", "achievement_label": "Text 24", "track": "Shape 25", "bar": "Shape 26", "actual_label": "Text 27", "actual": "Text 28", "target_label": "Text 29", "target": "Text 30", "direction": "Text 31"},
    {"card": "Shape 32", "rail": "Shape 33", "kpi": "Text 34", "achievement": "Text 35", "loss_bg": "Shape 36", "loss": "Text 37", "achievement_label": "Text 38", "track": "Shape 39", "bar": "Shape 40", "actual_label": "Text 41", "actual": "Text 42", "target_label": "Text 43", "target": "Text 44", "direction": "Text 45"},
    {"card": "Shape 46", "rail": "Shape 47", "kpi": "Text 48", "achievement": "Text 49", "loss_bg": "Shape 50", "loss": "Text 51", "achievement_label": "Text 52", "track": "Shape 53", "bar": "Shape 54", "actual_label": "Text 55", "actual": "Text 56", "target_label": "Text 57", "target": "Text 58", "direction": "Text 59"},
    {"card": "Shape 60", "rail": "Shape 61", "kpi": "Text 62", "achievement": "Text 63", "loss_bg": "Shape 64", "loss": "Text 65", "achievement_label": "Text 66", "track": "Shape 67", "bar": "Shape 68", "actual_label": "Text 69", "actual": "Text 70", "target_label": "Text 71", "target": "Text 72", "direction": "Text 73"},
    {"card": "Shape 74", "rail": "Shape 75", "kpi": "Text 76", "achievement": "Text 77", "loss_bg": "Shape 78", "loss": "Text 79", "achievement_label": "Text 80", "track": "Shape 81", "bar": "Shape 82", "actual_label": "Text 83", "actual": "Text 84", "target_label": "Text 85", "target": "Text 86", "direction": "Text 87"},
    {"card": "Shape 88", "rail": "Shape 89", "kpi": "Text 90", "achievement": "Text 91", "loss_bg": "Shape 92", "loss": "Text 93", "achievement_label": "Text 94", "track": "Shape 95", "bar": "Shape 96", "actual_label": "Text 97", "actual": "Text 98", "target_label": "Text 99", "target": "Text 100", "direction": "Text 101"},
]


def _clear_detail_card(slide: Any, spec: dict[str, str]) -> None:
    _clear_named(slide, [spec[key] for key in ("kpi", "achievement", "loss", "achievement_label", "actual_label", "actual", "target_label", "target", "direction")])
    for key in ("card", "rail", "loss_bg", "track", "bar"):
        shape = _find_shape(slide, spec[key])
        if shape:
            if key in {"track", "bar"}:
                shape.width = 1
            _set_fill(shape, WHITE, WHITE)


def _render_detail_card(slide: Any, spec: dict[str, str], kpi: dict[str, Any]) -> None:
    loss = _number(kpi.get("weighted_impact")) or 0
    achievement = _number(kpi.get("achievement_pct"))
    status = _clean(kpi.get("status"), "Data quality")
    color = RED if loss > 0 else _status_color(status)
    fill = PALE_RED if loss > 0 else _status_fill(status)
    _set_fill(_find_shape(slide, spec["card"]), WHITE, LINE)
    _set_fill(_find_shape(slide, spec["rail"]), color, color)
    kpi_label = _fit(kpi.get("label"), 25)
    _set_named_text(slide, spec["kpi"], kpi_label, color=NAVY, bold=True, size=6.4 if len(kpi_label) > 24 else 6.8 if len(kpi_label) > 18 else None)
    achievement_text = _fmt_percent(achievement)
    _set_named_text(slide, spec["achievement"], achievement_text, color=color, bold=True, size=12.5 if len(achievement_text) > 5 else None)
    _set_fill(_find_shape(slide, spec["loss_bg"]), fill, fill)
    gap = max(0.0, -(_number(kpi.get("gap")) or 0.0))
    loss_badge = f"Achievement gap {_fmt_percent(gap)}\nWeighted loss {_fmt_percent(loss)}"
    loss_bg = _find_shape(slide, spec["loss_bg"])
    loss_text = _find_shape(slide, spec["loss"])
    if loss_bg and loss_text:
        card = _find_shape(slide, spec["card"])
        if card:
            badge_width = int(1.36 * 914400)
            loss_bg.left = int(card.left + card.width - badge_width - 0.1 * 914400)
            loss_bg.width = badge_width
            loss_bg.top = int(card.top + 0.37 * 914400)
            loss_bg.height = int(0.28 * 914400)
            loss_text.left = loss_bg.left
            loss_text.top = loss_bg.top
            loss_text.width = loss_bg.width
            loss_text.height = loss_bg.height
    _set_named_text(slide, spec["loss"], loss_badge, color=color, bold=True, size=5.0)
    _set_named_text(slide, spec["achievement_label"], "Achievement", color=FAINT, bold=True, size=5.4)
    track = _find_shape(slide, spec["track"])
    if track:
        track.width = int(2.17 * 914400)
        _set_fill(track, TRACK, TRACK)
    _set_bar(_find_shape(slide, spec["bar"]), 2.17 * 914400 * max(0.0, min(1.0, (achievement or 0) / 100)), color)
    _set_named_text(slide, spec["actual_label"], "Actual", color=FAINT, bold=True, size=5.5)
    actual_text = _fmt_native(kpi.get("actual"), kpi.get("unit"))
    _set_named_text(slide, spec["actual"], actual_text, color=NAVY, size=5.2 if len(actual_text) > 11 else 6.0)
    _set_named_text(slide, spec["target_label"], "Target", color=FAINT, bold=True, size=5.5)
    target_text = _fmt_native(kpi.get("target"), kpi.get("unit"))
    _set_named_text(slide, spec["target"], target_text, color=NAVY, size=5.2 if len(target_text) > 11 else 6.0)
    direction = "Lower is better" if _clean(kpi.get("direction")).casefold() == "lower_better" else "Higher is better"
    _set_named_text(slide, spec["direction"], direction, color=MUTED, size=5.4)


def _detail_kpis(person: dict[str, Any], limit: int) -> list[dict[str, Any]]:
    kpis = sorted(person.get("kpis") or [], key=lambda row: (-(_number(row.get("weighted_impact")) or 0), _clean(row.get("label"))))
    return kpis[:limit]


def _clear_detail_page(slide: Any, *, max_cards: int = 6) -> None:
    for spec in _DETAIL_CARD_SPECS[:max_cards]:
        _clear_detail_card(slide, spec)
    for spec in _DETAIL_CARD_SPECS[max_cards:]:
        if _find_shape(slide, spec["card"]):
            _clear_detail_card(slide, spec)


def _detail_bottom_names(slide: Any) -> tuple[str, str]:
    return ("Shape 102", "Text 103") if _find_shape(slide, "Text 103") else ("Shape 74", "Text 75")


def _render_detail_profile(slide: Any, person: dict[str, Any] | None, period: str, title: str, page: int, *, max_cards: int = 6) -> None:
    if person is None:
        _set_header(slide, title, "Actual vs Target is shown for every KPI so the employee number is clear and auditable.", period)
        _set_fill(_find_shape(slide, "Shape 6"), PALE_BLUE, PALE_BLUE)
        _set_fill(_find_shape(slide, "Shape 7"), NEUTRAL, NEUTRAL)
        _set_named_text(slide, "Text 8", "DETAIL UNAVAILABLE", color=WHITE, bold=True, size=6.0)
        _set_named_text(slide, "Text 9", "Employee detail unavailable", color=NAVY, bold=True, size=8.0)
        _set_named_text(slide, "Text 10", "The selected role or filter does not expose person-level rows.", color=MUTED, size=6.4)
        _set_named_text(slide, "Text 11", "MAIN SCORE", color=FAINT, bold=True)
        _set_named_text(slide, "Text 12", "N/A", color=NEUTRAL, bold=True)
        _set_named_text(slide, "Text 13", "MOM CHANGE", color=FAINT, bold=True)
        _set_named_text(slide, "Text 14", "N/A", color=NEUTRAL, bold=True)
        _set_fill(_find_shape(slide, "Shape 15"), PALE_NEUTRAL, PALE_NEUTRAL)
        _set_named_text(slide, "Text 16", "Aggregate-only view — no employee-level action detail is included.", color=NAVY, bold=True, size=6.8)
        _set_named_text(slide, "Text 17", "KPI Drivers — Actual / Target / Achievement / Weighted Loss", color=FAINT, bold=True)
        _clear_detail_page(slide, max_cards=max_cards)
        bottom_shape, bottom_text = _detail_bottom_names(slide)
        _set_fill(_find_shape(slide, bottom_shape), PALE_NEUTRAL, PALE_NEUTRAL)
        _set_named_text(slide, bottom_text, "Review rule: person-level evidence becomes available when the selected role permits employee detail.", color=NAVY, bold=True, size=7.0)
        _footer(slide, period, page)
        return
    score = _number(person.get("score"))
    status = _status_key(person)
    status_label = "BELOW TARGET" if _status_label(person) in {"BELOW", "WATCH"} else "ON TRACK"
    movement = person.get("movement")
    _set_header(slide, title, "Actual vs Target is shown for every KPI so the employee number is clear and auditable.", period)
    _set_fill(_find_shape(slide, "Shape 6"), WHITE, LINE)
    _set_fill(_find_shape(slide, "Shape 7"), _status_color(status), _status_color(status))
    _set_named_text(slide, "Text 8", status_label, color=WHITE, bold=True, size=6.2)
    name = _clean(person.get("name"), "Unknown employee")
    position = _fit(person.get("position"), 34)
    _set_named_text(slide, "Text 9", name, color=NAVY, bold=True, size=8.2 if len(name) > 20 else None)
    _set_named_text(slide, "Text 10", position, color=MUTED, size=6.4 if len(position) > 26 else None)
    _set_named_text(slide, "Text 11", "MAIN SCORE", color=FAINT, bold=True)
    _set_named_text(slide, "Text 12", _fmt_percent(score), color=_status_color(status), bold=True)
    _set_named_text(slide, "Text 13", "MOM CHANGE", color=FAINT, bold=True)
    _set_named_text(slide, "Text 14", _movement_display(movement), color=_movement_display_color(movement), bold=True)
    _set_fill(_find_shape(slide, "Shape 15"), PALE_AMBER if (_number(person.get("weighted_loss")) or 0) > 0 else PALE_GREEN, PALE_AMBER if (_number(person.get("weighted_loss")) or 0) > 0 else PALE_GREEN)
    _set_named_text(slide, "Text 16", _person_action_focus(person), color=NAVY, bold=True, size=6.1 if len(_person_action_focus(person)) > 95 else 6.8)
    _set_named_text(slide, "Text 17", "KPI Drivers — Actual / Target / Achievement / Weighted Loss", color=FAINT, bold=True)
    _clear_detail_page(slide, max_cards=max_cards)
    for spec, kpi in zip(_DETAIL_CARD_SPECS[:max_cards], _detail_kpis(person, max_cards)):
        _render_detail_card(slide, spec, kpi)
    bottom_shape, bottom_text = _detail_bottom_names(slide)
    _set_fill(_find_shape(slide, bottom_shape), PALE_BLUE, PALE_BLUE)
    _set_named_text(slide, bottom_text, "Review rule: prioritize the highest weighted-loss KPI first; use Actual / Target to validate the gap before assigning corrective action.", color=NAVY, bold=True, size=7.0)
    _footer(slide, period, page)


def _detail_selection(payload: dict[str, Any]) -> tuple[dict[str, Any] | None, dict[str, Any] | None, list[dict[str, Any]], dict[str, Any] | None]:
    if not payload.get("people_visible"):
        return None, None, [], None
    affected = list(payload.get("people") or [])
    first = affected[0] if len(affected) > 0 else None
    second = affected[1] if len(affected) > 1 else None
    remainder = affected[2:]
    group: list[dict[str, Any]] = []
    if remainder:
        grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
        for person in remainder:
            grouped[_clean(person.get("position"), "Unassigned")].append(person)
        group = next((values[:2] for values in grouped.values() if len(values) >= 2), [])
    used = {id(person) for person in (first, second, *group) if person is not None}
    fourth = next((person for person in remainder if id(person) not in used), None)
    if fourth is None and group:
        fourth = next((person for person in remainder if id(person) not in used), None)
    return first, second, group, fourth


_GROUP_CARD_SPECS = [
    {"card": "Shape 16", "rail": "Shape 17", "kpi": "Text 18", "achievement": "Text 19", "loss_bg": "Shape 20", "loss": "Text 21", "achievement_label": "Text 22", "track": "Shape 23", "bar": "Shape 24", "actual_label": "Text 25", "actual": "Text 26", "target_label": "Text 27", "target": "Text 28", "direction": "Text 29"},
    {"card": "Shape 30", "rail": "Shape 31", "kpi": "Text 32", "achievement": "Text 33", "loss_bg": "Shape 34", "loss": "Text 35", "achievement_label": "Text 36", "track": "Shape 37", "bar": "Shape 38", "actual_label": "Text 39", "actual": "Text 40", "target_label": "Text 41", "target": "Text 42", "direction": "Text 43"},
    {"card": "Shape 44", "rail": "Shape 45", "kpi": "Text 46", "achievement": "Text 47", "loss_bg": "Shape 48", "loss": "Text 49", "achievement_label": "Text 50", "track": "Shape 51", "bar": "Shape 52", "actual_label": "Text 53", "actual": "Text 54", "target_label": "Text 55", "target": "Text 56", "direction": "Text 57"},
]

_GROUP_CARD_SPECS_SECOND = [
    {"card": "Shape 68", "rail": "Shape 69", "kpi": "Text 70", "achievement": "Text 71", "loss_bg": "Shape 72", "loss": "Text 73", "achievement_label": "Text 74", "track": "Shape 75", "bar": "Shape 76", "actual_label": "Text 77", "actual": "Text 78", "target_label": "Text 79", "target": "Text 80", "direction": "Text 81"},
    {"card": "Shape 82", "rail": "Shape 83", "kpi": "Text 84", "achievement": "Text 85", "loss_bg": "Shape 86", "loss": "Text 87", "achievement_label": "Text 88", "track": "Shape 89", "bar": "Shape 90", "actual_label": "Text 91", "actual": "Text 92", "target_label": "Text 93", "target": "Text 94", "direction": "Text 95"},
    {"card": "Shape 96", "rail": "Shape 97", "kpi": "Text 98", "achievement": "Text 99", "loss_bg": "Shape 100", "loss": "Text 101", "achievement_label": "Text 102", "track": "Shape 103", "bar": "Shape 104", "actual_label": "Text 105", "actual": "Text 106", "target_label": "Text 107", "target": "Text 108", "direction": "Text 109"},
]


def _render_group_profile(slide: Any, person: dict[str, Any] | None, prefix: int) -> None:
    if prefix == 1:
        names = {"card": "Shape 6", "badge": "Shape 9", "status": "Text 10", "name": "Text 7", "position": "Text 8", "score_label": "Text 11", "score": "Text 12", "movement": "Text 13", "action_bg": "Shape 14", "action": "Text 15"}
    else:
        names = {"card": "Shape 58", "badge": "Shape 61", "status": "Text 62", "name": "Text 59", "position": "Text 60", "score_label": "Text 63", "score": "Text 64", "movement": "Text 65", "action_bg": "Shape 66", "action": "Text 67"}
    if person is None:
        _set_fill(_find_shape(slide, names["card"]), PALE_BLUE, PALE_BLUE)
        _set_named_text(slide, names["name"], "Employee detail unavailable", color=NAVY, bold=True, size=7.2)
        _set_named_text(slide, names["position"], "No person-level rows are available for this filter.", color=MUTED, size=5.8)
        _set_named_text(slide, names["status"], "DETAIL UNAVAILABLE", color=NEUTRAL, bold=True, size=5.2)
        _set_named_text(slide, names["score_label"], "MAIN SCORE", color=FAINT, bold=True, size=5.4)
        _set_named_text(slide, names["score"], "N/A", color=NEUTRAL, bold=True, size=12)
        _set_named_text(slide, names["movement"], "N/A", color=NEUTRAL, bold=True, size=6.8)
        _set_fill(_find_shape(slide, names["action_bg"]), PALE_NEUTRAL, PALE_NEUTRAL)
        _set_named_text(slide, names["action"], "Aggregate-only view.", color=NAVY, bold=True, size=5.7)
        return
    status = _status_key(person)
    _set_fill(_find_shape(slide, names["card"]), WHITE, LINE)
    _set_fill(_find_shape(slide, names["badge"]), _status_fill(status), _status_fill(status))
    _set_named_text(slide, names["name"], _clean(person.get("name"), "Unknown employee"), color=NAVY, bold=True, size=6.8)
    _set_named_text(slide, names["position"], _fit(person.get("position"), 28), color=MUTED, size=5.8)
    _set_named_text(slide, names["status"], "BELOW TARGET" if _status_label(person) != "ON TRACK" else "ON TRACK", color=_status_color(status), bold=True, size=5.4)
    _set_named_text(slide, names["score_label"], "MAIN SCORE", color=FAINT, bold=True, size=5.4)
    _set_named_text(slide, names["score"], _fmt_percent(person.get("score")), color=_status_color(status), bold=True, size=12)
    _set_named_text(slide, names["movement"], _movement_display(person.get("movement")), color=_movement_display_color(person.get("movement")), bold=True, size=6.5)
    _set_fill(_find_shape(slide, names["action_bg"]), PALE_AMBER, PALE_AMBER)
    _set_named_text(slide, names["action"], _person_action_focus(person), color=NAVY, bold=True, size=5.7)


def _render_group_cards(slide: Any, cards: list[dict[str, Any]], kpis: list[dict[str, Any]]) -> None:
    for spec in cards:
        _clear_detail_card(slide, spec)
    for spec, kpi in zip(cards, kpis[:3]):
        _render_detail_card(slide, spec, kpi)


def _slide_detail_pages(prs: Presentation, payload: dict[str, Any]) -> None:
    period = _latest_period(payload)
    first, second, group, fourth = _detail_selection(payload)
    _render_detail_profile(prs.slides[5], first, period, f"6A. Employee Detail — {_clean(first.get('name'), 'No employee detail') if first else 'No employee detail'}", 6, max_cards=6)
    _render_detail_profile(prs.slides[6], second, period, f"6B. Employee Detail — {_clean(second.get('name'), 'No employee detail') if second else 'No employee detail'}", 7, max_cards=6)

    slide = prs.slides[7]
    group_label = _clean(group[0].get("position"), "Selected role") + "s" if group else "Selected role"
    _set_header(slide, f"6C. Employee Detail — {group_label}", "Actual vs Target is shown for every KPI so the employee number is clear and auditable.", period)
    _render_group_profile(slide, group[0] if group else None, 1)
    _render_group_profile(slide, group[1] if len(group) > 1 else None, 2)
    for person, cards in ((group[0] if group else None, _GROUP_CARD_SPECS), (group[1] if len(group) > 1 else None, _GROUP_CARD_SPECS_SECOND)):
        _render_group_cards(slide, cards, _detail_kpis(person, 3) if person else [])
    _footer(slide, period, 8)

    _render_detail_profile(prs.slides[8], fourth, period, f"6D. Employee Detail — {_clean(fourth.get('name'), 'No employee detail') if fourth else 'No employee detail'}", 9, max_cards=4)


def _slide_ten(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = prs.slides[9]
    period = _latest_period(payload)
    trend = list(payload.get("trend") or [])
    driver = payload.get("driver") or {}
    driver_trend = list(payload.get("driver_trend") or [])
    latest = trend[-1] if trend else {}
    first = trend[0] if trend else {}
    best = payload.get("best_period") or {}
    worst = payload.get("worst_period") or {}
    next_review = payload.get("next_review") or {}
    trend_status = _clean(payload.get("trend_status"), "No measured history")
    overall_score = _number(payload.get("overall_score"))
    if trend_status.casefold() == "improving":
        context_title = "Improving; target gap unresolved" if overall_score is None or overall_score < 100 else "Improving toward target"
    elif trend_status.casefold() == "declining":
        context_title = "Declining; intervention required"
    elif trend_status.casefold() == "stable":
        context_title = "Stable; target gap remains"
    elif trend_status.casefold() == "movement only":
        context_title = "Movement observed; sustained trend unconfirmed"
    elif trend_status.casefold() == "trend unavailable":
        context_title = "Trend unavailable; one measured period"
    else:
        context_title = "No measured history"
    title = f"6E. Performance Context — {context_title}"
    _set_header(slide, title, "Historical direction and leading KPI evidence are shown from the same filtered snapshot.", period)
    title_shape = _find_shape(slide, "Text 1")
    if title_shape:
        # The inherited title box was sized for the old, shorter appendix
        # title.  Give the conclusion room to stay on one line and keep it
        # clear of the period label at the right edge.
        title_shape.width = int(11.2 * 914400)
        title_shape.height = int(0.30 * 914400)
        _set_shape_text(title_shape, title, color=WHITE, bold=True, size=16.5)
    specs = [
        {"card": "Shape 6", "name": "Text 7", "position": "Text 8", "score": "Text 9", "movement": "Text 10", "watch_label": "Text 11", "watch": "Text 12"},
        {"card": "Shape 13", "name": "Text 14", "position": "Text 15", "score": "Text 16", "movement": "Text 17", "watch_label": "Text 18", "watch": "Text 19"},
        {"card": "Shape 20", "name": "Text 21", "position": "Text 22", "score": "Text 23", "movement": "Text 24", "watch_label": "Text 25", "watch": "Text 26"},
        {"card": "Shape 27", "name": "Text 28", "position": "Text 29", "score": "Text 30", "movement": "Text 31", "watch_label": "Text 32", "watch": "Text 33"},
        {"card": "Shape 34", "name": "Text 35", "position": "Text 36", "score": "Text 37", "movement": "Text 38", "watch_label": "Text 39", "watch": "Text 40"},
        {"card": "Shape 41", "name": "Text 42", "position": "Text 43", "score": "Text 44", "movement": "Text 45", "watch_label": "Text 46", "watch": "Text 47"},
        {"card": "Shape 48", "name": "Text 49", "position": "Text 50", "score": "Text 51", "movement": "Text 52", "watch_label": "Text 53", "watch": "Text 54"},
        {"card": "Shape 55", "name": "Text 56", "position": "Text 57", "score": "Text 58", "movement": "Text 59", "watch_label": "Text 60", "watch": "Text 61"},
    ]
    first_score = _number(first.get("score"))
    last_score = _number(latest.get("score"))
    net = _number(payload.get("net_movement"))
    driver_actual = _number(driver.get("actual"))
    driver_target = _number(driver.get("target"))
    driver_achievement = _fmt_percent(driver.get("achievement_pct"))
    driver_direction = "Lower is better" if _clean(driver.get("direction")).casefold() == "lower_better" else "Higher is better"
    decision_request = _clean(next_review.get("decision_request"))
    next_review_card = (
        "Assign owner + review date"
        if "assign an owner" in decision_request.casefold() or "review date" in decision_request.casefold()
        else "Confirm owner + success evidence"
    )
    ratio_text = ""
    if driver_actual is not None and driver_target not in {None, 0}:
        ratio_text = f"{driver_actual / driver_target:.1f}x target"
    actual_history = " → ".join(_fmt_native(row.get("actual"), driver.get("unit")) for row in driver_trend[-4:])
    history_scores = " → ".join(f"{_clean(row.get('label')).split()[0]} {_fmt_percent(row.get('score'))}" for row in trend[-4:])
    cards = [
        ("HISTORY COVERAGE", f"{len(trend)} measured periods", str(len(trend)), "PERIODS", "Window", f"{_clean(first.get('label'), 'History unavailable')} → {_clean(latest.get('label'), 'Selected period')}", BLUE, PALE_BLUE),
        ("BEST PERIOD", _clean(best.get("label"), "Unavailable"), _fmt_percent(best.get("score")), "BEST", "Compared with target", _fmt_percent((_number(best.get("score")) or 0) - 100), GREEN if _number(best.get("score")) is not None and _number(best.get("score")) >= 100 else BLUE, PALE_GREEN if _number(best.get("score")) is not None and _number(best.get("score")) >= 100 else PALE_BLUE),
        ("WORST PERIOD", _clean(worst.get("label"), "Unavailable"), _fmt_percent(worst.get("score")), "LOWEST", "Compared with target", _fmt_percent((_number(worst.get("score")) or 0) - 100), RED if _number(worst.get("score")) is not None and _number(worst.get("score")) < 70 else AMBER, PALE_RED if _number(worst.get("score")) is not None and _number(worst.get("score")) < 70 else PALE_AMBER),
        ("NET MOVEMENT", "First to latest", _fmt_percent(net, signed=True), "NET", "Trend status", _clean(payload.get("trend_status"), "History unavailable"), GREEN if net is not None and net > 0 else NEUTRAL, PALE_GREEN if net is not None and net > 0 else PALE_NEUTRAL),
        ("LEADING KPI", _clean(driver.get("label"), "No driver"), _fmt_percent(driver.get("weighted_impact")), "LOSS", "Direction", driver_direction, RED if driver else NEUTRAL, PALE_RED if driver else PALE_NEUTRAL),
        ("CURRENT ACTUAL", _clean(driver.get("label"), "No driver"), _fmt_native(driver_actual, driver.get("unit")), ratio_text or "Ratio unavailable", "Configured target", _fmt_native(driver_target, driver.get("unit")), RED if driver else NEUTRAL, PALE_RED if driver else PALE_NEUTRAL),
        ("DRIVER ACHIEVEMENT", _clean(driver.get("label"), "No driver"), driver_achievement, "ACHIEVEMENT", "Achievement gap", _fmt_percent(driver.get("gap")), RED if driver else NEUTRAL, PALE_RED if driver else PALE_NEUTRAL),
        ("NEXT REVIEW", "Decision checkpoint", _fmt_percent(payload.get("gap_to_target")), "GAP", "Commitment", next_review_card, AMBER, PALE_AMBER),
    ]
    for spec, card in zip(specs, cards):
        title, detail, value, movement, watch_label, watch, color, fill = card
        _set_fill(_find_shape(slide, spec["card"]), fill, LINE)
        _set_named_text(slide, spec["name"], title, color=NAVY, bold=True, size=5.3 if len(title) > 15 else 5.7)
        _set_named_text(slide, spec["position"], detail, color=MUTED, size=5.0 if len(detail) > 24 else 5.6)
        _set_named_text(slide, spec["score"], value, color=color, bold=True, size=10.5 if len(value) > 7 else 12)
        _set_named_text(slide, spec["movement"], movement, color=color, bold=True, size=4.8 if len(movement) > 10 else 5.3)
        _set_named_text(slide, spec["watch_label"], watch_label, color=FAINT, bold=True, size=5.0)
        _set_named_text(slide, spec["watch"], watch, color=MUTED if color != AMBER else NAVY, size=4.5 if len(watch) > 42 else 5.0)
    context_bar = _find_shape(slide, "Shape 62")
    context_text = _find_shape(slide, "Text 63")
    if context_bar:
        context_bar.top = int(5.7 * 914400)
        context_bar.height = int(0.82 * 914400)
        _set_fill(context_bar, PALE_BLUE, PALE_BLUE)
    if context_text:
        context_text.top = int(5.77 * 914400)
        context_text.height = int(0.68 * 914400)
    direction_text = driver_direction.casefold() if driver else "direction unavailable"
    context = (
        f"History: {history_scores or 'No measured periods'} | {_clean(payload.get('trend_status'), 'History unavailable')}.\n"
        f"{_clean(driver.get('label'), 'Leading KPI')}: {actual_history or 'No driver history'} vs {_fmt_native(driver_target, driver.get('unit'))} target | {driver_achievement} achievement | {direction_text}."
    )
    _set_named_text(slide, "Text 63", context, color=NAVY, bold=True, size=6.7)
    _footer(slide, period, 10)


def _slide_eleven(prs: Presentation, payload: dict[str, Any]) -> None:
    slide = prs.slides[10]
    period = _latest_period(payload)
    _set_header(slide, "6F. Employee Detail Appendix — Full Scorecard", "One-page reference: score, MoM, status, top KPI, weighted loss, and actual / target values.", period)
    headers = [("Text 9", "Employee"), ("Text 11", "Position"), ("Text 13", "Score"), ("Text 15", "MoM"), ("Text 17", "Status"), ("Text 19", "Top KPI | weighted loss"), ("Text 21", "Actual"), ("Text 23", "Target"), ("Text 25", "Focus")]
    for name, value in headers:
        _set_named_text(slide, name, value, color=FAINT, bold=True)
    people = _person_rows(payload)[:13]
    groups = _text_groups(slide, top_min=1.5, top_max=6.2, threshold=0.14)[:13]
    for index in range(13):
        group = groups[index] if index < len(groups) else []
        fields = sorted(group, key=lambda shape: _inches(shape.left))
        person = people[index] if index < len(people) else None
        if not person or len(fields) < 9:
            for field in fields:
                _set_shape_text(field, "")
            if group:
                _set_table_row_background(slide, _inches(group[0].top), WHITE)
            continue
        kpi = _primary_kpi(person)
        status = _status_label(person)
        status_key = _status_key(person)
        values = [
            _clean(person.get("name"), "Unknown employee"),
            _clean(person.get("position"), "Unassigned"),
            _fmt_percent(person.get("score")),
            _movement_display(person.get("movement")),
            status,
            _kpi_gap_text(kpi) if kpi else "No KPI gap",
            _fmt_native(kpi.get("actual"), kpi.get("unit")) if kpi else "N/A",
            _fmt_native(kpi.get("target"), kpi.get("unit")) if kpi else "N/A",
            "Review KPI gap" if (_number(person.get("weighted_loss")) or 0) > 0 else "Monitor gap",
        ]
        for field, value in zip(fields[:9], values):
            color = _status_color(status_key) if field in fields[2:5] else RED if field is fields[5] and (_number(person.get("weighted_loss")) or 0) > 0 else None
            size = 4.7 if field is fields[5] and len(value) > 36 else 5.0 if field is fields[1] and len(value) > 18 else 5.45 if len(value) > 24 else 5.8
            _set_shape_text(field, value, color=color, bold=field is fields[0] or field is fields[4], size=size)
        _set_table_row_background(slide, _inches(fields[0].top), PALE_RED if status == "BELOW" else PALE_NEUTRAL if index % 2 else WHITE, status=status_key)
    _footer(slide, period, 11)


def _add_commitment_text(slide: Any, x: float, y: float, width: float, height: float, value: Any, *, size: float, color: RGBColor, bold: bool = False) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = _clean(value)
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _slide_twelve(prs: Presentation, payload: dict[str, Any]) -> None:
    """Close the deck with measurable commitments instead of a generic summary."""

    # The approved editable template intentionally exposes only DEFAULT and
    # blank layouts.  Use the blank layout so this page works with the real
    # template as well as the lightweight test fixture.
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(247, 249, 252)
    period = _latest_period(payload)
    next_review = payload.get("next_review") or {}
    driver = payload.get("driver") or {}
    actions = list(payload.get("actions") or [])
    action = actions[0] if actions else {}
    owner = _clean(action.get("owner_display"), "Owner needed")
    due = _clean(action.get("due_date_display"), "Due date needed")
    status = _clean(action.get("status_display"), "Status not recorded")
    direction = "lower is better" if _clean(driver.get("direction")).casefold() == "lower_better" else "higher is better"
    action_detail = (
        f"Owner: {owner} | Due/cadence: {due} | Status: {status}. "
        f"{_clean(next_review.get('action_requirement'), 'Action ownership requirement unavailable')}"
    )
    commitments = [
        ("Overall target", _clean(next_review.get("overall_target"), "Overall performance target: 100.0%.")),
        ("Leading KPI target", f"{_clean(next_review.get('leading_kpi_target'), 'Leading KPI target needed')}. Direction: {direction}. {_clean(next_review.get('expected_movement'), '')}"),
        ("Root-cause validation", _clean(next_review.get("root_cause_requirement"), "Root-cause validation required.")),
        ("People review", _clean(next_review.get("people_requirement"), "Employee review requirement unavailable.")),
        ("Action ownership", action_detail),
        ("Success evidence", f"{_clean(next_review.get('success_evidence'), 'Compare the next period score, KPI actuals, and targets.') } {_clean(next_review.get('escalation_rule'), '')}"),
    ]

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_WIDTH), Inches(0.7))
    _set_fill(header, NAVY, NAVY)
    _add_commitment_text(slide, 0.38, 0.10, 10.3, 0.28, "7. Next Review Commitments", size=18, color=WHITE, bold=True)
    _add_commitment_text(slide, 0.39, 0.40, 10.7, 0.16, "The next review is successful only when the target, owner, evidence, and escalation condition are explicit.", size=7.2, color=RGBColor(220, 228, 242))
    _add_commitment_text(slide, 11.9, 0.16, 1.0, 0.18, period, size=6.5, color=WHITE, bold=True)

    for index, (label, detail) in enumerate(commitments):
        y = 1.12 + index * 0.70
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(12.23), Inches(0.57))
        _set_fill(card, WHITE, LINE)
        _add_commitment_text(slide, 0.72, y + 0.06, 2.15, 0.42, label, size=7.2, color=NAVY, bold=True)
        _add_commitment_text(slide, 2.95, y + 0.04, 9.55, 0.46, detail, size=6.4 if len(detail) > 120 else 6.8, color=MUTED)

    decision = _clean(next_review.get("decision_request"), "Management decision required: assign an owner and review date before the next reporting period.")
    decision_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.55), Inches(12.23), Inches(0.62))
    decision_fill = PALE_AMBER if "needed" in decision.casefold() or "assign" in decision.casefold() else PALE_BLUE
    _set_fill(decision_box, decision_fill, decision_fill)
    _add_commitment_text(slide, 0.72, 5.66, 11.85, 0.36, decision, size=7.6, color=NAVY, bold=True)

    question = _clean(next_review.get("question"), "Did the action plan reduce the leading KPI gap?")
    question_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.25), Inches(6.45), Inches(8.85), Inches(0.48))
    _set_fill(question_box, NAVY, NAVY)
    _add_commitment_text(slide, 2.42, 6.53, 8.5, 0.30, question, size=9.2, color=WHITE, bold=True)
    _add_commitment_text(slide, 0.42, 7.22, 5.9, 0.13, f"Marketing PMS | {period} | Next Review Commitments", size=5.5, color=FAINT)
    _add_commitment_text(slide, 12.72, 7.20, 0.25, 0.15, "12", size=5.5, color=FAINT)


def build_marketing_pptx(period_label: str = "June 2026", report_data: dict[str, Any] | None = None) -> bytes:
    """Return a filter-specific Marketing deck using the approved slide frame."""

    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Marketing report template is missing: {TEMPLATE_PATH}")
    payload = _prepare_payload(period_label, report_data)
    prs = Presentation(str(TEMPLATE_PATH))
    team_label = _report_team_label(payload)
    prs.core_properties.title = f"{team_label} PMS Performance Report"
    prs.core_properties.subject = _clean(payload.get("scope_label"), team_label)
    prs.core_properties.author = "SGH Hub"
    _slide_one(prs, payload)
    _slide_two(prs, payload)
    _slide_three(prs, payload)
    _slide_four(prs, payload)
    _slide_five(prs, payload)
    _slide_detail_pages(prs, payload)
    _slide_ten(prs, payload)
    _slide_eleven(prs, payload)
    _slide_twelve(prs, payload)
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


__all__ = ["TEMPLATE_PATH", "build_marketing_pptx"]
