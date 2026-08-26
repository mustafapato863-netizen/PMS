"""Dynamic UAE / Offshore EGY PowerPoint using the approved Departments deck.

The corrected Departments Performance Review deck is the source of every
visual frame in this module.  The builder keeps its masters, typography,
geometry, spacing, table system, and narrative order, then replaces only the
inherited content with the selected report snapshot.

The report service supplies the same filtered records, history, KPI
definitions, evidence, and actions used by the API.  The snapshot normalizer
from the Marketing report is reused for the calculation contract because the
corrected Offshore reference contains the same score / baseline / MoM / KPI
loss grammar in its Marketing section.
"""

from __future__ import annotations

import io
import textwrap
from copy import deepcopy
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from exports import marketing_legendary_pptx_builder as base


TEMPLATE_PATH = Path(__file__).resolve().parent / "templates" / "Offshore_Departments_Performance_Review.pptx"
REFERENCE_SLIDES = 28
MISSING = "\u2014"

# The values are the shared semantic colors used by the approved reference
# and the existing Legendary Marketing export.  Geometry and typography still
# come from the source deck; these colors are used only for data-dependent
# status fills and progress bars.
NAVY = base.NAVY
DARK_NAVY = base.DARK_NAVY
TEAL = base.TEAL
RED = base.RED
AMBER = base.AMBER
GREEN = base.GREEN
BLUE = base.BLUE
MUTED = base.MUTED
FAINT = base.FAINT
BODY = base.BODY
WHITE = base.WHITE
LINE = base.LINE
TRACK = base.TRACK
PALE_RED = base.PALE_RED
PALE_AMBER = base.PALE_AMBER
PALE_GREEN = base.PALE_GREEN
PALE_BLUE = base.PALE_BLUE
PALE_NEUTRAL = base.PALE_NEUTRAL


SECTION_SPECS = [
    {
        "key": "inbound",
        "label": "Call Center \u2013 Inbound",
        "divider": 3,
        "performance": 4,
        "actions": 5,
        "followup": None,
    },
    {
        "key": "outbound",
        "label": "Call Center \u2013 Outbound",
        "divider": 6,
        "performance": 7,
        "actions": 8,
        "followup": None,
    },
    {
        "key": "rcm",
        "label": "RCM – Approvals",
        "divider": 9,
        "performance": 10,
        "actions": 11,
        "followup": 12,
    },
    {
        "key": "marketing",
        "label": "Marketing",
        "divider": 13,
        "performance": 14,
        "actions": 15,
        "followup": None,
    },
]


def _number(value: Any) -> float | None:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return None
    return number if number == number else None


def _clean(value: Any, fallback: str = MISSING) -> str:
    if value is None:
        return fallback
    text = " ".join(str(value).replace("\n", " ").split()).strip()
    return text or fallback


def _fit(value: Any, limit: int, fallback: str = MISSING) -> str:
    text = _clean(value, "")
    if not text:
        return fallback
    return textwrap.shorten(text, width=max(8, limit), placeholder="...")


def _full(value: Any, fallback: str = MISSING) -> str:
    """Return the complete value for a single-line, auto-fit text frame."""

    return _clean(value, fallback)


def _fit_wrapped(value: Any, limit: int, *, max_lines: int = 2, fallback: str = MISSING) -> str:
    """Fit decision text without hiding truncation behind an ellipsis."""

    text = _clean(value, "")
    if not text:
        return fallback
    lines = textwrap.wrap(
        text,
        width=max(8, limit),
        break_long_words=False,
        break_on_hyphens=False,
    )
    if len(lines) <= max_lines:
        return "\n".join(lines)
    visible = lines[: max(1, max_lines)]
    continuation = "[continued]"
    if visible:
        visible[-1] = f"{visible[-1]} {continuation}"
    return "\n".join(visible)


def _pct(value: Any, *, signed: bool = False, fallback: str = MISSING) -> str:
    number = _number(value)
    if number is None:
        return fallback
    return f"{number:+.1f}%" if signed else f"{number:.1f}%"


def _fmt_pct(value: Any) -> str:
    """Compatibility alias used by the presentation formatting helpers."""

    return _pct(value)


def _mom(value: Any) -> str:
    return "NEW" if _number(value) is None else _pct(value, signed=True)


def _native(value: Any, unit: Any = "") -> str:
    return base._fmt_native(value, unit, MISSING)


def _score_status(score: Any) -> str:
    number = _number(score)
    if number is None:
        return "Data unavailable"
    if number < 70:
        return "Requires action"
    if number < 90:
        return "Watch"
    return "On track"


def _kpi_status(kpi: dict[str, Any] | None) -> str:
    if not kpi or _number(kpi.get("achievement_pct")) is None:
        return "Data unavailable"
    achievement = _number(kpi.get("achievement_pct")) or 0.0
    if achievement >= 100:
        return "On track"
    if achievement >= 70:
        return "Watch"
    return "Requires action"


def _status_code(status: Any) -> str:
    normalized = _clean(status, "Data unavailable").casefold()
    if normalized in {"on track", "meets expectations", "within target", "completed"}:
        return "On track"
    if normalized in {"watch", "at risk", "average", "below target"}:
        return "Watch"
    if normalized in {"requires action", "critical", "urgent", "below average"}:
        return "Requires action"
    return "Data unavailable"


def _status_color(status: Any) -> RGBColor:
    return {
        "On track": GREEN,
        "Watch": AMBER,
        "Requires action": RED,
        "Data unavailable": MUTED,
    }.get(_status_code(status), MUTED)


def _status_fill(status: Any) -> RGBColor:
    return {
        "On track": PALE_GREEN,
        "Watch": PALE_AMBER,
        "Requires action": PALE_RED,
        "Data unavailable": PALE_NEUTRAL,
    }.get(_status_code(status), PALE_NEUTRAL)


def _status_label(status: Any) -> str:
    return {
        "On track": "ON TRACK",
        "Watch": "WATCH",
        "Requires action": "REQUIRES ACTION",
        "Data unavailable": "NOT AVAILABLE",
    }.get(_status_code(status), "NOT AVAILABLE")


def _direction_label(value: Any) -> str:
    return "Lower is better" if _clean(value, "").casefold() in {"lower_better", "lower is better"} else "Higher is better"


def _shape(slide: Any, name: str) -> Any:
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"Reference shape not found: {name}")


def _text_shape(slide: Any, name: str) -> Any:
    shape = _shape(slide, name)
    if not shape.has_text_frame:
        raise TypeError(f"Reference shape is not a text frame: {name}")
    return shape


def _write_shape(shape: Any, value: Any, *, color: RGBColor | None = None, align: PP_ALIGN | None = None) -> None:
    text = MISSING if value is None or value == "" else str(value)
    if not shape.has_text_frame:
        return
    frame = shape.text_frame
    frame.word_wrap = True
    paragraphs = list(frame.paragraphs)
    if not paragraphs:
        frame.text = text
        paragraphs = list(frame.paragraphs)
    first = paragraphs[0]
    runs = list(first.runs)
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
    else:
        first.text = text
    for paragraph in paragraphs[1:]:
        for run in paragraph.runs:
            run.text = ""
    if align is not None:
        for paragraph in frame.paragraphs:
            paragraph.alignment = align
    if color is not None:
        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = color


def _write(slide: Any, name: str, value: Any, *, color: RGBColor | None = None, align: PP_ALIGN | None = None) -> None:
    _write_shape(_text_shape(slide, name), value, color=color, align=align)


def _write_single_line(
    slide: Any,
    name: str,
    value: Any,
    *,
    fallback: str = MISSING,
    color: RGBColor | None = None,
    align: PP_ALIGN | None = None,
) -> None:
    """Write complete data without ellipsizing it inside a narrow source cell.

    PowerPoint's native text-to-fit behavior keeps long names and KPI labels
    visible while preserving the reference geometry.  This is intentionally
    limited to one-line data cells; narrative cells continue to use the
    explicit wrapped/continued treatment above.
    """

    shape = _text_shape(slide, name)
    _write_shape(shape, _full(value, fallback), color=color, align=align)
    shape.text_frame.word_wrap = False
    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _fill(slide: Any, name: str, color: RGBColor, *, line: RGBColor | None = None) -> None:
    shape = _shape(slide, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is not None:
        shape.line.color.rgb = line


def _set_width(shape: Any, width: float, *, left: float | None = None) -> None:
    shape.width = Inches(max(0.01, width))
    if left is not None:
        shape.left = Inches(left)


def _resize_bar(slide: Any, name: str, ratio: Any, *, color: RGBColor | None = None) -> None:
    shape = _shape(slide, name)
    left = shape.left / 914400
    original_width = shape.width / 914400
    value = _number(ratio)
    value = 0.0 if value is None else max(0.0, min(1.0, value))
    _set_width(shape, original_width * value, left=left)
    if color is not None:
        _fill(slide, name, color)


def _clear_slides(prs: Presentation) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def _append_template_slide(prs: Presentation, source: Presentation, source_index: int) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    tree = slide.shapes._spTree
    for child in list(tree)[2:]:
        tree.remove(child)
    source_tree = source.slides[source_index].shapes._spTree
    for child in list(source_tree)[2:]:
        tree.insert_element_before(deepcopy(child), "p:extLst")
    return slide


def _period(snapshot: dict[str, Any], fallback: str) -> str:
    return _clean(snapshot.get("latest_period_label") or snapshot.get("period_label"), fallback)


def _branch_label(report_data: dict[str, Any]) -> str:
    region = _clean(report_data.get("region_label"), "Selected branch")
    return "SGH Offshore Branch" if "offshore" in region.casefold() or region.casefold() in {"egy", "egypt"} else region


def _set_common(slide: Any, title: str, subtitle: str, period: str, branch: str, page: int) -> None:
    _write(slide, "Text 2", title)
    _write(slide, "Text 3", subtitle)
    _write(slide, "Text 4", _fit(f"{period}  |  {branch}  |  {page:02d}", 96))
    _write(slide, "Text 6", "CONFIDENTIAL \u2014 INTERNAL USE ONLY")
    _write(slide, "Text 7", f"{page:02d}")


def _set_divider(slide: Any, title: str, subtitle: str, period: str, branch: str, page: int, bullets: list[str]) -> None:
    _write(slide, "Text 3", "SECTION")
    _write(slide, "Text 4", title)
    _write(slide, "Text 5", subtitle)
    _write(slide, "Text 7", "What we will review")
    for index, name in enumerate(("Text 9", "Text 11", "Text 13")):
        _write(slide, name, _fit(bullets[index] if index < len(bullets) else MISSING, 38))
    _write(slide, "Text 14", _fit(f"{period}  |  {branch}  |  {page:02d}", 96))
    _write(slide, "Text 15", "CONFIDENTIAL \u2014 INTERNAL USE ONLY")
    _write(slide, "Text 16", f"{page:02d}")


def _text_shapes(slide: Any, *, top_min: float, top_max: float, left_min: float = 0.0, left_max: float = 13.33) -> list[Any]:
    shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Source table/card backgrounds are auto-shapes with an empty text
        # frame.  They must not be mistaken for table value cells when rows
        # are discovered geometrically.
        if not any(paragraph.text for paragraph in shape.text_frame.paragraphs):
            continue
        top = shape.top / 914400
        left = shape.left / 914400
        if top_min <= top <= top_max and left_min <= left <= left_max:
            shapes.append(shape)
    return sorted(shapes, key=lambda item: (item.top, item.left))


def _row_groups(slide: Any, *, top_min: float, top_max: float, min_columns: int) -> list[list[Any]]:
    shapes = _text_shapes(slide, top_min=top_min, top_max=top_max)
    groups: list[list[Any]] = []
    for shape in shapes:
        top = shape.top / 914400
        if not groups or abs(top - groups[-1][0].top / 914400) > 0.07:
            groups.append([shape])
        else:
            groups[-1].append(shape)
    return [sorted(group, key=lambda item: item.left) for group in groups if len(group) >= min_columns]


def _write_role_people(shape: Any, role: Any, names: Iterable[Any], *, name_size: float = 5.4, height: float = 0.30) -> None:
    if not shape.has_text_frame:
        return
    frame = shape.text_frame
    original = frame.paragraphs[0].runs[0] if frame.paragraphs and frame.paragraphs[0].runs else None
    original_name = original.font.name if original else None
    original_size = original.font.size if original else None
    original_bold = original.font.bold if original else None
    original_color = original.font.color.rgb if original and original.font.color and original.font.color.type is not None else NAVY
    frame.clear()
    frame.word_wrap = True
    first = frame.paragraphs[0]
    first.text = _full(role)
    first.alignment = PP_ALIGN.LEFT
    first.space_before = 0
    first.space_after = 0
    for run in first.runs:
        if original_name:
            run.font.name = original_name
        if original_size:
            run.font.size = original_size
        if original_bold is not None:
            run.font.bold = original_bold
        run.font.color.rgb = original_color
    role_runs = list(first.runs)
    names_clean = []
    for value in names:
        name = _clean(value, "")
        if name and name.casefold() not in {candidate.casefold() for candidate in names_clean}:
            names_clean.append(name)
    if names_clean:
        detail = frame.add_paragraph()
        detail.text = ", ".join(names_clean)
        detail.alignment = PP_ALIGN.LEFT
        detail.space_before = 0
        detail.space_after = 0
        detail.line_spacing = 1.0
        for run in detail.runs:
            if role_runs:
                run.font.name = role_runs[0].font.name
            run.font.size = Pt(name_size)
            run.font.bold = False
            run.font.color.rgb = MUTED
    shape.height = Inches(height)
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _copy_data(data: dict[str, Any]) -> dict[str, Any]:
    # The payload consists of JSON-like values.  A shallow copy of each list
    # is enough and avoids retaining SQLAlchemy objects in the presentation.
    result = dict(data)
    result["filters"] = dict(data.get("filters") or {})
    result["records"] = [dict(row) for row in data.get("records") or []]
    result["selected_records"] = [dict(row) for row in data.get("selected_records") or []]
    result["history"] = [
        {**dict(item), "records": [dict(row) for row in item.get("records") or []]}
        for item in data.get("history") or []
    ]
    result["actions"] = [dict(row) for row in data.get("actions") or []]
    return result


def _team_kind(value: Any) -> str:
    text = _clean(value, "").casefold()
    if "marketing" in text:
        return "marketing"
    if "inbound" in text:
        return "inbound"
    if "outbound" in text:
        return "outbound"
    if "rcm" in text or "revenue cycle" in text or "approval" in text:
        return "rcm"
    return "other"


def _belongs_to_team(record: dict[str, Any], team: str) -> bool:
    actual = _clean(record.get("team"), "")
    return actual.casefold() == team.casefold() or _team_kind(actual) == _team_kind(team) and _team_kind(team) != "other"


def _team_snapshot(raw: dict[str, Any], team: str, period_label: str) -> dict[str, Any]:
    local = _copy_data(raw)
    current = [row for row in local.get("records") or [] if _belongs_to_team(row, team)]
    selected = [row for row in local.get("selected_records") or [] if _belongs_to_team(row, team)]
    history = []
    for item in local.get("history") or []:
        rows = [row for row in item.get("records") or [] if _belongs_to_team(row, team)]
        if rows:
            history.append({**item, "records": rows})
    local["records"] = current
    local["selected_records"] = selected
    local["history"] = history
    local["actions"] = [
        action
        for action in local.get("actions") or []
        if _clean(action.get("team"), "").casefold() in {team.casefold(), "selected scope", ""}
        or _team_kind(action.get("team")) == _team_kind(team) and _team_kind(team) != "other"
    ]
    local["filters"] = {**(local.get("filters") or {}), "team": team}
    return base._prepare_snapshot(period_label, local)


def _empty_snapshot(period_label: str, raw: dict[str, Any]) -> dict[str, Any]:
    local = _copy_data(raw)
    local["records"] = []
    local["selected_records"] = []
    local["history"] = []
    local["actions"] = []
    return base._prepare_snapshot(period_label, local)


def _team_view(name: str, snapshot: dict[str, Any]) -> dict[str, Any]:
    score = _number(snapshot.get("latest_score"))
    baseline = _number(snapshot.get("baseline_score"))
    people = list(snapshot.get("all_people") or [])
    if score is None:
        values = [_number(person.get("score")) for person in people]
        values = [value for value in values if value is not None]
        score = sum(values) / len(values) if values else None
    if baseline is None:
        values = [_number(person.get("baseline_score")) for person in people]
        values = [value for value in values if value is not None]
        baseline = sum(values) / len(values) if values else None
    movement = score - baseline if score is not None and baseline is not None else None
    return {
        "name": name,
        "score": score,
        "baseline": baseline,
        "movement": movement,
        "target": _number(snapshot.get("target_score")) or 100.0,
        "status": _score_status(score),
        "people": people,
        "affected_count": sum(1 for person in people if _score_status(person.get("score")) == "Requires action"),
        "kpis": [dict(row) for row in snapshot.get("kpis") or []],
        "snapshot": snapshot,
    }


def _role_rows(snapshot: dict[str, Any]) -> list[dict[str, Any]]:
    grouped: dict[str, list[dict[str, Any]]] = {}
    for person in snapshot.get("all_people") or []:
        role = _clean(person.get("position"), "Unassigned")
        grouped.setdefault(role, []).append(person)
    rows = []
    for role, people in grouped.items():
        def average(key: str) -> float | None:
            values = [_number(person.get(key)) for person in people]
            values = [value for value in values if value is not None]
            return sum(values) / len(values) if values else None

        score = average("score")
        baseline = average("baseline_score")
        rows.append({
            "role": role,
            "score": score,
            "baseline": baseline,
            "movement": score - baseline if score is not None and baseline is not None else None,
            "status": _score_status(score),
            "people": [person.get("name") for person in people],
            "count": len(people),
        })
    return sorted(rows, key=lambda row: (_number(row.get("score")) is None, _number(row.get("score")) or 0, _clean(row.get("role")).casefold()))


def _kpi_match(left: dict[str, Any], right: dict[str, Any]) -> bool:
    left_keys = {
        _clean(left.get(key), "").casefold()
        for key in ("key", "group_key", "label")
        if _clean(left.get(key), "")
    }
    left_keys.update(_clean(value, "").casefold() for value in left.get("source_keys") or [])
    right_keys = {
        _clean(right.get(key), "").casefold()
        for key in ("key", "group_key", "label", "kpi_key")
        if _clean(right.get(key), "")
    }
    return bool(left_keys & right_keys)


def _person_kpis(person: dict[str, Any]) -> list[dict[str, Any]]:
    return sorted(
        [dict(row) for row in person.get("kpis") or [] if isinstance(row, dict)],
        key=lambda row: (-(_number(row.get("weighted_impact")) or 0.0), _clean(row.get("label")).casefold()),
    )


def _driver_people(snapshot: dict[str, Any], kpi: dict[str, Any]) -> list[dict[str, Any]]:
    return [
        person
        for person in snapshot.get("all_people") or []
        if any(_kpi_match(kpi, employee_kpi) for employee_kpi in _person_kpis(person))
    ]


def _root_evidence(snapshot: dict[str, Any], person: dict[str, Any], kpi: dict[str, Any]) -> tuple[str, str]:
    row = next((candidate for candidate in snapshot.get("root_cause_rows") or [] if _kpi_match(candidate, kpi)), None)
    if row and "confirmed" in _clean(row.get("evidence_status"), "").casefold():
        return "Confirmed Root Cause", _clean(row.get("recorded_root_cause"), "Confirmed cause recorded")
    if "confirmed" in _clean(person.get("root_cause_status"), "").casefold() and _clean(person.get("root_cause"), ""):
        return "Confirmed Root Cause", _clean(person.get("root_cause"))
    if row:
        return "Investigation Required", _clean(row.get("required_validation"), "Validate the operational cause before assigning corrective action.")
    return "KPI Evidence", "The KPI identifies where to investigate; no operational root cause is recorded."


def _action_rows(snapshot: dict[str, Any]) -> list[dict[str, Any]]:
    actions = [dict(row) for row in snapshot.get("actions") or []]
    if not actions:
        # The snapshot normally includes transparent proposals for uncovered
        # drivers.  This fallback keeps the action section useful when a
        # filtered team has no recorded action and no proposal was produced.
        for kpi in (snapshot.get("kpis") or [])[:4]:
            if (_number(kpi.get("weighted_impact")) or 0.0) <= 0:
                continue
            actions.append({
                "is_proposed": True,
                "action_type": kpi.get("label") or "KPI workstream",
                "linked_kpi_key": kpi.get("key") or kpi.get("group_key"),
                "action_display": f"Validate {kpi.get('label', 'KPI')} workflow and focused coaching.",
                "owner_display": "Owner needed",
                "due_date_display": "Due date needed",
                "status_display": "Proposed",
                "success_metric_display": f"Move {_clean(kpi.get('label'), 'KPI')} toward {_native(kpi.get('target'), kpi.get('unit'))}.",
                "employee_name": "Selected team scope",
            })
    result = []
    kpis = [dict(row) for row in snapshot.get("kpis") or []]
    for action in actions:
        linked = next((kpi for kpi in kpis if _kpi_match(kpi, {"key": action.get("linked_kpi_key")})), None)
        people = _driver_people(snapshot, linked) if linked else []
        names = [person.get("name") for person in people if person.get("name")]
        scope = _clean(action.get("employee_name"), "")
        if names:
            scope = ", ".join(names[:4])
        priority = _clean(action.get("priority"), "")
        if not priority:
            priority = "High" if linked and (_number(linked.get("weighted_impact")) or 0) >= 10 else "Medium"
        result.append({
            "priority": priority,
            "workstream": _clean(action.get("action_type") or (linked or {}).get("label"), "Management workstream"),
            "owner": _clean(action.get("owner_display") or action.get("owner"), "Owner needed"),
            "due": _clean(action.get("due_date_display") or action.get("due_date"), "Due date needed"),
            "status": _clean(action.get("status_display") or action.get("status"), "Status not recorded"),
            "scope": scope or "Selected team scope",
            "execution": _clean(action.get("action_display") or action.get("action_text"), "Action text needed"),
            "success": _clean(action.get("success_metric_display") or action.get("success_metric"), "Success metric needed"),
            "linked": linked,
        })
    return sorted(result, key=lambda row: ({"high": 0, "medium": 1, "low": 2}.get(_clean(row.get("priority"), "").casefold(), 3), _clean(row.get("workstream")).casefold()))


def _kpi_gap(kpi: dict[str, Any]) -> float | None:
    achievement = _number(kpi.get("achievement_pct"))
    return max(0.0, 100.0 - achievement) if achievement is not None else None


def _kpi_loss(kpi: dict[str, Any]) -> float:
    return max(0.0, _number(kpi.get("weighted_impact")) or 0.0)


def _kpi_text(kpi: dict[str, Any]) -> str:
    if not kpi:
        return MISSING
    actual = _native(kpi.get("actual"), kpi.get("unit"))
    target = _native(kpi.get("target"), kpi.get("unit"))
    loss = _kpi_loss(kpi)
    return _fit_wrapped(
        f"{_clean(kpi.get('label'), 'KPI')} {_fmt_pct(kpi.get('achievement_pct'))} loss {_fmt_pct(loss)}",
        42,
        max_lines=2,
    )


def _action_focus(snapshot: dict[str, Any], *, limit: int = 2) -> str:
    kpis = [row for row in snapshot.get("kpis") or [] if _kpi_loss(row) > 0]
    if not kpis:
        return "No measurable KPI gap available"
    return ", ".join(_clean(row.get("label"), "KPI") for row in kpis[:limit])


def _employee_priority(snapshot: dict[str, Any]) -> list[dict[str, Any]]:
    people = [dict(row) for row in snapshot.get("all_people") or []]
    return sorted(
        people,
        key=lambda row: (_number(row.get("score")) is None, _number(row.get("score")) or 0.0, -sum(_kpi_loss(kpi) for kpi in _person_kpis(row)), _clean(row.get("name")).casefold()),
    )


def _ordered_sections(raw: dict[str, Any], period_label: str) -> list[dict[str, Any]]:
    current_names = []
    for record in raw.get("records") or []:
        name = _clean(record.get("team"), "Unassigned")
        if name.casefold() not in {item.casefold() for item in current_names}:
            current_names.append(name)
    used: set[str] = set()
    sections = []
    filters = raw.get("filters") or {}
    # A narrowed Reports-page scope should not leave unrelated empty
    # department dividers in the deck.  Region-only exports retain the full
    # approved branch rhythm; team/employee/position/level/KPI filters are
    # intentionally presented as the selected slice.
    narrowed_scope = any(
        filters.get(key)
        for key in ("team", "position", "employee_id", "performance_level", "grade", "status", "kpi")
    )
    for spec in SECTION_SPECS:
        actual = next((name for name in current_names if name.casefold() not in used and _team_kind(name) == spec["key"]), None)
        if actual:
            used.add(actual.casefold())
            snapshot = _team_snapshot(raw, actual, period_label)
            # Keep the approved reference vocabulary in the visual layer;
            # the underlying records remain identified by their live team.
            display = spec["label"]
        elif not narrowed_scope:
            snapshot = _empty_snapshot(period_label, raw)
            display = spec["label"]
        else:
            continue
        sections.append({**spec, "display": display, "snapshot": snapshot, "view": _team_view(display, snapshot)})
    for actual in current_names:
        if actual.casefold() in used:
            continue
        snapshot = _team_snapshot(raw, actual, period_label)
        sections.append({
            "key": f"other-{len(sections)}",
            "label": actual,
            "display": actual,
            "divider": 3,
            "performance": 4,
            "actions": 5,
            "followup": None,
            "snapshot": snapshot,
            "view": _team_view(actual, snapshot),
            "extra": True,
        })
    return sections


def _priority_views(sections: list[dict[str, Any]]) -> list[dict[str, Any]]:
    views = [section["view"] for section in sections]
    return sorted(views, key=lambda row: (_number(row.get("score")) is None, _number(row.get("score")) or 0.0, _clean(row.get("name")).casefold()))


def _overview_views(sections: list[dict[str, Any]]) -> list[dict[str, Any]]:
    preferred = {"inbound": 0, "marketing": 1, "outbound": 2, "rcm": 3}
    return sorted(
        [section["view"] for section in sections],
        key=lambda row: (preferred.get(_team_kind(row.get("name")), 99), _clean(row.get("name")).casefold()),
    )


def _card_status(
    slide: Any,
    card_shape: str,
    accent_shape: str,
    status: Any,
    *,
    card_fill: RGBColor | None = None,
) -> None:
    fill = card_fill if card_fill is not None else _status_fill(status)
    color = _status_color(status)
    _fill(slide, card_shape, fill)
    _fill(slide, accent_shape, color)


def _write_score_card(slide: Any, names: tuple[str, str, str, str, str, str, str], view: dict[str, Any], *, title: str | None = None) -> None:
    label, score_name, monthly_name, baseline_name, mom_name, target_name, status_name = names
    status = view.get("status")
    _write_single_line(slide, label, title or view.get("name"))
    _write(slide, score_name, _pct(view.get("score")))
    _write(slide, monthly_name, "Monthly performance")
    _write(slide, baseline_name, f"Baseline: {_pct(view.get('baseline'))}")
    _write(slide, mom_name, f"MoM: {_mom(view.get('movement'))}")
    _write(slide, target_name, f"Target: {_pct(view.get('target'))}")
    _write(slide, status_name, _status_label(status), color=_status_color(status))


def _write_kpi_card(slide: Any, names: tuple[str, str, str, str, str, str], kpi: dict[str, Any] | None) -> None:
    label, actual_name, baseline_name, mom_name, target_name, status_name = names
    if not kpi:
        for name in names:
            _write(slide, name, MISSING)
        return
    status = _kpi_status(kpi)
    _write_single_line(slide, label, kpi.get("label"))
    _write(slide, actual_name, _native(kpi.get("actual"), kpi.get("unit")))
    _write(slide, baseline_name, f"Baseline: {_native(kpi.get('baseline_actual'), kpi.get('unit'))}")
    _write(slide, mom_name, f"MoM: {_mom(kpi.get('mom'))}")
    _write(slide, target_name, f"Target: {_native(kpi.get('target'), kpi.get('unit'))}")
    _write(slide, status_name, _status_label(status), color=_status_color(status))


def _populate_cover(slide: Any, model: dict[str, Any], page: int) -> None:
    snapshot = model["snapshot"]
    period = model["period"]
    branch = model["branch"]
    sections = [section["display"] for section in model["sections"] if not section.get("extra")]
    actions = sum(len(_action_rows(section["snapshot"])) for section in model["sections"])
    _write(slide, "Text 3", "Offshore Departments")
    _write(slide, "Text 4", "Combined Performance Review")
    _write(slide, "Text 5", _fit(f"{period}  |  {branch}", 70))
    _write(slide, "Text 6", "  •  ".join(sections))
    _write(slide, "Text 8", "Narrative Flow")
    _write(slide, "Text 9", "Monthly performance summary, KPI cards with baseline and MoM change, then employee evidence and corrective action control.")
    _write(slide, "Text 11", f"{len(sections)} Departments")
    _write(slide, "Text 13", "Exceptions First")
    _write(slide, "Text 15", f"{actions} Action Focus")
    _write(slide, "Text 16", "CONFIDENTIAL \u2013 INTERNAL USE ONLY")


def _populate_overview(slide: Any, model: dict[str, Any], page: int) -> None:
    views = _overview_views(model["sections"])
    _set_common(slide, "Combined Monthly Performance Overview", "Selected-branch view across all department teams", model["period"], model["branch"], page)
    card_specs = [
        (("Text 10", "Text 11", "Text 12", "Text 13", "Text 14", "Text 15", "Text 17"), "Shape 8", "Shape 9"),
        (("Text 20", "Text 21", "Text 22", "Text 23", "Text 24", "Text 25", "Text 27"), "Shape 18", "Shape 19"),
        (("Text 30", "Text 31", "Text 32", "Text 33", "Text 34", "Text 35", "Text 37"), "Shape 28", "Shape 29"),
        (("Text 40", "Text 41", "Text 42", "Text 43", "Text 44", "Text 45", "Text 47"), "Shape 38", "Shape 39"),
    ]
    for index, (names, card, accent) in enumerate(card_specs):
        view = views[index] if index < len(views) else _team_view(MISSING, _empty_snapshot(model["period"], model["raw"]))
        _card_status(slide, card, accent, view.get("status"))
        _write_score_card(slide, names, view, title=_fit(view.get("name"), 25))
        # The overview source uses filled status pills; keep the inherited
        # white pill typography instead of applying the semantic color to the
        # text itself (which would disappear on an amber/red pill).
        _write(slide, names[-1], _status_label(view.get("status")), color=WHITE)
    score = model["snapshot"].get("overall_score")
    movement = model["snapshot"].get("movement")
    driver = (model["snapshot"].get("driver") or {})
    affected = len([person for person in model["snapshot"].get("all_people") or [] if _score_status(person.get("score")) == "Requires action"])
    top = max(model["snapshot"].get("all_people") or [], key=lambda person: _number(person.get("score")) or -1, default={})
    decline = min(
        [person for person in model["snapshot"].get("all_people") or [] if _number(person.get("movement")) is not None],
        key=lambda person: _number(person.get("movement")) or 0,
        default={},
    )
    takeaway = f"Executive takeaway: overall performance is {_pct(score)} with {_mom(movement)} MoM; {_clean(driver.get('label'), 'the leading KPI driver')} and {affected} employee(s) require the primary management focus."
    _write(slide, "Text 49", _fit(takeaway, 170))
    _write(slide, "Text 50", "Department-Level Action View")
    row_groups = _row_groups(slide, top_min=5.25, top_max=6.85, min_columns=6)
    for index, group in enumerate(row_groups[:4]):
        view = views[index] if index < len(views) else None
        if not view:
            values = [MISSING] * 6
        else:
            kpis = view.get("kpis") or []
            # The source cell is intentionally compact.  Show the highest
            # weighted-loss KPI here and keep the full KPI list in the
            # department detail / appendix slides.
            gap_text = _kpi_text(kpis[0]) if kpis else "No measured KPI gap"
            values = [
                _full(view.get("name")),
                _pct(view.get("score")),
                _pct(view.get("baseline")),
                _mom(view.get("movement")),
                _fit_wrapped(gap_text, 40, max_lines=2),
                _fit_wrapped(_action_focus(view["snapshot"]), 42, max_lines=2),
            ]
        for column, (shape, value) in enumerate(zip(group, values)):
            if column == 0:
                _write_single_line(slide, shape.name, value)
            else:
                _write_shape(shape, value)


def _populate_priority(slide: Any, model: dict[str, Any], page: int) -> None:
    _set_common(slide, "Combined Priority Map", "Action order across departments based on risk and weighted performance gaps", model["period"], model["branch"], page)
    views = _priority_views(model["sections"])
    slots = [
        (("Text 10", "Text 11", "Text 12", "Text 13", "Text 14", "Text 15"), "Shape 8", "Shape 9"),
        (("Text 18", "Text 19", "Text 20", "Text 21", "Text 22", "Text 23"), "Shape 16", "Shape 17"),
        (("Text 26", "Text 27", "Text 28", "Text 29", "Text 30", "Text 31"), "Shape 24", "Shape 25"),
        (("Text 34", "Text 35", "Text 36", "Text 37", "Text 38", "Text 39"), "Shape 32", "Shape 33"),
    ]
    for index, (names, card, accent) in enumerate(slots):
        view = views[index] if index < len(views) else None
        status = view.get("status") if view else "Data unavailable"
        _card_status(slide, card, accent, status)
        if not view:
            for name in names:
                _write(slide, name, MISSING)
            continue
        role = _clean(view.get("name"), "Department")
        kpis = view.get("kpis") or []
        bullets = [
            f"{_clean(kpi.get('label'), 'KPI')} {_native(kpi.get('actual'), kpi.get('unit'))} vs {_native(kpi.get('target'), kpi.get('unit'))} target"
            for kpi in kpis[:2]
        ]
        bullets.append(f"{view.get('affected_count', 0)} employee(s) require attention")
        _write_single_line(slide, names[0], f"{index + 1}. {_full(role)}")
        _write(slide, names[1], _pct(view.get("score")))
        _write(slide, names[2], _fit_wrapped(f"{_status_label(status).title()}: {_action_focus(view['snapshot'])}", 46, max_lines=2))
        for name, value in zip(names[3:], bullets):
            _write(slide, name, _fit_wrapped(f"• {value}", 46, max_lines=2))
    _write(slide, "Text 41", "Management reading: move from branch-level risk to department detail, then employee-level evidence and corrective action control.")


def _populate_section_divider(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int) -> None:
    view = section["view"]
    score = _pct(view.get("score"))
    subtitle = f"Current performance {_clean(score)}; review the largest KPI gaps and employee actions before the next cycle."
    bullets = [
        "Monthly KPI scorecards",
        f"{len(view.get('people') or [])} employee(s) in scope",
        "Employee details and action tracker",
    ]
    _set_divider(slide, section["display"], subtitle, model["period"], model["branch"], page, bullets)


def _populate_department_performance(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int, total: int) -> None:
    view = section["view"]
    title = f"{section['display']} | Performance" + (f" ({part}/{total})" if total > 1 else "")
    subtitle = "Current scorecard, baseline, MoM and employee-level evidence" + (" | continuation" if total > 1 and part > 1 else "")
    _set_common(slide, title, subtitle, model["period"], model["branch"], page)
    _write(slide, "Text 9", f"DEPARTMENT: {section['display']}")
    score_view = {
        "name": section["display"],
        "score": view.get("score"),
        "baseline": view.get("baseline"),
        "movement": view.get("movement"),
        "target": view.get("target"),
        "status": view.get("status"),
    }
    _card_status(slide, "Shape 10", "Shape 11", view.get("status"))
    _write_score_card(slide, ("Text 12", "Text 13", "Text 12", "Text 14", "Text 15", "Text 16", "Text 17"), score_view, title="Team Avg Performance")
    # The score card shares its source label frame with the monthly label;
    # restore the approved two-line composition explicitly.
    _write(slide, "Text 12", "Team Avg Performance")
    _write(slide, "Text 13", _pct(view.get("score")))
    _write(slide, "Text 14", f"Baseline {_pct(view.get('baseline'))}")
    _write(slide, "Text 15", f"MoM {_mom(view.get('movement'))}")
    _write(slide, "Text 16", f"Target {_pct(view.get('target'))}")
    _write(slide, "Text 17", _status_label(view.get("status")), color=_status_color(view.get("status")))
    kpis = (view.get("kpis") or [])[:3]
    kpi_card_groups = [
        (("Text 20", "Text 21", "Text 22", "Text 23", "Text 24", "Text 25"), "Shape 18", "Shape 19"),
        (("Text 28", "Text 29", "Text 30", "Text 31", "Text 32", "Text 33"), "Shape 26", "Shape 27"),
        (("Text 36", "Text 37", "Text 38", "Text 39", "Text 40", "Text 41"), "Shape 34", "Shape 35"),
    ]
    for index, (names, card, accent) in enumerate(kpi_card_groups):
        kpi = kpis[index] if index < len(kpis) else None
        _card_status(slide, card, accent, _kpi_status(kpi))
        _write_kpi_card(slide, names, kpi)
    if _shape_exists(slide, "Text 92"):
        _write(slide, "Text 92", f"Action Focus: {_action_focus(view.get('snapshot') or {})}.")
    _write(slide, "Text 42", "Employees Under Review" if section["key"] != "outbound" else "Below Average Employees")
    headers = _row_groups(slide, top_min=4.15, top_max=4.35, min_columns=6)
    if headers:
        kpi_labels = [(_clean(kpi.get("label"), "KPI") if kpi else "KPI") for kpi in (view.get("kpis") or [])[:3]]
        for shape, value in zip(headers[0], ["Employee", "Score", "Grade", *(kpi_labels + ["KPI", "KPI", "KPI"])[:3]]):
            _write_shape(shape, value)
    row_groups = _row_groups(slide, top_min=4.5, top_max=6.45, min_columns=6)
    top_kpis = (view.get("kpis") or [])[:3]
    for group, person in zip(row_groups[:5], rows[:5]):
        values = [
            _full(person.get("name")),
            _pct(person.get("score")),
            _full(person.get("grade")),
        ]
        for kpi in top_kpis:
            person_kpi = next((item for item in _person_kpis(person) if _kpi_match(kpi, item)), None)
            values.append(_native(person_kpi.get("actual"), person_kpi.get("unit")) if person_kpi else MISSING)
        values = (values + [MISSING] * len(group))[: len(group)]
        for column, (shape, value) in enumerate(zip(group, values)):
            if column in {0, 2}:
                _write_single_line(slide, shape.name, value)
            else:
                _write_shape(shape, value)
    for group in row_groups[len(rows[:5]):5]:
        for shape in group:
            _write_shape(shape, MISSING)


def _person_issue(person: dict[str, Any]) -> str:
    kpis = _person_kpis(person)
    if not kpis:
        return "No KPI evidence recorded"
    return "; ".join(_clean(kpi.get("label"), "KPI") for kpi in kpis[:3])


def _person_root_cause(snapshot: dict[str, Any], person: dict[str, Any]) -> str:
    kpi = _person_kpis(person)[0] if _person_kpis(person) else {}
    state, note = _root_evidence(snapshot, person, kpi)
    if state == "Confirmed Root Cause":
        return _fit_wrapped(note, 62, max_lines=3)
    return _fit_wrapped(f"{state} — validate before assigning cause", 62, max_lines=2)


def _person_action(person: dict[str, Any]) -> str:
    actions = person.get("actions") or []
    if actions:
        first = actions[0]
        return _fit_wrapped(first.get("action_display") or first.get("action_text"), 90, max_lines=3)
    return "Management action required; owner and due date are not recorded."


def _populate_department_actions(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int, total: int) -> None:
    title = f"{section['display']} | Corrective Actions" + (f" ({part}/{total})" if total > 1 else "")
    _set_common(slide, title, "Evidence-led issues, root-cause status and corrective action control", model["period"], model["branch"], page)
    row_groups = _row_groups(slide, top_min=1.75, top_max=5.7, min_columns=5)
    snapshot = section["snapshot"]
    for group, person in zip(row_groups[:5], rows[:5]):
        values = [
            _full(person.get("name")),
            f"{_pct(person.get('score'))} / {_full(person.get('grade'))}",
            _fit_wrapped(_person_issue(person), 38, max_lines=2),
            _person_root_cause(snapshot, person),
            _person_action(person),
        ]
        for column, (shape, value) in enumerate(zip(group, values)):
            if column == 0:
                _write_single_line(slide, shape.name, value)
            else:
                _write_shape(shape, value)
    for group in row_groups[len(rows[:5]):5]:
        for shape in group:
            _write_shape(shape, MISSING)


def _populate_rcm_followup(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int) -> None:
    _set_common(slide, "RCM Follow-up Control", "Named control areas for the next reporting cycle", model["period"], model["branch"], page)
    actions = _action_rows(section["snapshot"])
    people = _employee_priority(section["snapshot"])
    top_groups = _row_groups(slide, top_min=1.75, top_max=2.1, min_columns=5)
    if top_groups:
        person = people[0] if people else None
        if person:
            values = [
                _full(person.get("name")),
                f"{_pct(person.get('score'))} / {_full(person.get('grade'))}",
                _fit_wrapped(_person_issue(person), 38, max_lines=2),
                _person_root_cause(section["snapshot"], person),
                _person_action(person),
            ]
        else:
            values = [MISSING] * 5
        for column, (shape, value) in enumerate(zip(top_groups[0], values)):
            if column == 0:
                _write_single_line(slide, shape.name, value)
            else:
                _write_shape(shape, value)
    rows = actions[:3]
    fallback = [
        {"workstream": "Rejected cases review", "owner": "Owner needed", "due": "Weekly", "status": "Validate result"},
        {"workstream": "Pre-submission checklist", "owner": "Owner needed", "due": "Daily", "status": "Rejection rate moves toward target"},
        {"workstream": "Aging case queue", "owner": "Owner needed", "due": "Start / end of day", "status": "Submission SLA improves"},
    ]
    table_rows = []
    for index in range(3):
        row = rows[index] if index < len(rows) else fallback[index]
        table_rows.append([
            _fit_wrapped(row.get("workstream"), 32, max_lines=2),
            _fit_wrapped(row.get("owner"), 24, max_lines=2),
            _fit_wrapped(row.get("due"), 22, max_lines=2),
            _fit_wrapped(row.get("success") or row.get("status"), 52, max_lines=2),
        ])
    groups = _row_groups(slide, top_min=4.5, top_max=5.8, min_columns=4)
    for group, values in zip(groups[:3], table_rows):
        for shape, value in zip(group, values):
            _write_shape(shape, value)


def _populate_marketing_overview(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int) -> None:
    view = section["view"]
    snapshot = section["snapshot"]
    _set_common(slide, "Marketing | Monthly Performance Overview", "Marketing is included as a team in this combined review", model["period"], model["branch"], page)
    people = snapshot.get("all_people") or []
    below = [person for person in people if _number(person.get("score")) is not None and (_number(person.get("score")) or 0) < 70]
    top = max(people, key=lambda person: _number(person.get("score")) or -1, default={})
    decline = min([person for person in people if _number(person.get("movement")) is not None], key=lambda person: _number(person.get("movement")) or 0, default={})
    driver = (view.get("kpis") or [None])[0]
    cards = [
        (("Text 10", "Text 11", "Text 12", "Text 13", "Text 14", "Text 15", "Text 17"), "Shape 8", "Shape 9"),
    ]
    _card_status(slide, "Shape 8", "Shape 9", view.get("status"))
    _write(slide, "Text 10", "Marketing Main Score")
    _write(slide, "Text 11", _pct(view.get("score")))
    _write(slide, "Text 12", "Monthly performance")
    _write(slide, "Text 13", f"Baseline: {_pct(view.get('baseline'))}")
    _write(slide, "Text 14", f"MoM: {_mom(view.get('movement'))}")
    _write(slide, "Text 15", f"Target: {_pct(view.get('target'))}")
    _write(slide, "Text 17", _status_label(view.get("status")), color=WHITE)
    _card_status(slide, "Shape 18", "Shape 19", "Requires action" if below else "On track", card_fill=WHITE)
    _write(slide, "Text 20", "EMPLOYEES BELOW 70%")
    _write(slide, "Text 21", len(below))
    _write(slide, "Text 22", "require attention")
    _card_status(slide, "Shape 23", "Shape 24", "On track", card_fill=WHITE)
    _write(slide, "Text 25", "TOP SCORE")
    _write(slide, "Text 26", _pct(top.get("score")))
    _write_single_line(slide, "Text 27", top.get("name"))
    _card_status(
        slide,
        "Shape 28",
        "Shape 29",
        "Requires action" if _number(decline.get("movement")) is not None and (_number(decline.get("movement")) or 0) < 0 else "On track",
        card_fill=WHITE,
    )
    _write(slide, "Text 30", "LARGEST DECLINE")
    _write(slide, "Text 31", _mom(decline.get("movement")))
    _write_single_line(slide, "Text 32", decline.get("name"))
    _card_status(slide, "Shape 33", "Shape 34", view.get("status"), card_fill=WHITE)
    _write(slide, "Text 35", "GAP")
    _write(slide, "Text 36", _pct(max(0.0, (view.get("target") or 100) - (view.get("score") or 0))))
    _write(slide, "Text 37", "to target")
    takeaway = f"Executive takeaway: Marketing overall performance is {_pct(view.get('score'))} with {_mom(view.get('movement'))} MoM; {len(below)} employee(s) and {_clean((driver or {}).get('label'), 'the leading KPI driver')} are the primary management focus."
    _write(slide, "Text 39", _fit(takeaway, 170))
    _write(slide, "Text 40", "Department-Level Action View")
    roles = _role_rows(snapshot)
    groups = _row_groups(slide, top_min=5.22, top_max=6.95, min_columns=6)
    for group, role in zip(groups[:5], roles[:5]):
        values = [
            _full(role.get("role")),
            _pct(role.get("score")),
            _pct(role.get("baseline")),
            _mom(role.get("movement")),
            _status_label(role.get("status")),
            "Immediate corrective action" if role.get("status") == "Requires action" else "Monitor key gap" if role.get("status") == "Watch" else "Maintain performance",
        ]
        for column, (shape, value) in enumerate(zip(group, values)):
            if column == 0:
                _write_single_line(slide, shape.name, value)
            else:
                _write_shape(shape, value, color=_status_color(role.get("status")) if value == _status_label(role.get("status")) else None)
    for group in groups[len(roles[:5]):5]:
        for shape in group:
            _write_shape(shape, MISSING)


def _populate_marketing_roles(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int, total: int) -> None:
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"Marketing Performance by Role{suffix}", "Current performance, baseline, MoM and people assigned to each function", model["period"], model["branch"], page)
    _write(slide, "Text 8", "Role Performance \u2014 Main Score & MoM")
    _write(slide, "Text 55", "Attention Order")
    role_names = ["Text 10", "Text 15", "Text 20", "Text 25", "Text 30", "Text 35", "Text 40", "Text 45", "Text 50"]
    score_names = ["Text 13", "Text 18", "Text 23", "Text 28", "Text 33", "Text 38", "Text 43", "Text 48", "Text 53"]
    mom_names = ["Text 14", "Text 19", "Text 24", "Text 29", "Text 34", "Text 39", "Text 44", "Text 49", "Text 54"]
    track_names = ["Shape 11", "Shape 16", "Shape 21", "Shape 26", "Shape 31", "Shape 36", "Shape 41", "Shape 46", "Shape 51"]
    bar_names = ["Shape 12", "Shape 17", "Shape 22", "Shape 27", "Shape 32", "Shape 37", "Shape 42", "Shape 47", "Shape 52"]
    for index, name in enumerate(role_names):
        if index >= len(rows):
            _write(slide, name, MISSING)
            _write(slide, score_names[index], MISSING)
            _write(slide, mom_names[index], MISSING)
            _resize_bar(slide, bar_names[index], 0.0, color=MUTED)
            continue
        row = rows[index]
        _write_role_people(_text_shape(slide, name), row.get("role"), row.get("people") or [])
        _write(slide, score_names[index], _pct(row.get("score")))
        _write(slide, mom_names[index], _mom(row.get("movement")))
        _fill(slide, track_names[index], TRACK)
        _resize_bar(slide, bar_names[index], (_number(row.get("score")) or 0) / 100, color=_status_color(row.get("status")))
    attention = rows[:3]
    attention_names = [("Text 58", "Text 59", "Text 60"), ("Text 62", "Text 63", "Text 64"), ("Text 66", "Text 67", "Text 68")]
    for index, names in enumerate(attention_names):
        row = attention[index] if index < len(attention) else None
        if row:
            _write_single_line(slide, names[0], f"{index + 1}. {_full(row.get('role'))}")
            _write(slide, names[1], _pct(row.get("score")))
            _write(slide, names[2], f"MoM {_mom(row.get('movement'))} | {row.get('count', 0)} employee(s)")
        else:
            for name in names:
                _write(slide, name, MISSING)
    story = "Strong roles are masking lower scores in " + ", ".join(_clean(row.get("role"), "role") for row in rows[:3]) + "." if rows else "No role-level performance data is available."
    _write(slide, "Text 70", "Story signal")
    _write_single_line(slide, "Text 71", story)


def _populate_marketing_trend(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int, kpis: list[dict[str, Any]], part: int, total: int) -> None:
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"Marketing Trend & KPI Driver Impact{suffix}", "Historical direction before drilling into employee drivers", model["period"], model["branch"], page)
    snapshot = section["snapshot"]
    trend = snapshot.get("trend") or []
    history_count = snapshot.get("history_count") or len(trend)
    measured_rows = sum(_number(item.get("record_count")) or 0 for item in trend)
    best = snapshot.get("best_period") or {}
    worst = snapshot.get("worst_period") or {}
    _write(slide, "Text 10", "HISTORY COVERAGE")
    _write(slide, "Text 11", f"{history_count} PERIODS")
    _write(slide, "Text 12", f"{int(measured_rows)} measured rows")
    _write(slide, "Text 15", "BEST PERIOD")
    _write(slide, "Text 16", _pct(best.get("score")))
    _write_single_line(slide, "Text 17", best.get("label"))
    _write(slide, "Text 20", "WORST PERIOD")
    _write(slide, "Text 21", _pct(worst.get("score")))
    _write_single_line(slide, "Text 22", worst.get("label"))
    _write(slide, "Text 25", "NET MOVEMENT")
    _write(slide, "Text 26", _mom(snapshot.get("net_movement")))
    _write(slide, "Text 27", "First to latest")
    _write(slide, "Text 30", "TARGET GAP")
    _write(slide, "Text 31", _pct(max(0.0, (section["view"].get("target") or 100) - (section["view"].get("score") or 0))))
    _write(slide, "Text 32", f"{_clean(snapshot.get('latest_period_label'), 'Selected')} vs target")
    _write(slide, "Text 33", "Weighted score loss by KPI")
    rows = _row_groups(slide, top_min=3.0, top_max=6.0, min_columns=3)
    max_loss = max([_kpi_loss(kpi) for kpi in kpis] or [1.0])
    display_kpis = kpis[:8]
    for index, (group, kpi) in enumerate(zip(rows[:8], display_kpis)):
        label = _full(kpi.get("label"))
        loss = _kpi_loss(kpi)
        gap = _kpi_gap(kpi)
        affected = _number(kpi.get("affected_count"))
        values = [label, _fmt_pct(loss), f"{int(affected) if affected is not None else 0} gap(s) / {_fmt_pct(gap or 0)} gap"]
        # The row group contains label, loss, and gap/count text; progress
        # shapes are not text and are handled by their fixed name sequence.
        for column, (shape, value) in enumerate(zip(group, values)):
            if column == 0:
                _write_single_line(slide, shape.name, value)
            else:
                _write_shape(shape, value)
        row_index = index
        track_name = f"Shape {36 + row_index * 5}"
        bar_name = f"Shape {37 + row_index * 5}"
        _fill(slide, track_name, TRACK)
        _resize_bar(slide, bar_name, loss / max_loss if max_loss else 0, color=_status_color(_kpi_status(kpi)))
    for index in range(len(display_kpis), 8):
        label_name = f"Text {35 + index * 5}"
        loss_name = f"Text {38 + index * 5}"
        note_name = f"Text {39 + index * 5}"
        _write(slide, label_name, MISSING)
        _write(slide, loss_name, MISSING)
        _write(slide, note_name, MISSING)
        _resize_bar(slide, f"Shape {37 + index * 5}", 0.0, color=MUTED)
    _write(slide, "Text 76", "Reading rule: solve the highest weighted-loss drivers first; do not treat every KPI gap as equally material.")


def _populate_marketing_map(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int, kpis: list[dict[str, Any]], part: int, total: int) -> None:
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"Marketing Driver \u2192 Role \u2192 Employee Map{suffix}", "Connects overall KPI loss to affected people and functions", model["period"], model["branch"], page)
    _write(slide, "Text 8", "Priority driver map")
    row_groups = _row_groups(slide, top_min=2.27, top_max=5.1, min_columns=4)
    for group, kpi in zip(row_groups[:4], kpis[:4]):
        affected = _driver_people(section["snapshot"], kpi)
        roles = list(dict.fromkeys(_clean(person.get("position"), "Unassigned") for person in affected))
        names = list(dict.fromkeys(_clean(person.get("name"), "Unknown") for person in affected))
        action = next((row for row in _action_rows(section["snapshot"]) if row.get("linked") and _kpi_match(row["linked"], kpi)), None)
        action_text = action.get("execution") if action else f"Validate {_clean(kpi.get('label'), 'KPI')} operating cause and ownership."
        values = [
            f"{_clean(kpi.get('label'), 'KPI')} | {_fmt_pct(_kpi_loss(kpi))} loss",
            ", ".join(roles) or "Unassigned",
            ", ".join(names) or "No employee-level evidence",
            _fit_wrapped(action_text, 70, max_lines=3),
        ]
        for shape, value in zip(group, values):
            _write_shape(shape, value)
    for group in row_groups[len(kpis[:4]):4]:
        for shape in group:
            _write_shape(shape, MISSING)
    _write(slide, "Text 51", "Management interpretation: the score is the entry point; the intervention follows the driver, role and affected employee evidence.")


def _fill_priority_column(slide: Any, *, header: str, summary: str, people: list[dict[str, Any]], header_name: str, summary_name: str, card: str, accent: str, row_names: list[tuple[str, str, str]]) -> None:
    _card_status(
        slide,
        card,
        accent,
        {"Requires Action": "Requires action", "Watch": "Watch", "On Track": "On track"}.get(header, "Data unavailable"),
        card_fill=WHITE,
    )
    _write(slide, header_name, header)
    _write(slide, summary_name, summary)
    for index, names in enumerate(row_names):
        person = people[index] if index < len(people) else None
        if person:
            _write_single_line(slide, names[0], person.get("name"))
            _write(slide, names[1], _pct(person.get("score")))
            _write(slide, names[2], _mom(person.get("movement")))
        else:
            for name in names:
                _write(slide, name, MISSING)


def _populate_marketing_priority(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int) -> None:
    snapshot = section["snapshot"]
    people = _employee_priority(snapshot)
    action = [person for person in people if _score_status(person.get("score")) == "Requires action"]
    watch = [person for person in people if _score_status(person.get("score")) == "Watch"]
    on_track = [person for person in people if _score_status(person.get("score")) == "On track"]
    _set_common(slide, "Marketing Employee Performance Priority", "Exceptions first; on-track people stay visible but secondary", model["period"], model["branch"], page)
    _fill_priority_column(slide, header="Requires Action", summary=f"{len(action)} employee(s) | Score below 70%", people=action, header_name="Text 10", summary_name="Text 11", card="Shape 8", accent="Shape 9", row_names=[("Text 12", "Text 13", "Text 14"), ("Text 15", "Text 16", "Text 17"), ("Text 18", "Text 19", "Text 20"), ("Text 21", "Text 22", "Text 23"), ("Text 24", "Text 25", "Text 26")])
    _fill_priority_column(slide, header="Watch", summary=f"{len(watch)} employee(s) | 70\u201389%", people=watch, header_name="Text 29", summary_name="Text 30", card="Shape 27", accent="Shape 28", row_names=[("Text 31", "Text 32", "Text 33")])
    _fill_priority_column(slide, header="On Track", summary=f"{len(on_track)} employee(s) | 90%+", people=on_track, header_name="Text 36", summary_name="Text 37", card="Shape 34", accent="Shape 35", row_names=[("Text 38", "Text 39", "Text 40"), ("Text 41", "Text 42", "Text 43"), ("Text 44", "Text 45", "Text 46"), ("Text 47", "Text 48", "Text 49"), ("Text 50", "Text 51", "Text 52"), ("Text 53", "Text 54", "Text 55")])
    focus = ", ".join(_clean(person.get("name"), "Unknown") for person in action[:5]) or "No employee requires immediate action."
    _write(slide, "Text 57", _fit(f"Immediate focus: {len(action)} employee(s) below 70% \u2014 {focus}.", 170))


DETAIL_KPI_SLOTS = [
    ("Text 27", "Text 28", "Text 30", "Shape 31", "Shape 32", "Text 33", "Text 34", "Text 35"),
    ("Text 38", "Text 39", "Text 41", "Shape 42", "Shape 43", "Text 44", "Text 45", "Text 46"),
    ("Text 49", "Text 50", "Text 52", "Shape 53", "Shape 54", "Text 55", "Text 56", "Text 57"),
    ("Text 60", "Text 61", "Text 63", "Shape 64", "Shape 65", "Text 66", "Text 67", "Text 68"),
    ("Text 71", "Text 72", "Text 74", "Shape 75", "Shape 76", "Text 77", "Text 78", "Text 79"),
]


def _write_detail_kpi(slide: Any, slot: tuple[str, ...], kpi: dict[str, Any] | None) -> None:
    label_name, achievement_name, gap_name, track_name, bar_name, actual_name, target_name, direction_name = slot
    if not kpi:
        for name in (label_name, achievement_name, gap_name, actual_name, target_name, direction_name):
            if _shape_exists(slide, name):
                _write(slide, name, MISSING)
        if _shape_exists(slide, bar_name):
            _resize_bar(slide, bar_name, 0.0, color=MUTED)
        return
    loss = _kpi_loss(kpi)
    status = _kpi_status(kpi)
    if _shape_exists(slide, label_name):
        _write_single_line(slide, label_name, kpi.get("label"))
    if _shape_exists(slide, achievement_name):
        _write(slide, achievement_name, _fmt_pct(kpi.get("achievement_pct")))
    if _shape_exists(slide, gap_name):
        _write(slide, gap_name, f"Gap {_fmt_pct(_kpi_gap(kpi))} | Loss {_fmt_pct(loss)}")
    if _shape_exists(slide, track_name):
        _fill(slide, track_name, TRACK)
    if _shape_exists(slide, bar_name):
        _resize_bar(slide, bar_name, (_number(kpi.get("achievement_pct")) or 0) / 100, color=_status_color(status))
    if _shape_exists(slide, actual_name):
        _write(slide, actual_name, f"Actual: {_native(kpi.get('actual'), kpi.get('unit'))}")
    if _shape_exists(slide, target_name):
        _write(slide, target_name, f"Target: {_native(kpi.get('target'), kpi.get('unit'))}")
    if _shape_exists(slide, direction_name):
        _write(slide, direction_name, _direction_label(kpi.get("direction")))


def _populate_employee_detail(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int, person: dict[str, Any] | None, rank: int, source_slot: int) -> None:
    name = _clean((person or {}).get("name"), "Employee detail unavailable")
    position = _clean((person or {}).get("position"), "Position unavailable")
    title = f"Marketing Employee Detail \u2014 {name}"
    _set_common(slide, title, "Actual vs Target is shown for every KPI so the employee number is clear and auditable", model["period"], model["branch"], page)
    status = _score_status((person or {}).get("score"))
    _card_status(slide, "Shape 8", "Shape 9", status, card_fill=WHITE)
    if _shape_exists(slide, "Shape 12"):
        _fill(slide, "Shape 12", _status_color(status))
    _write_single_line(slide, "Text 10", name)
    _write_single_line(slide, "Text 11", position)
    _write(slide, "Text 13", _status_label(status), color=WHITE)
    _write(slide, "Text 15", _pct((person or {}).get("score")))
    _write(slide, "Text 17", _pct((person or {}).get("baseline_score")))
    _write(slide, "Text 19", _mom((person or {}).get("movement")))
    _write(slide, "Text 21", f"Action focus | Priority {rank}")
    _write(slide, "Text 22", _fit(_action_focus(section["snapshot"]), 62))
    _write(slide, "Text 23", "KPI Drivers \u2014 Actual / Target / Achievement / Weighted Loss")
    person_kpis = _person_kpis(person or {})[:5]
    detail_cards = [("Shape 25", "Shape 26"), ("Shape 36", "Shape 37"), ("Shape 47", "Shape 48"), ("Shape 58", "Shape 59"), ("Shape 69", "Shape 70")]
    for index, slot in enumerate(DETAIL_KPI_SLOTS):
        kpi = person_kpis[index] if index < len(person_kpis) else None
        if index < len(detail_cards):
            card, rail = detail_cards[index]
            if _shape_exists(slide, card):
                _fill(slide, card, WHITE if kpi else PALE_NEUTRAL)
            if _shape_exists(slide, rail):
                _fill(slide, rail, _status_color(_kpi_status(kpi)) if kpi else MUTED)
        _write_detail_kpi(slide, slot, kpi)
    review_name = "Text 81" if _shape_exists(slide, "Text 81") else "Text 70"
    if _shape_exists(slide, review_name):
        _write(slide, review_name, "Review rule: prioritize the highest weighted-loss KPI first; use Actual / Target to validate the gap before assigning corrective action.")


def _shape_exists(slide: Any, name: str) -> bool:
    return any(shape.name == name for shape in slide.shapes)


def _populate_employee_group(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int, people: list[dict[str, Any]]) -> None:
    title_people = " / ".join(_clean(person.get("name"), "Employee") for person in people[:2]) or "Employees requiring action"
    _set_common(slide, f"Marketing Employee Detail \u2014 {title_people}", "Employees sharing a common leading driver", model["period"], model["branch"], page)
    panels = [
        (("Shape 9", "Shape 10", "Shape 11", "Text 12", "Text 13", "Text 14", "Text 16", "Text 18", "Text 20"), ["Text 23", "Text 34", "Text 45"], ["Text 29", "Text 40", "Text 51"], ["Text 30", "Text 41", "Text 52"], ["Text 31", "Text 42", "Text 53"]),
        (("Shape 55", "Shape 56", "Shape 57", "Text 58", "Text 59", "Text 60", "Text 62", "Text 64", "Text 66"), ["Text 69", "Text 80", "Text 91"], ["Text 75", "Text 86", "Text 97"], ["Text 76", "Text 87", "Text 98"], ["Text 77", "Text 88", "Text 99"]),
    ]
    for index, (panel, labels, actuals, targets, directions) in enumerate(panels):
        person = people[index] if index < len(people) else None
        card_bg, rail, pill_bg, pill_text, name_name, role_name, score_name, mom_name, action_name = panel
        if person:
            status = _score_status(person.get("score"))
            _fill(slide, card_bg, WHITE)
            _fill(slide, rail, _status_color(status))
            _fill(slide, pill_bg, _status_color(status))
            _write(slide, pill_text, _status_label(status), color=WHITE)
            _write_single_line(slide, name_name, person.get("name"))
            _write_single_line(slide, role_name, person.get("position"))
            _write(slide, score_name, _pct(person.get("score")))
            _write(slide, mom_name, _mom(person.get("movement")))
            _write(slide, action_name, _fit_wrapped(f"Action focus: {_action_focus(section['snapshot'], limit=1)} gap review.", 48, max_lines=2))
            kpis = _person_kpis(person)[:3]
            for idx in range(3):
                kpi = kpis[idx] if idx < len(kpis) else None
                _write_single_line(slide, labels[idx], (kpi or {}).get("label"))
                _write(slide, actuals[idx], _fmt_pct((kpi or {}).get("achievement_pct")))
                _write(slide, targets[idx], f"Gap {_fmt_pct(_kpi_gap(kpi or {}))} | Loss {_fmt_pct(_kpi_loss(kpi or {}))}")
                _write(slide, directions[idx], _direction_label((kpi or {}).get("direction")))
        else:
            for name in [name_name, role_name, score_name, mom_name, action_name, *labels, *actuals, *targets, *directions]:
                if _shape_exists(slide, name):
                    _write(slide, name, MISSING)


def _evidence_rows(snapshot: dict[str, Any]) -> list[list[str]]:
    rows: list[tuple[float, list[str]]] = []
    for person in snapshot.get("all_people") or []:
        for kpi in _person_kpis(person):
            loss = _kpi_loss(kpi)
            if loss <= 0:
                continue
            state, note = _root_evidence(snapshot, person, kpi)
            evidence = f"{state} \u2014 {note}" if state != "KPI Evidence" else "KPI Evidence \u2014 Investigation Required"
            rows.append((loss, [
                _full(person.get("name")),
                _full(person.get("position")),
                _full(kpi.get("label")),
                _native(kpi.get("actual"), kpi.get("unit")),
                _native(kpi.get("target"), kpi.get("unit")),
                _fmt_pct(kpi.get("achievement_pct")),
                _fmt_pct(loss),
                _fit_wrapped(evidence, 76, max_lines=3),
            ]))
    rows.sort(key=lambda item: (-item[0], item[1][0].casefold()))
    return [row for _, row in rows]


def _populate_evidence(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int, rows: list[list[str]], part: int, total: int) -> None:
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"Marketing Root Cause Evidence{suffix}", "Evidence first: result, target, achievement, weighted loss and validation need", model["period"], model["branch"], page)
    groups = _row_groups(slide, top_min=1.57, top_max=5.25, min_columns=8)
    capacity = len(groups)
    displayed = rows[:capacity]
    for group, values in zip(groups, displayed):
        for shape, value in zip(group, values):
            _write_shape(shape, value)
    for group in groups[len(displayed):]:
        for shape in group:
            _write_shape(shape, MISSING)
    _write(slide, "Text 169", "Important: this slide separates evidence from assumption. A KPI identifies where to investigate; it does not invent an operational root cause.")


def _populate_marketing_actions(slide: Any, model: dict[str, Any], section: dict[str, Any], page: int, actions: list[dict[str, Any]], part: int, total: int) -> None:
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"Marketing Corrective Action Tracker{suffix}", "Data-driven workstreams created from the highest weighted-loss gaps", model["period"], model["branch"], page)
    groups = _row_groups(slide, top_min=1.78, top_max=5.0, min_columns=6)
    for group, action in zip(groups[:4], actions[:4]):
        owner = (
            f"{_fit(action.get('owner'), 16, 'Owner needed')}\n"
            f"Due: {_fit(action.get('due'), 16, 'Due date needed')}\n"
            f"Status: {_fit(action.get('status'), 16, 'Status not recorded')}"
        )
        values = [
            _clean(action.get("priority"), "Medium"),
            _fit_wrapped(action.get("workstream"), 26, max_lines=2),
            owner,
            _fit_wrapped(action.get("scope"), 34, max_lines=2),
            _fit_wrapped(action.get("execution"), 68, max_lines=3),
            _fit_wrapped(action.get("success"), 46, max_lines=2),
        ]
        for shape, value in zip(group, values):
            _write_shape(shape, value)
    for group in groups[len(actions[:4]):4]:
        for shape in group:
            _write_shape(shape, MISSING)
    _write(slide, "Text 69", "Decision needed before next review: confirm owner, due date, status and success evidence for every proposed workstream.")


def _populate_management_divider(slide: Any, model: dict[str, Any], page: int) -> None:
    _set_divider(slide, "Management Summary", "Close with cross-department priorities, next-review commitments and employee appendix.", model["period"], model["branch"], page, ["Next-review scorecard", "Owner / due date decisions", "Combined employee appendix"])


def _populate_management_summary(slide: Any, model: dict[str, Any], page: int) -> None:
    _set_common(slide, "Combined Management Summary / Next Review", "What should be reviewed before the next performance cycle", model["period"], model["branch"], page)
    views = _priority_views(model["sections"])
    kpis = [kpi for kpi in model["snapshot"].get("kpis") or [] if _kpi_loss(kpi) > 0]
    items = []
    for view in views[:3]:
        leading = (view.get("kpis") or [None])[0]
        if leading:
            items.append((_clean(leading.get("label"), "KPI").upper(), f"{_native(leading.get('actual'), leading.get('unit'))} → {_native(leading.get('target'), leading.get('unit'))}", f"{_fmt_pct(_kpi_loss(leading))} weighted loss"))
        else:
            items.append((_clean(view.get("name"), "Department").upper(), _pct(view.get("score")), f"{view.get('affected_count', 0)} employee(s) require review"))
    items.extend([
        ("EMPLOYEES REQUIRING ACTION", str(len([person for person in model["snapshot"].get("all_people") or [] if _score_status(person.get("score")) == "Requires action"])), "Priority people stay visible in the appendix"),
        ("OPEN ACTIONS", str(sum(len(_action_rows(section["snapshot"])) for section in model["sections"])), "Confirm owner, due date and success measure"),
        ("NEXT REVIEW", "Evidence → actual → target", "Recheck the same scope next cycle"),
    ])
    card_texts = [("Text 10", "Text 11", "Text 12"), ("Text 15", "Text 16", "Text 17"), ("Text 20", "Text 21", "Text 22"), ("Text 25", "Text 26", "Text 27"), ("Text 30", "Text 31", "Text 32"), ("Text 35", "Text 36", "Text 37")]
    card_shapes = [("Shape 8", "Shape 9"), ("Shape 13", "Shape 14"), ("Shape 18", "Shape 19"), ("Shape 23", "Shape 24"), ("Shape 28", "Shape 29"), ("Shape 33", "Shape 34")]
    for index, names in enumerate(card_texts):
        item = items[index] if index < len(items) else ("NOT AVAILABLE", MISSING, MISSING)
        status = "Requires action" if index < 5 else "Watch"
        _card_status(slide, *card_shapes[index], status)
        _write(slide, names[0], _fit(item[0], 28))
        _write(slide, names[1], _fit(item[1], 28))
        _write(slide, names[2], _fit(item[2], 44))
    _write(slide, "Text 38", "Final management decision")
    owners_missing = sum(1 for section in model["sections"] for action in _action_rows(section["snapshot"]) if _clean(action.get("owner"), "Owner needed") == "Owner needed")
    score = _pct(model["snapshot"].get("overall_score"))
    target = _pct(model["snapshot"].get("target_score") or 100)
    decision = f"Prioritize the highest weighted-loss gaps, confirm the {owners_missing} missing owner(s), then recheck score {score} against target {target} with actual KPI evidence at the next review."
    _write(slide, "Text 41", _fit(decision, 175))


def _appendix_rows(snapshot: dict[str, Any]) -> list[list[str]]:
    rows = []
    people = _employee_priority(snapshot)
    for person in people:
        primary = _person_kpis(person)[0] if _person_kpis(person) else {}
        rows.append([
            _clean(person.get("team"), "Unassigned"),
            _full(person.get("name")),
            _full(f"{_clean(person.get('position'), 'Unassigned')} / {_clean(person.get('grade'), 'N/A')}"),
            _pct(person.get("score")),
            _mom(person.get("movement")),
            _fit_wrapped(f"{_clean(primary.get('label'), 'No KPI')} | loss {_fmt_pct(_kpi_loss(primary))}", 46, max_lines=2),
            _full(person.get("action_status")),
        ])
    return rows


def _populate_appendix(slide: Any, model: dict[str, Any], page: int, rows: list[list[str]], part: int, total: int) -> None:
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"Combined Employee Appendix{suffix}", "One-page reference for employees requiring action or management follow-up", model["period"], model["branch"], page)
    groups = _row_groups(slide, top_min=1.46, top_max=5.9, min_columns=7)
    capacity = len(groups)
    displayed = rows[:capacity]
    for group, values in zip(groups, displayed):
        for shape, value in zip(group, values):
            _write_shape(shape, value)
    for group in groups[len(displayed):]:
        for shape in group:
            _write_shape(shape, MISSING)


def _marketing_detail_plan(section: dict[str, Any]) -> list[tuple[str, int, Any]]:
    people = [person for person in _employee_priority(section["snapshot"]) if _score_status(person.get("score")) in {"Requires action", "Watch"}]
    if not people:
        return []
    plan: list[tuple[str, int, Any]] = []
    consumed: set[str] = set()
    for index, person in enumerate(people):
        key = _clean(person.get("employee_id") or person.get("name"), str(index))
        if key in consumed:
            continue
        same_position = [candidate for candidate in people if _clean(candidate.get("position"), "").casefold() == _clean(person.get("position"), "").casefold() and _clean(candidate.get("employee_id") or candidate.get("name"), "") not in consumed]
        if len(same_position) >= 2:
            selected = same_position[:2]
            plan.append(("detail_group", 21, selected))
            consumed.update(_clean(candidate.get("employee_id") or candidate.get("name"), "") for candidate in selected)
            continue
        source_index = 19 if len(plan) % 3 == 0 else 20 if len(plan) % 3 == 1 else 22
        plan.append(("detail", source_index, person))
        consumed.add(key)
    return plan


def _build_plan(model: dict[str, Any]) -> list[tuple[str, int, Any]]:
    plan: list[tuple[str, int, Any]] = [("cover", 0, None), ("overview", 1, None), ("priority", 2, None)]
    for section in model["sections"]:
        if section.get("key") == "marketing":
            continue
        plan.append(("divider", section["divider"], section))
        people = _employee_priority(section["snapshot"])
        people_chunks = [people[index:index + 5] for index in range(0, len(people), 5)] or [[]]
        for part, chunk in enumerate(people_chunks, 1):
            plan.append(("department_performance", section["performance"], (section, chunk, part, len(people_chunks))))
        # The department corrective-action frames are employee evidence
        # tables in the approved source deck.  Use current people here; each
        # person carries their own KPI evidence and recorded/proposed action.
        action_chunks = [people[index:index + 5] for index in range(0, len(people), 5)] or [[]]
        for part, chunk in enumerate(action_chunks, 1):
            plan.append(("department_actions", section["actions"], (section, chunk, part, len(action_chunks))))
        if section.get("followup") is not None:
            plan.append(("rcm_followup", section["followup"], section))
    marketing = next((section for section in model["sections"] if section.get("key") == "marketing"), None)
    if marketing:
        plan.append(("divider", marketing["divider"], marketing))
        plan.append(("marketing_overview", marketing["performance"], marketing))
        roles = _role_rows(marketing["snapshot"])
        role_chunks = [roles[index:index + 9] for index in range(0, len(roles), 9)] or [[]]
        for part, chunk in enumerate(role_chunks, 1):
            plan.append(("marketing_roles", 15, (marketing, chunk, part, len(role_chunks))))
        kpis = [dict(row) for row in marketing["snapshot"].get("kpis") or []]
        kpi_chunks = [kpis[index:index + 8] for index in range(0, len(kpis), 8)] or [[]]
        for part, chunk in enumerate(kpi_chunks, 1):
            plan.append(("marketing_trend", 16, (marketing, chunk, part, len(kpi_chunks))))
        driver_chunks = [kpis[index:index + 4] for index in range(0, len(kpis), 4)] or [[]]
        for part, chunk in enumerate(driver_chunks, 1):
            plan.append(("marketing_map", 17, (marketing, chunk, part, len(driver_chunks))))
        plan.append(("marketing_priority", 18, marketing))
        plan.extend(_marketing_detail_plan(marketing))
        evidence = _evidence_rows(marketing["snapshot"])
        evidence_chunks = [evidence[index:index + 9] for index in range(0, len(evidence), 9)] or [[]]
        for part, chunk in enumerate(evidence_chunks, 1):
            plan.append(("evidence", 23, (marketing, chunk, part, len(evidence_chunks))))
        actions = _action_rows(marketing["snapshot"])
        action_chunks = [actions[index:index + 4] for index in range(0, len(actions), 4)] or [[]]
        for part, chunk in enumerate(action_chunks, 1):
            plan.append(("marketing_actions", 24, (marketing, chunk, part, len(action_chunks))))
    plan.append(("management_divider", 25, None))
    plan.append(("management_summary", 26, None))
    appendix = _appendix_rows(model["snapshot"])
    if appendix:
        appendix_chunks = [appendix[index:index + 12] for index in range(0, len(appendix), 12)]
        for part, chunk in enumerate(appendix_chunks, 1):
            plan.append(("appendix", 27, (chunk, part, len(appendix_chunks))))
    return plan


def _populate_slide(slide: Any, kind: str, context: Any, model: dict[str, Any], page: int) -> None:
    if kind == "cover":
        _populate_cover(slide, model, page)
    elif kind == "overview":
        _populate_overview(slide, model, page)
    elif kind == "priority":
        _populate_priority(slide, model, page)
    elif kind == "divider":
        _populate_section_divider(slide, model, context, page)
    elif kind == "department_performance":
        section, rows, part, total = context
        _populate_department_performance(slide, model, section, page, rows, part, total)
    elif kind == "department_actions":
        section, rows, part, total = context
        _populate_department_actions(slide, model, section, page, rows, part, total)
    elif kind == "rcm_followup":
        _populate_rcm_followup(slide, model, context, page)
    elif kind == "marketing_overview":
        _populate_marketing_overview(slide, model, context, page)
    elif kind == "marketing_roles":
        section, rows, part, total = context
        _populate_marketing_roles(slide, model, section, page, rows, part, total)
    elif kind == "marketing_trend":
        section, rows, part, total = context
        _populate_marketing_trend(slide, model, section, page, rows, part, total)
    elif kind == "marketing_map":
        section, rows, part, total = context
        _populate_marketing_map(slide, model, section, page, rows, part, total)
    elif kind == "marketing_priority":
        _populate_marketing_priority(slide, model, context, page)
    elif kind == "detail":
        person = context
        marketing = next(section for section in model["sections"] if section.get("key") == "marketing")
        rank = _employee_priority(marketing["snapshot"]).index(person) + 1 if person in _employee_priority(marketing["snapshot"]) else 1
        _populate_employee_detail(slide, model, marketing, page, person, rank, 0)
    elif kind == "detail_group":
        marketing = next(section for section in model["sections"] if section.get("key") == "marketing")
        _populate_employee_group(slide, model, marketing, page, context)
    elif kind == "evidence":
        section, rows, part, total = context
        _populate_evidence(slide, model, section, page, rows, part, total)
    elif kind == "marketing_actions":
        section, actions, part, total = context
        _populate_marketing_actions(slide, model, section, page, actions, part, total)
    elif kind == "management_divider":
        _populate_management_divider(slide, model, page)
    elif kind == "management_summary":
        _populate_management_summary(slide, model, page)
    elif kind == "appendix":
        rows, part, total = context
        _populate_appendix(slide, model, page, rows, part, total)


def _prepare_model(period_label: str, report_data: dict[str, Any] | None) -> dict[str, Any]:
    raw = _copy_data(report_data or {})
    raw["period_label"] = raw.get("period_label") or period_label
    snapshot = base._prepare_snapshot(period_label, raw)
    sections = _ordered_sections(raw, period_label)
    branch = _branch_label(raw)
    return {
        "raw": raw,
        "snapshot": snapshot,
        "sections": sections,
        "period": _period(snapshot, period_label),
        "branch": branch,
    }


def build_offshore_status_pptx(period_label: str = "Selected period", report_data: dict[str, Any] | None = None) -> bytes:
    """Generate a live UAE / Offshore EGY report using the corrected deck."""

    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Corrected Offshore Departments template is missing: {TEMPLATE_PATH}")
    template_bytes = TEMPLATE_PATH.read_bytes()
    source = Presentation(io.BytesIO(template_bytes))
    if len(source.slides) != REFERENCE_SLIDES:
        raise ValueError(f"Corrected Offshore Departments template must contain {REFERENCE_SLIDES} slides; found {len(source.slides)}")
    prs = Presentation(io.BytesIO(template_bytes))
    _clear_slides(prs)
    model = _prepare_model(period_label, report_data)
    plan = _build_plan(model)
    for page, (kind, source_index, context) in enumerate(plan, 1):
        slide = _append_template_slide(prs, source, source_index)
        _populate_slide(slide, kind, context, model, page)
    prs.core_properties.title = f"Offshore Departments Combined Performance Review - {model['period']}"
    prs.core_properties.subject = "Dynamic UAE / Offshore EGY departmental performance review"
    prs.core_properties.author = "PMS Dashboard"
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


__all__ = ["build_offshore_status_pptx", "TEMPLATE_PATH"]
