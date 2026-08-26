"""Source-bound UAE Executive Performance Summary PowerPoint builder.

The builder owns no scoring logic.  It consumes the canonical contract from
``services.uae_executive_summary`` and maps that contract onto the approved
Offshore source frames.  The detailed UAE/EGY builder remains untouched.
"""

from __future__ import annotations

import io
from copy import deepcopy
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation

from exports import offshore_status_pptx_builder as reference


TEMPLATE_PATH = reference.TEMPLATE_PATH
REFERENCE_SLIDES = reference.REFERENCE_SLIDES
MISSING = reference.MISSING


def _number(value: Any) -> float | None:
    try:
        result = float(value)
    except (TypeError, ValueError):
        return None
    return result if result == result else None


def _clean(value: Any, fallback: str = MISSING) -> str:
    if value is None:
        return fallback
    text = " ".join(str(value).replace("\n", " ").split()).strip()
    return text or fallback


def _pct(value: Any, *, signed: bool = False) -> str:
    number = _number(value)
    if number is None:
        return MISSING
    return f"{number:+.1f}%" if signed else f"{number:.1f}%"


def _mom(value: Any, *, current_exists: bool = True) -> str:
    if _number(value) is None:
        return "NEW" if current_exists else "No current data"
    return _pct(value, signed=True)


def _fit(value: Any, limit: int) -> str:
    return reference._fit(value, limit)


def _fit_wrapped(value: Any, limit: int, max_lines: int = 2) -> str:
    return reference._fit_wrapped(value, limit, max_lines=max_lines)


def _native(value: Any, unit: Any = "") -> str:
    return reference._native(value, unit)


def _status(value: Any) -> str:
    normalized = _clean(value, "Data Not Available").casefold()
    if normalized in {"requires action", "critical", "urgent"}:
        return "Requires action"
    if normalized in {"watch", "at risk", "below target"}:
        return "Watch"
    if normalized in {"on track", "meets expectations", "within target", "completed"}:
        return "On track"
    return "Data unavailable"


def _status_color(value: Any):
    return reference._status_color(_status(value))


def _status_label(value: Any) -> str:
    return reference._status_label(_status(value))


def _shape(slide: Any, name: str) -> Any:
    return reference._shape(slide, name)


def _write(slide: Any, name: str, value: Any, **kwargs: Any) -> None:
    reference._write(slide, name, value, **kwargs)


def _write_single_line(slide: Any, name: str, value: Any, **kwargs: Any) -> None:
    reference._write_single_line(slide, name, value, **kwargs)


def _write_shape(shape: Any, value: Any, **kwargs: Any) -> None:
    reference._write_shape(shape, value, **kwargs)


def _blank_shape(shape: Any) -> None:
    """Clear an inherited table cell without writing a filler em dash."""

    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True


def _blank_rows(groups: list[list[Any]], start: int) -> None:
    for group in groups[start:]:
        for shape in group:
            _blank_shape(shape)


def _set_common(slide: Any, title: str, subtitle: str, period: str, scope: str, page: int) -> None:
    reference._set_common(slide, title, subtitle, period, scope, page)


def _row_groups(slide: Any, *, top_min: float, top_max: float, min_columns: int) -> list[list[Any]]:
    return reference._row_groups(slide, top_min=top_min, top_max=top_max, min_columns=min_columns)


def _card_status(slide: Any, card: str, accent: str, status: Any) -> None:
    reference._card_status(slide, card, accent, _status(status))


def _resize_bar(slide: Any, name: str, ratio: Any, *, color: Any = None) -> None:
    reference._resize_bar(slide, name, ratio, color=color)


def _clear_slides(prs: Presentation) -> None:
    reference._clear_slides(prs)


def _append_template_slide(prs: Presentation, source: Presentation, source_index: int) -> Any:
    return reference._append_template_slide(prs, source, source_index)


def _copy_payload(payload: dict[str, Any]) -> dict[str, Any]:
    result = dict(payload)
    for key in (
        "overall",
        "cover",
        "largest_kpi_driver",
        "largest_department_decline",
        "best_department",
        "next_review",
    ):
        result[key] = dict(payload.get(key) or {})
    for key in (
        "departments",
        "department_ranking",
        "kpis",
        "drivers",
        "trend",
        "history",
        "actions",
        "recorded_actions",
        "action_rows",
        "decisions",
        "data_coverage",
        "employee_priority_rows",
        "kpi_evidence_rows",
        "driver_department_impact",
        "historical_rows",
    ):
        result[key] = [dict(item) if isinstance(item, dict) else item for item in (payload.get(key) or [])]
    return result


def _populate_cover(slide: Any, payload: dict[str, Any], page: int) -> None:
    period = _clean(payload.get("period_label"), "Selected period")
    cover = payload.get("cover") or {}
    _write(slide, "Text 3", "UAE Offshore Departments")
    _write(slide, "Text 4", "Executive Performance Summary")
    _write(slide, "Text 5", _fit(f"{period}  |  UAE", 70))
    _write(slide, "Text 6", _fit("PMS / Performance Review  •  CEO / Group / Regional Leadership", 90))
    _write(slide, "Text 8", "Executive narrative")
    _write(slide, "Text 9", "UAE performance first, then department risk, KPI drivers, execution ownership, decisions, and the next-review agenda.")
    _write(slide, "Text 11", f"{cover.get('active_departments', 0)} Active Departments")
    _write(slide, "Text 13", f"{cover.get('employee_count', 0)} Employees in scope")
    _write(slide, "Text 15", f"{cover.get('open_actions', 0)} Open Actions")
    _write(slide, "Text 16", "CONFIDENTIAL – INTERNAL USE ONLY")


def _card_names(index: int) -> tuple[tuple[str, ...], str, str]:
    specs = [
        (("Text 10", "Text 11", "Text 12", "Text 13", "Text 14", "Text 15", "Text 17"), "Shape 8", "Shape 9"),
        (("Text 20", "Text 21", "Text 22", "Text 23", "Text 24", "Text 25", "Text 27"), "Shape 18", "Shape 19"),
        (("Text 30", "Text 31", "Text 32", "Text 33", "Text 34", "Text 35", "Text 37"), "Shape 28", "Shape 29"),
        (("Text 40", "Text 41", "Text 42", "Text 43", "Text 44", "Text 45", "Text 47"), "Shape 38", "Shape 39"),
    ]
    return specs[index]


def _write_score_card(slide: Any, names: tuple[str, ...], row: dict[str, Any], title: str, status: Any) -> None:
    label, score_name, monthly_name, baseline_name, mom_name, target_name, status_name = names
    _write_single_line(slide, label, title)
    _write(slide, score_name, _pct(row.get("score")))
    _write(slide, monthly_name, "Current performance")
    _write(slide, baseline_name, f"Baseline: {_pct(row.get('baseline'))}")
    _write(slide, mom_name, f"MoM: {_mom(row.get('mom'), current_exists=row.get('score') is not None)}")
    _write(slide, target_name, f"Target: {_pct(row.get('target'))}")
    _write(slide, status_name, _status_label(status), color=reference.WHITE)


def _populate_overview(slide: Any, payload: dict[str, Any], page: int) -> None:
    period = _clean(payload.get("period_label"), "Selected period")
    scope = "UAE"
    overall = dict(payload.get("overall") or {})
    current_exists = overall.get("score") is not None
    _set_common(slide, "UAE Executive Overview", "How did UAE perform this month?", period, scope, page)
    driver = dict(payload.get("largest_kpi_driver") or {})
    decline = dict(payload.get("largest_department_decline") or {})
    best = dict(payload.get("best_department") or {})
    cards = [
        ({**overall, "name": "Current UAE Performance"}, "Current UAE Performance"),
        ({
            "name": "Baseline / Target",
            "score": overall.get("baseline"),
            "baseline": overall.get("baseline"),
            "mom": overall.get("mom"),
            "target": overall.get("target"),
            "status": overall.get("status"),
        }, "Comparable Baseline"),
        ({
            "name": "Execution risk",
            "score": payload.get("open_action_count", 0),
            "baseline": payload.get("missing_owner_count", 0),
            "mom": payload.get("missing_due_date_count", 0),
            "target": payload.get("employees_requiring_action", 0),
            "status": "Requires Action" if payload.get("missing_owner_count", 0) or payload.get("missing_due_date_count", 0) else "On Track",
        }, "Open Actions"),
        ({
            "name": driver.get("label") or "Largest KPI Driver",
            "score": driver.get("achievement_pct"),
            "baseline": driver.get("baseline_achievement_pct"),
            "mom": driver.get("mom"),
            "target": 100.0,
            "status": driver.get("status"),
        }, "Largest KPI Driver"),
    ]
    for index, (row, title) in enumerate(cards):
        names, card, accent = _card_names(index)
        _card_status(slide, card, accent, row.get("status"))
        _write_score_card(slide, names, row, title if index != 3 else _fit(driver.get("label") or title, 25), row.get("status"))
    takeaway = (
        f"Executive takeaway: UAE is {_pct(overall.get('score'))} against a {_pct(overall.get('target'))} target; "
        f"MoM {_mom(overall.get('mom'), current_exists=current_exists)}. "
        f"{len(payload.get('risk_map', {}).get('Requires Action', []))} department(s) require action, "
        f"{payload.get('employees_requiring_action', 0)} employee(s) are below threshold, and "
        f"{payload.get('open_action_count', 0)} open action(s) remain."
        if current_exists
        else "Executive takeaway: no current UAE score is available for the selected period; see Data Coverage for known departments."
    )
    _write(slide, "Text 49", _fit(takeaway, 170))
    _write(slide, "Text 50", "Executive metric view")
    headers = ["Metric", "Value", "Signal"]
    groups = _row_groups(slide, top_min=4.92, top_max=6.9, min_columns=6)
    # The inherited six-column rows are repurposed as a compact, readable
    # management strip.  The same values are repeated nowhere else as a
    # separate population; they are aliases of the contract.
    rows = [
        ["Active Departments", str(payload.get("active_departments", 0)), "valid current data"],
        ["Departments Requiring Action", str(len(payload.get("risk_map", {}).get("Requires Action", []))), "score below 70%"],
        ["Employees Requiring Action", str(payload.get("employees_requiring_action", 0)), "current employee score"],
        ["Missing Owners / Due Dates", f"{payload.get('missing_owner_count', 0)} / {payload.get('missing_due_date_count', 0)}", "resolve before next review"],
    ]
    for group, row in zip(groups, rows):
        values = [row[0], row[1], row[2], _fit(driver.get("label") or "No measured KPI gap", 32), _pct(decline.get("mom")), _fit(best.get("department") or "No current department", 30)]
        for shape, value in zip(group, values):
            _write_shape(shape, value)
    _blank_rows(groups, len(rows))


def _ranking_pages(payload: dict[str, Any], source: Presentation) -> list[tuple[int, dict[str, Any]]]:
    rows = list(payload.get("department_ranking") or [])
    capacity = max(1, len(_row_groups(source.slides[15], top_min=1.7, top_max=5.0, min_columns=3)))
    return [(15, {"_ranking_rows": rows[index:index + capacity], "_ranking_part": (index // capacity) + 1, "_ranking_total": (len(rows) + capacity - 1) // capacity}) for index in range(0, len(rows), capacity)] or [(15, {"_ranking_rows": [], "_ranking_part": 1, "_ranking_total": 1})]


def _populate_ranking(slide: Any, payload: dict[str, Any], page: int) -> None:
    part = payload.get("_ranking_part", 1)
    total = payload.get("_ranking_total", 1)
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"UAE Department Performance Ranking{suffix}", "All active departments with valid current-period data; priority score is transparent and deterministic", _clean(payload.get("period_label")), "UAE", page)
    _write(slide, "Text 8", "Department ranking — highest management priority first")
    role_names = ["Text 10", "Text 15", "Text 20", "Text 25", "Text 30", "Text 35", "Text 40", "Text 45", "Text 50"]
    score_names = ["Text 13", "Text 18", "Text 23", "Text 28", "Text 33", "Text 38", "Text 43", "Text 48", "Text 53"]
    mom_names = ["Text 14", "Text 19", "Text 24", "Text 29", "Text 34", "Text 39", "Text 44", "Text 49", "Text 54"]
    track_names = ["Shape 11", "Shape 16", "Shape 21", "Shape 26", "Shape 31", "Shape 36", "Shape 41", "Shape 46", "Shape 51"]
    bar_names = ["Shape 12", "Shape 17", "Shape 22", "Shape 27", "Shape 32", "Shape 37", "Shape 42", "Shape 47", "Shape 52"]
    rows = list(payload.get("_ranking_rows") or [])
    for index, name in enumerate(role_names):
        row = rows[index] if index < len(rows) else None
        if not row:
            for shape_name in (name, score_names[index], mom_names[index]):
                _blank_shape(_shape(slide, shape_name))
            _resize_bar(slide, bar_names[index], 0, color=reference.MUTED)
            continue
        _write_single_line(slide, name, f"{row.get('rank')}. {_clean(row.get('department'))}")
        _write(slide, score_names[index], _pct(row.get("current_score")))
        _write(slide, mom_names[index], _mom(row.get("mom"), current_exists=row.get("current_score") is not None))
        reference._fill(slide, track_names[index], reference.TRACK)
        _resize_bar(slide, bar_names[index], (_number(row.get("current_score")) or 0) / 100, color=_status_color(row.get("status")))
    attention_names = [("Text 58", "Text 59", "Text 60"), ("Text 62", "Text 63", "Text 64"), ("Text 66", "Text 67", "Text 68")]
    for index, names in enumerate(attention_names):
        row = rows[index] if index < len(rows) else None
        if row:
            _write_single_line(slide, names[0], f"{index + 1}. {_clean(row.get('department'))}")
            _write(slide, names[1], _pct(row.get("priority_score")))
            _write(slide, names[2], _fit_wrapped(
                f"{_status_label(row.get('status'))} | loss {_pct(row.get('weighted_loss'))} | {row.get('employees_requiring_action', 0)} employees | {row.get('open_actions', 0)} actions",
                42,
            ))
        else:
            for name in names:
                _blank_shape(_shape(slide, name))
    signal = payload.get("largest_department_decline") or {}
    _write(slide, "Text 70", "Priority signal")
    _write_single_line(slide, "Text 71", _fit(
        f"{_clean(signal.get('department'), 'No comparable decline')} has the largest negative movement ({_mom(signal.get('mom'))}); priority score includes target gap, KPI loss, affected employees, and unresolved actions.",
        170,
    ))


def _risk_pages(payload: dict[str, Any]) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for bucket in ("Requires Action", "Watch", "On Track", "Data Not Available"):
        for row in payload.get("risk_map", {}).get(bucket, []) or []:
            rows.append({**row, "risk_bucket": bucket})
    return [
        {"_risk_rows": rows[index:index + 4], "_risk_part": (index // 4) + 1, "_risk_total": (len(rows) + 3) // 4}
        for index in range(0, len(rows), 4)
    ] or [{"_risk_rows": [], "_risk_part": 1, "_risk_total": 1}]


def _populate_risk_map(slide: Any, payload: dict[str, Any], page: int) -> None:
    part = payload.get("_risk_part", 1)
    total = payload.get("_risk_total", 1)
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"UAE Risk and Opportunity Map{suffix}", "Requires Action, Watch, On Track, and Data Not Available departments", _clean(payload.get("period_label")), "UAE", page)
    slots = [
        (("Text 10", "Text 11", "Text 12", "Text 13", "Text 14", "Text 15"), "Shape 8", "Shape 9"),
        (("Text 18", "Text 19", "Text 20", "Text 21", "Text 22", "Text 23"), "Shape 16", "Shape 17"),
        (("Text 26", "Text 27", "Text 28", "Text 29", "Text 30", "Text 31"), "Shape 24", "Shape 25"),
        (("Text 34", "Text 35", "Text 36", "Text 37", "Text 38", "Text 39"), "Shape 32", "Shape 33"),
    ]
    rows = list(payload.get("_risk_rows") or [])
    for index, (names, card, accent) in enumerate(slots):
        row = rows[index] if index < len(rows) else None
        _card_status(slide, card, accent, row.get("status") if row else "Data Not Available")
        if not row:
            for name in names:
                _blank_shape(_shape(slide, name))
            continue
        bucket = row.get("risk_bucket")
        _write_single_line(slide, names[0], f"{bucket}: {_clean(row.get('department'))}")
        _write(slide, names[1], _pct(row.get("current_score")))
        _write(slide, names[2], _fit_wrapped(f"MoM {_mom(row.get('mom'), current_exists=row.get('score') is not None)} | {_status_label(row.get('status'))}", 46))
        driver = _clean(row.get("leading_kpi_driver"), "No measured KPI driver")
        values = [
            f"• Driver: {driver}",
            f"• Weighted loss: {_pct(row.get('weighted_loss'))}",
            f"• {row.get('employees_requiring_action', 0)} employees / {row.get('open_actions', 0)} open actions",
        ]
        for name, value in zip(names[3:], values):
            _write(slide, name, _fit_wrapped(value, 46))
    _write(slide, "Text 41", "Management reading: active departments stay visible in one risk map; departments without current data are shown once here and not expanded into empty sections.")


def _populate_trend(slide: Any, payload: dict[str, Any], page: int) -> None:
    snapshot = {
        "trend": payload.get("trend") or [],
        "history_count": payload.get("history_count") or 0,
        "best_period": payload.get("best_period") or {},
        "worst_period": payload.get("worst_period") or {},
        "net_movement": payload.get("net_movement"),
        "latest_period_label": payload.get("period_label"),
    }
    section = {
        "snapshot": snapshot,
        "view": {"score": payload.get("current_score"), "target": payload.get("target", 100.0)},
    }
    reference._populate_marketing_trend(
        slide,
        {"period": payload.get("period_label"), "branch": "UAE"},
        section,
        page,
        list(payload.get("kpis") or [])[:8],
        1,
        1,
    )
    _set_common(slide, "UAE Performance Trend", "Every valid historical period for the same comparable UAE scope", _clean(payload.get("period_label")), "UAE", page)
    current_exists = payload.get("current_score") is not None
    _write(slide, "Text 32", f"{_clean(payload.get('period_label'))} vs target" if current_exists else "No current UAE data")
    _write(slide, "Text 76", "Reading rule: trend uses all available comparable periods; if no comparable baseline exists, MoM is NEW rather than a fabricated movement.")


def _populate_driver_table(slide: Any, payload: dict[str, Any], page: int) -> None:
    _set_common(slide, "Largest KPI Performance Drivers", "Direction-aware weighted performance loss ranked from the live UAE scope", _clean(payload.get("period_label")), "UAE", page)
    headers = ["KPI", "Affected Department", "Current → Baseline", "MoM", "Target", "Achievement", "Direction / Loss", "Status / Employees"]
    for name, value in zip(("Text 9", "Text 11", "Text 13", "Text 15", "Text 17", "Text 19", "Text 21", "Text 23"), headers):
        _write(slide, name, value)
    groups = _row_groups(slide, top_min=1.57, top_max=5.25, min_columns=8)
    rows = list(payload.get("kpi_evidence_rows") or [])
    displayed = rows[:len(groups)]
    displayed_count = len(displayed)
    if not rows and groups:
        _write_shape(groups[0][0], "No measured KPI data")
        for shape in groups[0][1:]:
            _blank_shape(shape)
        displayed_count = 1
    for group, row in zip(groups, displayed):
        values = [
            _fit_wrapped(row.get("kpi"), 22),
            _fit_wrapped(row.get("affected_department"), 24),
            f"{_native(row.get('current_actual'))} → {_native(row.get('baseline'))}",
            _mom(row.get("mom"), current_exists=row.get("current_actual") is not None),
            _native(row.get("target")),
            _pct(row.get("achievement")),
            f"{reference._direction_label(row.get('direction'))} | {_pct(row.get('weighted_loss'))}",
            f"{_status_label(row.get('status'))} | {row.get('employees_affected', 0)}",
        ]
        for shape, value in zip(group, values):
            _write_shape(shape, value)
    _blank_rows(groups, displayed_count)
    _write(slide, "Text 169", "Direction is applied before achievement, target gap, status color, and weighted loss. Lower response time, rejection, or error rate is positive when the KPI is lower-is-better.")


def _driver_pages(payload: dict[str, Any]) -> list[dict[str, Any]]:
    rows = list(payload.get("driver_department_impact") or [])
    return [
        {"_driver_rows": rows[index:index + 4], "_driver_part": (index // 4) + 1, "_driver_total": (len(rows) + 3) // 4}
        for index in range(0, len(rows), 4)
    ] or [{"_driver_rows": [], "_driver_part": 1, "_driver_total": 1}]


def _populate_driver_map(slide: Any, payload: dict[str, Any], page: int) -> None:
    part = payload.get("_driver_part", 1)
    total = payload.get("_driver_total", 1)
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"Driver → Department → Function → Impact{suffix}", "Aggregate executive view: employee counts replace names for group-safe access", _clean(payload.get("period_label")), "UAE", page)
    _write(slide, "Text 8", "Priority driver map")
    _write(slide, "Text 11", "Driver")
    _write(slide, "Text 13", "Role / Function")
    _write(slide, "Text 15", "Employees affected")
    _write(slide, "Text 17", "Required management action")
    groups = _row_groups(slide, top_min=2.27, top_max=5.1, min_columns=4)
    rows = list(payload.get("_driver_rows") or [])
    for group, row in zip(groups[:4], rows):
        values = [
            f"{_clean(row.get('driver'))} | {row.get('department')}",
            _clean(row.get("role_function"), "Function not available"),
            str(row.get("employees_affected", 0)),
            "Confirm owner, intervention, and next-review evidence.",
        ]
        for shape, value in zip(group, values):
            _write_shape(shape, _fit_wrapped(value, 70, max_lines=3))
    _blank_rows(groups, min(len(rows), 4))
    _write(slide, "Text 51", "Management interpretation: the score is the entry point; the intervention follows the measured driver, department, function, affected employee count, and required action.")


def _action_pages(payload: dict[str, Any], source: Presentation) -> list[dict[str, Any]]:
    rows = list(payload.get("action_rows") or [])
    capacity = max(1, len(_row_groups(source.slides[24], top_min=1.78, top_max=5.0, min_columns=6)))
    return [
        {"_action_rows": rows[index:index + capacity], "_action_part": (index // capacity) + 1, "_action_total": (len(rows) + capacity - 1) // capacity}
        for index in range(0, len(rows), capacity)
    ] or [{"_action_rows": [], "_action_part": 1, "_action_total": 1}]


def _populate_actions(slide: Any, payload: dict[str, Any], page: int) -> None:
    part = payload.get("_action_part", 1)
    total = payload.get("_action_total", 1)
    suffix = f" ({part}/{total})" if total > 1 else ""
    _set_common(slide, f"Action Ownership and Execution Risk{suffix}", "Recorded actions are in scope; unresolved owners, dates, success measures, and evidence remain visible", _clean(payload.get("period_label")), "UAE", page)
    headers = ["Priority", "KPI / Issue", "Owner", "Department / Workstream", "Execution / Status", "Success / Evidence"]
    for name, value in zip(("Text 9", "Text 11", "Text 13", "Text 15", "Text 17", "Text 19"), headers):
        _write(slide, name, value)
    groups = _row_groups(slide, top_min=1.78, top_max=5.0, min_columns=6)
    rows = list(payload.get("_action_rows") or [])
    for group, row in zip(groups, rows):
        owner = _clean(row.get("owner"), "Owner needed")
        due = _clean(row.get("due_date"), "Due date needed")
        status = _clean(row.get("status"), "Status not recorded")
        success = _clean(row.get("success_measure"), "Success metric needed")
        evidence = _clean(row.get("evidence_reference"), "Evidence reference needed")
        values = [
            _clean(row.get("priority"), "Medium"),
            _clean(row.get("kpi"), "Issue not linked"),
            f"{owner}\nDue: {due}",
            _fit_wrapped(row.get("department"), 34),
            _fit_wrapped(f"{row.get('action')}\nStatus: {status}", 68, max_lines=3),
            _fit_wrapped(f"{success}\nEvidence: {evidence}", 46, max_lines=3),
        ]
        for shape, value in zip(group, values):
            _write_shape(shape, value)
    _blank_rows(groups, len(rows))
    summary = (
        f"Execution totals: {payload.get('open_action_count', 0)} open actions | {payload.get('missing_owner_count', 0)} missing owner(s) | {payload.get('missing_due_date_count', 0)} missing due date(s) | {payload.get('overdue_actions', 0)} overdue."
        if rows
        else "No recorded open actions are in the selected UAE scope. No empty action rows are carried into the appendix."
    )
    _write(slide, "Text 69", summary)


def _populate_decisions(slide: Any, payload: dict[str, Any], page: int) -> None:
    decisions = list(payload.get("decisions") or [])[:5]
    _set_common(slide, "Management Decisions Required", "Top decisions required before the next UAE review; evidence is separated from unconfirmed cause", _clean(payload.get("period_label")), "UAE", page)
    cards = [
        (("Text 10", "Text 11", "Text 12"), ("Shape 8", "Shape 9")),
        (("Text 15", "Text 16", "Text 17"), ("Shape 13", "Shape 14")),
        (("Text 20", "Text 21", "Text 22"), ("Shape 18", "Shape 19")),
        (("Text 25", "Text 26", "Text 27"), ("Shape 23", "Shape 24")),
        (("Text 30", "Text 31", "Text 32"), ("Shape 28", "Shape 29")),
        (("Text 35", "Text 36", "Text 37"), ("Shape 33", "Shape 34")),
    ]
    for index, (names, shapes) in enumerate(cards):
        decision = decisions[index] if index < len(decisions) else None
        _card_status(slide, shapes[0], shapes[1], "Requires Action" if decision else "Data Not Available")
        if not decision:
            for name in names:
                _blank_shape(_shape(slide, name))
            continue
        _write(slide, names[0], _fit(decision.get("issue"), 34))
        _write(slide, names[1], _fit_wrapped(f"{decision.get('affected_scope')} | {decision.get('evidence_state')}", 38))
        _write(slide, names[2], _fit_wrapped(f"Decision: {decision.get('required_decision')}\nOwner: {decision.get('accountable_owner')}\nCommitment: {decision.get('commitment_date')}", 44, max_lines=3))
    _write(slide, "Text 38", "Evidence rule")
    _write(slide, "Text 41", _fit(
        "KPI Evidence is measured actual versus target. Confirmed Root Cause appears only when recorded evidence supports it. Otherwise: Investigation Required — cause validation is not recorded.",
        175,
    ))


def _populate_next_review(slide: Any, payload: dict[str, Any], page: int) -> None:
    next_review = payload.get("next_review") or {}
    _set_common(slide, "Next Review Commitments", "Close the executive story with a measurable next-review agenda", _clean(payload.get("period_label")), "UAE", page)
    items = [
        ("CURRENT UAE SCORE", _pct(payload.get("current_score")), f"MoM {_mom(payload.get('mom'), current_exists=payload.get('current_score') is not None)}"),
        ("TARGET", _pct(payload.get("target")), f"Gap {_pct(payload.get('gap_to_target'))}"),
        ("EXPECTED IMPROVEMENT", _pct(next_review.get("expected_improvement")), "Recheck against target"),
        ("TOP KPI GAPS", ", ".join(next_review.get("top_kpi_gaps") or []) or "No measured gap", "Actual / target / direction"),
        ("DEPARTMENTS REQUIRING ACTION", str(len(next_review.get("departments_requiring_action") or [])), ", ".join(next_review.get("departments_requiring_action") or []) or "None"),
        ("OPEN ACTIONS", str(payload.get("open_action_count", 0)), f"Missing owners {payload.get('missing_owner_count', 0)} | due dates {payload.get('missing_due_date_count', 0)}"),
    ]
    cards = [
        (("Text 10", "Text 11", "Text 12"), ("Shape 8", "Shape 9")),
        (("Text 15", "Text 16", "Text 17"), ("Shape 13", "Shape 14")),
        (("Text 20", "Text 21", "Text 22"), ("Shape 18", "Shape 19")),
        (("Text 25", "Text 26", "Text 27"), ("Shape 23", "Shape 24")),
        (("Text 30", "Text 31", "Text 32"), ("Shape 28", "Shape 29")),
        (("Text 35", "Text 36", "Text 37"), ("Shape 33", "Shape 34")),
    ]
    for index, (names, shapes) in enumerate(cards):
        item = items[index]
        _card_status(slide, shapes[0], shapes[1], "Requires Action" if index in {2, 3, 4, 5} else "Watch")
        _write(slide, names[0], _fit(item[0], 30))
        _write(slide, names[1], _fit(item[1], 30))
        _write(slide, names[2], _fit_wrapped(item[2], 44))
    _write(slide, "Text 38", "Required management commitments")
    _write(slide, "Text 41", _fit(
        f"Next review: {next_review.get('period', 'next available reporting period')}. {next_review.get('commitments', 'Confirm accountable owners and dates for unresolved actions.')}",
        175,
    ))


def _appendix_pages(payload: dict[str, Any], source: Presentation) -> list[tuple[int, dict[str, Any], str]]:
    pages: list[tuple[int, dict[str, Any], str]] = []
    departments = list(payload.get("department_scorecard_rows") or [])
    coverage = list(payload.get("data_coverage") or [])
    employee_rows = list(payload.get("employee_priority_rows") or [])
    # Main ranking handles the first department page.  Only overflow or data
    # coverage gets an appendix, preventing empty filler pages.
    department_capacity = max(1, len(_row_groups(source.slides[27], top_min=1.46, top_max=5.9, min_columns=7)))
    if len(departments) > 9:
        for index in range(9, len(departments), department_capacity):
            pages.append((27, {"_appendix_rows": departments[index:index + department_capacity], "_appendix_kind": "Department scorecard", "_appendix_part": (index - 9) // department_capacity + 1}, "department"))
    if any(not row.get("current_data") for row in coverage):
        pages.append((27, {"_appendix_rows": coverage, "_appendix_kind": "Data Coverage", "_appendix_part": 1}, "coverage"))
    if employee_rows:
        pages.append((27, {"_appendix_rows": employee_rows, "_appendix_kind": "Priority population", "_appendix_part": 1}, "employee"))
    return pages


def _populate_appendix(slide: Any, payload: dict[str, Any], page: int) -> None:
    kind = payload.get("_appendix_kind", "Appendix")
    rows = list(payload.get("_appendix_rows") or [])
    _set_common(slide, f"UAE Executive Appendix — {kind}", "Dynamic reference rows only when measured data or coverage requires them", _clean(payload.get("period_label")), "UAE", page)
    headers = ["Area", "Employee count", "Role / Function", "Score", "MoM", "Primary Gap", "Action Status"]
    for name, value in zip(("Text 9", "Text 11", "Text 13", "Text 15", "Text 17", "Text 19", "Text 21"), headers):
        _write(slide, name, value)
    groups = _row_groups(slide, top_min=1.46, top_max=5.9, min_columns=7)
    values_rows: list[list[str]] = []
    for row in rows:
        if kind == "Data Coverage":
            values_rows.append([_clean(row.get("department")), "—", "—", "—", "—", _clean(row.get("status")), "Shown once"])
        elif kind == "Department scorecard":
            values_rows.append([
                _clean(row.get("department")),
                str(row.get("employees_requiring_action", 0)),
                _clean(row.get("leading_kpi_driver"), "No measured KPI"),
                _pct(row.get("current_score")),
                _mom(row.get("mom"), current_exists=row.get("current_score") is not None),
                f"Loss {_pct(row.get('priority_score'))} priority points",
                _clean(row.get("status")),
            ])
        else:
            values_rows.append([
                _clean(row.get("department")),
                str(row.get("employee_count", 0)),
                _clean(row.get("role_function")),
                _pct(row.get("score")),
                "—",
                "Employee below action threshold",
                "Action focus",
            ])
    for group, values in zip(groups, values_rows):
        for shape, value in zip(group, values):
            _write_shape(shape, value)
    _blank_rows(groups, len(values_rows))


def build_uae_executive_summary_pptx(period_label: str = "Selected period", report_data: dict[str, Any] | None = None) -> bytes:
    """Generate the concise, dynamic UAE CEO report from the approved source."""

    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Approved Offshore template is missing: {TEMPLATE_PATH}")
    template_bytes = TEMPLATE_PATH.read_bytes()
    source = Presentation(io.BytesIO(template_bytes))
    if len(source.slides) != REFERENCE_SLIDES:
        raise ValueError(f"Approved Offshore template must contain {REFERENCE_SLIDES} slides; found {len(source.slides)}")
    prs = Presentation(io.BytesIO(template_bytes))
    _clear_slides(prs)
    payload = _copy_payload(report_data or {})
    payload["period_label"] = payload.get("period_label") or period_label

    plan: list[tuple[int, Any, dict[str, Any]]] = [
        (0, _populate_cover, payload),
        (1, _populate_overview, payload),
    ]
    for source_index, page_payload in _ranking_pages(payload, source):
        merged = _copy_payload(payload)
        merged.update(page_payload)
        plan.append((source_index, _populate_ranking, merged))
    for page_payload in _risk_pages(payload):
        merged = _copy_payload(payload)
        merged.update(page_payload)
        plan.append((2, _populate_risk_map, merged))
    plan.extend([
        (16, _populate_trend, payload),
        (23, _populate_driver_table, payload),
    ])
    for page_payload in _driver_pages(payload):
        merged = _copy_payload(payload)
        merged.update(page_payload)
        plan.append((17, _populate_driver_map, merged))
    for page_payload in _action_pages(payload, source):
        merged = _copy_payload(payload)
        merged.update(page_payload)
        plan.append((24, _populate_actions, merged))
    plan.extend([
        (26, _populate_decisions, payload),
        (26, _populate_next_review, payload),
    ])
    for source_index, page_payload, _kind in _appendix_pages(payload, source):
        merged = _copy_payload(payload)
        merged.update(page_payload)
        plan.append((source_index, _populate_appendix, merged))

    for page, (source_index, populate, page_payload) in enumerate(plan, 1):
        slide = _append_template_slide(prs, source, source_index)
        populate(slide, page_payload, page)
    prs.core_properties.title = f"UAE Offshore Departments Executive Performance Summary - {payload['period_label']}"
    prs.core_properties.subject = "Dynamic UAE CEO / group / regional performance summary"
    prs.core_properties.author = "PMS Dashboard"
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


__all__ = ["TEMPLATE_PATH", "build_uae_executive_summary_pptx"]
