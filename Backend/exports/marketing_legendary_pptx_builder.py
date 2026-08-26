"""Dynamic Marketing PMS PowerPoint using the approved Legendary story.

The builder deliberately owns only the Marketing ``team_marketing`` export.
It consumes the same filtered input and canonical Insights snapshot as the
existing report service, but recreates the approved 16:9 visual system in code
so no July-specific file or absolute reference path is required at runtime.
"""

from __future__ import annotations

import io
import re
import textwrap
from collections import defaultdict
from statistics import mean
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from services.insights_report_service import build_insights_snapshot


SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
FONT = "Aptos"

# Exact or near-exact colors sampled from the approved reference deck.
NAVY = RGBColor(23, 59, 93)
DARK_NAVY = RGBColor(16, 41, 65)
TEAL = RGBColor(4, 149, 143)
RED = RGBColor(213, 63, 73)
AMBER = RGBColor(226, 161, 0)
GREEN = RGBColor(34, 166, 111)
BLUE = RGBColor(45, 108, 223)
MUTED = RGBColor(102, 119, 143)
FAINT = RGBColor(117, 132, 153)
BODY = RGBColor(243, 247, 251)
FOOTER = RGBColor(232, 238, 245)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(221, 230, 239)
TRACK = RGBColor(227, 234, 242)
PALE_RED = RGBColor(251, 232, 234)
PALE_AMBER = RGBColor(255, 247, 231)
PALE_GREEN = RGBColor(235, 249, 242)
PALE_BLUE = RGBColor(234, 241, 255)
PALE_NEUTRAL = RGBColor(247, 250, 253)
SHADOW = RGBColor(232, 238, 245)

MONTHS = {
    name: index
    for index, name in enumerate(
        (
            "January",
            "February",
            "March",
            "April",
            "May",
            "June",
            "July",
            "August",
            "September",
            "October",
            "November",
            "December",
        ),
        1,
    )
}
MONTH_NAMES = {value: key for key, value in MONTHS.items()}


def _number(value: Any) -> float | None:
    if value is None or value == "":
        return None
    try:
        number = float(value)
    except (TypeError, ValueError):
        return None
    return number if number == number else None


def _clean(value: Any, fallback: str = "") -> str:
    if value is None:
        return fallback
    text = " ".join(str(value).split())
    # Keep user data readable when a legacy source contains mojibake.
    text = (
        text.replace("Ã¢â‚¬â€", "—")
        .replace("Ã¢â‚¬â€œ", "–")
        .replace("Ã¢â‚¬Â¢", "•")
        .replace("Ã¢â‚¬â„¢", "'")
    )
    return text or fallback


def _fmt_pct(value: Any, *, signed: bool = False, fallback: str = "—") -> str:
    number = _number(value)
    if number is None:
        return fallback
    return f"{number:+.1f}%" if signed else f"{number:.1f}%"


def _fmt_mom(value: Any) -> str:
    return "NEW" if _number(value) is None else _fmt_pct(value, signed=True)


def _fmt_native(value: Any, unit: Any = "", fallback: str = "—") -> str:
    number = _number(value)
    if number is None:
        return fallback
    normalized = _clean(unit).casefold()
    if normalized in {"%", "percent", "percentage"}:
        if abs(number) <= 1:
            number *= 100
        return _fmt_pct(number)
    if normalized in {"minutes", "minute", "min"}:
        return f"{number:,.0f} min" if number.is_integer() else f"{number:,.1f} min"
    if normalized in {"seconds", "second", "sec"}:
        return f"{number:,.0f} sec" if number.is_integer() else f"{number:,.1f} sec"
    if normalized in {"count", "number", "visits"}:
        return f"{number:,.0f}" if number.is_integer() else f"{number:,.1f}"
    formatted = f"{number:,.0f}" if number.is_integer() else f"{number:,.2f}".rstrip("0").rstrip(".")
    return f"{formatted} {_clean(unit)}".strip()


def _direction(value: Any) -> str:
    return "lower_better" if _clean(value).casefold() in {"lower_better", "lower is better"} else "higher_better"


def _direction_label(value: Any) -> str:
    return "Lower is better" if _direction(value) == "lower_better" else "Higher is better"


def _status_from_achievement(achievement: Any) -> str:
    value = _number(achievement)
    if value is None:
        return "Data quality"
    if value >= 100:
        return "On track"
    if value >= 70:
        return "Watch"
    return "Requires action"


def _employee_status(score: Any) -> str:
    value = _number(score)
    if value is None:
        return "Data unavailable"
    if value < 70:
        return "Requires action"
    if value < 90:
        return "Watch"
    return "On track"


def _status_color(status: Any) -> RGBColor:
    value = _clean(status).casefold()
    if value in {"requires action", "critical", "below target", "below"}:
        return RED
    if value in {"watch", "at risk", "pending", "proposed"}:
        return AMBER
    if value in {"on track", "improving", "recorded"}:
        return GREEN
    return BLUE


def _status_fill(status: Any) -> RGBColor:
    value = _clean(status).casefold()
    if value in {"requires action", "critical", "below target", "below"}:
        return PALE_RED
    if value in {"watch", "at risk", "pending", "proposed"}:
        return PALE_AMBER
    if value in {"on track", "improving", "recorded"}:
        return PALE_GREEN
    return PALE_BLUE


def _movement_color(value: Any) -> RGBColor:
    number = _number(value)
    if number is None or number == 0:
        return MUTED
    return GREEN if number > 0 else RED


def _month_key(record: dict[str, Any]) -> tuple[int, int] | None:
    year = _number(record.get("year"))
    raw_month = record.get("month")
    month = _number(raw_month) or MONTHS.get(_clean(raw_month).title())
    return (int(year), int(month)) if year and month else None


def _period_label(key: tuple[int, int]) -> str:
    return f"{MONTH_NAMES.get(key[1], key[1])} {key[0]}"


def _history_fallback(records: list[dict[str, Any]], period_label: str) -> list[dict[str, Any]]:
    grouped: dict[tuple[int, int], list[dict[str, Any]]] = defaultdict(list)
    for record in records:
        key = _month_key(record)
        if key:
            grouped[key].append(record)
    if not grouped:
        return [{"key": period_label, "label": period_label, "records": records}] if records else []
    return [
        {"key": f"{year}-{month:02d}", "label": _period_label((year, month)), "records": rows}
        for (year, month), rows in sorted(grouped.items())
    ]


def _prepare_snapshot(period_label: str, report_data: dict[str, Any] | None) -> dict[str, Any]:
    data = dict(report_data or {})
    records = [dict(row) for row in (data.get("records") or [])]
    selected = [dict(row) for row in (data.get("selected_records") or records)]
    history = list(data.get("history") or [])
    if not history:
        history = _history_fallback(selected, period_label)
    return build_insights_snapshot(
        {
            "period_label": data.get("period_label") or period_label,
            "scope_label": data.get("scope_label") or "Marketing",
            "filters": dict(data.get("filters") or {}),
            "aggregate_only": bool(data.get("aggregate_only")),
            "kpi_definitions": list(data.get("kpi_definitions") or []),
            "records": records,
            "selected_records": selected,
            "history": history,
            "actions": list(data.get("actions") or []),
        }
    )


def _team_label(payload: dict[str, Any]) -> str:
    filters = payload.get("filters") or {}
    team = _clean(filters.get("team"))
    if team:
        return team
    records = payload.get("records") or payload.get("all_people") or []
    teams = sorted({_clean(row.get("team")) for row in records if _clean(row.get("team"))})
    return teams[0] if len(teams) == 1 else "Marketing"


def _period(payload: dict[str, Any]) -> str:
    return _clean(payload.get("latest_period_label") or payload.get("period_label"), "Selected period")


def _people(payload: dict[str, Any]) -> list[dict[str, Any]]:
    if not payload.get("people_visible"):
        return []
    return [dict(row) for row in (payload.get("all_people") or [])]


def _primary_kpi(person: dict[str, Any]) -> dict[str, Any]:
    kpis = [dict(row) for row in (person.get("kpis") or []) if isinstance(row, dict)]
    if not kpis:
        return {}
    positive = [row for row in kpis if (_number(row.get("weighted_impact")) or 0) > 0]
    return (positive or kpis)[0]


def _role_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    people = _people(payload)
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for person in people:
        grouped[_clean(person.get("position"), "Unassigned")].append(person)
    rows: list[dict[str, Any]] = []
    for role, members in grouped.items():
        def avg(key: str) -> float | None:
            values = [_number(member.get(key)) for member in members]
            values = [value for value in values if value is not None]
            return mean(values) if values else None

        rows.append(
            {
                "role": role,
                "score": avg("score"),
                "baseline": avg("baseline_score"),
                "movement": avg("movement"),
                "count": len(members),
            }
        )
    return sorted(rows, key=lambda row: (_number(row.get("score")) is None, -(_number(row.get("score")) or 0)))


def _fit_title(value: Any, limit: int = 72) -> str:
    text = _clean(value)
    if len(text) <= limit:
        return text
    return textwrap.shorten(text, width=limit, placeholder="…")


def _wrap(value: Any, width: int) -> str:
    text = _clean(value)
    if not text:
        return "—"
    return "\n".join(textwrap.wrap(text, width=max(8, width), break_long_words=False, break_on_hyphens=False))


def _safe_name(value: Any) -> str:
    return re.sub(r"[^A-Za-z0-9]+", "_", _clean(value)).strip("_") or "scope"


def _add_rect(slide: Any, x: float, y: float, width: float, height: float, fill: RGBColor, line: RGBColor | None = None, *, rounded: bool = False) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(max(0.01, width)),
        Inches(max(0.01, height)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.7)
    return shape


def _add_text(
    slide: Any,
    x: float,
    y: float,
    width: float,
    height: float,
    value: Any,
    *,
    size: float = 8.0,
    color: RGBColor = NAVY,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
    italic: bool = False,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(max(0.01, width)), Inches(max(0.01, height)))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = "" if value is None else str(value)
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _add_card(slide: Any, x: float, y: float, width: float, height: float, *, fill: RGBColor = WHITE, rail: RGBColor = TEAL, line: RGBColor = LINE, rounded: bool = True) -> Any:
    _add_rect(slide, x + 0.045, y + 0.045, width, height, SHADOW, None, rounded=rounded)
    card = _add_rect(slide, x, y, width, height, fill, line, rounded=rounded)
    _add_rect(slide, x, y, 0.075, height, rail, rail, rounded=False)
    return card


def _add_header(slide: Any, title: str, subtitle: str, period: str, page: int) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BODY
    _add_rect(slide, 0, 0, SLIDE_WIDTH, 0.82, NAVY, NAVY)
    _add_rect(slide, 0, 0.82, SLIDE_WIDTH, 0.08, TEAL, TEAL)
    _add_text(slide, 0.38, 0.14, 10.9, 0.32, _fit_title(title, 76), size=18.5 if len(_clean(title)) < 58 else 16.5, color=WHITE, bold=True)
    _add_text(slide, 0.39, 0.52, 10.9, 0.18, _wrap(subtitle, 125), size=7.3, color=RGBColor(220, 228, 242))
    _add_text(slide, 11.72, 0.16, 1.2, 0.22, f"{period} | {page:02d}", size=7.7, color=RGBColor(220, 228, 242), align=PP_ALIGN.RIGHT)


def _add_footer(slide: Any, period: str, page: int) -> None:
    _add_rect(slide, 0, 7.23, SLIDE_WIDTH, 0.27, FOOTER, FOOTER)
    _add_text(slide, 0.30, 7.29, 5.5, 0.12, f"Marketing PMS | {period}", size=5.8, color=FAINT)
    _add_text(slide, 12.42, 7.29, 0.55, 0.12, f"{page:02d}", size=5.8, color=FAINT, align=PP_ALIGN.RIGHT)


def _add_section_title(slide: Any, x: float, y: float, width: float, title: str, *, accent: RGBColor = TEAL) -> None:
    _add_text(slide, x, y, width, 0.25, title, size=13.5, color=NAVY, bold=True)
    _add_rect(slide, x, y + 0.36, width, 0.045, accent, accent)


def _add_bar(slide: Any, x: float, y: float, width: float, height: float, ratio: Any, color: RGBColor) -> None:
    _add_rect(slide, x, y, width, height, TRACK, TRACK, rounded=True)
    value = _number(ratio)
    if value is None:
        return
    _add_rect(slide, x, y, width * max(0.0, min(1.0, value)), height, color, color, rounded=True)


def _add_table(
    slide: Any,
    x: float,
    y: float,
    width: float,
    columns: list[tuple[str, str, float]],
    rows: list[dict[str, Any]],
    *,
    row_height: float = 0.38,
    header_height: float = 0.38,
    font_size: float = 6.7,
    header_size: float = 6.8,
    row_fill: Any = None,
    bold_keys: Iterable[str] = (),
    center_keys: Iterable[str] = (),
) -> float:
    total = sum(column[2] for column in columns)
    if total <= 0:
        return y
    cursor_x = x
    for key, label, fraction in columns:
        cell_width = width * fraction / total
        _add_rect(slide, cursor_x, y, cell_width, header_height, NAVY, WHITE)
        _add_text(slide, cursor_x + 0.02, y + 0.01, cell_width - 0.04, header_height - 0.02, label, size=header_size, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER if key in center_keys else PP_ALIGN.LEFT)
        cursor_x += cell_width
    for row_index, row in enumerate(rows):
        top = y + header_height + row_index * row_height
        fill = row_fill(row, row_index) if callable(row_fill) else (PALE_NEUTRAL if row_index % 2 else WHITE)
        cursor_x = x
        for key, _label, fraction in columns:
            cell_width = width * fraction / total
            value = row.get(key, "—")
            cell_fill = _status_fill(value) if key == "status" else fill
            _add_rect(slide, cursor_x, top, cell_width, row_height, cell_fill, LINE)
            color = _status_color(value) if key == "status" else RED if key in {"loss", "weighted_loss", "gap"} and _number(row.get(key)) is not None and (_number(row.get(key)) or 0) > 0 else NAVY
            display = _wrap(value, max(8, int(cell_width * 12)))
            _add_text(
                slide,
                cursor_x + 0.03,
                top + 0.02,
                cell_width - 0.06,
                row_height - 0.04,
                display,
                size=font_size,
                color=color,
                bold=key in set(bold_keys) or key == "status",
                valign=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.CENTER if key in center_keys else PP_ALIGN.LEFT,
            )
            cursor_x += cell_width
    return y + header_height + len(rows) * row_height


def _chunks(items: list[Any], size: int) -> list[list[Any]]:
    return [items[index : index + size] for index in range(0, len(items), size)] or [[]]


def _score_status(score: Any) -> str:
    value = _number(score)
    if value is None:
        return "Data unavailable"
    if value < 70:
        return "Requires action"
    if value < 90:
        return "Watch"
    return "On track"


def _score_baseline(payload: dict[str, Any]) -> Any:
    return payload.get("baseline_score")


def _score_mom(payload: dict[str, Any]) -> Any:
    return payload.get("movement")


def _score_gap(payload: dict[str, Any]) -> Any:
    score = _number(payload.get("overall_score"))
    target = _number(payload.get("target_score"))
    return score - target if score is not None and target is not None else None


def _top_performer(payload: dict[str, Any]) -> dict[str, Any]:
    people = [row for row in _people(payload) if _number(row.get("score")) is not None]
    return max(people, key=lambda row: (_number(row.get("score")) or 0, _clean(row.get("name")))) if people else {}


def _largest_decline(payload: dict[str, Any]) -> dict[str, Any]:
    people = [row for row in _people(payload) if _number(row.get("movement")) is not None]
    return min(people, key=lambda row: (_number(row.get("movement")) or 0, _clean(row.get("name")))) if people else {}


def _below_threshold(payload: dict[str, Any], threshold: float = 70.0) -> list[dict[str, Any]]:
    return [
        row
        for row in _people(payload)
        if _number(row.get("score")) is not None and (_number(row.get("score")) or 0) < threshold
    ]


def _kpi_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    return [dict(row) for row in (payload.get("kpis") or []) if isinstance(row, dict)]


def _kpi_loss(row: dict[str, Any]) -> float:
    return max(0.0, _number(row.get("weighted_impact")) or 0.0)


def _kpi_action(payload: dict[str, Any], kpi: dict[str, Any]) -> dict[str, Any]:
    source_keys = {
        _clean(kpi.get("key")).casefold(),
        _clean(kpi.get("group_key")).casefold(),
        *[_clean(value).casefold() for value in (kpi.get("source_keys") or [])],
    }
    for action in payload.get("actions") or []:
        linked = _clean(action.get("linked_kpi_key")).casefold()
        if linked and (linked in source_keys or any(linked == key for key in source_keys if key)):
            return dict(action)
    return {}


def _person_kpi_rows(person: dict[str, Any]) -> list[dict[str, Any]]:
    rows = [dict(row) for row in (person.get("kpis") or []) if isinstance(row, dict)]
    return sorted(rows, key=lambda row: (-_kpi_loss(row), _clean(row.get("label")).casefold()))


def _employee_priority_groups(payload: dict[str, Any]) -> dict[str, list[dict[str, Any]]]:
    groups = {"Requires Action": [], "Watch": [], "On Track": []}
    for person in _people(payload):
        status = _employee_status(person.get("score"))
        key = {
            "Requires action": "Requires Action",
            "Watch": "Watch",
            "On track": "On Track",
        }.get(status, "Watch")
        groups[key].append(person)
    for values in groups.values():
        values.sort(key=lambda row: (_number(row.get("score")) is None, _number(row.get("score")) or 0, -_kpi_loss(_primary_kpi(row))))
    return groups


def _driver_role_employee_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    people = _people(payload)
    rows: list[dict[str, Any]] = []
    for kpi in _kpi_rows(payload):
        affected = [
            person
            for person in people
            if any(
                _clean(item.get("group_key")).casefold() == _clean(kpi.get("group_key")).casefold()
                or _clean(item.get("key")).casefold() in {
                    _clean(kpi.get("key")).casefold(),
                    *[_clean(value).casefold() for value in (kpi.get("source_keys") or [])],
                }
                for item in _person_kpi_rows(person)
                if _kpi_loss(item) > 0
            )
        ]
        if not affected and _kpi_loss(kpi) <= 0:
            continue
        roles = sorted({_clean(person.get("position"), "Unassigned") for person in affected})
        names = sorted({_clean(person.get("name"), "Unknown") for person in affected})
        action = _kpi_action(payload, kpi)
        rows.append(
            {
                "driver": _clean(kpi.get("label"), "KPI"),
                "loss": _kpi_loss(kpi),
                "role": ", ".join(roles) or "Selected role(s)",
                "employees": ", ".join(names) or "Employee-level detail unavailable",
                "employee_count": len(names),
                "action": _clean(action.get("action_display"), "Validate the operational cause and assign corrective action"),
                "status": _clean(kpi.get("status"), "Data quality"),
            }
        )
    return sorted(rows, key=lambda row: (-(_number(row.get("loss")) or 0), _clean(row.get("driver")).casefold()))


def _evidence_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for item in payload.get("root_cause_rows") or []:
        state = _clean(item.get("evidence_status"), "Investigation Required")
        if state.casefold() == "kpi signal only":
            display_state = "Investigation Required"
        elif "confirmed" in state.casefold():
            display_state = "Confirmed Root Cause"
        else:
            display_state = "KPI Evidence"
        rows.append(
            {
                "driver": _clean(item.get("label"), "KPI"),
                "evidence": display_state,
                "recorded": _clean(item.get("recorded_root_cause"), "No operational cause recorded"),
                "validation": _clean(item.get("required_validation"), "Validate the operational cause before closing the gap"),
                "owner": _clean(item.get("owner_display"), "Owner needed"),
                "loss": _kpi_loss(item),
            }
        )
    return rows


def _detail_people(payload: dict[str, Any]) -> list[dict[str, Any]]:
    people = _people(payload)
    return sorted(
        [person for person in people if _employee_status(person.get("score")) in {"Requires action", "Watch"}],
        key=lambda row: (-_kpi_loss(_primary_kpi(row)), _number(row.get("score")) is None, _number(row.get("score")) or 0, _clean(row.get("name")).casefold()),
    )


def _priority_rank(payload: dict[str, Any], person: dict[str, Any]) -> int:
    people = _detail_people(payload)
    for index, row in enumerate(people, 1):
        if _clean(row.get("employee_id")) == _clean(person.get("employee_id")):
            return index
    return len(people) + 1


def _management_actions(payload: dict[str, Any]) -> list[dict[str, Any]]:
    return [dict(row) for row in (payload.get("actions") or []) if isinstance(row, dict)]


def _action_issue(action: dict[str, Any]) -> str:
    return _clean(action.get("action_type"), _clean(action.get("linked_kpi_key"), "Management workstream"))


def _action_status(action: dict[str, Any]) -> str:
    return _clean(action.get("status_display"), _clean(action.get("status"), "Status not recorded"))


def _action_owner(action: dict[str, Any]) -> str:
    return _clean(action.get("owner_display"), _clean(action.get("owner"), "Owner needed"))


def _action_due(action: dict[str, Any]) -> str:
    return _clean(action.get("due_date_display"), _clean(action.get("due_date"), "Due date needed"))


def _action_text(action: dict[str, Any]) -> str:
    return _clean(action.get("action_display"), _clean(action.get("action_text"), "Action text needed"))


def _action_success(action: dict[str, Any]) -> str:
    return _clean(action.get("success_metric_display"), _clean(action.get("success_metric"), "Success metric needed"))


def _history_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    rows = []
    for row in payload.get("trend") or []:
        rows.append(dict(row))
    return rows


def _new_slide(prs: Presentation) -> Any:
    return prs.slides.add_slide(prs.slide_layouts[6])


def _metric_card(
    slide: Any,
    x: float,
    y: float,
    width: float,
    height: float,
    label: str,
    value: Any,
    *,
    value_color: RGBColor = NAVY,
    rail: RGBColor = TEAL,
    subtext: str = "",
    value_size: float = 19.0,
) -> None:
    _add_card(slide, x, y, width, height, rail=rail)
    _add_text(slide, x + 0.18, y + 0.12, width - 0.30, 0.20, label.upper(), size=6.2, color=MUTED, bold=True)
    display_value = value if value not in {None, ""} else "Not available"
    multiline = "\n" in str(display_value)
    _add_text(slide, x + 0.18, y + 0.34, width - 0.30, 0.55 if multiline else 0.40, display_value, size=min(value_size, 10.5) if multiline else value_size, color=value_color, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if subtext:
        _add_text(slide, x + 0.18, y + height - 0.28, width - 0.30, 0.16, subtext, size=5.8, color=FAINT)


def _status_badge(slide: Any, x: float, y: float, width: float, height: float, status: str) -> None:
    _add_rect(slide, x, y, width, height, _status_fill(status), _status_fill(status), rounded=True)
    _add_text(slide, x + 0.04, y + 0.01, width - 0.08, height - 0.02, status, size=6.3, color=_status_color(status), bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def _status_label_for_kpi(kpi: dict[str, Any]) -> str:
    return _clean(kpi.get("status"), _status_from_achievement(kpi.get("achievement")))


def _kpi_value_line(kpi: dict[str, Any]) -> tuple[str, str, str, str, str]:
    return (
        _fmt_native(kpi.get("actual"), kpi.get("unit")),
        _fmt_native(kpi.get("baseline_actual"), kpi.get("unit")),
        _fmt_mom(kpi.get("mom")),
        _fmt_native(kpi.get("target"), kpi.get("unit")),
        _status_label_for_kpi(kpi),
    )


def _slide_cover(prs: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _new_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_NAVY
    _add_rect(slide, 0, 0, SLIDE_WIDTH, 0.12, TEAL, TEAL)
    _add_rect(slide, 0, 7.34, SLIDE_WIDTH, 0.16, TEAL, TEAL)
    _add_text(slide, 0.68, 0.70, 6.8, 0.22, "MARKETING PMS / PERFORMANCE REVIEW", size=8.2, color=RGBColor(170, 223, 220), bold=True)
    _add_text(slide, 0.68, 1.45, 7.6, 0.62, "Marketing Team\nPerformance Review", size=28.0, color=WHITE, bold=True)
    _add_text(slide, 0.70, 2.90, 4.0, 0.26, _period(payload), size=13.5, color=RGBColor(220, 228, 242), bold=True)
    _add_text(slide, 0.70, 3.35, 5.9, 0.46, "Executive performance story\nCurrent month, drivers, people, and actions", size=10.2, color=RGBColor(185, 204, 222))
    _add_rect(slide, 0.70, 4.55, 5.35, 0.04, TEAL, TEAL)
    _add_text(slide, 0.70, 4.78, 5.4, 0.25, "Prepared from the selected Marketing PMS dataset", size=7.0, color=RGBColor(185, 204, 222))
    current = _number(payload.get("overall_score"))
    driver = payload.get("driver") or {}
    lens_x = 8.45
    _add_card(slide, lens_x, 1.00, 4.15, 5.35, fill=RGBColor(21, 55, 84), rail=TEAL, line=RGBColor(54, 91, 122), rounded=True)
    _add_text(slide, lens_x + 0.34, 1.37, 3.4, 0.24, "EXECUTIVE LENS", size=7.2, color=RGBColor(170, 223, 220), bold=True)
    _add_text(slide, lens_x + 0.34, 1.78, 3.2, 0.26, "Current monthly performance", size=8.0, color=RGBColor(185, 204, 222))
    _add_text(slide, lens_x + 0.34, 2.07, 3.2, 0.62, _fmt_pct(current), size=31.0, color=WHITE, bold=True)
    _add_text(slide, lens_x + 0.34, 2.83, 3.2, 0.23, f"MoM: {_fmt_mom(payload.get('movement'))}", size=8.2, color=_movement_color(payload.get("movement")), bold=True)
    _add_rect(slide, lens_x + 0.34, 3.34, 3.35, 0.02, RGBColor(54, 91, 122), RGBColor(54, 91, 122))
    _add_text(slide, lens_x + 0.34, 3.62, 1.55, 0.18, "Below threshold", size=6.5, color=RGBColor(185, 204, 222))
    _add_text(slide, lens_x + 2.25, 3.56, 1.45, 0.28, str(len(_below_threshold(payload))), size=13.5, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    _add_text(slide, lens_x + 0.34, 4.15, 1.55, 0.18, "Leading driver", size=6.5, color=RGBColor(185, 204, 222))
    _add_text(slide, lens_x + 0.34, 4.40, 3.1, 0.48, _fit_title(driver.get("label"), 34) if driver else "No measurable gap", size=9.2, color=WHITE, bold=True)
    _add_text(slide, lens_x + 0.34, 5.22, 3.1, 0.20, f"Weighted loss {_fmt_pct(driver.get('weighted_impact'))}", size=6.7, color=RGBColor(255, 205, 145) if driver else RGBColor(185, 204, 222))
    _add_text(slide, 0.70, 6.82, 4.6, 0.14, f"Marketing PMS | {_period(payload)}", size=5.8, color=RGBColor(124, 157, 182))
    _add_text(slide, 12.35, 6.82, 0.45, 0.14, f"{page:02d}", size=5.8, color=RGBColor(124, 157, 182), align=PP_ALIGN.RIGHT)


def _slide_overview(prs: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _new_slide(prs)
    current = payload.get("overall_score")
    baseline = _score_baseline(payload)
    status = _score_status(current)
    _add_header(slide, "Monthly Performance Overview", "How did Marketing perform this month?", _period(payload), page)
    card_y, card_h, gap = 1.08, 1.13, 0.13
    card_w = (12.45 - gap * 4) / 5
    cards = [
        ("Current monthly performance", _fmt_pct(current), TEAL, NAVY),
        ("Baseline / previous comparable", _fmt_pct(baseline), BLUE, NAVY),
        ("MoM change", _fmt_mom(_score_mom(payload)), _movement_color(_score_mom(payload)), _movement_color(_score_mom(payload))),
        ("Target", _fmt_pct(payload.get("target_score")), NAVY, NAVY),
        ("Status", status, _status_color(status), _status_color(status)),
    ]
    for index, (label, value, rail, value_color) in enumerate(cards):
        _metric_card(slide, 0.43 + index * (card_w + gap), card_y, card_w, card_h, label, value, rail=rail, value_color=value_color, value_size=16.0 if label == "Status" else 18.5)

    below = _below_threshold(payload)
    top = _top_performer(payload)
    decline = _largest_decline(payload)
    driver = payload.get("driver") or {}
    second_y = 2.48
    second = [
        ("Gap to target", _fmt_pct(_score_gap(payload)), RED if (_number(_score_gap(payload)) or 0) < 0 else GREEN, RED if (_number(_score_gap(payload)) or 0) < 0 else GREEN),
        ("Employees below threshold", str(len(below)), RED if below else GREEN, NAVY),
        ("Top performer", _fit_title(top.get("name"), 20) if top else "Not available", GREEN, NAVY),
        ("Largest decline", f"{_clean(decline.get('name'), 'Not available')}\nMoM: {_fmt_mom(decline.get('movement'))}" if decline else "Not available", RED if decline else BLUE, NAVY),
        ("Leading performance driver", _fit_title(driver.get("label"), 24) if driver else "No measurable gap", RED if driver else GREEN, NAVY),
    ]
    for index, (label, value, rail, value_color) in enumerate(second):
        _metric_card(slide, 0.43 + index * (card_w + gap), second_y, card_w, 1.20, label, value, rail=rail, value_color=value_color, value_size=12.0 if label in {"Top performer", "Largest decline", "Leading performance driver"} else 18.0)

    _add_section_title(slide, 0.48, 4.05, 5.0, "Executive read")
    headline = _clean(payload.get("headline"), "No performance narrative is available.")
    driver_story = _clean(payload.get("driver_story"), "No measurable KPI gap is available.")
    _add_card(slide, 0.48, 4.50, 5.0, 1.70, fill=PALE_BLUE, rail=BLUE, line=LINE)
    _add_text(slide, 0.75, 4.78, 4.48, 0.55, headline, size=12.0, color=NAVY, bold=True)
    _add_text(slide, 0.75, 5.49, 4.43, 0.46, driver_story, size=7.4, color=MUTED)
    _add_text(slide, 0.75, 5.98, 4.3, 0.15, "Use the driver and employee pages to move from signal to action.", size=6.1, color=BLUE, italic=True)

    role_rows = _role_rows(payload)
    _add_section_title(slide, 5.78, 4.05, 7.05, "Role-level signal", accent=TEAL)
    table_rows = [
        {
            "role": _fit_title(row.get("role"), 28),
            "current": _fmt_pct(row.get("score")),
            "baseline": _fmt_pct(row.get("baseline")),
            "mom": _fmt_mom(row.get("movement")),
            "status": _employee_status(row.get("score")),
        }
        for row in role_rows[:5]
    ]
    _add_table(
        slide,
        5.78,
        4.50,
        7.05,
        [("role", "Role", 2.4), ("current", "Current", 1.0), ("baseline", "Baseline", 1.0), ("mom", "MoM", 0.85), ("status", "Status", 1.25)],
        table_rows,
        row_height=0.31,
        header_height=0.35,
        font_size=6.4,
        header_size=6.2,
        bold_keys=("role",),
        center_keys=("current", "baseline", "mom", "status"),
    )
    _add_footer(slide, _period(payload), page)


def _slide_roles(prs: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _new_slide(prs)
    _add_header(slide, "Performance by Role", "Prioritize roles with the largest current performance gap.", _period(payload), page)
    role_rows = _role_rows(payload)
    _add_section_title(slide, 0.48, 1.08, 7.65, "Role scorecard")
    table_rows = []
    for row in role_rows:
        score = _number(row.get("score"))
        table_rows.append(
            {
                "role": _fit_title(row.get("role"), 34),
                "members": str(row.get("count") or 0),
                "current": _fmt_pct(score),
                "baseline": _fmt_pct(row.get("baseline")),
                "mom": _fmt_mom(row.get("movement")),
                "status": _employee_status(score),
                "score_ratio": max(0.0, min(1.0, (score or 0) / 100.0)) if score is not None else None,
            }
        )
    _add_table(
        slide,
        0.48,
        1.53,
        7.65,
        [("role", "Role", 2.7), ("members", "People", 0.65), ("current", "Current", 1.0), ("baseline", "Baseline", 1.0), ("mom", "MoM", 0.85), ("status", "Status", 1.15)],
        table_rows,
        row_height=0.47,
        header_height=0.37,
        font_size=6.7,
        header_size=6.5,
        bold_keys=("role",),
        center_keys=("members", "current", "baseline", "mom", "status"),
        row_fill=lambda row, index: PALE_RED if _clean(row.get("status")).casefold() == "requires action" else PALE_AMBER if _clean(row.get("status")).casefold() == "watch" else (PALE_NEUTRAL if index % 2 else WHITE),
    )

    _add_section_title(slide, 8.48, 1.08, 4.35, "Management attention order", accent=RED)
    ordered_people = sorted(_people(payload), key=lambda row: (_number(row.get("score")) is None, _number(row.get("score")) or 0, -_kpi_loss(_primary_kpi(row))))[:6]
    if not ordered_people:
        _add_card(slide, 8.48, 1.53, 4.35, 1.05, fill=PALE_NEUTRAL, rail=BLUE)
        _add_text(slide, 8.77, 1.83, 3.75, 0.28, "Employee-level detail unavailable for this scope.", size=8.0, color=MUTED)
    else:
        for index, person in enumerate(ordered_people):
            y = 1.53 + index * 0.70
            status = _employee_status(person.get("score"))
            _add_card(slide, 8.48, y, 4.35, 0.55, fill=_status_fill(status), rail=_status_color(status), line=LINE)
            _add_text(slide, 8.72, y + 0.10, 2.1, 0.18, f"{index + 1}. {_fit_title(person.get('name'), 23)}", size=7.0, color=NAVY, bold=True)
            _add_text(slide, 10.85, y + 0.10, 0.85, 0.18, _fmt_pct(person.get("score")), size=7.1, color=_status_color(status), bold=True, align=PP_ALIGN.RIGHT)
            _add_text(slide, 11.73, y + 0.10, 0.88, 0.18, _fmt_mom(person.get("movement")), size=6.6, color=_movement_color(person.get("movement")), bold=True, align=PP_ALIGN.RIGHT)
            _add_text(slide, 8.72, y + 0.30, 3.65, 0.14, _fit_title((_primary_kpi(person) or {}).get("label"), 48) or "No leading KPI", size=5.7, color=MUTED)
    _add_card(slide, 8.48, 6.05, 4.35, 0.75, fill=PALE_BLUE, rail=TEAL, line=LINE)
    _add_text(slide, 8.73, 6.23, 3.85, 0.20, "Interpretation", size=6.5, color=TEAL, bold=True)
    _add_text(slide, 8.73, 6.45, 3.85, 0.22, "Low score plus high weighted loss moves a role to the front of the queue.", size=6.4, color=NAVY)
    _add_footer(slide, _period(payload), page)


def _slide_trend(prs: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _new_slide(prs)
    history = _history_rows(payload)
    history_count = len(history) or _number(payload.get("history_count")) or 0
    title = f"Performance Trend ({int(history_count)} period{'s' if history_count != 1 else ''})"
    _add_header(slide, title, "Use the available history to frame current performance and movement.", _period(payload), page)
    best = payload.get("best_period") or {}
    worst = payload.get("worst_period") or {}
    cards = [
        ("Periods available", str(int(history_count)), BLUE, NAVY),
        ("Best period", f"{_clean(best.get('label'), 'Not available')} {_fmt_pct(best.get('score'))}" if best else "Not available", GREEN, NAVY),
        ("Worst period", f"{_clean(worst.get('label'), 'Not available')} {_fmt_pct(worst.get('score'))}" if worst else "Not available", RED, NAVY),
        ("Net movement", _fmt_mom(payload.get("net_movement")), _movement_color(payload.get("net_movement")), _movement_color(payload.get("net_movement"))),
        ("Gap to target", _fmt_pct(_score_gap(payload)), RED if (_number(_score_gap(payload)) or 0) < 0 else GREEN, RED if (_number(_score_gap(payload)) or 0) < 0 else GREEN),
    ]
    gap = 0.13
    width = (12.45 - gap * 4) / 5
    for index, (label, value, rail, color) in enumerate(cards):
        _metric_card(slide, 0.43 + index * (width + gap), 1.08, width, 1.02, label, value, rail=rail, value_color=color, value_size=11.5 if label in {"Best period", "Worst period"} else 17.0)

    _add_section_title(slide, 0.48, 2.48, 7.90, "Historical performance")
    chart_x, chart_y, chart_w, chart_h = 0.78, 2.98, 7.25, 2.88
    _add_rect(slide, chart_x, chart_y + chart_h, chart_w, 0.02, LINE, LINE)
    if history:
        max_score = max([_number(row.get("score")) or 0 for row in history] + [100.0])
        min_score = min([_number(row.get("score")) or 100 for row in history] + [0.0])
        scale_min = max(0.0, min(70.0, min_score - 10.0))
        scale_max = max(100.0, max_score)
        bar_w = min(0.86, (chart_w - 0.30) / max(1, len(history)) - 0.18)
        for index, row in enumerate(history):
            score = _number(row.get("score"))
            ratio = ((score or 0) - scale_min) / max(1.0, scale_max - scale_min) if score is not None else 0
            x = chart_x + 0.28 + index * (chart_w - 0.35) / max(1, len(history)) + 0.05
            h = max(0.05, chart_h * max(0.0, min(1.0, ratio))) if score is not None else 0.05
            color = GREEN if (score or 0) >= 90 else AMBER if (score or 0) >= 70 else RED
            _add_rect(slide, x, chart_y + chart_h - h, bar_w, h, color, color, rounded=True)
            _add_text(slide, x - 0.08, chart_y + chart_h - h - 0.25, bar_w + 0.16, 0.18, _fmt_pct(score), size=6.3, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            _add_text(slide, x - 0.15, chart_y + chart_h + 0.10, bar_w + 0.30, 0.25, _fit_title(row.get("label"), 16), size=5.7, color=MUTED, align=PP_ALIGN.CENTER)
        target_y = chart_y + chart_h - chart_h * ((100.0 - scale_min) / max(1.0, scale_max - scale_min))
        if 0 <= target_y <= chart_y + chart_h:
            _add_rect(slide, chart_x, target_y, chart_w, 0.025, NAVY, NAVY)
            _add_text(slide, chart_x + chart_w - 0.70, target_y - 0.18, 0.70, 0.15, "Target", size=5.7, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    else:
        _add_text(slide, chart_x, chart_y + 1.15, chart_w, 0.26, "No valid historical periods are available.", size=9.0, color=MUTED, align=PP_ALIGN.CENTER)

    _add_section_title(slide, 8.55, 2.48, 4.28, "Trend context", accent=BLUE)
    _add_card(slide, 8.55, 2.98, 4.28, 2.88, fill=PALE_BLUE, rail=BLUE, line=LINE)
    trend_status = _clean(payload.get("trend_status"), "No measured history")
    _add_text(slide, 8.87, 3.28, 3.65, 0.25, trend_status, size=12.5, color=NAVY, bold=True)
    _add_text(slide, 8.87, 3.75, 3.62, 0.60, _clean(payload.get("context_headline"), "No trend context is available."), size=8.0, color=MUTED)
    _add_text(slide, 8.87, 4.63, 1.30, 0.16, "Current vs baseline", size=6.2, color=FAINT)
    _add_text(slide, 10.47, 4.60, 1.95, 0.22, f"{_fmt_pct(payload.get('overall_score'))} / {_fmt_pct(payload.get('baseline_score'))}", size=9.0, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    _add_text(slide, 8.87, 5.08, 1.30, 0.16, "Movement", size=6.2, color=FAINT)
    _add_text(slide, 10.47, 5.05, 1.95, 0.22, _fmt_mom(payload.get("movement")), size=9.0, color=_movement_color(payload.get("movement")), bold=True, align=PP_ALIGN.RIGHT)
    _add_text(slide, 8.87, 5.47, 3.62, 0.16, "The chart uses every valid period returned for the selected scope.", size=6.0, color=BLUE, italic=True)
    _add_footer(slide, _period(payload), page)


def _slide_kpi_health(prs: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]] | None = None, part: int = 1, total_parts: int = 1) -> None:
    slide = _new_slide(prs)
    kpis = rows if rows is not None else _kpi_rows(payload)
    suffix = f" | Part {part}/{total_parts}" if total_parts > 1 else ""
    _add_header(slide, f"KPI Health Overview{suffix}", "Every KPI uses the same Current -> Baseline -> MoM -> Target -> Status grammar.", _period(payload), page)
    if not kpis:
        _add_card(slide, 0.55, 1.25, 12.25, 1.0, fill=PALE_NEUTRAL, rail=BLUE)
        _add_text(slide, 0.88, 1.57, 11.5, 0.23, "No KPI rows are available for the selected Marketing scope.", size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
        _add_footer(slide, _period(payload), page)
        return
    columns, rows_per_column = 3, 3
    card_w = 3.98
    card_h = 1.67
    for index, kpi in enumerate(kpis[:9]):
        col = index % columns
        row = index // columns
        x = 0.43 + col * 4.20
        y = 1.12 + row * 1.88
        status = _status_label_for_kpi(kpi)
        rail = _status_color(status)
        _add_card(slide, x, y, card_w, card_h, rail=rail, line=LINE)
        _add_text(slide, x + 0.20, y + 0.12, 3.42, 0.25, _fit_title(kpi.get("label"), 42), size=8.0, color=NAVY, bold=True)
        _status_badge(slide, x + 2.77, y + 0.13, 1.00, 0.22, status)
        labels = ["Current", "Baseline", "MoM", "Target"]
        values = [_fmt_native(kpi.get("actual"), kpi.get("unit")), _fmt_native(kpi.get("baseline_actual"), kpi.get("unit")), _fmt_mom(kpi.get("mom")), _fmt_native(kpi.get("target"), kpi.get("unit"))]
        for value_index, (label, value) in enumerate(zip(labels, values)):
            bx = x + 0.20 + value_index * 0.88
            _add_text(slide, bx, y + 0.54, 0.82, 0.14, label, size=5.1, color=FAINT, bold=True, align=PP_ALIGN.CENTER)
            _add_text(slide, bx, y + 0.72, 0.82, 0.30, value, size=8.8 if value_index != 2 else 8.2, color=_movement_color(kpi.get("mom")) if value_index == 2 else NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_rect(slide, x + 0.20, y + 1.18, 3.40, 0.02, LINE, LINE)
        _add_text(slide, x + 0.20, y + 1.29, 1.72, 0.16, f"Achievement {_fmt_pct(kpi.get('achievement_pct'))}", size=5.8, color=MUTED)
        _add_text(slide, x + 1.96, y + 1.29, 1.64, 0.16, f"Loss {_fmt_pct(kpi.get('weighted_impact'))}", size=5.8, color=RED if _kpi_loss(kpi) > 0 else GREEN, bold=True, align=PP_ALIGN.RIGHT)
        _add_text(slide, x + 0.20, y + 1.48, 3.40, 0.13, _direction_label(kpi.get("direction")), size=5.6, color=FAINT, italic=True)
    if len(kpis) > 9:
        _add_text(slide, 0.50, 6.90, 12.2, 0.16, f"{len(kpis) - 9} additional KPI row(s) continue on the next health page.", size=6.0, color=AMBER, italic=True)
    _add_footer(slide, _period(payload), page)


def _slide_driver_impact(prs: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]] | None = None, part: int = 1, total_parts: int = 1) -> None:
    slide = _new_slide(prs)
    kpis = rows if rows is not None else _kpi_rows(payload)
    suffix = f" | Part {part}/{total_parts}" if total_parts > 1 else ""
    _add_header(slide, f"Performance Driver Impact{suffix}", "Rank KPI issues by weighted performance loss so management sees the biggest levers first.", _period(payload), page)
    if not kpis:
        _add_text(slide, 0.60, 1.55, 12.0, 0.25, "No measurable KPI drivers are available.", size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
        _add_footer(slide, _period(payload), page)
        return
    max_loss = max([_kpi_loss(kpi) for kpi in kpis] + [1.0])
    _add_section_title(slide, 0.50, 1.08, 12.3, "Weighted loss ranking", accent=RED)
    for index, kpi in enumerate(kpis[:10]):
        y = 1.56 + index * 0.48
        loss = _kpi_loss(kpi)
        status = _status_label_for_kpi(kpi)
        _add_text(slide, 0.58, y, 2.58, 0.22, f"{index + 1:02d}  {_clean(kpi.get('label'), 'KPI')}", size=6.7, color=NAVY, bold=True)
        _add_bar(slide, 3.28, y + 0.035, 5.65, 0.17, loss / max_loss if max_loss else 0, RED if loss > 0 else GREEN)
        _add_text(slide, 9.08, y, 0.85, 0.22, _fmt_pct(loss), size=7.0, color=RED if loss > 0 else GREEN, bold=True, align=PP_ALIGN.RIGHT)
        _add_text(slide, 10.03, y, 0.86, 0.22, str(kpi.get("affected_count") if kpi.get("affected_count") is not None else "-"), size=6.6, color=MUTED, align=PP_ALIGN.CENTER)
        _status_badge(slide, 10.98, y - 0.02, 1.62, 0.26, status)
    _add_card(slide, 0.55, 6.45, 12.2, 0.52, fill=PALE_RED, rail=RED, line=LINE)
    driver = payload.get("driver") or {}
    _add_text(slide, 0.83, 6.62, 11.6, 0.18, f"Read first: {_clean(driver.get('label'), 'No leading driver')} | {_clean(payload.get('driver_story'), 'No driver story is available.')}", size=6.8, color=NAVY, bold=True)
    _add_footer(slide, _period(payload), page)


def _slide_driver_map(prs: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]] | None = None, part: int = 1, total_parts: int = 1) -> None:
    slide = _new_slide(prs)
    driver_rows = rows if rows is not None else _driver_role_employee_rows(payload)
    suffix = f" | Part {part}/{total_parts}" if total_parts > 1 else ""
    _add_header(slide, f"Driver -> Role -> Employee{suffix}", "Connect the KPI signal to the role and people who need management attention.", _period(payload), page)
    _add_section_title(slide, 0.48, 1.08, 12.35, "Performance driver bridge", accent=TEAL)
    table_rows = [
        {
            "driver": _clean(row.get("driver"), "KPI"),
            "loss": _fmt_pct(row.get("loss")),
            "role": _clean(row.get("role"), "Selected role(s)"),
            "employees": _clean(row.get("employees"), "Employee-level detail unavailable"),
            "action": _clean(row.get("action"), "Validate the operational cause and assign corrective action"),
        }
        for row in driver_rows[:8]
    ]
    if table_rows:
        _add_table(
            slide,
            0.48,
            1.55,
            12.35,
            [("driver", "Performance driver", 1.55), ("loss", "Loss", 0.62), ("role", "Affected role", 1.45), ("employees", "Affected employees", 2.35), ("action", "Required management action", 2.65)],
            table_rows,
            row_height=0.57,
            header_height=0.40,
            font_size=6.2,
            header_size=6.1,
            bold_keys=("driver", "loss"),
            center_keys=("loss",),
            row_fill=lambda row, index: PALE_RED if index == 0 else PALE_NEUTRAL if index % 2 else WHITE,
        )
    else:
        _add_card(slide, 0.55, 1.55, 12.2, 1.0, fill=PALE_NEUTRAL, rail=BLUE)
        _add_text(slide, 0.85, 1.88, 11.6, 0.24, "No employee-level driver bridge is available for the selected scope.", size=9.0, color=MUTED, align=PP_ALIGN.CENTER)
    _add_card(slide, 0.55, 6.55, 12.2, 0.60, fill=PALE_BLUE, rail=BLUE, line=LINE)
    _add_text(slide, 0.84, 6.75, 11.6, 0.18, "The bridge preserves the distinction between a measured KPI signal and a confirmed operational cause.", size=6.7, color=NAVY)
    _add_footer(slide, _period(payload), page)


def _slide_employee_priority(prs: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _new_slide(prs)
    _add_header(slide, "Employee Performance Priority", "Separate critical intervention, watch-list coaching, and people who are on track.", _period(payload), page)
    groups = _employee_priority_groups(payload)
    titles = [("Requires Action", RED, PALE_RED), ("Watch", AMBER, PALE_AMBER), ("On Track", GREEN, PALE_GREEN)]
    for column, (title, color, fill) in enumerate(titles):
        x = 0.48 + column * 4.28
        _add_rect(slide, x, 1.12, 3.98, 0.38, color, color)
        _add_text(slide, x + 0.16, 1.22, 3.65, 0.17, f"{title.upper()}  ({len(groups[title])})", size=7.1, color=WHITE, bold=True)
        members = groups[title]
        if not members:
            _add_card(slide, x, 1.65, 3.98, 0.75, fill=PALE_NEUTRAL, rail=BLUE)
            _add_text(slide, x + 0.22, 1.90, 3.52, 0.20, "No employees in this group.", size=7.0, color=MUTED, align=PP_ALIGN.CENTER)
        for index, person in enumerate(members[:7]):
            y = 1.65 + index * 0.68
            primary = _primary_kpi(person)
            _add_card(slide, x, y, 3.98, 0.55, fill=fill if title != "On Track" else WHITE, rail=color, line=LINE)
            _add_text(slide, x + 0.21, y + 0.08, 2.35, 0.19, _fit_title(person.get("name"), 26), size=7.0, color=NAVY, bold=True)
            _add_text(slide, x + 2.65, y + 0.08, 1.04, 0.19, _fmt_pct(person.get("score")), size=7.4, color=color, bold=True, align=PP_ALIGN.RIGHT)
            _add_text(slide, x + 0.21, y + 0.31, 2.70, 0.13, _fit_title(primary.get("label"), 37) if primary else "No leading KPI loss", size=5.6, color=MUTED)
            _add_text(slide, x + 2.93, y + 0.31, 0.76, 0.13, f"Loss {_fmt_pct(_kpi_loss(primary))}", size=5.5, color=RED if _kpi_loss(primary) > 0 else GREEN, align=PP_ALIGN.RIGHT)
        if len(members) > 7:
            _add_text(slide, x + 0.2, 6.65, 3.6, 0.16, f"+{len(members) - 7} shown in appendix", size=5.8, color=FAINT, italic=True)
    _add_card(slide, 0.55, 6.91, 12.2, 0.22, fill=BODY, rail=TEAL, line=BODY, rounded=False)
    _add_text(slide, 0.80, 6.94, 11.8, 0.12, "Priority is ordered by score, then weighted loss. High performers remain visible without competing with intervention items.", size=5.8, color=MUTED)
    _add_footer(slide, _period(payload), page)


def _employee_evidence_label(person: dict[str, Any]) -> tuple[str, str]:
    state = _clean(person.get("root_cause_status"), "KPI signal only")
    if "confirmed" in state.casefold():
        return "Confirmed Root Cause", _clean(person.get("root_cause"), "Confirmed cause recorded")
    if "evidence" in state.casefold() or _clean(person.get("root_cause")):
        return "KPI Evidence", _clean(person.get("root_cause"), "Evidence recorded; cause requires confirmation")
    return "Investigation Required", "The KPI signal identifies where to investigate; an operational cause is not recorded."


def _add_profile_card(slide: Any, payload: dict[str, Any], person: dict[str, Any], x: float, y: float, width: float, height: float) -> None:
    status = _employee_status(person.get("score"))
    _add_card(slide, x, y, width, height, fill=WHITE, rail=_status_color(status), line=LINE)
    _add_text(slide, x + 0.22, y + 0.22, width - 0.42, 0.52, _clean(person.get("name"), "Unknown employee"), size=16.0, color=NAVY, bold=True)
    _add_text(slide, x + 0.22, y + 0.83, width - 0.42, 0.24, _clean(person.get("position"), "Position not available"), size=8.0, color=MUTED)
    _add_text(slide, x + 0.22, y + 1.32, width - 0.42, 0.16, "CURRENT SCORE", size=6.0, color=FAINT, bold=True)
    _add_text(slide, x + 0.22, y + 1.53, width - 0.42, 0.46, _fmt_pct(person.get("score")), size=24.0, color=_status_color(status), bold=True)
    _status_badge(slide, x + 0.22, y + 2.12, width - 0.44, 0.30, status)
    _add_text(slide, x + 0.22, y + 2.62, 1.18, 0.15, "Baseline", size=6.0, color=FAINT)
    _add_text(slide, x + 1.56, y + 2.59, width - 1.78, 0.19, _fmt_pct(person.get("baseline_score")), size=8.0, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    _add_text(slide, x + 0.22, y + 2.97, 1.18, 0.15, "MoM", size=6.0, color=FAINT)
    _add_text(slide, x + 1.56, y + 2.94, width - 1.78, 0.19, _fmt_mom(person.get("movement")), size=8.0, color=_movement_color(person.get("movement")), bold=True, align=PP_ALIGN.RIGHT)
    _add_text(slide, x + 0.22, y + 3.32, 1.18, 0.15, "Priority / rank", size=6.0, color=FAINT)
    _add_text(slide, x + 1.56, y + 3.29, width - 1.78, 0.19, f"P{_priority_rank(payload, person)}", size=8.0, color=RED if status == "Requires action" else AMBER, bold=True, align=PP_ALIGN.RIGHT)
    _add_rect(slide, x + 0.22, y + 3.70, width - 0.44, 0.02, LINE, LINE)
    evidence_label, evidence_text = _employee_evidence_label(person)
    _add_text(slide, x + 0.22, y + 3.92, width - 0.44, 0.18, evidence_label, size=6.2, color=_status_color(evidence_label), bold=True)
    _add_text(slide, x + 0.22, y + 4.19, width - 0.44, 0.75, _wrap(evidence_text, 40), size=6.3, color=MUTED)
    actions = person.get("actions") or []
    action_text = _action_text(actions[0]) if actions else "No action recorded; assign a workstream if intervention is required."
    _add_text(slide, x + 0.22, y + 5.15, width - 0.44, 0.18, "ACTION LINK", size=6.0, color=TEAL, bold=True)
    _add_text(slide, x + 0.22, y + 5.42, width - 0.44, 0.43, _wrap(action_text, 40), size=6.2, color=NAVY)


def _add_detail_kpi_card(slide: Any, kpi: dict[str, Any], x: float, y: float, width: float, height: float) -> None:
    status = _status_label_for_kpi(kpi)
    loss = _kpi_loss(kpi)
    rail = RED if loss > 0 else GREEN
    _add_card(slide, x, y, width, height, fill=PALE_RED if loss > 0 else WHITE, rail=rail, line=LINE)
    _add_text(slide, x + 0.17, y + 0.14, width - 1.30, 0.27, _fit_title(kpi.get("label"), 34), size=8.0, color=NAVY, bold=True)
    _status_badge(slide, x + width - 1.07, y + 0.13, 0.88, 0.24, status)
    rows = [
        ("Actual", _fmt_native(kpi.get("actual"), kpi.get("unit"))),
        ("Target", _fmt_native(kpi.get("target"), kpi.get("unit"))),
        ("Achievement", _fmt_pct(kpi.get("achievement_pct"))),
        ("Achievement gap", _fmt_pct(kpi.get("shortfall"))),
        ("Weighted loss", _fmt_pct(loss)),
    ]
    for index, (label, value) in enumerate(rows):
        col = index % 2
        row = index // 2
        bx = x + 0.17 + col * (width - 0.34) / 2
        by = y + 0.55 + row * 0.39
        _add_text(slide, bx, by, (width - 0.45) / 2, 0.13, label, size=5.1, color=FAINT, bold=True)
        _add_text(slide, bx, by + 0.14, (width - 0.45) / 2, 0.16, value, size=7.0, color=RED if label in {"Achievement gap", "Weighted loss"} and loss > 0 else NAVY, bold=True)
    _add_rect(slide, x + 0.17, y + height - 0.38, width - 0.34, 0.02, LINE, LINE)
    _add_text(slide, x + 0.17, y + height - 0.28, width - 0.34, 0.13, _direction_label(kpi.get("direction")), size=5.7, color=FAINT, italic=True)


def _slide_detail(prs: Presentation, payload: dict[str, Any], person: dict[str, Any], page: int, kpis: list[dict[str, Any]], part: int, total_parts: int) -> None:
    slide = _new_slide(prs)
    name = _clean(person.get("name"), "Employee detail")
    suffix = f" | Part {part}/{total_parts}" if total_parts > 1 else ""
    _add_header(slide, f"Employee Detail: {_fit_title(name, 45)}{suffix}", "Current score, baseline, movement, and the KPI losses behind the priority ranking.", _period(payload), page)
    _add_profile_card(slide, payload, person, 0.48, 1.08, 3.08, 5.85)
    _add_section_title(slide, 3.86, 1.08, 8.96, "KPI drivers", accent=RED)
    if not kpis:
        _add_card(slide, 3.86, 1.55, 8.96, 1.0, fill=PALE_NEUTRAL, rail=BLUE)
        _add_text(slide, 4.16, 1.87, 8.35, 0.20, "No employee KPI detail is available for this record.", size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    else:
        card_w, card_h = 2.83, 2.20
        for index, kpi in enumerate(kpis[:6]):
            col = index % 3
            row = index // 3
            _add_detail_kpi_card(slide, kpi, 3.86 + col * 3.04, 1.55 + row * 2.48, card_w, card_h)
        if len(kpis) > 6:
            _add_text(slide, 3.88, 6.62, 8.8, 0.15, f"{len(kpis) - 6} additional KPI driver(s) continue on the next detail page.", size=5.8, color=AMBER, italic=True)
    _add_footer(slide, _period(payload), page)


def _slide_evidence(prs: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]] | None = None, part: int = 1, total_parts: int = 1) -> None:
    slide = _new_slide(prs)
    evidence_rows = rows if rows is not None else _evidence_rows(payload)
    suffix = f" | Part {part}/{total_parts}" if total_parts > 1 else ""
    _add_header(slide, f"Root Cause / Evidence{suffix}", "Separate KPI evidence, confirmed root cause, and investigation required.", _period(payload), page)
    _add_card(slide, 0.50, 1.06, 12.30, 0.56, fill=PALE_BLUE, rail=BLUE, line=LINE)
    _add_text(slide, 0.80, 1.20, 11.7, 0.25, "A performance signal is not a root cause. States used: KPI Evidence | Confirmed Root Cause | Investigation Required.", size=6.4, color=NAVY)
    table_rows = [
        {
            "driver": _clean(row.get("driver"), "KPI"),
            "loss": _fmt_pct(row.get("loss")),
            "evidence": row.get("evidence"),
            "recorded": _clean(row.get("recorded"), "No operational cause recorded"),
            "validation": _clean(row.get("validation"), "Validate the operational cause before closing the gap"),
            "owner": _clean(row.get("owner"), "Owner needed"),
        }
        for row in evidence_rows[:7]
    ]
    if table_rows:
        _add_table(
            slide,
            0.50,
            1.86,
            12.30,
            [("driver", "KPI evidence", 1.35), ("loss", "Loss", 0.52), ("evidence", "State", 1.20), ("recorded", "Recorded root cause", 2.25), ("validation", "Investigation / validation", 2.70), ("owner", "Owner", 1.05)],
            table_rows,
            row_height=0.80,
            header_height=0.40,
            font_size=5.5,
            header_size=5.8,
            bold_keys=("driver", "evidence"),
            center_keys=("loss",),
            row_fill=lambda row, index: PALE_RED if _clean(row.get("evidence")).casefold() == "investigation required" else PALE_AMBER if _clean(row.get("evidence")).casefold() == "kpi evidence" else PALE_GREEN,
        )
    else:
        _add_card(slide, 0.55, 1.86, 12.2, 0.95, fill=PALE_NEUTRAL, rail=BLUE)
        _add_text(slide, 0.85, 2.17, 11.6, 0.20, "No root-cause evidence rows are available for the current KPI set.", size=8.4, color=MUTED, align=PP_ALIGN.CENTER)
    _add_footer(slide, _period(payload), page)


def _slide_actions(prs: Presentation, payload: dict[str, Any], page: int, actions: list[dict[str, Any]] | None = None, part: int = 1, total_parts: int = 1) -> None:
    slide = _new_slide(prs)
    action_rows = actions if actions is not None else _management_actions(payload)
    suffix = f" | Part {part}/{total_parts}" if total_parts > 1 else ""
    _add_header(slide, f"Corrective Action Tracker{suffix}", "Recorded and proposed workstreams remain explicit, including missing owners and due dates.", _period(payload), page)
    table_rows = [
        {
            "issue": _clean(_action_issue(action), "Management workstream"),
            "action": _clean(_action_text(action), "Action text needed"),
            "owner": _action_owner(action),
            "due": _action_due(action),
            "status": _action_status(action),
            "success": _clean(_action_success(action), "Success metric needed"),
        }
        for action in action_rows[:6]
    ]
    if table_rows:
        _add_table(
            slide,
            0.47,
            1.28,
            12.36,
            [("issue", "Issue / workstream", 1.25), ("action", "Action", 2.45), ("owner", "Owner", 1.00), ("due", "Due date", 0.95), ("status", "Status", 0.95), ("success", "Success measure", 2.20)],
            table_rows,
            row_height=0.78,
            header_height=0.42,
            font_size=5.9,
            header_size=5.9,
            bold_keys=("issue", "status"),
            center_keys=("owner", "due", "status"),
            row_fill=lambda row, index: PALE_AMBER if _clean(row.get("status")).casefold() == "proposed" else PALE_RED if "needed" in f"{row.get('owner')} {row.get('due')}".casefold() else (PALE_NEUTRAL if index % 2 else WHITE),
        )
    else:
        _add_card(slide, 0.55, 1.28, 12.2, 0.95, fill=PALE_NEUTRAL, rail=BLUE)
        _add_text(slide, 0.85, 1.60, 11.6, 0.20, "No recorded or proposed corrective actions are available.", size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    missing_owner = sum(1 for action in action_rows if _action_owner(action).casefold() == "owner needed")
    missing_due = sum(1 for action in action_rows if _action_due(action).casefold() == "due date needed")
    _add_card(slide, 0.55, 6.52, 12.2, 0.43, fill=PALE_RED if missing_owner or missing_due else PALE_GREEN, rail=RED if missing_owner or missing_due else GREEN, line=LINE)
    _add_text(slide, 0.82, 6.65, 11.7, 0.16, f"Open management items: {len(action_rows)} | Missing owners: {missing_owner} | Missing due dates: {missing_due}", size=6.6, color=NAVY, bold=True)
    _add_footer(slide, _period(payload), page)


def _slide_management(prs: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _new_slide(prs)
    _add_header(slide, "Management Summary / Next Review", "Close the story with the score, the decisions required, and the evidence expected next cycle.", _period(payload), page)
    actions = _management_actions(payload)
    missing_owner = sum(1 for action in actions if _action_owner(action).casefold() == "owner needed")
    cards = [
        ("Current overall score", _fmt_pct(payload.get("overall_score")), TEAL, NAVY),
        ("Target", _fmt_pct(payload.get("target_score")), NAVY, NAVY),
        ("Important KPI gaps", str(sum(1 for kpi in _kpi_rows(payload) if _kpi_loss(kpi) > 0)), RED, RED),
        ("Employees requiring action", str(len(_employee_priority_groups(payload)["Requires Action"])), RED, RED),
        ("Open actions", str(len(actions)), AMBER if actions else GREEN, NAVY),
        ("Missing owners", str(missing_owner), RED if missing_owner else GREEN, RED if missing_owner else GREEN),
    ]
    gap = 0.12
    width = (12.45 - gap * 5) / 6
    for index, (label, value, rail, color) in enumerate(cards):
        _metric_card(slide, 0.43 + index * (width + gap), 1.08, width, 1.03, label, value, rail=rail, value_color=color, value_size=13.0 if len(str(value)) > 7 else 17.0)
    review = payload.get("next_review") or {}
    _add_section_title(slide, 0.50, 2.48, 6.0, "Next review commitments", accent=TEAL)
    _add_card(slide, 0.50, 2.92, 6.0, 3.58, fill=PALE_BLUE, rail=TEAL, line=LINE)
    commitment_rows = [
        ("Expected movement", review.get("expected_movement")),
        ("Root-cause requirement", review.get("root_cause_requirement")),
        ("People requirement", review.get("people_requirement")),
        ("Action requirement", review.get("action_requirement")),
        ("Due cadence", review.get("due_cadence")),
    ]
    for index, (label, value) in enumerate(commitment_rows):
        y = 3.22 + index * 0.58
        _add_text(slide, 0.78, y, 1.62, 0.16, label, size=5.9, color=FAINT, bold=True)
        _add_text(slide, 2.48, y - 0.02, 3.72, 0.30, _wrap(value, 47), size=6.3, color=NAVY)
        if index < len(commitment_rows) - 1:
            _add_rect(slide, 0.78, y + 0.38, 5.36, 0.015, LINE, LINE)

    _add_section_title(slide, 6.84, 2.48, 5.98, "Success criteria and decision request", accent=RED)
    _add_card(slide, 6.84, 2.92, 5.98, 3.58, fill=PALE_RED, rail=RED, line=LINE)
    _add_text(slide, 7.15, 3.24, 5.35, 0.16, "Success evidence", size=6.1, color=RED, bold=True)
    _add_text(slide, 7.15, 3.49, 5.30, 0.68, _wrap(review.get("success_evidence"), 60), size=8.0, color=NAVY, bold=True)
    _add_text(slide, 7.15, 4.48, 5.35, 0.16, "Management decision", size=6.1, color=RED, bold=True)
    _add_text(slide, 7.15, 4.74, 5.30, 0.65, _wrap(review.get("decision_request"), 60), size=8.0, color=NAVY, bold=True)
    _add_text(slide, 7.15, 5.70, 5.30, 0.16, "Next-cycle question", size=6.1, color=RED, bold=True)
    _add_text(slide, 7.15, 5.95, 5.30, 0.32, _wrap(review.get("question"), 60), size=7.0, color=MUTED)
    _add_footer(slide, _period(payload), page)


def _slide_employee_appendix(prs: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]] | None = None, part: int = 1, total_parts: int = 1) -> None:
    slide = _new_slide(prs)
    people = rows if rows is not None else _people(payload)
    suffix = f" | Part {part}/{total_parts}" if total_parts > 1 else ""
    _add_header(slide, f"Appendix: Full Employee Scorecard{suffix}", "Reference table for questions during the presentation; the primary storyline appears earlier.", _period(payload), page)
    table_rows = []
    for person in people[:12]:
        primary = _primary_kpi(person)
        actions = person.get("actions") or []
        table_rows.append(
            {
                "employee": _fit_title(person.get("name"), 26),
                "role": _fit_title(person.get("position"), 24),
                "current": _fmt_pct(person.get("score")),
                "baseline": _fmt_pct(person.get("baseline_score")),
                "mom": _fmt_mom(person.get("movement")),
                "status": _employee_status(person.get("score")),
                "driver": _clean(primary.get("label"), "No leading KPI") if primary else "No leading KPI",
                "actual": _fmt_native(primary.get("actual"), primary.get("unit")) if primary else "Not available",
                "target": _fmt_native(primary.get("target"), primary.get("unit")) if primary else "Not available",
                "loss": _fmt_pct(_kpi_loss(primary)),
                "action": _clean((actions[0] if actions else _kpi_action(payload, primary)).get("source_display"), "No action") if (actions or primary) else "No action",
            }
        )
    if table_rows:
        _add_table(
            slide,
            0.38,
            1.22,
            12.58,
            [("employee", "Employee", 1.45), ("role", "Position", 1.35), ("current", "Current", 0.65), ("baseline", "Baseline", 0.68), ("mom", "MoM", 0.60), ("status", "Status", 0.90), ("driver", "Leading KPI", 1.45), ("actual", "Actual", 0.66), ("target", "Target", 0.66), ("loss", "Loss", 0.58), ("action", "Action ref", 0.85)],
            table_rows,
            row_height=0.42,
            header_height=0.41,
            font_size=5.5,
            header_size=5.5,
            bold_keys=("employee", "current", "loss"),
            center_keys=("current", "baseline", "mom", "status", "actual", "target", "loss", "action"),
            row_fill=lambda row, index: PALE_RED if _clean(row.get("status")).casefold() == "requires action" else PALE_AMBER if _clean(row.get("status")).casefold() == "watch" else (PALE_NEUTRAL if index % 2 else WHITE),
        )
    else:
        _add_card(slide, 0.55, 1.22, 12.2, 0.95, fill=PALE_NEUTRAL, rail=BLUE)
        _add_text(slide, 0.85, 1.54, 11.6, 0.20, "No employee scorecard rows are available for the selected scope.", size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    _add_footer(slide, _period(payload), page)


def _slide_kpi_appendix(prs: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]] | None = None, part: int = 1, total_parts: int = 1) -> None:
    slide = _new_slide(prs)
    kpis = rows if rows is not None else _kpi_rows(payload)
    suffix = f" | Part {part}/{total_parts}" if total_parts > 1 else ""
    _add_header(slide, f"Appendix: Full KPI Reference{suffix}", "Actual, target, MoM, direction, status, weighted loss, and action reference.", _period(payload), page)
    table_rows = []
    for kpi in kpis[:12]:
        action = _kpi_action(payload, kpi)
        table_rows.append(
            {
                "kpi": _clean(kpi.get("label"), "KPI"),
                "actual": _fmt_native(kpi.get("actual"), kpi.get("unit")),
                "baseline": _fmt_native(kpi.get("baseline_actual"), kpi.get("unit")),
                "mom": _fmt_mom(kpi.get("mom")),
                "target": _fmt_native(kpi.get("target"), kpi.get("unit")),
                "achievement": _fmt_pct(kpi.get("achievement_pct")),
                "gap": _fmt_pct(kpi.get("shortfall")),
                "loss": _fmt_pct(_kpi_loss(kpi)),
                "direction": _direction_label(kpi.get("direction")),
                "status": _status_label_for_kpi(kpi),
                "action": _clean(action.get("source_display"), "No action"),
            }
        )
    if table_rows:
        _add_table(
            slide,
            0.32,
            1.20,
            12.66,
            [("kpi", "KPI", 1.55), ("actual", "Actual", 0.70), ("baseline", "Baseline", 0.70), ("mom", "MoM", 0.58), ("target", "Target", 0.70), ("achievement", "Achiev.", 0.68), ("gap", "Gap", 0.58), ("loss", "Weighted loss", 0.75), ("direction", "Direction", 1.15), ("status", "Status", 0.88), ("action", "Action", 0.75)],
            table_rows,
            row_height=0.43,
            header_height=0.41,
            font_size=5.25,
            header_size=5.3,
            bold_keys=("kpi", "loss"),
            center_keys=("actual", "baseline", "mom", "target", "achievement", "gap", "loss", "status", "action"),
            row_fill=lambda row, index: PALE_RED if _clean(row.get("status")).casefold() in {"critical", "requires action"} else PALE_AMBER if _clean(row.get("status")).casefold() in {"at risk", "watch"} else (PALE_NEUTRAL if index % 2 else WHITE),
        )
    else:
        _add_card(slide, 0.55, 1.20, 12.2, 0.95, fill=PALE_NEUTRAL, rail=BLUE)
        _add_text(slide, 0.85, 1.52, 11.6, 0.20, "No KPI reference rows are available for the selected scope.", size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    _add_footer(slide, _period(payload), page)


def build_marketing_legendary_pptx(period_label: str = "June 2026", report_data: dict[str, Any] | None = None) -> bytes:
    """Build the approved dynamic Marketing report without a July template dependency."""

    payload = _prepare_snapshot(period_label, report_data)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    prs.core_properties.title = f"Marketing Team Performance Review - {_period(payload)}"
    prs.core_properties.subject = "Marketing PMS performance review"
    prs.core_properties.author = "PMS Dashboard"

    page = 1
    _slide_cover(prs, payload, page)
    page += 1
    _slide_overview(prs, payload, page)
    page += 1
    _slide_roles(prs, payload, page)
    page += 1
    _slide_trend(prs, payload, page)
    page += 1

    kpis = _kpi_rows(payload)
    kpi_health_parts = _chunks(kpis, 9)
    for index, chunk in enumerate(kpi_health_parts, 1):
        _slide_kpi_health(prs, payload, page, chunk, index, len(kpi_health_parts))
        page += 1

    driver_parts = _chunks(kpis, 10)
    for index, chunk in enumerate(driver_parts, 1):
        _slide_driver_impact(prs, payload, page, chunk, index, len(driver_parts))
        page += 1

    bridge_rows = _driver_role_employee_rows(payload)
    bridge_parts = _chunks(bridge_rows, 8)
    for index, chunk in enumerate(bridge_parts, 1):
        _slide_driver_map(prs, payload, page, chunk, index, len(bridge_parts))
        page += 1

    _slide_employee_priority(prs, payload, page)
    page += 1

    detail_people = _detail_people(payload)
    for person in detail_people:
        detail_kpis = _person_kpi_rows(person)
        detail_parts = _chunks(detail_kpis, 6)
        for index, chunk in enumerate(detail_parts, 1):
            _slide_detail(prs, payload, person, page, chunk, index, len(detail_parts))
            page += 1

    evidence_rows = _evidence_rows(payload)
    for index, chunk in enumerate(_chunks(evidence_rows, 7), 1):
        _slide_evidence(prs, payload, page, chunk, index, len(_chunks(evidence_rows, 7)))
        page += 1

    actions = _management_actions(payload)
    action_parts = _chunks(actions, 6)
    for index, chunk in enumerate(action_parts, 1):
        _slide_actions(prs, payload, page, chunk, index, len(action_parts))
        page += 1

    _slide_management(prs, payload, page)
    page += 1

    people = _people(payload)
    employee_parts = _chunks(people, 12)
    for index, chunk in enumerate(employee_parts, 1):
        _slide_employee_appendix(prs, payload, page, chunk, index, len(employee_parts))
        page += 1

    kpi_parts = _chunks(kpis, 12)
    for index, chunk in enumerate(kpi_parts, 1):
        _slide_kpi_appendix(prs, payload, page, chunk, index, len(kpi_parts))
        page += 1

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


__all__ = ["build_marketing_legendary_pptx"]


# Keep the existing public import stable for the report service and callers,
# while the Marketing export itself is now filled from the approved reference
# deck.  The import is intentionally lazy so the snapshot helpers above stay
# importable without creating a module cycle.
_legacy_marketing_legendary_builder = build_marketing_legendary_pptx


def build_marketing_legendary_pptx(period_label: str = "June 2026", report_data: dict[str, Any] | None = None) -> bytes:
    from exports.marketing_reference_pptx_builder import build_marketing_reference_pptx

    return build_marketing_reference_pptx(period_label, report_data)
