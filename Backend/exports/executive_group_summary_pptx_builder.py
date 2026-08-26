"""CEO / group-level performance summary using the approved Offshore visual system.

This builder is intentionally separate from the Marketing and Offshore monthly
builders.  It reuses the approved Offshore source frames and semantic colors,
but its payload is an aggregate story: authorized regions, teams, KPI drivers,
evidence state, commitments, and next-review decisions.  It never needs to
render employee names for the executive audience.
"""

from __future__ import annotations

import io
from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from exports import offshore_status_pptx_builder as reference


TEMPLATE_PATH = reference.TEMPLATE_PATH
REFERENCE_SLIDES = reference.REFERENCE_SLIDES
MISSING = reference.MISSING


def _number(value: Any) -> float | None:
    try:
        result = float(value)
    except (TypeError, ValueError):
        return None
    return result if result == result else None


def _clean(value: Any, fallback: str = MISSING) -> str:
    if value is None:
        return fallback
    text = " ".join(str(value).replace("\n", " ").split()).strip()
    return text or fallback


def _pct(value: Any, *, signed: bool = False) -> str:
    number = _number(value)
    if number is None:
        return MISSING
    return f"{number:+.1f}%" if signed else f"{number:.1f}%"


def _mom(value: Any) -> str:
    return "NEW" if _number(value) is None else _pct(value, signed=True)


def _score_status(score: Any) -> str:
    value = _number(score)
    if value is None:
        return "Data unavailable"
    if value < 70:
        return "Requires action"
    if value < 90:
        return "Watch"
    return "On track"


def _kpi_status(kpi: dict[str, Any] | None) -> str:
    achievement = _number((kpi or {}).get("achievement_pct"))
    if achievement is None:
        return "Data unavailable"
    if achievement >= 100:
        return "On track"
    if achievement >= 70:
        return "Watch"
    return "Requires action"


def _kpi_loss(kpi: dict[str, Any] | None) -> float:
    return max(0.0, _number((kpi or {}).get("weighted_impact")) or 0.0)


def _fit(value: Any, limit: int) -> str:
    return reference._fit(value, limit)


def _fit_wrapped(value: Any, limit: int, max_lines: int = 2) -> str:
    return reference._fit_wrapped(value, limit, max_lines=max_lines)


def _native(value: Any, unit: Any = "") -> str:
    return reference._native(value, unit)


def _set_common(slide: Any, title: str, subtitle: str, period: str, scope: str, page: int) -> None:
    reference._set_common(slide, title, subtitle, period, scope, page)


def _status_color(status: Any):
    return reference._status_color(status)


def _status_label(status: Any) -> str:
    return reference._status_label(status)


def _write(slide: Any, name: str, value: Any, **kwargs: Any) -> None:
    reference._write(slide, name, value, **kwargs)


def _write_single_line(slide: Any, name: str, value: Any) -> None:
    reference._write_single_line(slide, name, value)


def _write_shape(shape: Any, value: Any, **kwargs: Any) -> None:
    reference._write_shape(shape, value, **kwargs)


def _shape(slide: Any, name: str) -> Any:
    return reference._shape(slide, name)


def _text_shape(slide: Any, name: str) -> Any:
    return reference._text_shape(slide, name)


def _row_groups(slide: Any, *, top_min: float, top_max: float, min_columns: int) -> list[list[Any]]:
    return reference._row_groups(slide, top_min=top_min, top_max=top_max, min_columns=min_columns)


def _card_status(slide: Any, card: str, accent: str, status: Any, *, card_fill: Any = None) -> None:
    reference._card_status(slide, card, accent, status, card_fill=card_fill)


def _resize_bar(slide: Any, name: str, ratio: Any, *, color: Any = None) -> None:
    reference._resize_bar(slide, name, ratio, color=color)


def _append_template_slide(prs: Presentation, source: Presentation, source_index: int) -> Any:
    return reference._append_template_slide(prs, source, source_index)


def _clear_slides(prs: Presentation) -> None:
    reference._clear_slides(prs)


def _group_view(row: dict[str, Any]) -> dict[str, Any]:
    return {
        "name": row.get("name"),
        "score": row.get("score"),
        "baseline": row.get("baseline"),
        "movement": row.get("movement"),
        "target": row.get("target", 100.0),
        "status": row.get("status") or _score_status(row.get("score")),
        "people": [row.get("scope_label") or row.get("name")],
        "affected_count": row.get("below_threshold", 0),
        "kpis": list(row.get("kpis") or []),
        "snapshot": {"all_people": [row], "kpis": list(row.get("kpis") or [])},
    }


def _view_rows(payload: dict[str, Any]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    regions = [dict(row) for row in payload.get("regions") or []]
    teams = [dict(row) for row in payload.get("teams") or []]
    return regions, teams


def _current_view(payload: dict[str, Any]) -> dict[str, Any]:
    overall = dict(payload.get("overall") or {})
    overall.setdefault("name", "Authorized Group Scope")
    overall.setdefault("scope_label", payload.get("scope_label") or "Authorized scope")
    overall.setdefault("status", _score_status(overall.get("score")))
    overall.setdefault("target", 100.0)
    return overall


def _write_score_card(slide: Any, names: tuple[str, ...], row: dict[str, Any], *, title: str | None = None) -> None:
    label, score_name, monthly_name, baseline_name, mom_name, target_name, status_name = names
    status = row.get("status") or _score_status(row.get("score"))
    _write_single_line(slide, label, title or row.get("name"))
    _write(slide, score_name, _pct(row.get("score")))
    _write(slide, monthly_name, "Current performance")
    _write(slide, baseline_name, f"Baseline: {_pct(row.get('baseline'))}")
    _write(slide, mom_name, f"MoM: {_mom(row.get('movement'))}")
    _write(slide, target_name, f"Target: {_pct(row.get('target', 100.0))}")
    _write(slide, status_name, _status_label(status), color=reference.WHITE)


def _leading_kpi(row: dict[str, Any]) -> dict[str, Any]:
    kpis = [dict(kpi) for kpi in row.get("kpis") or [] if isinstance(kpi, dict)]
    return max(kpis, key=lambda item: (_kpi_loss(item), _number(item.get("shortfall")) or 0.0), default={})


def _risk_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    regions, teams = _view_rows(payload)
    rows = []
    # Team risk is the most actionable CEO drilldown.  Region rows are used as
    # a fallback when the authorized dataset contains no team dimension.
    rows.extend(teams or regions)
    return sorted(
        rows,
        key=lambda row: (
            _number(row.get("score")) is None,
            _number(row.get("score")) if _number(row.get("score")) is not None else 999.0,
            -(_number(row.get("weighted_loss")) or 0.0),
            _clean(row.get("name"), "").casefold(),
        ),
    )


def _populate_cover(slide: Any, payload: dict[str, Any], page: int) -> None:
    period = _clean(payload.get("period_label"), "Selected period")
    scope = _clean(payload.get("scope_label"), "Authorized group scope")
    regions, teams = _view_rows(payload)
    actions = payload.get("actions") or []
    _write(slide, "Text 3", "GROUP EXECUTIVE")
    _write(slide, "Text 4", "Performance Review")
    _write(slide, "Text 5", _fit(f"{period}  |  {scope}", 70))
    _write(slide, "Text 6", _fit("Authorized regions and teams  •  CEO / group leadership view", 90))
    _write(slide, "Text 8", "Executive Narrative")
    _write(slide, "Text 9", "Overall performance first, then region and team risk, KPI drivers, evidence, commitments, and the next review decision.")
    _write(slide, "Text 11", f"{len(regions)} Regions")
    _write(slide, "Text 13", f"{len(teams)} Teams")
    _write(slide, "Text 15", f"{len(actions)} Commitments")
    _write(slide, "Text 16", "CONFIDENTIAL – INTERNAL USE ONLY")


def _populate_overview(slide: Any, payload: dict[str, Any], page: int) -> None:
    period = _clean(payload.get("period_label"), "Selected period")
    scope = _clean(payload.get("scope_label"), "Authorized group scope")
    regions, teams = _view_rows(payload)
    overall = _current_view(payload)
    cards = [overall]
    seen_group_names = {str(overall.get("name") or "").casefold()}
    for row in [*regions, *teams]:
        key = str(row.get("name") or "").casefold()
        if key and key not in seen_group_names:
            cards.append(row)
            seen_group_names.add(key)
        if len(cards) >= 4:
            break
    _set_common(slide, "Executive Group Performance Overview", "How the authorized group performed this month", period, scope, page)
    card_specs = [
        (("Text 10", "Text 11", "Text 12", "Text 13", "Text 14", "Text 15", "Text 17"), "Shape 8", "Shape 9"),
        (("Text 20", "Text 21", "Text 22", "Text 23", "Text 24", "Text 25", "Text 27"), "Shape 18", "Shape 19"),
        (("Text 30", "Text 31", "Text 32", "Text 33", "Text 34", "Text 35", "Text 37"), "Shape 28", "Shape 29"),
        (("Text 40", "Text 41", "Text 42", "Text 43", "Text 44", "Text 45", "Text 47"), "Shape 38", "Shape 39"),
    ]
    for index, (names, card, accent) in enumerate(card_specs):
        row = cards[index] if index < len(cards) else {}
        status = row.get("status") or _score_status(row.get("score"))
        _card_status(slide, card, accent, status)
        _write_score_card(slide, names, row, title=_fit(row.get("name"), 25))
    overall_score = overall.get("score")
    leading = dict(payload.get("drivers", [{}])[0] or {}) if payload.get("drivers") else {}
    risk_count = sum(1 for row in teams or regions if _score_status(row.get("score")) == "Requires action")
    takeaway = (
        f"CEO takeaway: the authorized group is {_pct(overall_score)} with {_mom(overall.get('movement'))} MoM; "
        f"{risk_count} group(s) require action and {_clean(leading.get('label'), 'no leading KPI driver')} carries the largest weighted loss."
    )
    _write(slide, "Text 49", _fit(takeaway, 170))
    _write(slide, "Text 50", "Team / Region Action View")
    groups = _row_groups(slide, top_min=5.25, top_max=6.85, min_columns=6)
    for index, group in enumerate(groups[:4]):
        row = (teams or regions)[index] if index < len(teams or regions) else None
        if not row:
            values = [MISSING] * 6
        else:
            driver = _leading_kpi(row)
            gap_text = (
                f"{_clean(driver.get('label'), 'KPI')} {_pct(driver.get('achievement_pct'))} | loss {_pct(driver.get('weighted_impact'))}"
                if driver else "No measured KPI gap"
            )
            values = [
                _clean(row.get("name")),
                _pct(row.get("score")),
                _pct(row.get("baseline")),
                _mom(row.get("movement")),
                _fit_wrapped(gap_text, 40),
                _fit_wrapped(row.get("action_focus") or "Review the highest weighted-loss driver", 42),
            ]
        for column, (shape, value) in enumerate(zip(group, values)):
            if column == 0:
                _write_single_line(slide, shape.name, value)
            else:
                _write_shape(shape, value)


def _populate_priority(slide: Any, payload: dict[str, Any], page: int) -> None:
    period = _clean(payload.get("period_label"), "Selected period")
    scope = _clean(payload.get("scope_label"), "Authorized group scope")
    _set_common(slide, "Region / Team Priority Map", "Action order based on score gap, movement, and weighted KPI loss", period, scope, page)
    rows = _risk_rows(payload)
    slots = [
        (("Text 10", "Text 11", "Text 12", "Text 13", "Text 14", "Text 15"), "Shape 8", "Shape 9"),
        (("Text 18", "Text 19", "Text 20", "Text 21", "Text 22", "Text 23"), "Shape 16", "Shape 17"),
        (("Text 26", "Text 27", "Text 28", "Text 29", "Text 30", "Text 31"), "Shape 24", "Shape 25"),
        (("Text 34", "Text 35", "Text 36", "Text 37", "Text 38", "Text 39"), "Shape 32", "Shape 33"),
    ]
    for index, (names, card, accent) in enumerate(slots):
        row = rows[index] if index < len(rows) else None
        status = row.get("status") if row else "Data unavailable"
        _card_status(slide, card, accent, status)
        if not row:
            for name in names:
                _write(slide, name, MISSING)
            continue
        driver = _leading_kpi(row)
        _write_single_line(slide, names[0], f"{index + 1}. {_clean(row.get('name'))}")
        _write(slide, names[1], _pct(row.get("score")))
        _write(slide, names[2], _fit_wrapped(f"{_status_label(status).title()}: {_clean(row.get('scope_label'))}", 46))
        bullets = [
            f"Gap to target {_pct(row.get('gap_to_target'))}",
            f"{row.get('below_threshold', 0)} below threshold / {row.get('record_count', 0)} measured rows",
            f"{_clean(driver.get('label'), 'No leading KPI')} loss {_pct(driver.get('weighted_impact'))}",
        ]
        for name, value in zip(names[3:], bullets):
            _write(slide, name, _fit_wrapped(f"• {value}", 46))
    _write(slide, "Text 41", "Management reading: resolve the lowest-performing group and its highest weighted-loss KPI first, then confirm ownership and evidence.")


def _populate_group_performance(slide: Any, payload: dict[str, Any], page: int) -> None:
    period = _clean(payload.get("period_label"), "Selected period")
    scope = _clean(payload.get("scope_label"), "Authorized group scope")
    _set_common(slide, "Performance by Region / Team", "Current score, baseline, MoM, and group scope", period, scope, page)
    _write(slide, "Text 8", "Region / Team Performance — Main Score & MoM")
    _write(slide, "Text 55", "Lowest Performing Groups")
    role_names = ["Text 10", "Text 15", "Text 20", "Text 25", "Text 30", "Text 35", "Text 40", "Text 45", "Text 50"]
    score_names = ["Text 13", "Text 18", "Text 23", "Text 28", "Text 33", "Text 38", "Text 43", "Text 48", "Text 53"]
    mom_names = ["Text 14", "Text 19", "Text 24", "Text 29", "Text 34", "Text 39", "Text 44", "Text 49", "Text 54"]
    track_names = ["Shape 11", "Shape 16", "Shape 21", "Shape 26", "Shape 31", "Shape 36", "Shape 41", "Shape 46", "Shape 51"]
    bar_names = ["Shape 12", "Shape 17", "Shape 22", "Shape 27", "Shape 32", "Shape 37", "Shape 42", "Shape 47", "Shape 52"]
    rows = list((payload.get("teams") or [])) or list((payload.get("regions") or []))
    for index, name in enumerate(role_names):
        if index >= len(rows):
            _write(slide, name, MISSING)
            _write(slide, score_names[index], MISSING)
            _write(slide, mom_names[index], MISSING)
            _resize_bar(slide, bar_names[index], 0.0, color=reference.MUTED)
            continue
        row = rows[index]
        label = _text_shape(slide, name)
        reference._write_role_people(label, row.get("name"), [row.get("scope_label")], name_size=5.4)
        _write(slide, score_names[index], _pct(row.get("score")))
        _write(slide, mom_names[index], _mom(row.get("movement")))
        reference._fill(slide, track_names[index], reference.TRACK)
        _resize_bar(slide, bar_names[index], (_number(row.get("score")) or 0) / 100, color=_status_color(row.get("status")))
    attention_names = [("Text 58", "Text 59", "Text 60"), ("Text 62", "Text 63", "Text 64"), ("Text 66", "Text 67", "Text 68")]
    for index, names in enumerate(attention_names):
        row = rows[index] if index < len(rows) else None
        if row:
            _write_single_line(slide, names[0], f"{index + 1}. {_clean(row.get('name'))}")
            _write(slide, names[1], _pct(row.get("score")))
            _write(slide, names[2], f"MoM {_mom(row.get('movement'))} | {row.get('below_threshold', 0)} below")
        else:
            for name in names:
                _write(slide, name, MISSING)
    story = (
        f"The current priority is {_clean(rows[0].get('name'))} at {_pct(rows[0].get('score'))}. "
        f"Compare each group with the same baseline before escalating."
        if rows else "No region or team performance data is available."
    )
    _write(slide, "Text 70", "Story signal")
    _write_single_line(slide, "Text 71", story)


def _visual_section(payload: dict[str, Any]) -> dict[str, Any]:
    rows = []
    for source_row in list(payload.get("teams") or []) or list(payload.get("regions") or []):
        row = dict(source_row)
        row["position"] = row.get("scope_type") or "Group"
        row["name"] = row.get("name") or row.get("scope_label")
        rows.append(row)
    snapshot = {
        "all_people": rows,
        "kpis": [dict(kpi) for kpi in payload.get("kpis") or []],
        "trend": list(payload.get("trend") or []),
        "history_count": payload.get("history_count"),
        "best_period": payload.get("best_period") or {},
        "worst_period": payload.get("worst_period") or {},
        "net_movement": payload.get("net_movement"),
        "latest_period_label": payload.get("period_label"),
        "actions": list(payload.get("raw_actions") or []),
        "root_cause_rows": list(payload.get("evidence") or []),
    }
    overall = _current_view(payload)
    view = {
        "name": "Authorized Group Scope",
        "score": overall.get("score"),
        "baseline": overall.get("baseline"),
        "movement": overall.get("movement"),
        "target": overall.get("target", 100.0),
        "status": overall.get("status"),
        "people": rows,
        "kpis": list(payload.get("kpis") or []),
        "snapshot": snapshot,
    }
    return {"display": "Authorized Group Scope", "view": view, "snapshot": snapshot}


def _populate_trend(slide: Any, payload: dict[str, Any], page: int) -> None:
    section = _visual_section(payload)
    _set_common(slide, "Performance Trend & KPI Health", "Available history and direction-aware KPI driver impact", _clean(payload.get("period_label")), _clean(payload.get("scope_label")), page)
    # The approved trend frame already contains the complete history / loss
    # grammar.  Its inherited labels are then replaced with group-level copy.
    reference._populate_marketing_trend(slide, {"period": payload.get("period_label"), "branch": payload.get("scope_label")}, section, page, list(payload.get("kpis") or []), 1, 1)
    _set_common(slide, "Performance Trend & KPI Health", "Available history and direction-aware KPI driver impact", _clean(payload.get("period_label")), _clean(payload.get("scope_label")), page)
    directions = sorted({
        reference._direction_label(kpi.get("direction"))
        for kpi in payload.get("kpis") or []
        if kpi.get("direction")
    })
    _write(slide, "Text 76", f"Direction-aware scoring: {', '.join(directions) or 'Higher is better / Lower is better'}; solve the highest weighted-loss drivers first.")


def _populate_driver_map(slide: Any, payload: dict[str, Any], page: int) -> None:
    section = _visual_section(payload)
    model = {"period": payload.get("period_label"), "branch": payload.get("scope_label")}
    reference._populate_marketing_map(slide, model, section, page, list(payload.get("drivers") or [])[:4], 1, 1)
    _set_common(slide, "Driver → Region / Team", "Connects KPI loss to the affected group and the required management response", _clean(payload.get("period_label")), _clean(payload.get("scope_label")), page)
    _write(slide, "Text 8", "Priority driver map")
    _write(slide, "Text 13", "Group type")
    _write(slide, "Text 15", "Affected region / team")
    _write(slide, "Text 51", "Management interpretation: the score is the entry point; the intervention follows the driver, affected group, evidence state, and ownership.")


def _populate_evidence(slide: Any, payload: dict[str, Any], page: int) -> None:
    rows = [list(row) for row in payload.get("_evidence_page_rows", payload.get("evidence_rows")) or []]
    suffix = _clean(payload.get("_evidence_suffix"), "") if payload.get("_evidence_suffix") else ""
    title = payload.get("_evidence_title") or f"KPI Evidence & Root-Cause State{suffix}"
    _set_common(slide, title, "Evidence first: actual, target, achievement, weighted loss, and validation status", _clean(payload.get("period_label")), _clean(payload.get("scope_label")), page)
    for name, value in zip(("Text 9", "Text 11", "Text 13", "Text 15", "Text 17", "Text 19", "Text 21", "Text 23"), ("Scope", "Region / Team", "KPI", "Actual", "Target", "Achievement", "Weighted loss", "Evidence / validation")):
        _write(slide, name, value)
    groups = _row_groups(slide, top_min=1.57, top_max=5.25, min_columns=8)
    displayed = rows[: len(groups)]
    for group, values in zip(groups, displayed):
        for column, (shape, value) in enumerate(zip(group, values)):
            color = None
            if column == 5:
                try:
                    achievement = float(str(value).replace("%", "").strip())
                except (TypeError, ValueError):
                    achievement = None
                color = _status_color("On track" if achievement is not None and achievement >= 100 else "Watch" if achievement is not None and achievement >= 70 else "Requires action" if achievement is not None else "Data unavailable")
            elif column == 6:
                try:
                    loss = float(str(value).replace("%", "").strip())
                except (TypeError, ValueError):
                    loss = None
                color = reference.GREEN if loss is not None and loss <= 0 else reference.RED if loss is not None else reference.MUTED
            _write_shape(shape, value, color=color)
    for group in groups[len(displayed):]:
        for shape in group:
            _write_shape(shape, MISSING)
    _write(slide, "Text 169", "Important: a KPI identifies where to investigate. A confirmed root cause requires recorded evidence; otherwise the item remains Investigation Required.")


def _populate_actions(slide: Any, payload: dict[str, Any], page: int) -> None:
    section = _visual_section(payload)
    model = {"period": payload.get("period_label"), "branch": payload.get("scope_label")}
    normalized_actions = payload.get("_action_page_actions")
    if normalized_actions is None:
        normalized_actions = reference._action_rows(section["snapshot"])
    suffix = _clean(payload.get("_action_suffix"), "") if payload.get("_action_suffix") else ""
    reference._populate_marketing_actions(slide, model, section, page, list(normalized_actions)[:4], 1, 1)
    _set_common(slide, f"Executive Commitment Tracker{suffix}", "Recorded and proposed workstreams requiring owner, due date, status, and success evidence", _clean(payload.get("period_label")), _clean(payload.get("scope_label")), page)
    _write(slide, "Text 69", "Decision needed: confirm owner, due date, status, and success evidence for every recorded or proposed workstream.")


def _populate_management(slide: Any, payload: dict[str, Any], page: int) -> None:
    period = _clean(payload.get("period_label"), "Selected period")
    scope = _clean(payload.get("scope_label"), "Authorized group scope")
    overall = _current_view(payload)
    risk_rows = _risk_rows(payload)
    drivers = list(payload.get("drivers") or [])
    actions = list(payload.get("actions") or [])
    unresolved_owner = sum(1 for action in actions if _clean(action.get("owner"), "Owner needed") == "Owner needed")
    unresolved_due = sum(1 for action in actions if _clean(action.get("due"), "Due date needed") == "Due date needed")
    _set_common(slide, "CEO Management Summary / Next Review", "The decisions and evidence to carry into the next performance cycle", period, scope, page)
    lowest = risk_rows[0] if risk_rows else {}
    items = [
        ("CURRENT GROUP SCORE", _pct(overall.get("score")), f"Target {_pct(overall.get('target', 100.0))} | MoM {_mom(overall.get('movement'))}"),
        ("LOWEST GROUP", _clean(lowest.get("name")), f"Score {_pct(lowest.get('score'))} | Gap {_pct(lowest.get('gap_to_target'))}"),
        ("TOP KPI GAPS", f"{_clean(drivers[0].get('label'), 'KPI')} + {max(0, len(drivers) - 1)} more" if drivers else MISSING, f"Largest loss {_pct(drivers[0].get('weighted_impact')) if drivers else MISSING}"),
        ("OPEN COMMITMENTS", str(len(actions)), f"{unresolved_owner} owner(s) / {unresolved_due} due date(s) unresolved"),
        ("GROUPS REQUIRING ACTION", str(sum(1 for row in risk_rows if _score_status(row.get('score')) == "Requires action")), "Prioritize exceptions before on-track groups"),
        ("NEXT REVIEW", "Score + KPI evidence", "Recheck actual / target / closure evidence"),
    ]
    card_texts = [("Text 10", "Text 11", "Text 12"), ("Text 15", "Text 16", "Text 17"), ("Text 20", "Text 21", "Text 22"), ("Text 25", "Text 26", "Text 27"), ("Text 30", "Text 31", "Text 32"), ("Text 35", "Text 36", "Text 37")]
    card_shapes = [("Shape 8", "Shape 9"), ("Shape 13", "Shape 14"), ("Shape 18", "Shape 19"), ("Shape 23", "Shape 24"), ("Shape 28", "Shape 29"), ("Shape 33", "Shape 34")]
    for index, names in enumerate(card_texts):
        item = items[index]
        status = "Requires action" if index in {1, 2, 3, 4} else "Watch"
        _card_status(slide, *card_shapes[index], status)
        _write(slide, names[0], _fit(item[0], 28))
        _write(slide, names[1], _fit(item[1], 28))
        _write(slide, names[2], _fit(item[2], 44))
    _write(slide, "Text 38", "Final management decision")
    decision = (
        f"Prioritize {_clean(lowest.get('name'), 'the highest-risk group')}; close {_clean(drivers[0].get('label'), 'the leading KPI gap') if drivers else 'the leading KPI gap'}; "
        f"assign {unresolved_owner} owner(s) and {unresolved_due} due date(s), then recheck score and KPI evidence next cycle."
    )
    _write(slide, "Text 41", decision)


def _populate_appendix(slide: Any, payload: dict[str, Any], page: int) -> None:
    _set_common(slide, "Authorized Group Appendix", "Reference view for region/team score, leading KPI gap, and commitment status", _clean(payload.get("period_label")), _clean(payload.get("scope_label")), page)
    headers = ("Scope", "Group", "Region / Team", "Score", "MoM", "Leading KPI gap", "Action status")
    for name, value in zip(("Text 9", "Text 11", "Text 13", "Text 15", "Text 17", "Text 19", "Text 21"), headers):
        _write(slide, name, value)
    rows: list[list[str]] = []
    for source_row in list(payload.get("teams") or []) or list(payload.get("regions") or []):
        driver = _leading_kpi(source_row)
        rows.append([
            _clean(source_row.get("scope_type"), "Group"),
            _clean(source_row.get("name")),
            _clean(source_row.get("scope_label")),
            _pct(source_row.get("score")),
            _mom(source_row.get("movement")),
            _fit_wrapped(f"{_clean(driver.get('label'), 'No KPI')} | loss {_pct(driver.get('weighted_impact'))}", 46),
            _clean(source_row.get("action_status"), "No action recorded"),
        ])
    groups = _row_groups(slide, top_min=1.46, top_max=5.9, min_columns=7)
    displayed = rows[: len(groups)]
    for group, values in zip(groups, displayed):
        for shape, value in zip(group, values):
            _write_shape(shape, value)
    for group in groups[len(displayed):]:
        for shape in group:
            _write_shape(shape, MISSING)


def _copy_payload(payload: dict[str, Any]) -> dict[str, Any]:
    result = dict(payload)
    for key in ("overall", "filters"):
        result[key] = dict(payload.get(key) or {})
    for key in ("regions", "teams", "kpis", "drivers", "actions", "evidence", "evidence_rows", "kpi_appendix_rows", "trend", "raw_actions", "_evidence_page_rows", "_action_page_actions"):
        result[key] = [dict(item) if isinstance(item, dict) else item for item in (payload.get(key) or [])]
    return result


def build_executive_group_summary_pptx(period_label: str = "Selected period", report_data: dict[str, Any] | None = None) -> bytes:
    """Generate a live CEO/group performance summary using the approved frames."""

    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Approved Offshore template is missing: {TEMPLATE_PATH}")
    template_bytes = TEMPLATE_PATH.read_bytes()
    source = Presentation(io.BytesIO(template_bytes))
    if len(source.slides) != REFERENCE_SLIDES:
        raise ValueError(f"Approved Offshore template must contain {REFERENCE_SLIDES} slides; found {len(source.slides)}")
    prs = Presentation(io.BytesIO(template_bytes))
    _clear_slides(prs)
    payload = _copy_payload(report_data or {})
    payload["period_label"] = payload.get("period_label") or period_label
    plan: list[tuple[int, Any, dict[str, Any]]] = [
        (0, _populate_cover, payload),
        (1, _populate_overview, payload),
        (2, _populate_priority, payload),
        (15, _populate_group_performance, payload),
        (16, _populate_trend, payload),
        (17, _populate_driver_map, payload),
    ]
    evidence_capacity = max(1, len(_row_groups(source.slides[23], top_min=1.57, top_max=5.25, min_columns=8)))
    evidence_rows = list(payload.get("evidence_rows") or [])
    evidence_chunks = [evidence_rows[index:index + evidence_capacity] for index in range(0, len(evidence_rows), evidence_capacity)] or [[]]
    for index, chunk in enumerate(evidence_chunks, 1):
        page_payload = _copy_payload(payload)
        page_payload["_evidence_page_rows"] = chunk
        page_payload["_evidence_suffix"] = f" ({index}/{len(evidence_chunks)})" if len(evidence_chunks) > 1 else ""
        plan.append((23, _populate_evidence, page_payload))

    normalized_actions = reference._action_rows(_visual_section(payload)["snapshot"])
    action_capacity = max(1, len(_row_groups(source.slides[24], top_min=1.78, top_max=5.0, min_columns=6)))
    action_chunks = [normalized_actions[index:index + action_capacity] for index in range(0, len(normalized_actions), action_capacity)] or [[]]
    for index, chunk in enumerate(action_chunks, 1):
        page_payload = _copy_payload(payload)
        page_payload["_action_page_actions"] = chunk
        page_payload["_action_suffix"] = f" ({index}/{len(action_chunks)})" if len(action_chunks) > 1 else ""
        plan.append((24, _populate_actions, page_payload))

    plan.extend([
        (26, _populate_management, payload),
        (27, _populate_appendix, payload),
    ])
    appendix_rows = list(payload.get("kpi_appendix_rows") or [])
    appendix_chunks = [appendix_rows[index:index + evidence_capacity] for index in range(0, len(appendix_rows), evidence_capacity)]
    for index, chunk in enumerate(appendix_chunks, 1):
        page_payload = _copy_payload(payload)
        page_payload["_evidence_page_rows"] = chunk
        page_payload["_evidence_suffix"] = ""
        page_payload["_evidence_title"] = f"KPI Evidence Appendix ({index}/{len(appendix_chunks)})" if len(appendix_chunks) > 1 else "KPI Evidence Appendix"
        plan.append((23, _populate_evidence, page_payload))
    for page, (source_index, populate, page_payload) in enumerate(plan, 1):
        slide = _append_template_slide(prs, source, source_index)
        populate(slide, page_payload, page)
    prs.core_properties.title = f"Executive Group Performance Summary - {payload['period_label']}"
    prs.core_properties.subject = "Dynamic CEO / group-level performance summary"
    prs.core_properties.author = "PMS Dashboard"
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


__all__ = ["TEMPLATE_PATH", "build_executive_group_summary_pptx"]
