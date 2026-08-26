"""Marketing PowerPoint export built by following the approved reference deck.

The reference presentation is intentionally used as the slide source.  This
module only replaces text, fills, progress widths, and row values; it does not
recreate the visual system with a new blank presentation.  That keeps the
approved master, layout, theme, typography, geometry, and stacking order
intact while the report remains fully data-driven.
"""

from __future__ import annotations

import io
from copy import deepcopy
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from exports import marketing_legendary_pptx_builder as base


TEMPLATE_PATH = Path(__file__).resolve().parent / "templates" / "Marketing_PMS_Legendary_Template.pptx"
TEMPLATE_SLIDES = 16
SLIDE_WIDTH = base.SLIDE_WIDTH
SLIDE_HEIGHT = base.SLIDE_HEIGHT

NAVY = base.NAVY
DARK_NAVY = base.DARK_NAVY
TEAL = base.TEAL
RED = base.RED
AMBER = base.AMBER
GREEN = base.GREEN
BLUE = base.BLUE
MUTED = base.MUTED
FAINT = base.FAINT
BODY = base.BODY
WHITE = base.WHITE
LINE = base.LINE
TRACK = base.TRACK
PALE_RED = base.PALE_RED
PALE_AMBER = base.PALE_AMBER
PALE_GREEN = base.PALE_GREEN
PALE_BLUE = base.PALE_BLUE
PALE_NEUTRAL = base.PALE_NEUTRAL

MISSING = "—"


def _text(slide: Any, name: str) -> Any:
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"Template shape not found: {name}")


def _write(slide: Any, name: str, value: Any, *, color: RGBColor | None = None) -> None:
    """Replace text while retaining the reference shape's direct formatting."""

    shape = _text(slide, name)
    text = MISSING if value is None or value == "" else str(value)
    frame = shape.text_frame
    paragraphs = list(frame.paragraphs)
    if not paragraphs:
        frame.text = text
        return
    first = paragraphs[0]
    runs = list(first.runs)
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
    else:
        first.text = text
    for paragraph in paragraphs[1:]:
        for run in paragraph.runs:
            run.text = ""
    if color is not None:
        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = color


def _fill(slide: Any, name: str, color: RGBColor, *, line: RGBColor | None = None) -> None:
    shape = _text(slide, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is not None:
        shape.line.color.rgb = line


def _set_width(shape: Any, width: float, *, left: float | None = None) -> None:
    shape.width = Inches(max(0.01, width))
    if left is not None:
        shape.left = Inches(left)


def _shape(slide: Any, name: str) -> Any:
    return _text(slide, name)


def _clear_slides(prs: Presentation) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def _append_template_slide(prs: Presentation, source: Presentation, source_index: int) -> Any:
    """Append a source slide's shape tree to a slide using the same template layout."""

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    tree = slide.shapes._spTree
    for child in list(tree)[2:]:
        tree.remove(child)
    source_tree = source.slides[source_index].shapes._spTree
    for child in list(source_tree)[2:]:
        tree.insert_element_before(deepcopy(child), "p:extLst")
    return slide


def _period(payload: dict[str, Any]) -> str:
    return base._period(payload)


def _short_period(period: str) -> str:
    parts = str(period).split()
    if len(parts) >= 2 and parts[0].casefold() in {name.casefold() for name in base.MONTHS}:
        return f"{parts[0][:3]} {parts[1]}"
    return str(period)


def _pct(value: Any, *, signed: bool = False) -> str:
    return base._fmt_pct(value, signed=signed, fallback=MISSING)


def _mom(value: Any) -> str:
    return base._fmt_mom(value)


def _native(value: Any, unit: Any = "") -> str:
    return base._fmt_native(value, unit, MISSING)


def _num(value: Any) -> float | None:
    return base._number(value)


def _clean(value: Any, fallback: str = "") -> str:
    return base._clean(value, fallback)


def _fit(value: Any, limit: int, fallback: str = MISSING) -> str:
    cleaned = _clean(value)
    return base._fit_title(cleaned, limit) if cleaned else fallback


def _status(score: Any) -> str:
    return base._score_status(score)


def _employee_status(score: Any) -> str:
    return base._employee_status(score)


def _loss(kpi: dict[str, Any] | None) -> float:
    return base._kpi_loss(kpi or {})


def _kpis(payload: dict[str, Any]) -> list[dict[str, Any]]:
    return base._kpi_rows(payload)


def _people(payload: dict[str, Any]) -> list[dict[str, Any]]:
    return base._people(payload)


def _detail_order(payload: dict[str, Any]) -> list[dict[str, Any]]:
    """Approved priority order: lowest employee score first, then KPI loss."""

    return sorted(
        base._detail_people(payload),
        key=lambda row: (
            _num(row.get("score")) is None,
            _num(row.get("score")) or 0,
            -_loss(_person_kpis(row)[0]) if _person_kpis(row) else 0,
            _clean(row.get("name")).casefold(),
        ),
    )


def _priority_rank(payload: dict[str, Any], person: dict[str, Any]) -> int:
    person_id = _clean(person.get("employee_id"))
    for index, row in enumerate(_detail_order(payload), 1):
        if person_id and person_id == _clean(row.get("employee_id")):
            return index
    return len(_detail_order(payload)) + 1


def _appendix_order(payload: dict[str, Any]) -> list[dict[str, Any]]:
    status_order = {"Requires action": 0, "Watch": 1, "On track": 2, "Data unavailable": 3}
    return sorted(
        _people(payload),
        key=lambda row: (
            status_order.get(_employee_status(row.get("score")), 3),
            _num(row.get("score")) is None,
            _num(row.get("score")) or 0,
            _clean(row.get("name")).casefold(),
        ),
    )


def _role_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    return base._role_rows(payload)


def _status_code(status: str) -> str:
    return {
        "Requires action": "ACTION",
        "Watch": "WATCH",
        "On track": "TRACK",
        "Data quality": "DATA",
        "Data unavailable": "DATA",
    }.get(status, _fit(status.upper(), 10, "DATA"))


def _status_color(status: Any) -> RGBColor:
    return base._status_color(status)


def _status_fill(status: Any) -> RGBColor:
    return base._status_fill(status)


def _set_common(slide: Any, title: str, subtitle: str, period: str, page: int) -> None:
    _write(slide, "Text 2", title)
    _write(slide, "Text 3", subtitle)
    _write(slide, "Text 4", f"{period} | {page:02d}")
    _write(slide, "Text 6", f"Marketing PMS | {_short_period(period)}")
    _write(slide, "Text 7", f"{page:02d}")


def _set_row_fills(slide: Any, row_shape_groups: list[list[str]], status: str) -> None:
    fill = _status_fill(status)
    for group in row_shape_groups:
        for name in group:
            _fill(slide, name, fill)


def _role_action(status: str) -> str:
    return {
        "Requires action": "Immediate corrective action",
        "Watch": "Monitor key gap",
        "On track": "Maintain performance",
        "Data unavailable": "Validate data",
    }.get(status, "Validate data")


def _role_rows_low_first(payload: dict[str, Any]) -> list[dict[str, Any]]:
    return sorted(
        _role_rows_with_people(payload),
        key=lambda row: (_num(row.get("score")) is None, _num(row.get("score")) or 0, _clean(row.get("role")).casefold()),
    )


def _role_rows_with_people(payload: dict[str, Any]) -> list[dict[str, Any]]:
    """Attach the selected employees to each role row without changing scoring."""

    people_by_role: dict[str, list[str]] = {}
    for person in _people(payload):
        role = _clean(person.get("position"), "Unassigned")
        name = _clean(person.get("name"), "Unknown")
        key = role.casefold()
        names = people_by_role.setdefault(key, [])
        if name not in names:
            names.append(name)

    rows: list[dict[str, Any]] = []
    for source_row in _role_rows(payload):
        row = dict(source_row)
        role = _clean(row.get("role"), "Unassigned")
        row["employee_names"] = sorted(people_by_role.get(role.casefold(), []), key=str.casefold)
        rows.append(row)
    return rows


def _wrapped_names(names: list[str], max_chars: int = 38) -> str:
    """Wrap names at employee boundaries so the role row stays readable."""

    lines: list[str] = []
    current = ""
    for name in names:
        candidate = f"{current}, {name}" if current else name
        if current and len(candidate) > max_chars:
            lines.append(current)
            current = name
        else:
            current = candidate
    if current:
        lines.append(current)
    return "\n".join(lines)


def _write_role_with_people(
    slide: Any,
    shape_name: str,
    role: Any,
    names: list[str],
    *,
    height: float = 0.30,
    name_size: float = 5.7,
) -> None:
    """Write a source role label plus a muted employee-name line beneath it."""

    shape = _text(slide, shape_name)
    role_text = _clean(role, MISSING)
    people_text = _wrapped_names(names)
    frame = shape.text_frame
    first = frame.paragraphs[0]
    runs = list(first.runs)
    if runs:
        runs[0].text = role_text
        for run in runs[1:]:
            run.text = ""
        role_run = runs[0]
    else:
        first.text = role_text
        role_run = first.runs[0]
    first.alignment = PP_ALIGN.LEFT
    first.space_before = 0
    first.space_after = 0

    for paragraph in list(frame.paragraphs)[1:]:
        frame._element.remove(paragraph._p)
    if people_text:
        detail = frame.add_paragraph()
        detail.text = people_text
        detail.alignment = PP_ALIGN.LEFT
        detail.space_before = 0
        detail.space_after = 0
        detail.line_spacing = 1.0
        for run in detail.runs:
            if role_run.font.name:
                run.font.name = role_run.font.name
            run.font.size = Pt(name_size)
            run.font.bold = False
            run.font.color.rgb = MUTED
        shape.height = Inches(height)


def _top_driver(payload: dict[str, Any]) -> dict[str, Any]:
    rows = _kpis(payload)
    return rows[0] if rows else {}


def _kpi_status(kpi: dict[str, Any]) -> str:
    return _clean(kpi.get("status"), base._status_from_achievement(kpi.get("achievement")))


def _kpi_values(kpi: dict[str, Any]) -> tuple[str, str, str, str, str]:
    return (
        _native(kpi.get("actual"), kpi.get("unit")),
        _native(kpi.get("baseline_actual"), kpi.get("unit")),
        _mom(kpi.get("mom")),
        _native(kpi.get("target"), kpi.get("unit")),
        _kpi_status(kpi),
    )


def _kpi_key_matches(left: dict[str, Any], right: dict[str, Any]) -> bool:
    left_keys = {
        _clean(left.get("key")).casefold(),
        _clean(left.get("group_key")).casefold(),
        *[_clean(value).casefold() for value in (left.get("source_keys") or [])],
    }
    right_keys = {
        _clean(right.get("key")).casefold(),
        _clean(right.get("group_key")).casefold(),
        *[_clean(value).casefold() for value in (right.get("source_keys") or [])],
    }
    left_keys.discard("")
    right_keys.discard("")
    return bool(left_keys & right_keys) or _clean(left.get("label")).casefold() == _clean(right.get("label")).casefold()


def _person_kpis(person: dict[str, Any]) -> list[dict[str, Any]]:
    return sorted(base._person_kpi_rows(person), key=lambda row: (-_loss(row), _clean(row.get("label")).casefold()))


def _employee_driver_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    """Create auditable evidence rows from employee KPI data, without guessing causes."""

    rows: list[dict[str, Any]] = []
    for person in base._detail_people(payload) or _people(payload):
        for kpi in _person_kpis(person):
            if _loss(kpi) <= 0:
                continue
            root = next((candidate for candidate in payload.get("root_cause_rows") or [] if _kpi_key_matches(kpi, candidate)), {})
            root_state = _clean(root.get("evidence_status"))
            if "confirmed" in root_state.casefold():
                state, note = "Confirmed Root Cause", _clean(root.get("recorded_root_cause"), "Confirmed cause recorded")
            elif root:
                state, note = "Investigation Required", _clean(root.get("required_validation"), "No operational cause recorded")
            else:
                state, note = base._employee_evidence_label(person)
            rows.append(
                {
                    "employee": _clean(person.get("name"), "Unknown employee"),
                    "role": _clean(person.get("position"), "Unassigned"),
                    "kpi": _clean(kpi.get("label"), "KPI"),
                    "actual": _native(kpi.get("actual"), kpi.get("unit")),
                    "target": _native(kpi.get("target"), kpi.get("unit")),
                    "achievement": _pct(kpi.get("achievement_pct")),
                    "loss": _pct(_loss(kpi)),
                    "evidence": f"{state} — {_fit(note, 80, 'No operational cause recorded')}",
                    "state": state,
                    "sort_loss": _loss(kpi),
                }
            )
    # Preserve KPI-level investigation signals even when employee-level data
    # is not available for that KPI.  This is deliberately labelled as a
    # team-level signal rather than assigning it to an employee or inventing a
    # root cause.
    represented = {_clean(row.get("kpi")).casefold() for row in rows}
    for root in payload.get("root_cause_rows") or []:
        label = _clean(root.get("label"), "KPI")
        state = _clean(root.get("evidence_status"), "KPI signal only")
        if "confirmed" in state.casefold() or label.casefold() in represented:
            continue
        rows.append(
            {
                "employee": "Team-level signal",
                "role": "—",
                "kpi": label,
                "actual": _native(root.get("actual"), root.get("unit")),
                "target": _native(root.get("target"), root.get("unit")),
                "achievement": _pct(root.get("achievement_pct")),
                "loss": _pct(_loss(root)),
                "evidence": f"Investigation Required — {_fit(root.get('required_validation'), 80, 'No operational cause recorded')}",
                "state": "Investigation Required",
                "sort_loss": _loss(root),
            }
        )
    return sorted(rows, key=lambda row: (-row["sort_loss"], row["employee"].casefold(), row["kpi"].casefold()))


def _action_linked_kpi(payload: dict[str, Any], action: dict[str, Any]) -> dict[str, Any]:
    linked = _clean(action.get("linked_kpi_key")).casefold()
    if not linked:
        return {}
    for kpi in _kpis(payload):
        keys = {
            _clean(kpi.get("key")).casefold(),
            _clean(kpi.get("group_key")).casefold(),
            *[_clean(value).casefold() for value in (kpi.get("source_keys") or [])],
        }
        if linked in keys:
            return kpi
    return {}


def _action_people(payload: dict[str, Any], action: dict[str, Any], kpi: dict[str, Any]) -> str:
    employee_id = _clean(action.get("employee_id"))
    matching = []
    for person in _people(payload):
        if employee_id and _clean(person.get("employee_id")) == employee_id:
            matching.append(_clean(person.get("name"), "Unknown"))
        elif kpi and any(_kpi_key_matches(item, kpi) and _loss(item) > 0 for item in _person_kpis(person)):
            matching.append(_clean(person.get("name"), "Unknown"))
    return ", ".join(dict.fromkeys(matching)) or "Selected employees"


def _action_priority(action: dict[str, Any], index: int, linked_kpi: dict[str, Any]) -> str:
    if _loss(linked_kpi) >= 20 or index < 2:
        return "High"
    return "Medium"


def _action_rows(payload: dict[str, Any]) -> list[dict[str, Any]]:
    actions = base._management_actions(payload)
    rows: list[dict[str, Any]] = []
    for index, action in enumerate(actions):
        kpi = _action_linked_kpi(payload, action)
        owner = base._action_owner(action)
        due = base._action_due(action)
        status = base._action_status(action)
        workstream = _clean(kpi.get("label"), base._action_issue(action))
        rows.append(
            {
                "priority": _action_priority(action, index, kpi),
                "workstream": workstream,
                "owner_due": f"{owner}\nDue: {due}",
                "scope": _action_people(payload, action, kpi),
                "execution": f"{base._action_text(action)}\nStatus: {status}",
                "success": base._action_success(action),
                "status": status,
                "missing": owner.casefold() == "owner needed" or due.casefold() == "due date needed",
            }
        )
    return rows


def _action_rows_or_placeholder(payload: dict[str, Any]) -> list[dict[str, Any]]:
    rows = _action_rows(payload)
    if rows:
        return rows
    return [
        {
            "priority": "—",
            "workstream": "No recorded workstream",
            "owner_due": "Owner needed\nDue: Due date needed",
            "scope": "—",
            "execution": "No corrective action is recorded\nStatus: Unresolved",
            "success": "Success measure needed",
            "status": "Unresolved",
            "missing": True,
        }
    ]


def _history(payload: dict[str, Any]) -> list[dict[str, Any]]:
    return base._history_rows(payload)


def _history_score(row: dict[str, Any]) -> float | None:
    return _num(row.get("score"))


def _append_cover(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _append_template_slide(prs, source, 0)
    period = _period(payload)
    driver = _top_driver(payload)
    _write(slide, "Text 3", "Marketing Team")
    _write(slide, "Text 4", "Performance Review")
    _write(slide, "Text 5", f"{period} | Data-led performance story from Excel")
    _write(slide, "Text 6", "Overall performance, KPI drivers, individual impact and next-review commitments.")
    _write(slide, "Text 9", _pct(payload.get("overall_score")))
    _write(slide, "Text 10", "Marketing main score")
    _write(slide, "Text 12", f"MoM {_mom(payload.get('movement'))}", color=base._movement_color(payload.get("movement")))
    below = base._below_threshold(payload)
    _write(slide, "Text 14", f"{len(below)} employees below 70%")
    driver_text = f"Top driver: {_fit(driver.get('label'), 28)} {_pct(driver.get('weighted_impact'))} loss" if driver else "Top driver: —"
    _write(slide, "Text 16", driver_text)
    _write(slide, "Text 17", "CONFIDENTIAL — INTERNAL USE ONLY")


def _append_overview(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _append_template_slide(prs, source, 1)
    period = _period(payload)
    _set_common(slide, "Monthly Performance Overview", "How did Marketing perform this month?", period, page)
    current = _num(payload.get("overall_score"))
    target = _num(payload.get("target_score"))
    baseline = payload.get("baseline_score")
    overall_status = _status(current)
    driver = _top_driver(payload)
    top = base._top_performer(payload)
    decline = base._largest_decline(payload)
    below = base._below_threshold(payload)

    _write(slide, "Text 10", "CURRENT MONTHLY PERFORMANCE")
    _write(slide, "Text 11", _pct(current))
    _write(slide, "Text 12", "Marketing Main Score")
    _write(slide, "Text 13", f"Baseline: {_pct(baseline)}")
    _write(slide, "Text 14", f"MoM: {_mom(payload.get('movement'))}", color=base._movement_color(payload.get("movement")))
    _write(slide, "Text 16", overall_status.upper(), color=WHITE)
    _write(slide, "Text 17", f"Target: {_pct(target)}")
    _fill(slide, "Shape 8", base.PALE_AMBER if overall_status == "Watch" else base.PALE_RED if overall_status == "Requires action" else base.PALE_GREEN)
    _fill(slide, "Shape 9", _status_color(overall_status))
    _fill(slide, "Shape 15", _status_color(overall_status))

    _write(slide, "Text 20", "EMPLOYEES BELOW 70%")
    _write(slide, "Text 21", str(len(below)))
    _write(slide, "Text 22", f"{len(below)} employee(s) require attention")
    _write(slide, "Text 25", "TOP CURRENT SCORE")
    _write(slide, "Text 26", _pct(top.get("score")))
    _write(slide, "Text 27", _fit(top.get("name"), 22))
    _write(slide, "Text 30", "LARGEST DECLINE")
    _write(slide, "Text 31", _mom(decline.get("movement")))
    _write(slide, "Text 32", _fit(decline.get("name"), 22))
    gap = abs(current - target) if current is not None and target is not None else None
    _write(slide, "Text 35", "GAP")
    _write(slide, "Text 36", _pct(gap))
    _write(slide, "Text 37", "to target")
    _write(slide, "Text 39", f"Executive takeaway: Overall performance is {_pct(current)} with {_mom(payload.get('movement'))} MoM; {len(below)} employees and {_fit(driver.get('label'), 24) if driver else 'the KPI set'} are the primary management focus.")

    role_rows = _role_rows_low_first(payload)[:5]
    _write(slide, "Text 40", "Department-Level Action View")
    row_text_groups = [[f"Text {55 + row * 12 + column * 2}" for column in range(6)] for row in range(5)]
    row_shape_groups = [[f"Shape {54 + row * 12 + column * 2}" for column in range(6)] for row in range(5)]
    for row_index in range(5):
        row = role_rows[row_index] if row_index < len(role_rows) else {}
        score = row.get("score")
        state = _status(score)
        values = [
            _fit(row.get("role"), 30),
            _pct(score),
            _pct(row.get("baseline")),
            _mom(row.get("movement")),
            state.upper(),
            _role_action(state),
        ]
        for name, value in zip(row_text_groups[row_index], values):
            _write(slide, name, value, color=_status_color(state) if name.endswith("63") else None)
        _set_row_fills(slide, [row_shape_groups[row_index]], state)


def _append_roles(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int = 1, total_parts: int = 1) -> None:
    slide = _append_template_slide(prs, source, 2)
    period = _period(payload)
    suffix = f" — continued {part}/{total_parts}" if total_parts > 1 else ""
    _set_common(slide, f"Performance by Role{suffix}", "Current performance, baseline and MoM by function", period, page)
    _write(slide, "Text 8", "Role Performance — Main Score & MoM")
    _write(slide, "Text 55", "Attention Order")
    role_text_rows = [(10 + index * 5, 13 + index * 5, 14 + index * 5) for index in range(9)]
    role_bar_pairs = [(11 + index * 5, 12 + index * 5) for index in range(9)]
    for index in range(9):
        row = rows[index] if index < len(rows) else {}
        score = _num(row.get("score"))
        state = _status(score)
        name_name, score_name, mom_name = role_text_rows[index]
        if row.get("role") and row.get("employee_names"):
            _write_role_with_people(slide, f"Text {name_name}", row.get("role"), row.get("employee_names") or [])
        else:
            _write(slide, f"Text {name_name}", _fit(row.get("role"), 28))
        _write(slide, f"Text {score_name}", _pct(score), color=_status_color(state))
        _write(slide, f"Text {mom_name}", _mom(row.get("movement")) if row else MISSING, color=base._movement_color(row.get("movement")) if row else None)
        fill_name, track_name = role_bar_pairs[index]
        ratio = max(0.0, min(1.0, (score or 0) / 100.0)) if score is not None else 0.0
        fill_shape = _shape(slide, f"Shape {fill_name}")
        track_shape = _shape(slide, f"Shape {track_name}")
        start = fill_shape.left / 914400
        total_width = (fill_shape.width + track_shape.width) / 914400
        _fill(slide, f"Shape {fill_name}", _status_color(state))
        _fill(slide, f"Shape {track_name}", TRACK)
        _set_width(fill_shape, total_width * ratio)
        _set_width(track_shape, total_width * (1 - ratio), left=start + total_width * ratio)

    attention = sorted(
        rows,
        key=lambda row: (_num(row.get("score")) is None, _num(row.get("score")) or 0, _clean(row.get("role")).casefold()),
    )[:3]
    attention_text = [(58, 59, 60), (62, 63, 64), (66, 67, 68)]
    attention_shapes = ["Shape 57", "Shape 61", "Shape 65"]
    for index, (name_name, score_name, detail_name) in enumerate(attention_text):
        row = attention[index] if index < len(attention) else {}
        state = _status(row.get("score"))
        role = _fit(row.get("role"), 28)
        if row and row.get("employee_names"):
            _write_role_with_people(
                slide,
                f"Text {name_name}",
                f"{index + 1}. {role}",
                row.get("employee_names") or [],
                height=0.28,
                name_size=5.5,
            )
        else:
            _write(slide, f"Text {name_name}", f"{index + 1}. {role}")
        _write(slide, f"Text {score_name}", _pct(row.get("score")), color=_status_color(state))
        _write(slide, f"Text {detail_name}", f"MoM {_mom(row.get('movement'))} | {int(row.get('count') or 0)} employee(s)" if row else "No role-level signal")
        _fill(slide, attention_shapes[index], _status_fill(state) if row else PALE_NEUTRAL)
    low_roles = ", ".join(_fit(row.get("role"), 24) for row in attention if row)
    _write(slide, "Text 71", f"Strong roles can mask low scores in {low_roles}." if low_roles else "No role-level performance signal is available.")
    _fill(slide, "Shape 69", DARK_NAVY)


def _trend_summary(payload: dict[str, Any], history: list[dict[str, Any]]) -> tuple[dict[str, Any], dict[str, Any], Any, Any]:
    valid = [row for row in history if _history_score(row) is not None]
    best = max(valid, key=lambda row: _history_score(row) or 0) if valid else {}
    worst = min(valid, key=lambda row: _history_score(row) or 0) if valid else {}
    if valid:
        movement = (_history_score(valid[-1]) or 0) - (_history_score(valid[0]) or 0)
    else:
        movement = None
    target = _num(payload.get("target_score"))
    current = _num(payload.get("overall_score"))
    return best, worst, movement, abs(current - target) if current is not None and target is not None else None


def _append_trend(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int, history_slice: list[dict[str, Any]], history_count: int, part: int = 1, total_parts: int = 1) -> None:
    slide = _append_template_slide(prs, source, 3)
    period = _period(payload)
    title = f"Performance Trend ({history_count} period{'s' if history_count != 1 else ''})"
    if total_parts > 1:
        title += f" — {part}/{total_parts}"
    _set_common(slide, title, "Historical direction before drilling into drivers", period, page)
    best, worst, net_movement, target_gap = _trend_summary(payload, _history(payload))
    history = _history(payload)
    _write(slide, "Text 10", "HISTORY COVERAGE")
    _write(slide, "Text 11", f"{history_count} PERIOD{'S' if history_count != 1 else ''}")
    _write(slide, "Text 12", f"{sum(_num(row.get('record_count')) or 0 for row in history)} measured rows" if history else "No measured rows")
    _write(slide, "Text 15", "BEST PERIOD")
    _write(slide, "Text 16", _pct(best.get("score")))
    _write(slide, "Text 17", _fit(best.get("label"), 18))
    _write(slide, "Text 20", "WORST PERIOD")
    _write(slide, "Text 21", _pct(worst.get("score")))
    _write(slide, "Text 22", _fit(worst.get("label"), 18))
    _write(slide, "Text 25", "NET MOVEMENT")
    _write(slide, "Text 26", _pct(net_movement, signed=True), color=base._movement_color(net_movement))
    _write(slide, "Text 27", "First to latest")
    _write(slide, "Text 30", "TARGET GAP")
    _write(slide, "Text 31", _pct(target_gap))
    _write(slide, "Text 32", f"{_short_period(period)} vs {_pct(payload.get('target_score'))}")

    _write(slide, "Text 33", "Performance history")
    _write(slide, "Text 47", "Context summary")
    row_slots = [(35, 38, 36, 37), (39, 42, 40, 41), (43, 46, 44, 45)]
    max_score = max([_history_score(row) or 0 for row in history_slice] + [100.0])
    for index, (label_name, score_name, fill_name, track_name) in enumerate(row_slots):
        row = history_slice[index] if index < len(history_slice) else {}
        score = _history_score(row)
        _write(slide, f"Text {label_name}", _fit(row.get("label"), 14))
        _write(slide, f"Text {score_name}", _pct(score))
        fill_shape = _shape(slide, f"Shape {fill_name}")
        track_shape = _shape(slide, f"Shape {track_name}")
        start = fill_shape.left / 914400
        total_width = (fill_shape.width + track_shape.width) / 914400
        ratio = max(0.0, min(1.0, (score or 0) / max_score)) if score is not None else 0.0
        color = GREEN if (score or 0) >= 90 else AMBER if (score or 0) >= 70 else RED
        _fill(slide, f"Shape {fill_name}", color)
        _fill(slide, f"Shape {track_name}", TRACK)
        _set_width(fill_shape, total_width * ratio)
        _set_width(track_shape, total_width * (1 - ratio), left=start + total_width * ratio)
    context = _clean(payload.get("context_headline"), "No trend context is available.")
    _write(slide, "Text 50", context)
    _write(slide, "Text 51", f"Current {_pct(payload.get('overall_score'))} vs baseline {_pct(payload.get('baseline_score'))}; net movement {_pct(net_movement, signed=True)}.")


def _kpi_card_slots() -> list[tuple[int, int, int, int, int]]:
    return [(12, 13, 14, 15, 16), (19, 20, 21, 22, 23), (26, 27, 28, 29, 30), (33, 34, 35, 36, 37), (40, 41, 42, 43, 44), (47, 48, 49, 50, 51)]


def _append_kpi_health(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int = 1, total_parts: int = 1) -> None:
    slide = _append_template_slide(prs, source, 4)
    period = _period(payload)
    suffix = f" — {part}/{total_parts}" if total_parts > 1 else ""
    _set_common(slide, f"KPI Health Overview{suffix}", "Current / baseline / MoM / target / status for every primary KPI", period, page)
    _write(slide, "Text 8", "Top KPI health signals")
    card_shapes = [("Shape 10", "Shape 11"), ("Shape 17", "Shape 18"), ("Shape 24", "Shape 25"), ("Shape 31", "Shape 32"), ("Shape 38", "Shape 39"), ("Shape 45", "Shape 46")]
    slots = _kpi_card_slots()
    for index, (label_name, value_name, weighted_label_name, current_name, detail_name) in enumerate(slots):
        kpi = rows[index] if index < len(rows) else {}
        if kpi:
            status = _kpi_status(kpi)
            loss = _loss(kpi)
            actual, baseline, mom, target, _ = _kpi_values(kpi)
            card_fill = _status_fill(status) if loss > 0 else PALE_GREEN
            rail = _status_color(status) if loss > 0 else GREEN
            _write(slide, f"Text {label_name}", _fit(kpi.get("label"), 26))
            _write(slide, f"Text {value_name}", _pct(loss), color=rail)
            _write(slide, f"Text {weighted_label_name}", f"Weighted loss · {_status_code(status)}")
            _write(slide, f"Text {current_name}", f"C {actual} · B {baseline}")
            _write(slide, f"Text {detail_name}", f"M {mom} · T {target}")
        else:
            status = "Data unavailable"
            card_fill = PALE_NEUTRAL
            rail = BLUE
            _write(slide, f"Text {label_name}", "KPI not available")
            _write(slide, f"Text {value_name}", MISSING)
            _write(slide, f"Text {weighted_label_name}", "No signal")
            _write(slide, f"Text {current_name}", "C — · B —")
            _write(slide, f"Text {detail_name}", "M — · T —")
        card_name, rail_name = card_shapes[index]
        _fill(slide, card_name, card_fill)
        _fill(slide, rail_name, rail)
    driver = _top_driver(payload)
    affected = [
        _clean(row.get("name"))
        for row in _people(payload)
        if driver and any(_kpi_key_matches(item, driver) and _loss(item) > 0 for item in _person_kpis(row))
    ]
    affected = list(dict.fromkeys(affected))
    if driver:
        _write(slide, "Text 55", f"{_clean(driver.get('label'), 'KPI')}: weighted loss {_pct(_loss(driver))} across {', '.join(affected) if affected else 'employee-level detail unavailable'}.")
    else:
        _write(slide, "Text 55", "No measurable KPI evidence is available for the selected Marketing dataset.")
    _fill(slide, "Shape 54", DARK_NAVY)


def _append_driver_impact(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int = 1, total_parts: int = 1) -> None:
    slide = _append_template_slide(prs, source, 5)
    period = _period(payload)
    suffix = f" — {part}/{total_parts}" if total_parts > 1 else ""
    _set_common(slide, f"Performance Driver Impact{suffix}", "Weighted performance loss ranks the management action order", period, page)
    _write(slide, "Text 8", "Weighted score loss by KPI")
    max_loss = max([_loss(kpi) for kpi in rows] + [1.0])
    text_rows = [(10, 12, 13), (14, 17, 18), (19, 22, 23), (24, 27, 28), (29, 32, 33), (34, 37, 38), (39, 42, 43), (44, 47, 48), (49, 52, 53), (54, 57, 58)]
    shape_rows = [(11, None), (15, 16), (20, 21), (25, 26), (30, 31), (35, 36), (40, 41), (45, 46), (50, 51), (55, 56)]
    for index, (label_name, loss_name, context_name) in enumerate(text_rows):
        kpi = rows[index] if index < len(rows) else {}
        fill_shape_name, track_shape_name = shape_rows[index]
        if kpi:
            status = _kpi_status(kpi)
            loss = _loss(kpi)
            affected_count = kpi.get("affected_count")
            if affected_count is None:
                affected_count = sum(1 for person in _people(payload) if any(_kpi_key_matches(item, kpi) and _loss(item) > 0 for item in _person_kpis(person)))
            _write(slide, f"Text {label_name}", _fit(kpi.get("label"), 28))
            _write(slide, f"Text {loss_name}", _pct(loss), color=_status_color(status))
            _write(slide, f"Text {context_name}", f"{int(affected_count)} affected | {_status_code(status)}")
            ratio = max(0.0, min(1.0, loss / max_loss))
            fill_shape = _shape(slide, f"Shape {fill_shape_name}")
            start = fill_shape.left / 914400
            total_width = fill_shape.width / 914400
            color = _status_color(status)
            _fill(slide, f"Shape {fill_shape_name}", color)
            if track_shape_name is not None:
                track_shape = _shape(slide, f"Shape {track_shape_name}")
                total_width = (fill_shape.width + track_shape.width) / 914400
                _fill(slide, f"Shape {track_shape_name}", TRACK)
                _set_width(fill_shape, total_width * ratio)
                _set_width(track_shape, total_width * (1 - ratio), left=start + total_width * ratio)
            else:
                _set_width(fill_shape, total_width * ratio)
        else:
            _write(slide, f"Text {label_name}", "KPI not available")
            _write(slide, f"Text {loss_name}", MISSING)
            _write(slide, f"Text {context_name}", "—")
            _fill(slide, f"Shape {fill_shape_name}", TRACK)
            if track_shape_name is not None:
                _fill(slide, f"Shape {track_shape_name}", TRACK)
    driver = _top_driver(payload)
    _write(slide, "Text 60", f"Reading rule: solve the highest weighted-loss drivers first; { _clean(driver.get('label'), 'no leading KPI') if driver else 'no measurable driver'} is currently ranked first.")
    _fill(slide, "Shape 59", PALE_AMBER)


def _append_driver_map(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int = 1, total_parts: int = 1) -> None:
    slide = _append_template_slide(prs, source, 6)
    period = _period(payload)
    suffix = f" — {part}/{total_parts}" if total_parts > 1 else ""
    _set_common(slide, f"Driver → Role → Employee Map{suffix}", "Connects overall KPI loss to affected people and functions", period, page)
    _write(slide, "Text 8", "Priority driver map")
    text_rows = [(19, 21, 23, 25), (27, 29, 31, 33), (35, 37, 39, 41), (43, 45, 47, 49)]
    shape_rows = [[f"Shape {18 + row * 8 + col * 2}" for col in range(4)] for row in range(4)]
    for index, (driver_name, role_name, employee_name, action_name) in enumerate(text_rows):
        row = rows[index] if index < len(rows) else {}
        if row:
            _write(slide, f"Text {driver_name}", f"{_fit(row.get('driver'), 24)} / {_pct(row.get('loss'))} loss")
            _write(slide, f"Text {role_name}", _fit(row.get("role"), 24))
            _write(slide, f"Text {employee_name}", _fit(row.get("employees"), 44))
            _write(slide, f"Text {action_name}", _fit(row.get("action"), 64))
            state = "Requires action" if _loss(row) > 0 else "On track"
        else:
            for name in (driver_name, role_name, employee_name, action_name):
                _write(slide, f"Text {name}", MISSING)
            state = "Data unavailable"
        for shape_name in shape_rows[index]:
            _fill(slide, shape_name, PALE_RED if index == 0 and row else (PALE_NEUTRAL if index % 2 else WHITE) if row else PALE_NEUTRAL)
    _write(slide, "Text 51", "Management interpretation: overall score is not the story; the story is which driver touches which role, then which employees require intervention.")
    _fill(slide, "Shape 50", DARK_NAVY)


def _priority_row_names() -> tuple[list[tuple[int, int, int]], list[tuple[int, int, int]]]:
    action = [(12 + i * 3, 13 + i * 3, 14 + i * 3) for i in range(5)]
    watch = [(31, 32, 33)]
    track = [(38 + i * 3, 39 + i * 3, 40 + i * 3) for i in range(6)]
    return action + watch, track


def _write_priority_column(slide: Any, members: list[dict[str, Any]], title: str, column: str, start: int, max_rows: int) -> None:
    for index in range(max_rows):
        row = members[index] if index < len(members) else {}
        text_names = (start + index * 3, start + index * 3 + 1, start + index * 3 + 2)
        if row:
            primary = _person_kpis(row)[0] if _person_kpis(row) else {}
            _write(slide, f"Text {text_names[0]}", _fit(row.get("name"), 24))
            _write(slide, f"Text {text_names[1]}", _pct(row.get("score")), color=_status_color(_employee_status(row.get("score"))))
            _write(slide, f"Text {text_names[2]}", _mom(row.get("movement")), color=base._movement_color(row.get("movement")))
        else:
            for name in text_names:
                _write(slide, f"Text {name}", MISSING)


def _append_priority(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _append_template_slide(prs, source, 7)
    period = _period(payload)
    _set_common(slide, "Employee Performance Priority", "Exceptions first; on-track people stay visible but secondary", period, page)
    groups = base._employee_priority_groups(payload)
    action = groups["Requires Action"]
    watch = groups["Watch"]
    track = groups["On Track"]
    _write(slide, "Text 10", "Requires Action")
    _write(slide, "Text 11", f"{len(action)} employee(s) | Score below 70%")
    _write(slide, "Text 29", "Watch")
    _write(slide, "Text 30", f"{len(watch)} employee(s) | 70–89%")
    _write(slide, "Text 36", "On Track")
    _write(slide, "Text 37", f"{len(track)} employee(s) | 90%+")
    _write_priority_column(slide, action, "Requires Action", "action", 12, 5)
    _write_priority_column(slide, watch, "Watch", "watch", 31, 1)
    _write_priority_column(slide, track, "On Track", "track", 38, 6)
    # Restore the semantic panel colors even when the live data has different counts.
    _fill(slide, "Shape 9", RED)
    _fill(slide, "Shape 28", AMBER)
    _fill(slide, "Shape 35", GREEN)
    names = ", ".join(_fit(person.get("name"), 22) for person in action[:5])
    _write(slide, "Text 57", f"Immediate focus: {len(action)} employees below 70% — {names or 'none recorded'}.")


def _detail_slots(source_index: int) -> list[tuple[int, int, int, int, int, int, int, int]]:
    """Return the text/progress slots for an individual detail source slide."""

    count = 5 if source_index == 8 else 4
    return [
        (27 + index * 11, 28 + index * 11, 29 + index * 11, 30 + index * 11, 33 + index * 11, 34 + index * 11, 35 + index * 11, 31 + index * 11)
        for index in range(count)
    ]


def _detail_progress_shapes(source_index: int) -> list[tuple[int, int]]:
    count = 5 if source_index == 8 else 4
    return [(31 + index * 11, 32 + index * 11) for index in range(count)]


def _detail_card_shapes(source_index: int) -> list[tuple[int, int, int]]:
    count = 5 if source_index == 8 else 4
    return [(25 + index * 11, 26 + index * 11, 25 + index * 11) for index in range(count)]


def _detail_action_text(person: dict[str, Any]) -> str:
    actions = person.get("actions") or []
    if actions:
        return base._action_text(actions[0])
    primary = _person_kpis(person)[0] if _person_kpis(person) else {}
    if primary:
        return f"{_clean(primary.get('label'), 'KPI')} gap review; owner and due date require validation."
    return "No action recorded; assign a workstream if intervention is required."


def _set_progress(slide: Any, fill_name: str, track_name: str, ratio: float, color: RGBColor) -> None:
    fill_shape = _shape(slide, f"Shape {fill_name}")
    track_shape = _shape(slide, f"Shape {track_name}")
    start = fill_shape.left / 914400
    total_width = (fill_shape.width + track_shape.width) / 914400
    ratio = max(0.0, min(1.0, ratio))
    _fill(slide, f"Shape {fill_name}", color)
    _fill(slide, f"Shape {track_name}", TRACK)
    _set_width(fill_shape, total_width * ratio)
    _set_width(track_shape, total_width * (1 - ratio), left=start + total_width * ratio)


def _write_detail_kpi(slide: Any, kpi: dict[str, Any], slot: tuple[int, int, int, int, int, int, int, int], progress: tuple[int, int], card_shapes: tuple[int, int, int]) -> None:
    title_name, achievement_name, status_name, gap_name, actual_name, target_name, direction_name, _ = slot
    status = _kpi_status(kpi)
    loss = _loss(kpi)
    rail = _status_color(status) if loss > 0 else GREEN
    fill = _status_fill(status) if loss > 0 else PALE_GREEN
    _write(slide, f"Text {title_name}", _fit(kpi.get("label"), 24))
    _write(slide, f"Text {achievement_name}", _pct(kpi.get("achievement_pct")), color=rail)
    _write(slide, f"Text {status_name}", f"{_status_code(status)} | Achiev.")
    _write(slide, f"Text {gap_name}", f"Gap {_pct(kpi.get('shortfall'))} / Loss {_pct(loss)}", color=RED if loss > 0 else GREEN)
    _write(slide, f"Text {actual_name}", f"Actual: {_native(kpi.get('actual'), kpi.get('unit'))}")
    _write(slide, f"Text {target_name}", f"Target: {_native(kpi.get('target'), kpi.get('unit'))}")
    _write(slide, f"Text {direction_name}", base._direction_label(kpi.get("direction")))
    _fill(slide, f"Shape {card_shapes[0]}", fill)
    _fill(slide, f"Shape {card_shapes[1]}", rail)
    achievement = _num(kpi.get("achievement_pct"))
    _set_progress(slide, str(progress[0]), str(progress[1]), (achievement or 0) / 100.0 if achievement is not None else 0, rail)


def _append_detail(prs: Presentation, source: Presentation, payload: dict[str, Any], person: dict[str, Any], page: int, sequence: int, source_index: int) -> None:
    slide = _append_template_slide(prs, source, source_index)
    period = _period(payload)
    status = _employee_status(person.get("score"))
    name = _clean(person.get("name"), "Employee detail")
    position = _clean(person.get("position"), "Position not available")
    letter = chr(64 + sequence) if sequence <= 26 else str(sequence)
    _set_common(slide, f"6{letter}. Employee Detail: {_fit(name, 45)}", "Actual vs Target is shown for every KPI so the employee number is clear and auditable", period, page)
    _write(slide, "Text 11", status.upper(), color=WHITE)
    _write(slide, "Text 12", _fit(name, 24))
    _write(slide, "Text 13", _fit(f"{position} · Priority P{_priority_rank(payload, person)}", 42))
    _write(slide, "Text 14", "CURRENT SCORE")
    _write(slide, "Text 15", _pct(person.get("score")), color=_status_color(status))
    _write(slide, "Text 16", "BASELINE")
    _write(slide, "Text 17", _pct(person.get("baseline_score")))
    _write(slide, "Text 18", "MOM")
    _write(slide, "Text 19", _mom(person.get("movement")), color=base._movement_color(person.get("movement")))
    _write(slide, "Text 21", "Action focus")
    _write(slide, "Text 22", _fit(_detail_action_text(person), 58))
    _fill(slide, "Shape 9", _status_color(status))
    _fill(slide, "Shape 10", _status_color(status))
    _fill(slide, "Shape 20", PALE_RED if status == "Requires action" else PALE_AMBER)

    kpis = _person_kpis(person)
    _write(slide, "Text 23", "KPI Drivers — Actual / Target / Achievement / Weighted Loss")
    slots = _detail_slots(source_index)
    progress = _detail_progress_shapes(source_index)
    cards = _detail_card_shapes(source_index)
    for index, slot in enumerate(slots):
        if index < len(kpis):
            _write_detail_kpi(slide, kpis[index], slot, progress[index], cards[index])
        else:
            for text_id in slot[:7]:
                _write(slide, f"Text {text_id}", MISSING)
            _fill(slide, f"Shape {cards[index][0]}", PALE_NEUTRAL)
            _fill(slide, f"Shape {cards[index][1]}", BLUE)
            _set_progress(slide, str(progress[index][0]), str(progress[index][1]), 0, BLUE)
    _write(slide, "Text 81" if source_index == 8 else "Text 70", "Review rule: prioritize the highest weighted-loss KPI first; use Actual / Target to validate the gap before assigning corrective action.")


def _group_person_text(slide: Any, person: dict[str, Any], *, top: bool, payload: dict[str, Any]) -> None:
    prefix = "" if top else ""  # The source uses a second identical profile frame.
    names = (11, 12, 13, 14, 15, 16, 17) if top else (54, 55, 56, 57, 58, 59, 60)
    name_name, position_name, score_label, score_name, mom_label, mom_name, action_name = names
    _write(slide, f"Text {name_name}", _fit(person.get("name"), 24))
    _write(slide, f"Text {position_name}", _fit(person.get("position"), 27))
    _write(slide, f"Text {score_label}", "MAIN SCORE")
    _write(slide, f"Text {score_name}", _pct(person.get("score")), color=_status_color(_employee_status(person.get("score"))))
    _write(slide, f"Text {mom_label}", "MOM")
    _write(slide, f"Text {mom_name}", _mom(person.get("movement")), color=base._movement_color(person.get("movement")))
    primary = _person_kpis(person)[0] if _person_kpis(person) else {}
    _write(slide, f"Text {action_name}", f"Action focus: {_fit(primary.get('label'), 25) if primary else 'KPI'} gap review.")


def _write_group_kpi(slide: Any, kpi: dict[str, Any], names: tuple[int, int, int, int, int, int, int], card_shape: tuple[int, int], progress: tuple[int, int], people: list[dict[str, Any]]) -> None:
    label_name, achievement_name, status_label_name, gap_name, actual_name, target_name, direction_name = names
    status = _kpi_status(kpi)
    loss = _loss(kpi)
    rail = _status_color(status) if loss > 0 else GREEN
    _write(slide, f"Text {label_name}", _fit(kpi.get("label"), 24))
    _write(slide, f"Text {achievement_name}", _pct(kpi.get("achievement_pct")), color=rail)
    _write(slide, f"Text {status_label_name}", "Achievement")
    _write(slide, f"Text {gap_name}", f"Gap {_pct(kpi.get('shortfall'))} / Loss {_pct(loss)}", color=RED if loss > 0 else GREEN)
    values = []
    for person in people:
        match = next((item for item in _person_kpis(person) if _kpi_key_matches(item, kpi)), {})
        values.append(_native(match.get("actual"), match.get("unit")))
    actual = " / ".join(values) if values else _native(kpi.get("actual"), kpi.get("unit"))
    _write(slide, f"Text {actual_name}", f"Actual: {actual}")
    _write(slide, f"Text {target_name}", f"Target: {_native(kpi.get('target'), kpi.get('unit'))}")
    _write(slide, f"Text {direction_name}", base._direction_label(kpi.get("direction")))
    _fill(slide, f"Shape {card_shape[0]}", _status_fill(status) if loss > 0 else PALE_GREEN)
    _fill(slide, f"Shape {card_shape[1]}", rail)
    _set_progress(slide, str(progress[0]), str(progress[1]), (_num(kpi.get("achievement_pct")) or 0) / 100.0, rail)


def _append_group_detail(prs: Presentation, source: Presentation, payload: dict[str, Any], people: list[dict[str, Any]], page: int, sequence: int) -> None:
    slide = _append_template_slide(prs, source, 10)
    period = _period(payload)
    position = _fit(people[0].get("position"), 30) if people else "Shared role"
    _set_common(slide, f"6{chr(64 + sequence)}. Employee Detail: {position}s", "Grouped view for employees sharing the same leading KPI pattern", period, page)
    _group_person_text(slide, people[0], top=True, payload=payload)
    if len(people) > 1:
        _group_person_text(slide, people[1], top=False, payload=payload)
    else:
        for name in (54, 55, 56, 57, 58, 59, 60):
            _write(slide, f"Text {name}", MISSING)
    kpi_candidates = _person_kpis(people[0]) if people else []
    selected: list[dict[str, Any]] = []
    for kpi in kpi_candidates:
        if _loss(kpi) > 0 or len(selected) < 3:
            selected.append(kpi)
        if len(selected) == 3:
            break
    top_slots = [(20, 21, 22, 23, 26, 27, 28), (31, 32, 33, 34, 37, 38, 39), (42, 43, 44, 45, 48, 49, 50)]
    bottom_slots = [(63, 64, 65, 66, 69, 70, 71), (74, 75, 76, 77, 80, 81, 82), (85, 86, 87, 88, 91, 92, 93)]
    top_cards = [(18, 19), (29, 30), (40, 41)]
    bottom_cards = [(61, 62), (72, 73), (83, 84)]
    top_progress = [(24, 25), (35, 36), (46, 47)]
    bottom_progress = [(67, 68), (78, 79), (89, 90)]
    for index in range(3):
        kpi = selected[index] if index < len(selected) else {}
        if kpi:
            _write_group_kpi(slide, kpi, top_slots[index], top_cards[index], top_progress[index], people)
            _write_group_kpi(slide, kpi, bottom_slots[index], bottom_cards[index], bottom_progress[index], people)
        else:
            for name in (*top_slots[index], *bottom_slots[index]):
                _write(slide, f"Text {name}", MISSING)
    _write(slide, "Text 3", "Grouped view keeps shared KPI drivers together while showing each employee's score and MoM separately.")


def _append_evidence(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int, total_parts: int) -> None:
    slide = _append_template_slide(prs, source, 12)
    period = _period(payload)
    suffix = f" — {part}/{total_parts}" if total_parts > 1 else ""
    _set_common(slide, f"Root Cause / Evidence{suffix}", "Evidence first: result, target, achievement, weighted loss and validation need", period, page)
    headers = [(9, "Employee"), (11, "Role"), (13, "KPI"), (15, "Actual"), (17, "Target"), (19, "Ach."), (21, "Loss"), (23, "Evidence / operating issue")]
    for name, value in headers:
        _write(slide, f"Text {name}", value)
    row_text_start = 25
    row_shape_start = 24
    for index in range(9):
        row = rows[index] if index < len(rows) else {}
        text_names = [row_text_start + index * 16 + offset for offset in (0, 2, 4, 6, 8, 10, 12, 14)]
        shape_names = [row_shape_start + index * 16 + offset for offset in (0, 2, 4, 6, 8, 10, 12, 14)]
        if row:
            values = [row.get("employee"), row.get("role"), row.get("kpi"), row.get("actual"), row.get("target"), row.get("achievement"), row.get("loss"), row.get("evidence")]
            state = row.get("state", "Investigation Required")
            fill = PALE_RED if "investigation" in state.casefold() else PALE_AMBER if "evidence" in state.casefold() else PALE_GREEN
        else:
            values = [MISSING] * 8
            fill = PALE_NEUTRAL
        for text_name, value in zip(text_names, values):
            _write(slide, f"Text {text_name}", _fit(value, 60) if text_name in {text_names[-1], text_names[0], text_names[1], text_names[2]} else value)
        for shape_name in shape_names:
            _fill(slide, f"Shape {shape_name}", fill)
    _write(slide, "Text 169", "Important: this slide separates evidence from assumption. The KPI identifies where to investigate; it does not invent an operational root cause.")
    _fill(slide, "Shape 168", PALE_AMBER)


def _append_actions(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int, total_parts: int) -> None:
    slide = _append_template_slide(prs, source, 13)
    period = _period(payload)
    suffix = f" — {part}/{total_parts}" if total_parts > 1 else ""
    _set_common(slide, f"Corrective Action Tracker{suffix}", "Data-driven workstreams created from the highest weighted-loss gaps", period, page)
    for name, value in ((9, "Priority"), (11, "Workstream"), (13, "Owner / Due"), (15, "Scope"), (17, "Execution / status"), (19, "Success measure")):
        _write(slide, f"Text {name}", value)
    for index in range(4):
        row = rows[index] if index < len(rows) else {}
        text_names = [21 + index * 12 + offset for offset in (0, 2, 4, 6, 8, 10)]
        shape_names = [20 + index * 12 + offset for offset in (0, 2, 4, 6, 8, 10)]
        if row:
            values = [row.get("priority"), row.get("workstream"), row.get("owner_due"), row.get("scope"), row.get("execution"), row.get("success")]
            fill = PALE_RED if row.get("missing") else PALE_AMBER if row.get("priority") == "High" else PALE_NEUTRAL
        else:
            values = [MISSING] * 6
            fill = PALE_NEUTRAL
        for text_name, value in zip(text_names, values):
            _write(slide, f"Text {text_name}", _fit(value, 68))
        for shape_name in shape_names:
            _fill(slide, f"Shape {shape_name}", fill)
    missing_owner = sum(1 for row in _action_rows(payload) if "Owner needed" in row.get("owner_due", ""))
    missing_due = sum(1 for row in _action_rows(payload) if "Due: Due date needed" in row.get("owner_due", ""))
    _write(slide, "Text 69", f"Decision needed before next review: confirm owners and due dates for every workstream; missing owners {missing_owner}, missing due dates {missing_due}.")
    _fill(slide, "Shape 68", DARK_NAVY)


def _append_next_review(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int) -> None:
    slide = _append_template_slide(prs, source, 14)
    period = _period(payload)
    _set_common(slide, "Management Summary / Next Review", "Next Review Commitments | The next review is successful only when targets, owners and evidence are explicit", period, page)
    kpis = _kpis(payload)
    driver = kpis[0] if kpis else {}
    second = kpis[1] if len(kpis) > 1 else {}
    actions = _action_rows(payload)
    missing_owner = sum(1 for row in actions if "Owner needed" in row.get("owner_due", ""))
    missing_due = sum(1 for row in actions if "Due: Due date needed" in row.get("owner_due", ""))
    below = base._below_threshold(payload)
    evidence = _employee_driver_rows(payload)
    confirmed = sum(1 for row in evidence if row.get("state") == "Confirmed Root Cause")
    investigation = sum(1 for row in evidence if row.get("state") == "Investigation Required")
    cards = [
        ("OVERALL PERFORMANCE", f"{_pct(payload.get('overall_score'))} | Target {_pct(payload.get('target_score'))}", f"Current vs target gap: {_pct(abs((_num(payload.get('overall_score')) or 0) - (_num(payload.get('target_score')) or 0)))}", BLUE, WHITE),
        (_fit(driver.get("label"), 22, "LEADING KPI" ).upper(), f"{_native(driver.get('actual'), driver.get('unit'))} | Target {_native(driver.get('target'), driver.get('unit'))}", f"Weighted loss {_pct(_loss(driver))} | {base._direction_label(driver.get('direction'))}", RED, PALE_RED),
        (_fit(second.get("label"), 22, "SECOND KPI" ).upper(), f"{_native(second.get('actual'), second.get('unit'))} | Target {_native(second.get('target'), second.get('unit'))}", f"Weighted loss {_pct(_loss(second))}", RED if _loss(second) else GREEN, PALE_RED if _loss(second) else PALE_GREEN),
        ("EMPLOYEES BELOW 70%", f"{len(below)} below threshold", "Priority list is carried into the employee detail and appendix views.", RED, PALE_RED),
        ("OPEN ACTIONS / OWNERS", f"{len(actions)} open | {missing_owner} owner gaps", f"Due dates missing: {missing_due}", AMBER if actions else GREEN, PALE_AMBER),
        ("EVIDENCE QUALITY", f"{confirmed} confirmed / {investigation} validate", "Unknown causes remain explicitly flagged for investigation.", GREEN if investigation == 0 else AMBER, PALE_GREEN if investigation == 0 else PALE_AMBER),
    ]
    starts = [(10, 11, 12), (15, 16, 17), (20, 21, 22), (25, 26, 27), (30, 31, 32), (35, 36, 37)]
    shapes = [("Shape 8", "Shape 9"), ("Shape 13", "Shape 14"), ("Shape 18", "Shape 19"), ("Shape 23", "Shape 24"), ("Shape 28", "Shape 29"), ("Shape 33", "Shape 34")]
    for (label, value, detail, rail, card_fill), (label_name, value_name, detail_name), (card_name, rail_name) in zip(cards, starts, shapes):
        _write(slide, f"Text {label_name}", label)
        _write(slide, f"Text {value_name}", value)
        _write(slide, f"Text {detail_name}", detail)
        _fill(slide, card_name, card_fill)
        _fill(slide, rail_name, rail)
    workstreams = ", ".join(dict.fromkeys(row.get("workstream", "workstream") for row in actions[:4])) or "the recorded workstreams"
    _write(slide, "Text 38", "Management decision summary")
    _write(slide, "Text 41", f"Before the next review: confirm owners and due dates for {workstreams}; validate each KPI against target and close only with recorded success evidence.")
    _fill(slide, "Shape 40", PALE_AMBER if missing_owner or missing_due else PALE_GREEN)


def _employee_appendix_rows(payload: dict[str, Any], people: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    rows = []
    for person in people:
        person_kpis = _person_kpis(person)
        primary = person_kpis[0] if person_kpis else {}
        status = _employee_status(person.get("score"))
        rows.append(
            {
                "employee": _fit(person.get("name"), 28),
                "role": _fit(person.get("position"), 24),
                "score": _pct(person.get("score")),
                "mom": _mom(person.get("movement")),
                "status": status.upper(),
                "status_key": status,
                "driver": f"{_fit(primary.get('label'), 30)} | {_pct(_loss(primary))}" if primary else "No leading KPI",
                "actual": _native(primary.get("actual"), primary.get("unit")) if primary else MISSING,
                "target": _native(primary.get("target"), primary.get("unit")) if primary else MISSING,
                "focus": "Review KPI gap" if status in {"Requires action", "Watch"} else "Monitor gap",
            }
        )
    return rows


def _append_employee_appendix(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int, total_parts: int) -> None:
    slide = _append_template_slide(prs, source, 15)
    period = _period(payload)
    suffix = f" — {part}/{total_parts}" if total_parts > 1 else ""
    _set_common(slide, f"Appendix: Full Employee Scorecard{suffix}", "Employee Detail Appendix — Full Scorecard | score, MoM, status, top KPI, weighted loss and actual / target", period, page)
    for name, value in ((9, "Employee"), (11, "Role"), (13, "Score"), (15, "MoM"), (17, "Status"), (19, "Top KPI | loss"), (21, "Actual"), (23, "Target"), (25, "Focus")):
        _write(slide, f"Text {name}", value)
    for index in range(12):
        row = rows[index] if index < len(rows) else {}
        text_names = [27 + index * 18 + offset for offset in (0, 2, 4, 6, 8, 10, 12, 14, 16)]
        shape_names = [26 + index * 18 + offset for offset in (0, 2, 4, 6, 8, 10, 12, 14, 16)]
        if row:
            values = [row.get("employee"), row.get("role"), row.get("score"), row.get("mom"), row.get("status"), row.get("driver"), row.get("actual"), row.get("target"), row.get("focus")]
            status = row.get("status_key", "Data unavailable")
            fill = PALE_RED if status == "Requires action" else PALE_AMBER if status == "Watch" else PALE_NEUTRAL if index % 2 == 0 else WHITE
        else:
            values = [MISSING] * 9
            fill = PALE_NEUTRAL
        for text_name, value in zip(text_names, values):
            _write(slide, f"Text {text_name}", value)
        for shape_name in shape_names:
            _fill(slide, f"Shape {shape_name}", fill)


def _kpi_appendix_rows(payload: dict[str, Any], rows: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    result = []
    for kpi in rows:
        action = base._kpi_action(payload, kpi)
        status = _kpi_status(kpi)
        result.append(
            {
                "kpi": _fit(kpi.get("label"), 30),
                "direction": base._direction_label(kpi.get("direction")),
                "actual": _native(kpi.get("actual"), kpi.get("unit")),
                "baseline": _native(kpi.get("baseline_actual"), kpi.get("unit")),
                "mom": _mom(kpi.get("mom")),
                "target": _native(kpi.get("target"), kpi.get("unit")),
                "achievement": f"{_pct(kpi.get('achievement_pct'))}\nGap {_pct(kpi.get('shortfall'))}",
                "loss": f"{_pct(_loss(kpi))}\n{_status_code(status)}",
                "action": "Open" if action else MISSING,
                "status_key": status,
            }
        )
    return result


def _append_kpi_appendix(prs: Presentation, source: Presentation, payload: dict[str, Any], page: int, rows: list[dict[str, Any]], part: int, total_parts: int) -> None:
    slide = _append_template_slide(prs, source, 15)
    period = _period(payload)
    suffix = f" — {part}/{total_parts}" if total_parts > 1 else ""
    _set_common(slide, f"Appendix: Full KPI Reference{suffix}", "KPI Reference Appendix — Full KPI List | Actual, baseline, MoM, target, achievement gap, weighted loss, direction, status and action reference", period, page)
    headers = [(9, "KPI"), (11, "Direction"), (13, "Actual"), (15, "Baseline"), (17, "MoM"), (19, "Target"), (21, "Achiev. / gap"), (23, "Loss / status"), (25, "Action")]
    for name, value in headers:
        _write(slide, f"Text {name}", value)
    for index in range(12):
        row = rows[index] if index < len(rows) else {}
        text_names = [27 + index * 18 + offset for offset in (0, 2, 4, 6, 8, 10, 12, 14, 16)]
        shape_names = [26 + index * 18 + offset for offset in (0, 2, 4, 6, 8, 10, 12, 14, 16)]
        if row:
            values = [row.get("kpi"), row.get("direction"), row.get("actual"), row.get("baseline"), row.get("mom"), row.get("target"), row.get("achievement"), row.get("loss"), row.get("action")]
            status = row.get("status_key", "Data unavailable")
            fill = PALE_RED if status in {"Requires action", "Critical"} else PALE_AMBER if status in {"Watch", "At risk"} else PALE_NEUTRAL if index % 2 == 0 else WHITE
        else:
            values = [MISSING] * 9
            fill = PALE_NEUTRAL
        for text_name, value in zip(text_names, values):
            _write(slide, f"Text {text_name}", value)
        for shape_name in shape_names:
            _fill(slide, f"Shape {shape_name}", fill)


def build_marketing_reference_pptx(period_label: str = "June 2026", report_data: dict[str, Any] | None = None) -> bytes:
    """Build the Marketing report by filling the approved reference frames."""

    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Marketing reference template is missing: {TEMPLATE_PATH}")

    payload = base._prepare_snapshot(period_label, report_data)
    template_bytes = TEMPLATE_PATH.read_bytes()
    source = Presentation(io.BytesIO(template_bytes))
    prs = Presentation(io.BytesIO(template_bytes))
    _clear_slides(prs)
    prs.core_properties.title = f"Marketing Team Performance Review - {_period(payload)}"
    prs.core_properties.subject = "Marketing PMS performance review"
    prs.core_properties.author = "PMS Dashboard"

    page = 1
    _append_cover(prs, source, payload, page)
    page += 1
    _append_overview(prs, source, payload, page)
    page += 1

    role_rows = _role_rows_with_people(payload)
    role_chunks = [role_rows[index : index + 9] for index in range(0, len(role_rows), 9)] or [[]]
    for index, chunk in enumerate(role_chunks, 1):
        _append_roles(prs, source, payload, page, chunk, index, len(role_chunks))
        page += 1

    all_history = _history(payload)
    history_count = len(all_history)
    primary_history = all_history[-3:] if len(all_history) > 3 else all_history
    trend_slices = [primary_history]
    if len(all_history) > 3:
        earlier = all_history[:-3]
        trend_slices.extend(earlier[index : index + 3] for index in range(0, len(earlier), 3))
    for index, chunk in enumerate(trend_slices, 1):
        _append_trend(prs, source, payload, page, chunk, history_count, index, len(trend_slices))
        page += 1

    kpis = _kpis(payload)
    # The reference story uses one primary page for each of these sections.
    # The highest-impact rows receive the executive treatment; the complete
    # list remains available in the appendix so a larger live KPI catalogue
    # does not push the main narrative out of the approved sequence.
    _append_kpi_health(prs, source, payload, page, kpis[:6], 1, 1)
    page += 1

    _append_driver_impact(prs, source, payload, page, kpis[:10], 1, 1)
    page += 1

    bridge_rows = base._driver_role_employee_rows(payload)
    _append_driver_map(prs, source, payload, page, bridge_rows[:4], 1, 1)
    page += 1

    _append_priority(prs, source, payload, page)
    page += 1

    detail_people = _detail_order(payload)
    consumed: set[str] = set()
    detail_sequence = 1
    for index, person in enumerate(detail_people):
        person_id = _clean(person.get("employee_id"), f"person-{index}")
        if person_id in consumed:
            continue
        same_position = [
            candidate
            for candidate in detail_people[index:]
            if _clean(candidate.get("position"), "Unassigned").casefold() == _clean(person.get("position"), "Unassigned").casefold()
            and _clean(candidate.get("employee_id"), "") not in consumed
        ]
        if len(same_position) == 2 and "social media specialist" in _clean(person.get("position"), "").casefold():
            _append_group_detail(prs, source, payload, same_position, page, detail_sequence)
            consumed.update(_clean(item.get("employee_id"), "") for item in same_position)
            page += 1
            detail_sequence += 1
            continue
        # Keep the approved detail composition: the first attention profile
        # uses the five-card frame, the second uses the four-card frame, and
        # later individual profiles use the final four-card frame.  Additional
        # employee KPIs remain available in the scorecard/KPI appendices.
        source_index = 8 if detail_sequence == 1 else 9 if detail_sequence == 2 else 11
        _append_detail(prs, source, payload, person, page, detail_sequence, source_index)
        consumed.add(person_id)
        page += 1
        detail_sequence += 1

    evidence_rows = _employee_driver_rows(payload)
    _append_evidence(prs, source, payload, page, evidence_rows[:9], 1, 1)
    page += 1

    action_rows = _action_rows_or_placeholder(payload)
    action_chunks = [action_rows[index : index + 4] for index in range(0, len(action_rows), 4)] or [[]]
    for index, chunk in enumerate(action_chunks, 1):
        _append_actions(prs, source, payload, page, chunk, index, len(action_chunks))
        page += 1

    _append_next_review(prs, source, payload, page)
    page += 1

    people_rows = _employee_appendix_rows(payload, _appendix_order(payload))
    people_chunks = [people_rows[index : index + 12] for index in range(0, len(people_rows), 12)] or [[]]
    for index, chunk in enumerate(people_chunks, 1):
        _append_employee_appendix(prs, source, payload, page, chunk, index, len(people_chunks))
        page += 1

    kpi_rows = _kpi_appendix_rows(payload, kpis)
    kpi_appendix_chunks = [kpi_rows[index : index + 12] for index in range(0, len(kpi_rows), 12)] or [[]]
    for index, chunk in enumerate(kpi_appendix_chunks, 1):
        _append_kpi_appendix(prs, source, payload, page, chunk, index, len(kpi_appendix_chunks))
        page += 1

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


__all__ = ["build_marketing_reference_pptx", "TEMPLATE_PATH"]
