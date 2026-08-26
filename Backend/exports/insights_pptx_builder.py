"""Executive Insights PowerPoint export.

The builder is intentionally presentation-only.  ``build_insights_snapshot``
creates the one filtered, direction-aware data snapshot used by every slide;
this module only turns that snapshot into an 8-page decision narrative.
"""

from __future__ import annotations

import io
from typing import Any, Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from services.insights_report_service import build_insights_snapshot


SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5

# SGH Hub light-mode palette from the supplied story deck.
BG = RGBColor(248, 251, 255)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(11, 18, 42)
MUTED = RGBColor(71, 85, 105)
FAINT = RGBColor(100, 116, 139)
BLUE = RGBColor(36, 84, 244)
PURPLE = RGBColor(109, 94, 247)
GREEN = RGBColor(22, 163, 74)
RED = RGBColor(239, 68, 68)
PINK = RGBColor(225, 29, 72)
AMBER = RGBColor(234, 145, 24)
LINE = RGBColor(224, 232, 242)
PALE_BLUE = RGBColor(238, 245, 255)
PALE_PURPLE = RGBColor(241, 238, 255)
PALE_GREEN = RGBColor(234, 251, 241)
PALE_RED = RGBColor(255, 241, 242)
PALE_AMBER = RGBColor(255, 247, 230)
PALE_GREY = RGBColor(243, 246, 249)


def _number(value: Any) -> float | None:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return None
    return number if number == number else None


def _clean(value: Any, fallback: str = "") -> str:
    text = " ".join(str(value or fallback).split())
    return text.replace("â€¢", "|").replace("â€“", "-").replace("â€”", "-")


def _short(value: Any, limit: int = 58) -> str:
    text = _clean(value)
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "…"


def _fmt_percent(value: Any, signed: bool = False) -> str:
    number = _number(value)
    if number is None:
        return "N/A"
    return f"{number:+.1f}%" if signed else f"{number:.1f}%"


def _fmt_unavailable(value: Any, *, signed: bool = False) -> str:
    return _fmt_percent(value, signed=signed) if _number(value) is not None else "Unavailable"


def _fmt_value(value: Any, unit: Any = "") -> str:
    number = _number(value)
    if number is None:
        return "N/A"
    normalized_unit = _clean(unit)
    if normalized_unit == "%" and abs(number) <= 1:
        number *= 100
    if normalized_unit == "%":
        return f"{number:.1f}%"
    if number.is_integer():
        rendered = f"{number:,.0f}"
    else:
        rendered = f"{number:,.1f}"
    return f"{rendered} {normalized_unit}".strip()


def _shape(slide, geometry, left: float, top: float, width: float, height: float, fill=WHITE, line=LINE):
    shape = slide.shapes.add_shape(geometry, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def _text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    value: Any,
    *,
    size: float = 13,
    color=NAVY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = "Aptos",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    # Preserve intentional line breaks used by evidence, action, and people
    # tables while still normalizing whitespace within each line.
    raw_value = "" if value is None else str(value)
    cleaned_value = "\n".join(" ".join(line.split()) for line in raw_value.replace("\r\n", "\n").split("\n"))
    paragraph = frame.paragraphs[0]
    paragraph.text = cleaned_value
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _card(slide, left: float, top: float, width: float, height: float, fill=WHITE, line=LINE):
    # A very restrained shadow-like underlay keeps the rounded cards readable
    # in presentation mode without adding gradients or decoration.
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + 0.025, top + 0.025, width, height, RGBColor(239, 244, 250), RGBColor(239, 244, 250))
    return _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line)


def _header(slide, kicker: str, title: str, subtitle: str, accent=BLUE):
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, 0.10, accent, accent)
    _text(slide, 0.66, 0.28, 11.9, 0.20, kicker.upper(), size=9, color=accent, bold=True)
    # Conclusion titles should remain on one readable line.  Long dynamic
    # headlines are reduced slightly rather than allowed to collide with the
    # subtitle and the first content row.
    title_length = len(_clean(title))
    title_size = 25 if title_length <= 62 else 22 if title_length <= 72 else 19.5
    _text(slide, 0.66, 0.55, 11.95, 0.34, title, size=title_size, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _text(slide, 0.66, 1.02, 11.95, 0.24, subtitle, size=11, color=MUTED)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.66, 1.37, 12.00, 0.012, LINE, LINE)


def _footer(slide, payload: dict[str, Any], page: int):
    period = _clean(payload.get("period_label"), "Selected period")
    scope = _short(payload.get("scope_label"), 72)
    _text(slide, 0.66, 7.18, 11.1, 0.14, f"SGH Hub Insights | {period} | {scope} | Confidential", size=7.5, color=FAINT)
    page_count = payload.get("page_count")
    page_text = f"{page:02d} / {int(page_count):02d}" if page_count else f"{page:02d}"
    _text(slide, 11.75, 7.18, 0.90, 0.14, page_text, size=8, color=FAINT, bold=True, align=PP_ALIGN.RIGHT)


def _metric_card(slide, left, top, width, label, value, detail, accent=BLUE, fill=WHITE, height=1.0):
    _card(slide, left, top, width, height, fill, LINE)
    _shape(slide, MSO_SHAPE.RECTANGLE, left, top, 0.06, height, accent, accent)
    _text(slide, left + 0.18, top + 0.13, width - 0.30, 0.18, label.upper(), size=8.5, color=FAINT, bold=True)
    _text(slide, left + 0.18, top + 0.36, width - 0.30, 0.32, value, size=20, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _text(slide, left + 0.18, top + height - 0.25, width - 0.30, 0.15, detail, size=8.5, color=accent, bold=True)


def _status_color(status: str):
    normalized = _clean(status).casefold()
    if normalized in {"critical", "confirmed", "confirmed cause"}:
        return RED
    if normalized in {"at risk", "watch", "evidence recorded — cause pending confirmation", "kpi signal only", "no evidence recorded"}:
        return AMBER
    if normalized in {"on track", "completed"}:
        return GREEN
    if normalized == "data quality":
        return PURPLE
    return BLUE


def _draw_line_chart(
    slide,
    series: list[tuple[str, list[float | None], Any]],
    labels: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    minimum: float | None = None,
    maximum: float | None = None,
    target: float | None = None,
):
    values = [value for _, points, _ in series for value in points if value is not None]
    if target is not None:
        values.append(target)
    if not values:
        _text(slide, left, top + height / 2 - 0.15, width, 0.3, "No measured trend is available.", size=13, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        return
    low = minimum if minimum is not None else min(values)
    high = maximum if maximum is not None else max(values)
    if low == high:
        low -= 1
        high += 1
    if minimum is None and maximum is None:
        padding = max((high - low) * 0.12, 1.0)
        low -= padding
        high += padding
    left_axis = left + 0.55
    chart_width = width - 0.67
    bottom = top + height - 0.38
    chart_height = height - 0.58
    for tick in range(3):
        ratio = tick / 2
        y = bottom - ratio * chart_height
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left_axis), Inches(y), Inches(left_axis + chart_width), Inches(y))
        line.line.color.rgb = LINE
        line.line.width = Pt(0.7)
        tick_value = low + (high - low) * ratio
        _text(slide, left, y - 0.08, 0.50, 0.16, _fmt_value(tick_value), size=7.5, color=FAINT, align=PP_ALIGN.RIGHT)
    if target is not None:
        target_y = bottom - ((target - low) / (high - low)) * chart_height
        target_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left_axis), Inches(target_y), Inches(left_axis + chart_width), Inches(target_y))
        target_line.line.color.rgb = AMBER
        target_line.line.width = Pt(1.2)
        _text(slide, left_axis + chart_width - 0.85, target_y - 0.19, 0.85, 0.16, "Target", size=7.5, color=AMBER, bold=True, align=PP_ALIGN.RIGHT)
    for name, points, color in series:
        coords = []
        for index, point in enumerate(points):
            if point is None:
                coords.append(None)
                continue
            x = left_axis + index * chart_width / max(len(labels) - 1, 1)
            y = bottom - ((point - low) / (high - low)) * chart_height
            coords.append((x, y))
        previous = None
        for coordinate in coords:
            if coordinate is not None and previous is not None:
                connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(previous[0]), Inches(previous[1]), Inches(coordinate[0]), Inches(coordinate[1]))
                connector.line.color.rgb = color
                connector.line.width = Pt(2.1)
            if coordinate is not None:
                _shape(slide, MSO_SHAPE.OVAL, coordinate[0] - 0.045, coordinate[1] - 0.045, 0.09, 0.09, color, color)
                previous = coordinate
        _text(slide, left_axis + 0.08, top + 0.02 + 0.20 * series.index((name, points, color)), 1.7, 0.18, name, size=8.5, color=color, bold=True)
    for index, label in enumerate(labels):
        x = left_axis + index * chart_width / max(len(labels) - 1, 1)
        _text(slide, x - 0.5, bottom + 0.10, 1.0, 0.18, _short(label, 13), size=7.5, color=FAINT, align=PP_ALIGN.CENTER)


def _draw_actual_target_comparison(
    slide,
    actual: Any,
    target: Any,
    unit: Any,
    direction: str,
    left: float,
    top: float,
    width: float,
    height: float,
):
    """Render a useful current-period comparison when a line is impossible."""

    actual_number = _number(actual)
    target_number = _number(target)
    if actual_number is None or target_number is None:
        _empty_state(slide, left, top, width, height, "Actual-versus-target comparison unavailable", "The current snapshot does not contain both a measured actual and configured target.", PURPLE)
        return
    maximum = max(actual_number, target_number, 1.0) * 1.18
    bar_left = left + 1.65
    bar_width = width - 2.55
    row_height = 0.72
    actual_good = actual_number <= target_number if direction == "lower_better" else actual_number >= target_number
    actual_color = GREEN if actual_good else RED
    _text(slide, left, top + 0.06, 1.35, 0.22, "Actual", size=10.5, color=NAVY, bold=True)
    _text(slide, left, top + 0.78, 1.35, 0.22, "Target", size=10.5, color=NAVY, bold=True)
    for index, (label, value, color) in enumerate((("Actual", actual_number, actual_color), ("Target", target_number, AMBER))):
        y = top + index * row_height
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y + 0.07, bar_width, 0.22, PALE_GREY, PALE_GREY)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y + 0.07, max(0.08, bar_width * value / maximum), 0.22, color, color)
        _text(slide, bar_left + bar_width + 0.10, y + 0.02, 0.82, 0.26, _fmt_value(value, unit), size=10.5, color=color, bold=True, align=PP_ALIGN.RIGHT)
    _text(slide, bar_left, top + 1.62, bar_width, 0.20, f"0 {unit or ''}".strip(), size=8, color=FAINT)
    _text(slide, bar_left + bar_width - 1.10, top + 1.62, 1.10, 0.20, _fmt_value(maximum, unit), size=8, color=FAINT, align=PP_ALIGN.RIGHT)
    ratio = actual_number / target_number if target_number else None
    if direction == "lower_better" and ratio and ratio > 1:
        comparison = f"Current actual is {ratio:.1f}× the configured target. Lower is better."
    else:
        comparison = "Actual meets or exceeds the configured target." if actual_good else "Actual remains below the configured target."
    _text(slide, left, top + 2.02, width, 0.28, comparison, size=11, color=actual_color if not actual_good else GREEN, bold=True)


def _draw_bars(slide, rows: list[dict[str, Any]], left: float, top: float, width: float, height: float):
    if not rows:
        _text(slide, left, top + height / 2 - 0.15, width, 0.3, "No KPI loss is measurable in this scope.", size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
        return
    maximum = max((_number(row.get("weighted_impact")) or 0 for row in rows), default=1) or 1
    row_height = min(0.43, height / max(len(rows), 1))
    bar_left = left + 2.10
    bar_width = max(0.06, width - 2.72)
    axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(bar_left), Inches(top - 0.07), Inches(bar_left + bar_width), Inches(top - 0.07))
    axis.line.color.rgb = LINE
    axis.line.width = Pt(0.8)
    _text(slide, bar_left, top - 0.22, 0.55, 0.14, "0%", size=7.5, color=FAINT)
    _text(slide, bar_left + bar_width - 1.15, top - 0.22, 1.15, 0.14, f"{maximum:.1f}% max", size=7.5, color=FAINT, align=PP_ALIGN.RIGHT)
    for index, row in enumerate(rows):
        y = top + index * row_height
        label = f"P{index + 1}  {_short(row.get('label'), 27)}"
        _text(slide, left, y, 2.05, row_height, label, size=9.5, color=NAVY, bold=index < 3, valign=MSO_ANCHOR.MIDDLE)
        value = _number(row.get("weighted_impact")) or 0
        fill_width = bar_width * value / maximum
        color = RED if index == 0 else AMBER if value > 0 else FAINT
        bar_shape = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y + 0.09, fill_width, 0.18, color, color)
        bar_shape.name = f"weighted-loss-bar-{value:.4f}"
        _text(slide, left + width - 0.55, y + 0.01, 0.55, row_height, _fmt_percent(value), size=9.5, color=color, bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def _table(
    slide,
    left: float,
    top: float,
    width: float,
    row_height: float,
    columns: list[tuple[str, float, Callable[[dict[str, Any]], Any]]],
    rows: list[dict[str, Any]],
    *,
    max_rows: int = 6,
    header_fill=PALE_BLUE,
    font_size: float = 9.2,
    body_fill=WHITE,
):
    total_height = row_height * (1 + min(len(rows), max_rows))
    _card(slide, left, top, width, total_height, body_fill, LINE)
    cursor = left
    for label, column_width, _ in columns:
        _shape(slide, MSO_SHAPE.RECTANGLE, cursor, top, column_width, row_height, header_fill, header_fill)
        _text(slide, cursor + 0.06, top + 0.02, column_width - 0.12, row_height - 0.04, label.upper(), size=7.7, color=FAINT, bold=True, valign=MSO_ANCHOR.MIDDLE)
        cursor += column_width
    for index, row in enumerate(rows[:max_rows]):
        y = top + row_height * (index + 1)
        if index % 2:
            _shape(slide, MSO_SHAPE.RECTANGLE, left, y, width, row_height, RGBColor(252, 253, 255), RGBColor(252, 253, 255))
        cursor = left
        for _, column_width, value_fn in columns:
            value = value_fn(row)
            _text(slide, cursor + 0.06, y + 0.03, column_width - 0.12, row_height - 0.06, value, size=font_size, color=NAVY, bold=False, valign=MSO_ANCHOR.MIDDLE)
            cursor += column_width


def _empty_state(slide, left: float, top: float, width: float, height: float, title: str, detail: str, accent=BLUE):
    _card(slide, left, top, width, height, PALE_BLUE, accent)
    _text(slide, left + 0.30, top + 0.32, width - 0.60, 0.30, title, size=16, color=accent, bold=True)
    _text(slide, left + 0.30, top + 0.78, width - 0.60, height - 1.0, detail, size=11, color=MUTED)


def _driver_story(payload: dict[str, Any]) -> str:
    driver = payload.get("driver")
    trend = payload.get("driver_trend") or []
    if not driver:
        return "No leading KPI driver is available for the selected filters."
    actual = _fmt_value(driver.get("actual"), driver.get("unit"))
    target = _fmt_value(driver.get("target"), driver.get("unit"))
    direction = driver.get("direction")
    if direction == "lower_better":
        if trend and trend[0].get("actual") is not None and trend[-1].get("actual") is not None:
            start = _fmt_value(trend[0].get("actual"), driver.get("unit"))
            end = _fmt_value(trend[-1].get("actual"), driver.get("unit"))
            target_number = _number(driver.get("target"))
            ratio = _number(driver.get("actual")) / target_number if target_number and _number(driver.get("actual")) is not None else None
            multiplier = f" ({ratio:.1f}× target)" if ratio and ratio > 1 else ""
            history_note = "" if len(trend) >= 3 else " Historical direction is not yet established."
            return f"{driver['label']} moved from {start} to {end}, but remains above the {target} target{multiplier}.{history_note}"
        target_number = _number(driver.get("target"))
        actual_number = _number(driver.get("actual"))
        ratio = actual_number / target_number if target_number and actual_number is not None else None
        multiplier = f" Current actual is {ratio:.1f}× the configured target." if ratio and ratio > 1 else ""
        return f"{driver['label']} is {actual} against a {target} target; lower values are better.{multiplier} Historical direction is unavailable."
    if len(trend) < 2:
        return f"{driver['label']} is {actual} against a {target} target; higher values are better. Historical direction is unavailable."
    return f"{driver['label']} is {actual} against a {target} target; higher values are better."


def _management_ask(payload: dict[str, Any]) -> str:
    driver = payload.get("driver")
    if not driver:
        return "Confirm that a measurable KPI gap exists before assigning a corrective owner."
    evidence = next((row for row in payload.get("root_cause_rows", []) if row.get("key") == driver.get("key")), None)
    evidence_text = "confirm the operational cause" if not evidence or evidence.get("evidence_status") != "Confirmed cause" else "recheck the confirmed cause"
    action_text = "assign an owner and review date" if any(
        action.get("owner_display") == "Owner needed" or action.get("due_date_display") == "Due date needed"
        for action in payload.get("actions", [])
    ) else "confirm the accountable action"
    return f"Decision required: {action_text} for {driver['label']}, {evidence_text}, and measure the configured target at the next review."


def _page_one(prs: Presentation, payload: dict[str, Any], page: int = 1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    driver = payload.get("driver") or {}
    headline = payload.get("headline") or "No performance data is available for this filtered scope"
    subtitle = f"Latest available period: {payload.get('latest_period_label') or 'N/A'} | Comparison: {payload.get('comparison_period_label') or 'insufficient history'} | {_short(payload.get('scope_label'), 64)}"
    score = payload.get("overall_score")
    gap = payload.get("gap_to_target")
    movement = payload.get("movement")
    score_available = _number(score) is not None
    movement_available = _number(movement) is not None
    header_accent = RED if score_available and (gap or 0) < 0 else GREEN if score_available else MUTED
    _header(slide, "Executive decision snapshot", headline, subtitle, header_accent)
    score_color = BLUE if score_available else MUTED
    score_fill = PALE_BLUE if score_available else PALE_GREY
    gap_color = RED if score_available and (gap or 0) < 0 else GREEN if score_available else MUTED
    gap_fill = PALE_RED if score_available and (gap or 0) < 0 else PALE_GREEN if score_available else PALE_GREY
    movement_color = GREEN if movement_available and movement >= 0 else RED if movement_available else MUTED
    movement_fill = PALE_GREEN if movement_available and movement >= 0 else PALE_RED if movement_available else PALE_GREY
    population_label = "Records included" if payload.get("aggregate_only") else "Employees included"
    affected_label = "Affected records" if payload.get("aggregate_only") else "Affected employees"
    _metric_card(slide, 0.66, 1.67, 2.78, "Overall performance", _fmt_percent(score) if score_available else "Unavailable", "Target 100.0%", score_color, score_fill, 0.98)
    _metric_card(slide, 3.61, 1.67, 2.78, "Achievement gap", _fmt_percent(gap, signed=True) if score_available else "Unavailable", "Versus 100.0% target", gap_color, gap_fill, 0.98)
    _metric_card(slide, 6.56, 1.67, 2.78, "Movement vs comparison", _fmt_percent(movement, signed=True) if movement_available else "Unavailable", "No comparison period" if not movement_available else "Latest versus previous", movement_color, movement_fill, 0.98)
    _metric_card(slide, 9.51, 1.67, 3.15, population_label, str(payload.get("population_size") or 0), f"{payload.get('team_count') or 0} team(s) in scope", PURPLE, PALE_PURPLE, 0.98)

    _card(slide, 0.66, 2.91, 7.45, 1.28, WHITE, LINE)
    _text(slide, 0.94, 3.12, 2.2, 0.18, "LEADING KPI DRIVER", size=8.5, color=RED if driver else MUTED, bold=True)
    _text(slide, 0.94, 3.39, 4.35, 0.28, _short(driver.get("label") or "No measurable KPI loss", 42), size=19, bold=True)
    _text(slide, 5.55, 3.37, 2.2, 0.23, f"Weighted performance loss {_fmt_unavailable(driver.get('weighted_impact'))}", size=10.8, color=RED if driver else MUTED, bold=True, align=PP_ALIGN.RIGHT)
    driver_affected = driver.get("affected_count") if driver else None
    _text(slide, 0.94, 3.76, 6.8, 0.18, f"{affected_label} for leading KPI: {driver_affected if driver_affected is not None else 'Aggregate view'}", size=9.2, color=MUTED)

    _card(slide, 8.37, 2.91, 4.29, 1.28, PALE_AMBER, AMBER)
    _text(slide, 8.65, 3.12, 3.7, 0.18, "OPEN RECORDED ACTIONS", size=8.5, color=AMBER, bold=True)
    _text(slide, 8.65, 3.39, 3.55, 0.30, str(payload.get("open_actions") or 0), size=20, bold=True)
    _text(slide, 8.65, 3.76, 3.55, 0.18, f"{payload.get('proposed_action_count') or 0} proposed next step(s) shown", size=9.2, color=MUTED)

    _card(slide, 0.66, 4.52, 7.45, 1.80, WHITE, LINE)
    _text(slide, 0.94, 4.78, 2.7, 0.18, "DECISION STATEMENT", size=8.5, color=BLUE, bold=True)
    decision_text = (
        f"{payload.get('trend_headline') or 'Trend unavailable — only one measured period.'} "
        f"{driver.get('label') or 'The leading KPI'} remains the largest controllable gap and requires immediate ownership."
        if driver else payload.get("trend_headline") or "No performance data is available for the selected scope."
    )
    _text(slide, 0.94, 5.10, 6.75, 0.56, decision_text, size=12.3, bold=True)
    _text(slide, 0.94, 5.82, 6.75, 0.26, _short(payload.get("driver_story"), 118), size=10, color=MUTED)

    _card(slide, 8.37, 4.52, 4.29, 1.80, PALE_BLUE, BLUE)
    _text(slide, 8.65, 4.78, 3.5, 0.18, "MANAGEMENT ASK", size=8.5, color=BLUE, bold=True)
    _text(slide, 8.65, 5.10, 3.55, 0.90, _management_ask(payload), size=12.5, bold=True)
    _footer(slide, payload, page)


def _page_two(prs: Presentation, payload: dict[str, Any], page: int = 2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    _header(slide, "Performance context", payload.get("context_headline") or "Performance context is unavailable", f"Chronological score trend | {payload.get('scope_label')}", BLUE)
    _card(slide, 0.66, 1.63, 8.08, 3.63, WHITE, LINE)
    trend = payload.get("trend") or []
    if len(trend) >= 2:
        _text(slide, 0.95, 1.88, 3.3, 0.18, "OVERALL PERFORMANCE TREND", size=8.5, color=BLUE, bold=True)
        _draw_line_chart(
            slide,
            [("Actual score", [row.get("score") for row in trend], BLUE)],
            [row.get("label", "") for row in trend],
            1.06,
            2.14,
            7.15,
            2.55,
            minimum=0,
            maximum=110,
            target=100,
        )
    else:
        _text(slide, 0.95, 1.88, 4.0, 0.18, "CURRENT-PERIOD SNAPSHOT", size=8.5, color=BLUE, bold=True)
        _draw_actual_target_comparison(slide, payload.get("latest_score"), 100, "%", "higher_better", 1.00, 2.14, 7.22, 2.20)
        _text(slide, 1.00, 4.70, 2.10, 0.18, "KPI coverage", size=9.4, color=MUTED)
        _text(slide, 3.05, 4.70, 1.20, 0.18, str(len(payload.get("kpis") or [])), size=10.5, color=NAVY, bold=True)
        _text(slide, 4.55, 4.70, 1.90, 0.18, "History availability", size=9.4, color=MUTED)
        _text(slide, 6.55, 4.70, 1.25, 0.18, f"{payload.get('history_count') or 0} period(s)", size=10.5, color=FAINT, bold=True, align=PP_ALIGN.RIGHT)
    _card(slide, 8.99, 1.63, 3.67, 3.63, PALE_PURPLE, PURPLE)
    team_health = payload.get("team_health")
    if team_health:
        _text(slide, 9.28, 1.88, 2.8, 0.18, "TEAM HEALTH SNAPSHOT", size=8.5, color=PURPLE, bold=True)
        facts = [
            ("Team score", _fmt_unavailable(team_health.get("score"))),
            ("Employees", team_health.get("employee_count") or 0),
            ("Meeting target", team_health.get("employees_meeting_target") or 0),
            ("Below target", team_health.get("employees_below_target") or 0),
            ("Critical", team_health.get("critical_employees") or 0),
        ]
    else:
        _text(slide, 9.28, 1.88, 2.8, 0.18, "WHAT THE HISTORY SUPPORTS", size=8.5, color=PURPLE, bold=True)
        facts = [
            ("Current score", _fmt_unavailable(payload.get("latest_score"))),
            ("Best period", _fmt_unavailable((payload.get("best_period") or {}).get("score"))),
            ("Worst period", _fmt_unavailable((payload.get("worst_period") or {}).get("score"))),
            ("Net movement", _fmt_unavailable(payload.get("net_movement"), signed=True)),
            ("Status", payload.get("trend_status") or "No measured history"),
        ]
    for index, (label, value) in enumerate(facts):
        y = 2.28 + index * 0.52
        _text(slide, 9.28, y, 1.55, 0.18, label, size=9.5, color=MUTED)
        _text(slide, 10.80, y, 1.55, 0.18, value, size=10.5, color=NAVY if value not in {"Unavailable", "No measured history"} else FAINT, bold=True, align=PP_ALIGN.RIGHT)
    if team_health:
        _text(slide, 9.28, 4.70, 3.05, 0.22, f"Leading KPI: {team_health.get('leading_kpi')}", size=9.2, color=PURPLE, bold=True)
        _text(slide, 9.28, 4.91, 3.05, 0.16, f"Weighted performance loss: {_fmt_unavailable(team_health.get('leading_kpi_impact'))}", size=8.1, color=MUTED)
        status_summary = ", ".join(f"{key}: {value}" for key, value in (team_health.get("kpi_status_distribution") or {}).items()) or "No KPI status data"
        _text(slide, 9.28, 5.05, 3.05, 0.14, f"KPI status: {status_summary}", size=7.1, color=FAINT)
    trend_rows = []
    previous = None
    for row in trend:
        score = row.get("score")
        trend_rows.append({"period": row.get("label"), "score": score, "movement": score - previous if score is not None and previous is not None else None, "records": row.get("record_count")})
        if score is not None:
            previous = score
    _table(
        slide,
        0.66,
        5.42,
        12.00,
        0.24,
        [
            ("Period", 3.2, lambda row: row.get("period")),
            ("Score", 2.0, lambda row: _fmt_percent(row.get("score"))),
            ("Target", 2.0, lambda row: "100.0%"),
            ("Movement", 2.1, lambda row: _fmt_unavailable(row.get("movement"), signed=True)),
            ("Records", 2.7, lambda row: row.get("records")),
        ],
        trend_rows,
        max_rows=6,
        font_size=8.5,
    )
    _footer(slide, payload, page)


def _page_team_comparison(prs: Presentation, payload: dict[str, Any], page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    teams = payload.get("teams") or []
    priority = teams[0] if teams else {}
    title = f"{priority.get('name', 'The priority team')} has the largest team-level achievement gap"
    _header(slide, "Team comparison", title, "Only included when two or more authorized teams are present in the selected scope.", BLUE)
    _card(slide, 0.66, 1.63, 12.00, 2.32, WHITE, LINE)
    _text(slide, 0.95, 1.88, 4.5, 0.18, "RANKED TEAM SCORE", size=8.5, color=BLUE, bold=True)
    maximum = 100.0
    row_height = min(0.42, 1.55 / max(len(teams[:6]), 1))
    for index, team in enumerate(teams[:6]):
        y = 2.24 + index * row_height
        score = _number(team.get("score")) or 0
        _text(slide, 0.95, y, 2.35, row_height, f"P{index + 1}  {_clean(team.get('name'), 'Unassigned')}", size=9.5, bold=index == 0, valign=MSO_ANCHOR.MIDDLE)
        bar_left = 3.42
        bar_width = 6.35
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y + 0.10, bar_width, 0.18, PALE_GREY, PALE_GREY)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y + 0.10, max(0.06, bar_width * score / maximum), 0.18, GREEN if score >= 100 else BLUE, GREEN if score >= 100 else BLUE)
        _text(slide, 9.96, y + 0.02, 1.25, row_height, _fmt_unavailable(team.get("score")), size=9.8, color=NAVY, bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)
    _text(slide, 3.42, 3.65, 6.35, 0.16, "0%", size=7.5, color=FAINT)
    _text(slide, 8.70, 3.65, 1.07, 0.16, "100% target", size=7.5, color=FAINT, align=PP_ALIGN.RIGHT)
    rows = []
    for team in teams:
        rows.append({
            "team": team.get("name"),
            "score": _fmt_unavailable(team.get("score")),
            "movement": _fmt_unavailable(team.get("movement"), signed=True),
            "gap": _fmt_unavailable(team.get("gap"), signed=True),
            "employees": team.get("headcount"),
            "affected": team.get("affected_count"),
            "at_risk": team.get("at_risk_count"),
            "critical": team.get("critical_count"),
        })
    _table(
        slide,
        0.66,
        4.28,
        12.00,
        0.38,
        [
            ("Team", 2.30, lambda row: row.get("team")),
            ("Score", 1.25, lambda row: row.get("score")),
            ("Movement", 1.35, lambda row: row.get("movement")),
            ("Achievement gap", 1.55, lambda row: row.get("gap")),
            ("Employees", 1.45, lambda row: row.get("employees")),
            ("Affected", 1.25, lambda row: row.get("affected")),
            ("At risk", 1.40, lambda row: row.get("at_risk")),
            ("Critical", 1.25, lambda row: row.get("critical")),
        ],
        rows,
        max_rows=6,
        font_size=8.5,
        header_fill=PALE_BLUE,
    )
    _footer(slide, payload, page)


def _page_three(prs: Presentation, payload: dict[str, Any], page: int = 3):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    kpis = payload.get("kpis") or []
    driver = payload.get("driver") or {}
    title = f"{driver.get('label')} carries the largest weighted loss" if driver else "No weighted KPI loss is measurable"
    _header(slide, "Performance loss breakdown", title, "Ranked by weighted negative contribution to the selected scope", RED)
    _card(slide, 0.66, 1.63, 6.10, 2.48, WHITE, LINE)
    _text(slide, 0.94, 1.88, 4.2, 0.18, "RANKED WEIGHTED IMPACT", size=8.5, color=RED, bold=True)
    _draw_bars(slide, kpis[:6], 0.94, 2.25, 5.52, 1.55)
    _card(slide, 7.05, 1.63, 5.61, 2.48, PALE_RED, RED)
    _text(slide, 7.34, 1.88, 3.8, 0.18, "PRIORITY ORDER", size=8.5, color=RED, bold=True)
    for index, row in enumerate(kpis[:3]):
        y = 2.28 + index * 0.52
        _text(slide, 7.34, y, 0.44, 0.20, f"P{index + 1}", size=12, color=RED if index == 0 else AMBER, bold=True)
        _text(slide, 7.88, y, 2.55, 0.20, _short(row.get("label"), 30), size=10.5, bold=True)
        _text(slide, 10.58, y, 1.75, 0.20, f"Impact {_fmt_percent(row.get('weighted_impact'))}", size=9.5, color=RED if index == 0 else AMBER, bold=True, align=PP_ALIGN.RIGHT)
    if not kpis:
        _text(slide, 7.34, 2.45, 4.6, 0.45, "No KPI can be ranked because the filtered data contains no valid target-direction pair.", size=12, color=MUTED, bold=True)
    _text(slide, 7.34, 3.60, 4.8, 0.20, "Weighted performance loss = achievement shortfall × configured KPI weight.", size=9.2, color=MUTED)

    _table(
        slide,
        0.66,
        4.42,
        12.00,
        0.30,
        [
            ("KPI", 2.2, lambda row: row.get("label")),
            ("Actual", 1.42, lambda row: _fmt_value(row.get("actual"), row.get("unit"))),
            ("Target", 1.42, lambda row: _fmt_value(row.get("target"), row.get("unit"))),
            ("Achievement", 1.42, lambda row: _fmt_percent(row.get("achievement_pct"))),
            ("Achievement gap", 1.15, lambda row: _fmt_percent(row.get("gap"), signed=True)),
            ("Configured weight", 1.15, lambda row: _fmt_percent(row.get("weight_pct"))),
            ("Weighted loss", 1.30, lambda row: _fmt_percent(row.get("weighted_impact"))),
            ("Status", 1.54, lambda row: row.get("status")),
        ],
        kpis,
        max_rows=6,
        font_size=8.6,
    )
    _footer(slide, payload, page)


def _page_four(prs: Presentation, payload: dict[str, Any], page: int = 4):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    driver = payload.get("driver")
    if not driver:
        _header(slide, "Leading KPI deep dive", "No leading KPI is available", "The filtered scope has no valid measurable KPI gap.", PURPLE)
        _empty_state(slide, 0.66, 1.75, 12.0, 3.25, "No driver to investigate", "Check target completeness, KPI direction, and the selected filters before exporting again.", PURPLE)
        _footer(slide, payload, page)
        return
    direction = "Lower is better" if driver.get("direction") == "lower_better" else "Higher is better"
    _header(slide, "Leading KPI deep dive", _driver_story(payload), f"{driver.get('label')} | {direction} | Configured scoring direction", PURPLE)
    _metric_card(slide, 0.66, 1.63, 2.85, "Current actual", _fmt_value(driver.get("actual"), driver.get("unit")), f"Target {_fmt_value(driver.get('target'), driver.get('unit'))}", BLUE, PALE_BLUE, 0.95)
    _metric_card(slide, 3.70, 1.63, 2.85, "Achievement", _fmt_percent(driver.get("achievement_pct")), f"Achievement gap {_fmt_percent(driver.get('gap'), signed=True)}", RED if (driver.get("gap") or 0) < 0 else GREEN, PALE_RED if (driver.get("gap") or 0) < 0 else PALE_GREEN, 0.95)
    _metric_card(slide, 6.74, 1.63, 2.85, "Configured weight", _fmt_percent(driver.get("weight_pct")), f"Weighted performance loss {_fmt_percent(driver.get('weighted_impact'))}", RED, PALE_RED, 0.95)
    _metric_card(slide, 9.78, 1.63, 2.88, "Affected records" if payload.get("aggregate_only") else "Affected employees", str(driver.get("affected_count") if driver.get("affected_count") is not None else "Aggregate"), f"History: {payload.get('history_count') or 0} measured period(s)", AMBER, PALE_AMBER, 0.95)
    _card(slide, 0.66, 2.92, 12.00, 3.28, WHITE, LINE)
    chart_title = f"Actual versus target | {direction}"
    _text(slide, 0.95, 3.18, 4.2, 0.18, chart_title.upper(), size=8.5, color=PURPLE, bold=True)
    if driver.get("direction") == "lower_better":
        _text(slide, 8.0, 3.18, 4.25, 0.18, "Lower actual values are better; movement toward target moves down.", size=8.8, color=MUTED, align=PP_ALIGN.RIGHT)
    trend = payload.get("driver_trend") or []
    if len(trend) >= 2:
        _draw_line_chart(
            slide,
            [("Actual", [row.get("actual") for row in trend], RED if driver.get("direction") == "lower_better" else BLUE), ("Target", [row.get("target") for row in trend], AMBER)],
            [row.get("label", "") for row in trend],
            1.02,
            3.55,
            11.25,
            2.18,
        )
    else:
        _draw_actual_target_comparison(slide, driver.get("actual"), driver.get("target"), driver.get("unit"), driver.get("direction") or "higher_better", 1.02, 3.55, 7.25, 2.20)
        _text(slide, 8.55, 4.02, 3.40, 0.72, "Historical direction unavailable — only one measured period is available for this KPI.", size=13, color=MUTED, bold=True)
    status = driver.get("status") or "Data quality"
    status_color = _status_color(status)
    _text(slide, 0.95, 5.88, 4.0, 0.20, f"STATUS: {status.upper()}", size=9, color=status_color, bold=True)
    _text(slide, 5.15, 5.88, 7.0, 0.20, _short(payload.get("driver_story"), 100), size=9.3, color=MUTED, align=PP_ALIGN.RIGHT)
    _footer(slide, payload, page)


def _page_five(prs: Presentation, payload: dict[str, Any], page: int = 5):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    _header(slide, "Root cause evidence", "Measured KPI signals are not confirmed operational causes", "A poor KPI result identifies where to investigate; it does not prove why the result occurred.", AMBER)
    rows = payload.get("root_cause_rows") or []
    table_rows = []
    for row in rows:
        cause = row.get("recorded_root_cause") or "No operational cause recorded"
        if row.get("evidence_status") != "Confirmed cause" and row.get("required_validation"):
            cause = f"{cause}\nValidation: Confirm the operational cause before action closure."
        table_rows.append({
            **row,
            "cause_display": cause,
            "affected_display": row.get("affected_count") if row.get("affected_count") is not None else "Aggregate",
            "owner_display": row.get("owner_display") or "Owner needed",
        })
    if table_rows:
        _table(
            slide,
            0.66,
            1.73,
            12.00,
            0.72,
            [
                ("KPI", 1.25, lambda row: row.get("label")),
                ("Actual / target", 1.45, lambda row: f"{_fmt_value(row.get('actual'), row.get('unit'))} / {_fmt_value(row.get('target'), row.get('unit'))}"),
                ("Achievement", 0.95, lambda row: _fmt_percent(row.get("achievement_pct"))),
                ("Weighted loss", 1.00, lambda row: _fmt_percent(row.get("weighted_impact"))),
                ("Affected", 0.75, lambda row: row.get("affected_display")),
                ("Evidence status", 2.00, lambda row: row.get("evidence_status")),
                ("Owner", 1.10, lambda row: row.get("owner_display")),
                ("Cause / validation", 3.50, lambda row: row.get("cause_display")),
            ],
            table_rows,
            max_rows=4,
            font_size=7.1,
            header_fill=PALE_AMBER,
        )
    else:
        _empty_state(slide, 0.66, 1.85, 12.0, 2.4, "No KPI evidence is available", "Root cause analysis cannot proceed until the filtered snapshot contains a measured KPI with a valid target and direction.", AMBER)
    _card(slide, 0.66, 5.62, 12.00, 0.64, PALE_AMBER, AMBER)
    _text(slide, 0.96, 5.83, 11.35, 0.22, "Evidence rule: the KPI identifies where to investigate, but the operational cause has not been confirmed.", size=10.5, color=NAVY, bold=True)
    _footer(slide, payload, page)


def _page_six(prs: Presentation, payload: dict[str, Any], page: int = 6):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    if not payload.get("periods"):
        _header(slide, "People focus", "No performance data is available for the selected scope", "Person-level contribution cannot be ranked without a measured snapshot.", PURPLE)
        _empty_state(slide, 0.66, 1.78, 12.0, 3.10, "No affected employees identified", "Load a measured period before using employee-level performance loss ranking.", PURPLE)
        _footer(slide, payload, page)
        return
    count = payload.get("affected_count")
    title = f"{count} employee(s) are materially affecting the selected scope" if payload.get("people_visible") and count is not None else "Person-level detail is restricted for this role"
    people = payload.get("people") or []
    overflow_count = max(0, len(people) - 4)
    subtitle = "Ranked by weighted performance loss, then score | names and employee IDs are omitted from aggregate views"
    if overflow_count:
        subtitle = f"Top 4 weighted-loss contributors shown here; {overflow_count} additional affected employee(s) continue in the appendix."
    _header(slide, "People focus", title, subtitle, PURPLE)
    if not payload.get("people_visible"):
        _empty_state(slide, 0.66, 1.78, 12.0, 3.10, "Aggregate view", "This report contains team and KPI evidence only. Person-level employee names, IDs, root-cause details, and action details are not included for this role.", PURPLE)
        _footer(slide, payload, page)
        return
    rows = []
    for person in people:
        kpi = person.get("leading_kpi") or {}
        rows.append({
            "employee": f"{_clean(person.get('name'), 'Unknown employee')}\n{_clean(person.get('employee_id'), 'ID unavailable')}",
            "team": f"{_clean(person.get('team'), 'Unassigned')}\n{_clean(person.get('position'), 'Role unavailable')}",
            "score": _fmt_percent(person.get("score")),
            "movement": _fmt_unavailable(person.get("movement"), signed=True),
            "kpi": f"{_short(kpi.get('label'), 28)}\nActual / target: {_fmt_value(kpi.get('actual'), kpi.get('unit'))} / {_fmt_value(kpi.get('target'), kpi.get('unit'))}\nAchievement {_fmt_percent(kpi.get('achievement_pct'))} | Achievement gap {_fmt_percent(kpi.get('gap'), signed=True)} | Weight {_fmt_percent(kpi.get('weight_pct'))}",
            "loss": f"{_fmt_percent(person.get('weighted_loss_pct'))}\n{person.get('severity')}",
            "evidence": f"{person.get('root_cause_status')}\n{_short(person.get('action_status'), 22)}",
        })
    if rows:
        _table(
            slide,
            0.66,
            1.72,
            12.00,
            0.68,
            [
                ("Employee / ID", 2.00, lambda row: row.get("employee")),
                ("Team / role", 2.00, lambda row: row.get("team")),
                ("Score", 0.75, lambda row: row.get("score")),
                ("Movement", 0.95, lambda row: row.get("movement")),
                ("Leading KPI | actual / target | achievement | achievement gap | weight", 3.25, lambda row: row.get("kpi")),
                ("Loss / severity", 1.25, lambda row: row.get("loss")),
                ("Root cause / action", 1.80, lambda row: row.get("evidence")),
            ],
            rows,
            max_rows=4,
            font_size=7.0,
            header_fill=PALE_PURPLE,
        )
    if not rows:
        _empty_state(slide, 0.66, 1.85, 12.0, 2.20, "No affected employees identified", "Every employee in the selected snapshot is on track, or person-level KPI evidence is incomplete.", GREEN)
    snapshot = payload.get("score_snapshot") or []
    if snapshot:
        _text(slide, 0.70, 5.52, 3.2, 0.18, "LOWEST RELEVANT SCORE SNAPSHOT", size=8.5, color=PURPLE, bold=True)
        snapshot_rows = [{"employee": _clean(item.get("name"), "Unknown employee"), "score": _fmt_percent(item.get("score")), "severity": item.get("severity")} for item in snapshot]
        _table(
            slide,
            0.66,
            5.78,
            6.20,
            0.25,
            [("Employee", 3.6, lambda row: row.get("employee")), ("Score", 1.2, lambda row: row.get("score")), ("Severity", 1.4, lambda row: row.get("severity"))],
            snapshot_rows,
            max_rows=3,
            font_size=8.2,
            header_fill=PALE_PURPLE,
        )
    _footer(slide, payload, page)


def _page_people_appendix(prs: Presentation, payload: dict[str, Any], page: int):
    """Show overflow employee impact rows only when the scope needs them."""

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    people = (payload.get("people") or [])[4:]
    _header(
        slide,
        "Employee detail appendix",
        "Additional weighted-loss contributors require follow-through",
        "Continuation of the People Focus ranking; rows remain tied to the same filtered snapshot.",
        PURPLE,
    )
    rows = []
    for person in people:
        kpi = person.get("leading_kpi") or {}
        rows.append({
            "employee": f"{_clean(person.get('name'), 'Unknown employee')}\n{_clean(person.get('employee_id'), 'ID unavailable')}",
            "team": f"{_clean(person.get('team'), 'Unassigned')}\n{_clean(person.get('position'), 'Role unavailable')}",
            "score": _fmt_percent(person.get("score")),
            "movement": _fmt_unavailable(person.get("movement"), signed=True),
            "kpi": f"{_clean(kpi.get('label'), 'KPI')}\nActual / target: {_fmt_value(kpi.get('actual'), kpi.get('unit'))} / {_fmt_value(kpi.get('target'), kpi.get('unit'))}\nAchievement gap {_fmt_percent(kpi.get('gap'), signed=True)} | Weight {_fmt_percent(kpi.get('weight_pct'))}",
            "loss": f"{_fmt_percent(person.get('weighted_loss_pct'))}\n{person.get('severity')}",
            "evidence": f"{person.get('root_cause_status')}\n{_short(person.get('action_status'), 24)}",
        })
    _table(
        slide,
        0.66,
        1.73,
        12.00,
        0.55,
        [
            ("Employee / ID", 2.00, lambda row: row.get("employee")),
            ("Team / role", 2.00, lambda row: row.get("team")),
            ("Score", 0.75, lambda row: row.get("score")),
            ("Movement", 0.95, lambda row: row.get("movement")),
            ("Leading KPI | actual / target | achievement gap | weight", 3.25, lambda row: row.get("kpi")),
            ("Loss / severity", 1.25, lambda row: row.get("loss")),
            ("Root cause / action", 1.80, lambda row: row.get("evidence")),
        ],
        rows,
        max_rows=8,
        font_size=7.0,
        header_fill=PALE_PURPLE,
    )
    table_bottom = 1.73 + 0.55 * (1 + min(len(rows), 8))
    note_top = table_bottom + 0.25
    if note_top + 0.72 <= 7.05:
        if len(rows) > 8:
            _card(slide, 0.66, note_top, 12.00, 0.52, PALE_AMBER, AMBER)
            _text(slide, 0.96, note_top + 0.17, 11.35, 0.18, f"{len(rows) - 8} additional employee row(s) are outside this appendix page; export the underlying detail for the full authorized list.", size=9.5, color=NAVY, bold=True)
        else:
            _card(slide, 0.66, note_top, 12.00, 0.72, PALE_BLUE, BLUE)
            _text(slide, 0.96, note_top + 0.24, 11.35, 0.20, "These rows are continuation detail, not a second ranking; use the primary People Focus page for the top contributors.", size=9.5, color=NAVY, bold=True)
    _footer(slide, payload, page)


def _page_seven(prs: Presentation, payload: dict[str, Any], page: int = 7):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    _header(slide, "Action plan", "Every action needs complete accountability", "Recorded actions and transparent proposals are shown without inventing owners, deadlines, success metrics, or completion states.", GREEN)
    actions = payload.get("actions") or []
    rows = []
    for action in actions:
        related = f"KPI: {_clean(action.get('linked_kpi_key'), 'Scope KPI')}\nScope: {_short(action.get('employee_name') or action.get('team') or 'Selected scope', 34)}"
        rows.append({
            "action": _clean(action.get("action_display"), "Action text needed"),
            "kpi": _clean(action.get("linked_kpi_key"), "Scope KPI"),
            "owner": action.get("owner_display"),
            "due": action.get("due_date_display"),
            "metric": action.get("success_metric_display"),
            "status": action.get("status_display") if action.get("is_proposed") else f"Recorded\n{action.get('status_display')}",
            "evidence": action.get("evidence_display"),
        })
    if rows:
        _table(
            slide,
            0.66,
            1.73,
            12.00,
            0.60,
            [
                ("Action", 3.05, lambda row: row.get("action")),
                ("Related KPI", 1.25, lambda row: row.get("kpi")),
                ("Owner", 1.15, lambda row: row.get("owner")),
                ("Due / cadence", 1.20, lambda row: row.get("due")),
                ("Success metric", 2.15, lambda row: row.get("metric")),
                ("Status / source", 1.40, lambda row: row.get("status")),
                ("Evidence source", 1.80, lambda row: row.get("evidence")),
            ],
            rows,
            max_rows=5,
            font_size=7.8,
            header_fill=PALE_GREEN,
        )
    if not rows:
        _empty_state(slide, 0.66, 1.86, 12.0, 2.30, "No recorded or proposed action is available", "Confirm the leading KPI gap and assign an owner, review date, and measurable success metric before the next review.", GREEN)
    else:
        missing = sum(
            1
            for action in actions
            if action.get("owner_display") == "Owner needed"
            or action.get("due_date_display") == "Due date needed"
            or action.get("success_metric_display") == "Success metric needed"
        )
        table_bottom = 1.73 + 0.60 * (1 + min(len(rows), 5))
        status_top = 5.55
        if len(rows) <= 2:
            # A one-action report still needs a decision-useful readout; do
            # not leave most of the slide as empty whitespace.
            action = actions[0]
            _card(slide, 0.66, table_bottom + 0.30, 12.00, 1.18, PALE_BLUE, BLUE)
            _text(slide, 0.96, table_bottom + 0.52, 3.2, 0.18, "ACCOUNTABILITY CHECK", size=8.5, color=BLUE, bold=True)
            owner = action.get("owner_display") if action.get("owner_display") != "Owner needed" else "an owner (Owner needed)"
            due = action.get("due_date_display") if action.get("due_date_display") != "Due date needed" else "a due date or cadence (Due date needed)"
            metric = action.get("success_metric_display") if action.get("success_metric_display") != "Success metric needed" else "a measurable success metric (Success metric needed)"
            readout = f"Before the next review, assign {owner}, set {due}, and confirm {str(metric).rstrip('.')}."
            _text(slide, 0.96, table_bottom + 0.80, 11.35, 0.32, readout, size=10.5, color=NAVY, bold=True)
            status_top = table_bottom + 1.70
        _card(slide, 0.66, status_top, 12.00, 0.62, PALE_GREEN if not missing else PALE_AMBER, GREEN if not missing else AMBER)
        message = "All shown actions have an owner, deadline, and success metric." if not missing else f"{missing} action(s) still need an owner, due date, or success metric. Proposed actions are explicitly labelled."
        _text(slide, 0.96, status_top + 0.20, 11.35, 0.20, message, size=11, color=NAVY, bold=True)
    _footer(slide, payload, page)


def _page_eight(prs: Presentation, payload: dict[str, Any], page: int = 8):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    _header(slide, "Next review commitments", "Make the next review measurable", "Close the story with conditions that can be checked in the next available period.", BLUE)
    review = payload.get("next_review") or {}
    commitments = [
        ("Overall performance target", review.get("overall_target")),
        ("Leading KPI target", review.get("leading_kpi_target")),
        ("Expected movement", review.get("expected_movement")),
        ("Root cause confirmation", review.get("root_cause_requirement")),
        ("Employee review", review.get("people_requirement")),
        ("Action ownership", review.get("action_requirement")),
        ("Due / cadence", review.get("due_cadence")),
        ("Success evidence", review.get("success_evidence")),
        ("Escalation rule", review.get("escalation_rule")),
    ]
    _card(slide, 0.66, 1.72, 7.35, 4.58, WHITE, LINE)
    _text(slide, 0.96, 2.00, 3.5, 0.18, "COMMITMENTS TO CHECK", size=8.5, color=BLUE, bold=True)
    for index, (label, value) in enumerate(commitments):
        y = 2.30 + index * 0.43
        _shape(slide, MSO_SHAPE.OVAL, 0.98, y + 0.05, 0.18, 0.18, PALE_BLUE, BLUE)
        _text(slide, 1.02, y + 0.04, 0.10, 0.16, str(index + 1), size=7.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, 1.32, y, 2.2, 0.18, label, size=8.6, color=FAINT, bold=True)
        _text(slide, 3.62, y, 3.96, 0.32, value, size=9.0, color=NAVY, bold=True)
    _card(slide, 8.43, 1.72, 4.23, 4.58, PALE_BLUE, BLUE)
    _text(slide, 8.76, 2.05, 3.2, 0.18, "NEXT REVIEW QUESTION", size=8.5, color=BLUE, bold=True)
    _text(slide, 8.76, 2.47, 3.45, 1.05, review.get("question") or "Did the action plan reduce the leading KPI gap?", size=20, bold=True)
    _text(slide, 8.76, 3.82, 3.45, 0.18, "MANAGEMENT DECISION", size=8.5, color=AMBER, bold=True)
    _text(slide, 8.76, 4.08, 3.45, 0.85, review.get("decision_request") or "Management decision required: assign an owner and review date before the next reporting period.", size=11.5, color=NAVY, bold=True)
    _text(slide, 8.76, 5.23, 3.25, 0.55, "Bring the latest KPI actual, target, root-cause evidence, and action status to the review.", size=10.5, color=MUTED)
    _text(slide, 0.66, 6.58, 11.85, 0.24, "A successful next review shows measurable movement, confirmed evidence, and accountable follow-through in the same filtered scope.", size=11, color=MUTED, bold=True)
    _footer(slide, payload, page)


def _prepare_payload(report_data: dict[str, Any]) -> dict[str, Any]:
    """Compatibility wrapper for callers that imported the old helper."""

    return build_insights_snapshot(report_data)


def build_insights_pptx(period_label: str = "Selected period", report_data: dict[str, Any] | None = None) -> bytes:
    data = dict(report_data or {})
    data.setdefault("period_label", period_label)
    payload = _prepare_payload(data)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    prs.core_properties.title = "SGH Hub Insights Executive Report"
    prs.core_properties.subject = _clean(payload.get("scope_label"), "Authorized scope")
    prs.core_properties.author = "SGH Hub"

    renderers = [_page_one, _page_two]
    if len(payload.get("teams") or []) >= 2:
        renderers.append(_page_team_comparison)
    renderers.extend([_page_three, _page_four, _page_five, _page_six])
    if len(payload.get("people") or []) > 4:
        renderers.append(_page_people_appendix)
    renderers.extend([_page_seven, _page_eight])
    payload["page_count"] = len(renderers)
    for page_number, renderer in enumerate(renderers, 1):
        renderer(prs, payload, page_number)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


__all__ = ["build_insights_pptx"]
